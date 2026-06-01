"""Shared helpers for runtime-level PPT ADK HITL smoke tests."""

from __future__ import annotations

import contextlib
from typing import Any, AsyncGenerator
from unittest.mock import patch

from google.adk.models import BaseLlm
from google.adk.models.llm_request import LlmRequest
from google.adk.models.llm_response import LlmResponse
from google.genai.types import Content, FunctionCall, Part
from pptx import Presentation
from pydantic import PrivateAttr

from src.productions.ppt.ppt_product_manager import PptProductManager
from src.productions.ppt.schemas import (
    HtmlRouteBuildPackage,
    HtmlTemplatePackage,
    PptAssetResolutionResult,
    PptContentPlanningResult,
    PptContentPlanRevisionResult,
    PptRequirementAnalysisResult,
    PptRequirementRevisionResult,
    PptRouteExecutionResult,
    PptSourcePreparationResult,
    PptSystemSelectionResult,
    SourceUnderstanding,
)
from src.runtime.workspace import workspace_relative_path


class RuntimePptSmokeFakeLlm(BaseLlm):
    """Fake orchestrator model used to drive deterministic PPT runtime smoke tests."""

    _function_calls: list[FunctionCall] = PrivateAttr(default_factory=list)
    _final_text: str = PrivateAttr()
    _requests: list[LlmRequest] = PrivateAttr(default_factory=list)

    def __init__(
        self,
        *,
        function_calls: list[FunctionCall] | None = None,
        final_text: str,
    ) -> None:
        super().__init__(model="fake-runtime-ppt-smoke")
        self._function_calls = list(function_calls or [])
        self._final_text = final_text

    @property
    def requests(self) -> list[LlmRequest]:
        """Return captured LLM requests for assertions."""
        return self._requests

    async def generate_content_async(
        self,
        llm_request: LlmRequest,
        stream: bool = False,
    ) -> AsyncGenerator[LlmResponse, None]:
        """Yield one planned function call, then a harmless final fallback."""
        self._requests.append(llm_request)
        call_index = len(self._requests) - 1
        if call_index < len(self._function_calls):
            yield LlmResponse(
                content=Content(
                    role="model",
                    parts=[Part(function_call=self._function_calls[call_index])],
                )
            )
            return
        yield LlmResponse(
            content=Content(role="model", parts=[Part(text=self._final_text)])
        )


class RuntimePptSmokePatch:
    """Install deterministic Orchestrator and PPT phase stubs for runtime smokes."""

    def __init__(
        self,
        *,
        task: str,
        slide_count: int = 3,
        continue_responses: dict[int, str] | None = None,
    ) -> None:
        self.task = task
        self.slide_count = slide_count
        self.continue_responses = dict(continue_responses or {})
        self.fake_llms: list[RuntimePptSmokeFakeLlm] = []

    @contextlib.contextmanager
    def install(self):
        """Patch the runtime-level LLM and PPT provider phases."""
        with (
            patch("src.agents.orchestrator.orchestrator_agent.build_llm", side_effect=self._build_llm_stub),
            self.install_phase_stubs(),
        ):
            yield self

    @contextlib.contextmanager
    def install_phase_stubs(self):
        """Patch only the heavyweight PPT phases while leaving the Orchestrator LLM live."""
        with (
            patch.object(
                PptProductManager,
                "_prepare_initial_requirement_phase",
                self._prepare_initial_requirement_phase_stub,
            ),
            patch.object(
                PptProductManager,
                "_select_ppt_system_phase",
                self._select_ppt_system_phase_stub,
            ),
            patch.object(
                PptProductManager,
                "_run_source_preparation_phase",
                self._run_source_preparation_phase_stub,
            ),
            patch.object(
                PptProductManager,
                "_run_content_planning_phase",
                self._run_content_planning_phase_stub,
            ),
            patch.object(
                PptProductManager,
                "_revise_requirement_phase",
                self._revise_requirement_phase_stub,
            ),
            patch.object(
                PptProductManager,
                "_revise_content_plan_phase",
                self._revise_content_plan_phase_stub,
            ),
            patch.object(
                PptProductManager,
                "_run_asset_resolution_phase",
                self._run_asset_resolution_phase_stub,
            ),
            patch.object(
                PptProductManager,
                "_run_route_execution_phase",
                self._run_route_execution_phase_stub,
            ),
        ):
            yield self

    def _build_llm_stub(self, *_args: Any, **_kwargs: Any) -> RuntimePptSmokeFakeLlm:
        """Return fake Orchestrator models for initial, resumed, and final turns."""
        turn_index = len(self.fake_llms)
        if turn_index == 0:
            llm = RuntimePptSmokeFakeLlm(
                function_calls=[
                    FunctionCall(
                        name="run_ppt_product",
                        id="runtime_ppt_smoke_run_call",
                        args={
                            "task": self.task,
                            "inputs": [],
                            "output": {
                                "format": "pptx",
                                "route": "html",
                                "slide_count": self.slide_count,
                            },
                        },
                    )
                ],
                final_text='{"reply_text":"runtime ppt smoke initial fallback","final_file_paths":[]}',
            )
        elif turn_index >= 2:
            llm = RuntimePptSmokeFakeLlm(
                function_calls=[
                    FunctionCall(
                        name="continue_ppt_product",
                        id="runtime_ppt_smoke_continue_call",
                        args={"user_response": self.continue_responses.get(turn_index, "确认")},
                    )
                ],
                final_text='{"reply_text":"runtime ppt smoke final fallback","final_file_paths":[]}',
            )
        else:
            llm = RuntimePptSmokeFakeLlm(
                final_text='{"reply_text":"runtime ppt smoke resumed","final_file_paths":[]}'
            )
        self.fake_llms.append(llm)
        return llm

    @staticmethod
    async def _prepare_initial_requirement_phase_stub(
        manager: PptProductManager,
        **kwargs: Any,
    ) -> PptRequirementAnalysisResult:
        """Prepare a deterministic requirement while preserving product state writes."""
        requirement = manager.prepare_confirmed_requirement(
            task=kwargs["task"],
            inputs=kwargs["raw_inputs"],
            output=kwargs["output"],
            source_understanding=kwargs["source_understanding"],
        )
        result = PptRequirementAnalysisResult(
            confirmed_requirement=requirement,
            analysis_output={"source": "runtime_smoke_stub"},
            agent_message="runtime smoke requirement analysis",
        )
        manager._persist_requirement_analysis_result(kwargs["tool_context"], result)
        return result

    @staticmethod
    async def _select_ppt_system_phase_stub(
        manager: PptProductManager,
        **kwargs: Any,
    ) -> PptSystemSelectionResult:
        """Select the built-in HTML route deterministically."""
        result = PptSystemSelectionResult(
            system_selection=manager._build_default_system_selection(kwargs["requirement"]),
            selection_output={"source": "runtime_smoke_stub"},
            agent_message="runtime smoke system selection",
        )
        manager._persist_system_selection_result(kwargs["tool_context"], result)
        return result

    @staticmethod
    async def _run_source_preparation_phase_stub(
        manager: PptProductManager,
        **kwargs: Any,
    ) -> PptSourcePreparationResult:
        """Return an empty prepared-source result for source-free smoke tasks."""
        result = PptSourcePreparationResult(
            source_inputs=[],
            source_materials=SourceUnderstanding(document_type="brief"),
            input_signature="runtime-smoke",
        )
        manager._persist_source_preparation_result(kwargs["tool_context"], result)
        return result

    @staticmethod
    async def _run_content_planning_phase_stub(
        manager: PptProductManager,
        requirement: Any,
        **kwargs: Any,
    ) -> PptContentPlanningResult:
        """Build a deterministic deck plan without external provider calls."""
        content_plan = manager.build_initial_deck_content_plan(requirement)
        result = PptContentPlanningResult(
            content_plan=content_plan,
            deck_content_plan_markdown="runtime smoke content plan",
            planning_output={"source": "runtime_smoke_stub"},
        )
        manager._persist_content_planning_result(kwargs["tool_context"], result)
        return result

    @staticmethod
    async def _revise_requirement_phase_stub(
        manager: PptProductManager,
        **kwargs: Any,
    ) -> PptRequirementRevisionResult:
        """Apply deterministic requirement edits without provider calls."""
        revised_requirement = manager._revise_confirmed_requirement_deterministically(
            kwargs["existing_requirement"],
            user_response=kwargs["user_response"],
            raw_inputs=kwargs["raw_inputs"],
            output=kwargs["output"],
            source_understanding=kwargs["source_understanding"],
        )
        result = PptRequirementRevisionResult(
            confirmed_requirement=revised_requirement,
            revision_output={"source": "runtime_smoke_stub"},
            agent_message="runtime smoke requirement revision",
            user_revision=kwargs["user_response"],
        )
        manager._persist_requirement_revision_result(kwargs["tool_context"], result)
        return result

    @staticmethod
    async def _revise_content_plan_phase_stub(
        manager: PptProductManager,
        *,
        requirement: Any,
        user_response: str,
        **kwargs: Any,
    ) -> PptContentPlanRevisionResult:
        """Regenerate a deterministic content plan after a content-plan revision."""
        revised_requirement = requirement.model_copy(
            update={
                "request_brief": manager._append_user_revision(
                    requirement.request_brief,
                    user_response,
                    label="Content plan revision",
                )
            }
        )
        content_plan = manager.build_initial_deck_content_plan(revised_requirement)
        result = PptContentPlanRevisionResult(
            confirmed_requirement=revised_requirement,
            content_plan=content_plan,
            deck_content_plan_markdown="runtime smoke revised content plan",
            revision_output={"source": "runtime_smoke_stub"},
            user_revision=user_response,
        )
        manager._persist_content_plan_revision_result(kwargs["tool_context"], result)
        return result

    @staticmethod
    async def _run_asset_resolution_phase_stub(
        manager: PptProductManager,
        content_plan: Any,
        requirement: Any,
        **kwargs: Any,
    ) -> PptAssetResolutionResult:
        """Keep planned placeholder assets unchanged for the route smoke."""
        result = PptAssetResolutionResult(
            content_plan=content_plan,
            input_signature="runtime-smoke-assets",
            resolved_asset_manifest={"source": "runtime_smoke_stub"},
        )
        manager._persist_asset_resolution_result(kwargs["tool_context"], result)
        return result

    @staticmethod
    async def _run_route_execution_phase_stub(
        manager: PptProductManager,
        *,
        requirement: Any,
        content_plan: Any,
        tool_context: Any,
        **_kwargs: Any,
    ) -> PptRouteExecutionResult:
        """Create minimal real route artifacts and persist the route phase result."""
        output_dir = manager._build_route_output_dir(tool_context.state, route=requirement.route)
        output_dir_ref = manager._route_output_dir_reference(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)

        pptx_path = output_dir / "runtime_smoke.pptx"
        html_path = output_dir / "runtime_smoke.html"
        quality_path = output_dir / "quality_report.json"
        build_log_path = output_dir / "build_log.json"

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(914400, 914400, 7315200, 914400)
        textbox.text = "Runtime PPT smoke"
        presentation.save(pptx_path)

        html_path.write_text("<!doctype html><html><body>Runtime PPT smoke</body></html>", encoding="utf-8")
        quality_path.write_text('{"status":"pass"}', encoding="utf-8")
        build_log_path.write_text('{"source":"runtime_smoke_stub"}', encoding="utf-8")

        route_build = HtmlRouteBuildPackage(
            template=HtmlTemplatePackage(template_id="runtime-smoke", label="Runtime Smoke"),
            html_deck_path=workspace_relative_path(html_path),
            pptx_path=workspace_relative_path(pptx_path),
            quality_report_path=workspace_relative_path(quality_path),
            build_log_path=workspace_relative_path(build_log_path),
        )
        result = PptRouteExecutionResult(
            route=requirement.route,
            output_dir=output_dir_ref,
            input_signature=manager._route_execution_input_signature(requirement, content_plan),
            route_build=route_build,
        )
        manager._persist_route_execution_result(tool_context, result)
        return result
