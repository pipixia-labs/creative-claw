import json
import tempfile
import tomllib
import unittest
import zipfile
from pathlib import Path
from types import SimpleNamespace
from typing import AsyncGenerator
from unittest.mock import patch

from google.adk import Context, Workflow
from google.adk.artifacts import InMemoryArtifactService
from google.adk.models import BaseLlm
from google.adk.models.llm_request import LlmRequest
from google.adk.models.llm_response import LlmResponse
from google.adk.runners import Runner
from google.adk.sessions import InMemorySessionService
from google.adk.workflow import node
from google.genai.types import Content, FunctionCall, Part
from PIL import Image
from pptx import Presentation
from pydantic import PrivateAttr

from src.productions.ppt.ppt_product_manager import PptProductManager
from src.productions.ppt.routes.svg import (
    PPT_DESIGN_STRATEGY_STATE_KEY,
    PPT_SVG_EXECUTION_PLAN_STATE_KEY,
    PPT_SVG_ROUTE_GENERATED_PAGES_KEY,
    SVG_DELIVERY_STAGE,
    SVG_ROUTE_STAGE_SEQUENCE,
    build_default_svg_design_strategy,
    build_ppt_design_strategy_expert,
    build_ppt_svg_deck_executor_expert,
    build_svg_design_strategy_with_agent,
    build_svg_route,
    check_svg_pages_quality,
    export_svg_pages_to_pptx,
    generate_svg_pages_with_agent,
    prepare_svg_route_paths,
)
from src.productions.ppt.routes.svg import route as svg_route
from src.productions.ppt.schemas import ConfirmedRequirement, PptSvgExecutionPlan
from src.productions.ppt.schemas import DeckContentPlan, DeckPagePlan
from src.productions.ppt.templates.svg import (
    list_svg_layout_templates,
    load_svg_layout_template_package,
    select_svg_layout_template_match,
)
from src.runtime.workspace import resolve_workspace_path, workspace_root


class _SvgRouteToolCallingFakeLlm(BaseLlm):
    """Fake SVG-route model that drives a planned sequence of route tools."""

    _function_calls: list[FunctionCall] = PrivateAttr()
    _requests: list[LlmRequest] = PrivateAttr(default_factory=list)

    def __init__(self, *, function_calls: list[FunctionCall]) -> None:
        super().__init__(model="fake-ppt-svg-route")
        self._function_calls = function_calls

    @property
    def requests(self) -> list[LlmRequest]:
        return self._requests

    async def generate_content_async(
        self,
        llm_request: LlmRequest,
        stream: bool = False,
    ) -> AsyncGenerator[LlmResponse, None]:
        self._requests.append(llm_request)
        call_index = len(self._requests) - 1
        if call_index < len(self._function_calls):
            yield LlmResponse(
                content=Content(
                    role="model",
                    parts=[Part(function_call=self._function_calls[call_index])],
                )
            )
            return
        yield LlmResponse(
            content=Content(
                role="model",
                parts=[Part(text="SVG route agent completed.")],
            )
        )


def _function_declaration_names(llm_request: LlmRequest) -> list[str]:
    names: list[str] = []
    for tool in llm_request.config.tools or []:
        for declaration in tool.function_declarations or []:
            if declaration.name:
                names.append(declaration.name)
    return names


def _simple_svg(title: str, *, width: int = 1280, height: int = 720) -> str:
    return (
        f"<svg xmlns='http://www.w3.org/2000/svg' width='{width}' height='{height}' "
        f"viewBox='0 0 {width} {height}'>"
        "<rect x='0' y='0' width='100%' height='100%' fill='#FFFFFF'/>"
        f"<text x='72' y='120' font-size='40' fill='#172033'>{title}</text>"
        "</svg>"
    )


class PptSvgRouteTests(unittest.TestCase):
    def test_svg_layout_templates_are_included_in_package_data(self) -> None:
        pyproject_path = Path(__file__).resolve().parents[1] / "pyproject.toml"
        pyproject = tomllib.loads(pyproject_path.read_text(encoding="utf-8"))
        include_patterns = pyproject["tool"]["poetry"]["include"]

        self.assertIn("src/productions/ppt/templates/svg/**/*", include_patterns)

    def test_svg_layout_template_registry_loads_ppt_master_layouts(self) -> None:
        templates = list_svg_layout_templates()
        template = load_svg_layout_template_package("academic_defense")

        self.assertGreaterEqual(len(templates), 21)
        self.assertEqual(template.template_id, "academic_defense")
        self.assertIn("Academic Defense", template.label)
        self.assertIn("cover", template.page_svgs)
        self.assertIn("content", template.page_svgs)
        self.assertTrue(template.design_spec_path.exists())
        self.assertGreaterEqual(len(template.palette), 1)

    def test_svg_layout_template_selector_matches_task_and_explicit_template(self) -> None:
        auto_match = select_svg_layout_template_match(task="做一个毕业论文答辩 PPT，走 svg route。")
        explicit_match = select_svg_layout_template_match(task="做一个普通汇报。", template_id="mckinsey")
        weak_match = select_svg_layout_template_match(task="做一个普通产品介绍 PPT。")

        self.assertTrue(auto_match.use_template)
        self.assertEqual(auto_match.template_id, "academic_defense")
        self.assertGreaterEqual(auto_match.score, 35)
        self.assertTrue(explicit_match.use_template)
        self.assertTrue(explicit_match.explicit)
        self.assertEqual(explicit_match.template_id, "mckinsey")
        self.assertFalse(weak_match.use_template)

    def test_svg_layout_template_selector_covers_core_scenarios(self) -> None:
        cases = {
            "给董事会做一份麦肯锡风格的五年战略咨询汇报": "mckinsey",
            "做一个 Google 风格的技术分享和数据展示 PPT": "google_style",
            "做一个党建红色政府工作报告 PPT": "government_red",
            "做一个智慧城市数字政府治理方案汇报": "government_blue",
            "做一个医院病例研究和医学科研汇报": "medical_university",
        }

        for task, expected_template_id in cases.items():
            with self.subTest(task=task):
                match = select_svg_layout_template_match(task=task)
                self.assertTrue(match.use_template)
                self.assertEqual(match.template_id, expected_template_id)

    def test_svg_layout_template_selector_reports_unknown_explicit_template(self) -> None:
        match = select_svg_layout_template_match(task="做一个普通汇报。", template_id="missing_template")

        self.assertFalse(match.use_template)
        self.assertTrue(match.explicit)
        self.assertEqual(match.template_id, "missing_template")
        self.assertIn("Unknown SVG layout template", match.fallback_reason)

    def test_svg_layout_template_selector_requires_primary_task_signal(self) -> None:
        requirement = ConfirmedRequirement(
            route="svg",
            request_brief="做一个普通产品介绍 PPT。",
            topic="产品介绍",
            aspect_ratio="16:9",
        )
        content_plan = DeckContentPlan(
            title="毕业论文答辩",
            core_narrative="研究背景、论文方法和学术答辩结构都很完整。",
            pages=[
                DeckPagePlan(
                    slide_number=1,
                    page_type="content",
                    title="论文研究背景",
                    purpose="Explain the academic research context.",
                    key_takeaway="研究问题具有明确学术价值。",
                )
            ],
        )

        match = select_svg_layout_template_match(
            requirement=requirement,
            content_plan=content_plan,
        )

        self.assertFalse(match.use_template)
        self.assertIn("primary task signal", match.fallback_reason)

    def test_svg_route_build_log_records_auto_selected_layout_template(self) -> None:
        content_plan = DeckContentPlan(
            title="毕业论文答辩",
            core_narrative="用于学术论文答辩，说明研究背景、方法和结论。",
            pages=[
                DeckPagePlan(
                    slide_number=1,
                    page_type="cover",
                    title="毕业论文答辩",
                    purpose="Open the academic defense.",
                    key_takeaway="研究主题和答辩人信息清晰呈现。",
                ),
                DeckPagePlan(
                    slide_number=2,
                    page_type="content",
                    title="研究背景",
                    purpose="Explain the research context.",
                    key_takeaway="研究问题具有明确价值。",
                ),
            ],
        )

        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            route_build = build_svg_route(
                content_plan=content_plan,
                output_dir=Path(tmpdir),
                aspect_ratio="16:9",
            )

            build_log = json.loads(resolve_workspace_path(route_build.build_log_path).read_text(encoding="utf-8"))
            quality_report = json.loads(resolve_workspace_path(route_build.quality_report_path).read_text(encoding="utf-8"))

            self.assertTrue(route_build.svg_layout_template_selection["use_template"])
            self.assertEqual(route_build.svg_layout_template_selection["template_id"], "academic_defense")
            self.assertEqual(build_log["svg_layout_template_selection"]["template_id"], "academic_defense")
            self.assertEqual(quality_report["svg_layout_template_selection"]["template_id"], "academic_defense")

    def test_svg_route_skips_layout_template_for_4_3_requests(self) -> None:
        content_plan = DeckContentPlan(
            title="毕业论文答辩",
            core_narrative="用于学术论文答辩，说明研究背景、方法和结论。",
            pages=[
                DeckPagePlan(
                    slide_number=1,
                    page_type="content",
                    title="研究背景",
                    purpose="Explain the research context.",
                    key_takeaway="研究问题具有明确价值。",
                ),
                DeckPagePlan(
                    slide_number=2,
                    page_type="content",
                    title="研究方法",
                    purpose="Explain the research method.",
                    key_takeaway="方法设计体现可控创意生成。",
                ),
            ],
        )

        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            route_build = build_svg_route(
                content_plan=content_plan,
                output_dir=Path(tmpdir),
                aspect_ratio="4:3",
            )

            selection = route_build.svg_layout_template_selection
            self.assertFalse(selection["use_template"])
            self.assertEqual(selection["template_id"], "academic_defense")
            self.assertIn("16:9", selection["fallback_reason"])

    def test_svg_execution_plan_adds_rhythm_and_page_type_guidance(self) -> None:
        content_plan = DeckContentPlan(
            title="战略增长汇报",
            core_narrative="用于董事会战略讨论，突出增长路径和关键选择。",
            pages=[
                DeckPagePlan(
                    slide_number=1,
                    page_type="cover",
                    title="战略增长汇报",
                    purpose="Open the strategy deck.",
                    key_takeaway="明确增长主题和汇报对象。",
                ),
                DeckPagePlan(
                    slide_number=2,
                    page_type="comparison",
                    title="两条增长路径对比",
                    purpose="Compare strategic options.",
                    key_takeaway="不同路径的资源要求和收益结构不同。",
                ),
                DeckPagePlan(
                    slide_number=3,
                    page_type="quote",
                    title="核心判断",
                    purpose="Emphasize one executive message.",
                    key_takeaway="增长质量比增长速度更重要。",
                ),
            ],
        )
        requirement = ConfirmedRequirement(
            route="svg",
            request_brief="给董事会做一份麦肯锡风格的五年战略咨询汇报",
            topic=content_plan.title,
            aspect_ratio="16:9",
        )
        template_match = select_svg_layout_template_match(
            requirement=requirement,
            content_plan=content_plan,
        )

        design_stage = build_default_svg_design_strategy(
            requirement=requirement,
            content_plan=content_plan,
            template_match=template_match,
        )
        execution_plan = design_stage.execution_plan
        page_generation_plan = svg_route._svg_page_generation_plan(
            content_plan=content_plan,
            design_stage=design_stage,
        )

        self.assertEqual(execution_plan.page_rhythm_by_slide["P01"], "anchor")
        self.assertEqual(execution_plan.page_rhythm_by_slide["P02"], "dense")
        self.assertEqual(execution_plan.page_rhythm_by_slide["P03"], "breathing")
        self.assertGreaterEqual(execution_plan.typography_ramp["cover_title"], 60)
        self.assertIn("comparison", execution_plan.page_type_layout_guidance)
        self.assertIn("breathing", execution_plan.page_rhythm_guidance)
        self.assertTrue(template_match.use_template)
        self.assertEqual(page_generation_plan[0]["template_reference"]["reference_page_type"], "cover")
        self.assertEqual(page_generation_plan[1]["template_reference"]["reference_page_type"], "content")
        self.assertIn("03_content.svg", page_generation_plan[1]["template_reference"]["template_file"])
        self.assertTrue(page_generation_plan[1]["template_reference"]["svg_excerpt"])
        self.assertEqual(page_generation_plan[2]["page_rhythm"], "breathing")

    def test_svg_route_builds_svg_pages_quality_report_and_editable_pptx(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 3 页 PPTX 产品介绍，使用 SVG route。",
            inputs=[],
            output={"format": "pptx", "route": "svg", "slide_count": 3},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)

        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            route_build = build_svg_route(
                content_plan=content_plan,
                output_dir=Path(tmpdir),
                aspect_ratio="16:9",
            )

            pptx_path = resolve_workspace_path(route_build.pptx_path)
            quality_path = resolve_workspace_path(route_build.quality_report_path)
            build_log_path = resolve_workspace_path(route_build.build_log_path)
            svg_paths = [resolve_workspace_path(path) for path in route_build.svg_page_paths]
            prs = Presentation(str(pptx_path))
            pptx_text = "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
            )
            quality_report = json.loads(quality_path.read_text(encoding="utf-8"))
            build_log = json.loads(build_log_path.read_text(encoding="utf-8"))

            self.assertTrue(pptx_path.exists())
            self.assertTrue(all(path.exists() for path in svg_paths))
            self.assertEqual(len(route_build.svg_page_paths), len(content_plan.pages))
            self.assertEqual(len(prs.slides), len(content_plan.pages))
            self.assertIn(content_plan.pages[0].title, pptx_text)
            self.assertEqual(quality_report["status"], "pass")
            self.assertEqual(quality_report["route_stages"], list(SVG_ROUTE_STAGE_SEQUENCE))
            self.assertEqual(quality_report["delivery_stage"], SVG_DELIVERY_STAGE)
            self.assertEqual(build_log["workflow_name"], "SvgRouteSequentialAgent")
            self.assertEqual(build_log["route_stages"], list(SVG_ROUTE_STAGE_SEQUENCE))
            self.assertEqual(build_log["pptx_conversion"]["engine"], "native_drawingml_svg_converter")
            self.assertEqual(build_log["pptx_conversion"]["editable_level"], "high")

    def test_svg_quality_checker_reports_invalid_svg(self) -> None:
        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            svg_path = Path(tmpdir) / "bad.svg"
            svg_path.write_text("<svg viewBox='0 0 10 10'><script /></svg>", encoding="utf-8")

            report = check_svg_pages_quality(
                svg_page_paths=[str(svg_path)],
                expected_page_count=1,
                execution_plan=PptSvgExecutionPlan(),
            )

            self.assertEqual(report.status, "failed")
            self.assertFalse(report.checks["all_viewboxes_match_canvas"])
            self.assertFalse(report.checks["no_unsupported_tags"])

    def test_svg_quality_checker_rejects_malformed_path(self) -> None:
        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            svg_path = Path(tmpdir) / "bad_path.svg"
            svg_path.write_text(
                "<svg xmlns='http://www.w3.org/2000/svg' width='1280' height='720' viewBox='0 0 1280 720'><path d='M 10'/></svg>",
                encoding="utf-8",
            )

            report = check_svg_pages_quality(
                svg_page_paths=[str(svg_path)],
                expected_page_count=1,
                execution_plan=PptSvgExecutionPlan(),
            )

            self.assertEqual(report.status, "failed")
            self.assertFalse(report.checks["all_visual_elements_convertible"])
            self.assertIn("malformed_svg_path", {issue["code"] for issue in report.issues})

    def test_svg_expert_builders_use_ppt_svg_tools(self) -> None:
        manager = PptProductManager()
        design_agent = build_ppt_design_strategy_expert(
            save_design_strategy_tool=manager.save_ppt_design_strategy,
            save_svg_execution_plan_tool=manager.save_ppt_svg_execution_plan,
        )
        executor_agent = build_ppt_svg_deck_executor_expert(
            read_svg_execution_plan_tool=manager.read_ppt_svg_execution_plan,
            save_svg_page_tool=manager.save_ppt_svg_page,
        )

        self.assertEqual(design_agent.name, "PptDesignStrategyExpert")
        self.assertEqual(executor_agent.name, "PptSvgDeckExecutorExpert")
        self.assertEqual(
            {tool.__name__ for tool in design_agent.tools},
            {"save_ppt_design_strategy", "save_ppt_svg_execution_plan"},
        )
        self.assertEqual(
            {tool.__name__ for tool in executor_agent.tools},
            {"read_ppt_svg_execution_plan", "save_ppt_svg_page"},
        )

    def test_save_ppt_svg_page_tool_writes_route_page(self) -> None:
        manager = PptProductManager()
        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            paths = prepare_svg_route_paths(Path(tmpdir))
            tool_context = SimpleNamespace(
                state={
                    "sid": "svg-page-tool-test",
                    "turn_index": 1,
                    "step": 1,
                    "ppt_svg_route_output_dir": str(paths.output_dir),
                }
            )
            result = manager.save_ppt_svg_execution_plan(
                {"aspect_ratio": "16:9", "canvas_width": 1280, "canvas_height": 720},
                tool_context,
            )
            self.assertEqual(result["status"], "success")

            save_result = manager.save_ppt_svg_page(
                slide_number=1,
                file_name="slide_001.svg",
                title="Cover",
                page_type="cover",
                page_rhythm="anchor",
                svg_content="<svg xmlns='http://www.w3.org/2000/svg' width='1280' height='720' viewBox='0 0 1280 720'><text x='72' y='120'>Cover</text></svg>",
                tool_context=tool_context,
            )

            self.assertEqual(save_result["status"], "success")
            saved_page = tool_context.state[PPT_SVG_ROUTE_GENERATED_PAGES_KEY][0]
            self.assertEqual(saved_page["slide_number"], 1)
            self.assertTrue(resolve_workspace_path(saved_page["svg_path"]).exists())

    def test_native_svg_export_supports_paths_and_images(self) -> None:
        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            output_dir = Path(tmpdir)
            svg_path = output_dir / "native.svg"
            pptx_path = output_dir / "native.pptx"
            svg_path.write_text(
                """<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
  <rect x="0" y="0" width="1280" height="720" fill="#FFFFFF"/>
  <path d="M 80 120 L 260 120 L 220 220 Z" fill="#2457D6" stroke="#172033" stroke-width="2"/>
  <polyline points="320,120 390,170 460,120" fill="none" stroke="#43A6FF" stroke-width="4"/>
  <image x="80" y="280" width="80" height="80" href="data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="/>
  <text x="200" y="330" data-width="600" font-family="Aptos" font-size="32" fill="#172033"><tspan x="200" dy="0">Native SVG Text</tspan><tspan x="200" dy="42">Second line</tspan></text>
</svg>""",
                encoding="utf-8",
            )

            report = check_svg_pages_quality(
                svg_page_paths=[str(svg_path)],
                expected_page_count=1,
                execution_plan=PptSvgExecutionPlan(),
            )
            export = export_svg_pages_to_pptx(
                svg_page_paths=[str(svg_path)],
                pptx_path=pptx_path,
                execution_plan=PptSvgExecutionPlan(),
            )
            prs = Presentation(str(pptx_path))
            pptx_text = "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
            )

            self.assertEqual(report.status, "pass")
            self.assertTrue(pptx_path.exists())
            self.assertTrue(export.conversion_report["ok"])
            self.assertEqual(export.conversion_report["engine"], "native_drawingml_svg_converter")
            self.assertEqual(export.conversion_report["pages"][0]["image_count"], 1)
            self.assertIn("Native SVG Text", pptx_text)

    def test_native_svg_export_supports_ppt_master_baseline_features(self) -> None:
        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            output_dir = Path(tmpdir)
            image_path = output_dir / "wide.png"
            svg_path = output_dir / "baseline.svg"
            pptx_path = output_dir / "baseline.pptx"
            Image.new("RGB", (4, 2), "#2457D6").save(image_path)
            svg_path.write_text(
                f"""<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
  <defs>
    <linearGradient id="accentGrad" x1="0%" y1="0%" x2="100%" y2="0%">
      <stop offset="0%" stop-color="#2457D6"/>
      <stop offset="100%" stop-color="#43A6FF"/>
    </linearGradient>
    <marker id="arrow" markerWidth="8" markerHeight="8"><path d="M 0 0 L 8 4 L 0 8 Z"/></marker>
    <clipPath id="imageClip"><circle cx="980" cy="260" r="70"/></clipPath>
    <filter id="softShadow"><feGaussianBlur stdDeviation="4"/><feOffset dx="3" dy="4"/><feFlood flood-color="#000000" flood-opacity="0.25"/></filter>
  </defs>
  <g id="editableGroup" filter="url(#softShadow)">
    <rect x="64" y="64" width="500" height="150" rx="16" fill="url(#accentGrad)" stroke="#172033" stroke-width="2"/>
    <path d="M 96 360 C 160 300 220 420 280 360 S 400 300 460 360 Q 520 430 580 360 T 700 360 A 42 42 0 0 1 784 360" fill="none" stroke="#172033" stroke-width="4" marker-end="url(#arrow)"/>
    <image x="900" y="190" width="160" height="140" href="{image_path.name}" preserveAspectRatio="xMidYMid slice" clip-path="url(#imageClip)"/>
    <text x="96" y="150" data-width="420" font-family="Aptos" font-size="34" fill="#FFFFFF"><tspan font-weight="bold">Rich</tspan><tspan fill="#F4C542" font-style="italic"> Text</tspan></text>
  </g>
</svg>""",
                encoding="utf-8",
            )

            report = check_svg_pages_quality(
                svg_page_paths=[str(svg_path)],
                expected_page_count=1,
                execution_plan=PptSvgExecutionPlan(),
            )
            export = export_svg_pages_to_pptx(
                svg_page_paths=[str(svg_path)],
                pptx_path=pptx_path,
                execution_plan=PptSvgExecutionPlan(),
            )
            prs = Presentation(str(pptx_path))
            with zipfile.ZipFile(pptx_path) as pptx_zip:
                slide_xml = pptx_zip.read("ppt/slides/slide1.xml").decode("utf-8")

            self.assertEqual(report.status, "pass")
            self.assertTrue(export.conversion_report["ok"])
            self.assertEqual(len(prs.slides), 1)
            self.assertIn("<a:gradFill", slide_xml)
            self.assertIn("<a:tailEnd", slide_xml)
            self.assertIn("<a:srcRect", slide_xml)
            self.assertIn("<p:grpSp>", slide_xml)
            self.assertIn("<a:outerShdw", slide_xml)
            self.assertIn("<a:cubicBezTo>", slide_xml)

    def test_native_svg_export_failure_does_not_replace_existing_pptx(self) -> None:
        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            output_dir = Path(tmpdir)
            svg_path = output_dir / "bad.svg"
            pptx_path = output_dir / "deck.pptx"
            original_bytes = b"existing-pptx-placeholder"
            pptx_path.write_bytes(original_bytes)
            svg_path.write_text(
                "<svg xmlns='http://www.w3.org/2000/svg' width='1280' height='720' viewBox='0 0 1280 720'><path d='M 0 0 R 20 20 40 40'/></svg>",
                encoding="utf-8",
            )

            export = export_svg_pages_to_pptx(
                svg_page_paths=[str(svg_path)],
                pptx_path=pptx_path,
                execution_plan=PptSvgExecutionPlan(),
            )

            self.assertFalse(export.conversion_report["ok"])
            self.assertEqual(export.conversion_report["final_strategy"], "failed")
            self.assertEqual(pptx_path.read_bytes(), original_bytes)

    def test_save_ppt_svg_page_rejects_unsupported_native_subset(self) -> None:
        manager = PptProductManager()
        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            paths = prepare_svg_route_paths(Path(tmpdir))
            tool_context = SimpleNamespace(
                state={
                    "sid": "svg-page-tool-test",
                    "turn_index": 1,
                    "step": 1,
                    "ppt_svg_route_output_dir": str(paths.output_dir),
                }
            )
            manager.save_ppt_svg_execution_plan(
                {"aspect_ratio": "16:9", "canvas_width": 1280, "canvas_height": 720},
                tool_context,
            )

            with self.assertRaisesRegex(ValueError, "native PPTX converter subset"):
                manager.save_ppt_svg_page(
                    slide_number=1,
                    file_name="bad.svg",
                    title="Bad",
                    page_type="content",
                    page_rhythm="body",
                    svg_content="<svg xmlns='http://www.w3.org/2000/svg' width='1280' height='720' viewBox='0 0 1280 720'><foreignObject /></svg>",
                    tool_context=tool_context,
                )


class PptSvgRouteAgentToolTests(unittest.IsolatedAsyncioTestCase):
    async def test_svg_design_strategy_agenttool_main_path_saves_strategy(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 2 页 SVG route 产品介绍 PPT。",
            inputs=[],
            output={"format": "pptx", "route": "svg", "slide_count": 2},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        content_plan.pages = content_plan.pages[:2]
        template_match = select_svg_layout_template_match(
            requirement=requirement,
            content_plan=content_plan,
        )
        fallback = build_default_svg_design_strategy(
            requirement=requirement,
            content_plan=content_plan,
            template_match=template_match,
        )
        fake_llm = _SvgRouteToolCallingFakeLlm(
            function_calls=[
                FunctionCall(
                    name="save_ppt_design_strategy",
                    args={
                        "strategy_json": fallback.strategy.model_dump(mode="json"),
                        "confirmation_json": fallback.confirmation.model_dump(mode="json"),
                    },
                ),
                FunctionCall(
                    name="save_ppt_svg_execution_plan",
                    args={
                        "execution_plan_json": fallback.execution_plan.model_dump(mode="json"),
                    },
                ),
            ]
        )
        session_service = InMemorySessionService()
        artifact_service = InMemoryArtifactService()

        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            paths = prepare_svg_route_paths(Path(tmpdir) / "svg_route")
            with patch.object(svg_route, "build_llm", return_value=fake_llm):
                design_agent = build_ppt_design_strategy_expert(
                    save_design_strategy_tool=manager.save_ppt_design_strategy,
                    save_svg_execution_plan_tool=manager.save_ppt_svg_execution_plan,
                )

            @node(name="PptSvgDesignAgentToolHarnessNode", rerun_on_resume=True)
            async def design_harness(ctx: Context, node_input: str) -> dict:
                result = await build_svg_design_strategy_with_agent(
                    requirement=requirement,
                    content_plan=content_plan,
                    template_match=template_match,
                    paths=paths,
                    tool_context=ctx,
                    app_name="creative_claw",
                    artifact_service=artifact_service,
                    design_strategy_agent=design_agent,
                )
                ctx.state["svg_design_generation_mode_for_test"] = result.generation_mode
                return {"generation_mode": result.generation_mode}

            workflow = Workflow(
                name="PptSvgDesignAgentToolHarness",
                edges=[("START", design_harness)],
            )
            runner = Runner(
                node=workflow,
                session_service=session_service,
                artifact_service=artifact_service,
            )
            user_id = "user-ppt-svg-design"
            session_id = "session-ppt-svg-design"

            try:
                await session_service.create_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                    state={
                        "sid": "ppt-svg-design-agenttool-test",
                        "turn_index": 1,
                        "step": 1,
                    },
                )
                async for _ in runner.run_async(
                    user_id=user_id,
                    session_id=session_id,
                    new_message=Content(role="user", parts=[Part(text="Build design")]),
                ):
                    pass
                session = await session_service.get_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                )
            finally:
                await runner.close()

        self.assertIsNotNone(session)
        self.assertEqual(session.state["svg_design_generation_mode_for_test"], "llm_agent_design_strategy")
        self.assertEqual(session.state[PPT_DESIGN_STRATEGY_STATE_KEY]["style_name"], fallback.strategy.style_name)
        self.assertEqual(
            session.state[PPT_SVG_EXECUTION_PLAN_STATE_KEY]["canvas_width"],
            fallback.execution_plan.canvas_width,
        )
        first_request_tools = _function_declaration_names(fake_llm.requests[0])
        self.assertIn("save_ppt_design_strategy", first_request_tools)
        self.assertIn("save_ppt_svg_execution_plan", first_request_tools)

    async def test_svg_deck_executor_agenttool_main_path_saves_pages(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 2 页 SVG route 产品介绍 PPT。",
            inputs=[],
            output={"format": "pptx", "route": "svg", "slide_count": 2},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        content_plan.pages = content_plan.pages[:2]
        template_match = select_svg_layout_template_match(
            requirement=requirement,
            content_plan=content_plan,
        )
        design_stage = build_default_svg_design_strategy(
            requirement=requirement,
            content_plan=content_plan,
            template_match=template_match,
        )
        width = design_stage.execution_plan.canvas_width
        height = design_stage.execution_plan.canvas_height
        function_calls: list[FunctionCall] = []
        for page in content_plan.pages:
            function_calls.append(FunctionCall(name="read_ppt_svg_execution_plan", args={}))
            function_calls.append(
                FunctionCall(
                    name="save_ppt_svg_page",
                    args={
                        "slide_number": page.slide_number,
                        "svg_content": _simple_svg(page.title, width=width, height=height),
                        "file_name": f"slide_{page.slide_number:03d}.svg",
                        "title": page.title,
                        "page_type": page.page_type,
                        "page_rhythm": "anchor" if page.slide_number == 1 else "dense",
                    },
                )
            )
        fake_llm = _SvgRouteToolCallingFakeLlm(function_calls=function_calls)
        session_service = InMemorySessionService()
        artifact_service = InMemoryArtifactService()

        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            paths = prepare_svg_route_paths(Path(tmpdir) / "svg_route")
            with patch.object(svg_route, "build_llm", return_value=fake_llm):
                executor_agent = build_ppt_svg_deck_executor_expert(
                    read_svg_execution_plan_tool=manager.read_ppt_svg_execution_plan,
                    save_svg_page_tool=manager.save_ppt_svg_page,
                )

            @node(name="PptSvgExecutorAgentToolHarnessNode", rerun_on_resume=True)
            async def executor_harness(ctx: Context, node_input: str) -> dict:
                result = await generate_svg_pages_with_agent(
                    requirement=requirement,
                    content_plan=content_plan,
                    design_stage=design_stage,
                    paths=paths,
                    tool_context=ctx,
                    app_name="creative_claw",
                    artifact_service=artifact_service,
                    svg_executor_agent=executor_agent,
                )
                ctx.state["svg_executor_generation_mode_for_test"] = result.generation_mode
                ctx.state["svg_executor_saved_paths_exist_for_test"] = [
                    resolve_workspace_path(page.svg_path).exists()
                    for page in result.svg_pages
                ]
                return {
                    "generation_mode": result.generation_mode,
                    "page_count": len(result.svg_pages),
                }

            workflow = Workflow(
                name="PptSvgExecutorAgentToolHarness",
                edges=[("START", executor_harness)],
            )
            runner = Runner(
                node=workflow,
                session_service=session_service,
                artifact_service=artifact_service,
            )
            user_id = "user-ppt-svg-executor"
            session_id = "session-ppt-svg-executor"

            try:
                await session_service.create_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                    state={
                        "sid": "ppt-svg-executor-agenttool-test",
                        "turn_index": 1,
                        "step": 1,
                    },
                )
                async for _ in runner.run_async(
                    user_id=user_id,
                    session_id=session_id,
                    new_message=Content(role="user", parts=[Part(text="Generate SVG pages")]),
                ):
                    pass
                session = await session_service.get_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                )
            finally:
                await runner.close()

        self.assertIsNotNone(session)
        self.assertEqual(session.state["svg_executor_generation_mode_for_test"], "llm_agent_svg")
        saved_pages = session.state[PPT_SVG_ROUTE_GENERATED_PAGES_KEY]
        self.assertEqual(len(saved_pages), len(content_plan.pages))
        self.assertTrue(all(session.state["svg_executor_saved_paths_exist_for_test"]))
        first_request_tools = _function_declaration_names(fake_llm.requests[0])
        self.assertIn("read_ppt_svg_execution_plan", first_request_tools)
        self.assertIn("save_ppt_svg_page", first_request_tools)


if __name__ == "__main__":
    unittest.main()
