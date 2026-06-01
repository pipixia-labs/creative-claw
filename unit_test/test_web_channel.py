import asyncio
import base64
import contextlib
import json
import shutil
import subprocess
import tempfile
import unittest
import uuid
import zipfile
from pathlib import Path
from unittest.mock import patch
from urllib.parse import quote
from urllib.request import urlopen

from pptx import Presentation
from pptx.util import Inches
import websockets

from conf.channel import WebChannelConfig
from conf.system import SYS_CONFIG
from src.channels.events import OutboundMessage
from src.channels.manager import ChannelManager
from src.channels import web as web_channel_module
from src.channels.web import WebChannel
from src.runtime import InboundMessage
from src.runtime.workflow_service import CreativeClawRuntime
from src.runtime.workspace import generated_root, resolve_workspace_path, workspace_relative_path
from unit_test.ppt_runtime_smoke_helpers import RuntimePptSmokePatch


class WebchatStaticAssetTests(unittest.TestCase):
    def test_webchat_shell_uses_workspace_focused_desktop_split(self) -> None:
        styles_css = Path("src/webchat/static/styles.css").read_text(encoding="utf-8")

        self.assertIn("grid-template-columns: minmax(0, 1fr) clamp(360px, 30vw, 480px);", styles_css)
        self.assertNotIn("grid-template-columns: minmax(0, 3fr) minmax(380px, 2fr);", styles_css)
        self.assertIn("@media (max-width: 1080px)", styles_css)
        self.assertIn("grid-template-columns: 1fr;", styles_css)

    def test_pptx_preview_browser_bundle_is_built(self) -> None:
        bundle_path = Path("src/webchat/static/pptx-preview-assets/creative-claw-pptx-preview.js")
        source_path = Path("src/webchat/pptx_preview_app/main.js")

        self.assertTrue(bundle_path.exists())
        self.assertIn('from "pptx-preview"', source_path.read_text(encoding="utf-8"))

    def test_webchat_protects_selected_chat_text_from_tldraw_copy_handlers(self) -> None:
        app_js = Path("src/webchat/static/app.js").read_text(encoding="utf-8")

        self.assertIn('document.addEventListener("copy", preserveChatTextSelectionClipboard, true);', app_js)
        self.assertIn('document.addEventListener("cut", preserveChatTextSelectionClipboard, true);', app_js)
        self.assertIn("function preserveChatTextSelectionClipboard(event)", app_js)
        self.assertIn("selectionIntersectsChatPanel()", app_js)
        self.assertIn("event.stopImmediatePropagation();", app_js)
        self.assertNotIn("preserveChatTextSelectionClipboard(event) {\n  event.preventDefault", app_js)

    def test_webchat_does_not_render_hidden_preview_tabs_on_artifact_updates(self) -> None:
        app_js = Path("src/webchat/static/app.js").read_text(encoding="utf-8")

        self.assertNotIn("renderAllPreviewViews", app_js)
        activate_function = app_js[app_js.index("function activatePreviewTab") : app_js.index("function renderPreviewView")]
        self.assertIn("view.hidden = !isActive;", activate_function)
        self.assertIn("renderPreviewView(tabName);", activate_function)
        preview_update_function = app_js[app_js.index("function previewArtifactSet") : app_js.index("function groupPreviewArtifacts")]
        self.assertIn("activatePreviewTab(nextTab);", preview_update_function)
        self.assertNotIn("renderPreviewView(\"tldraw\")", preview_update_function)

    def test_webchat_replaces_thinking_placeholder_before_normal_delta(self) -> None:
        app_js = Path("src/webchat/static/app.js").read_text(encoding="utf-8")

        self.assertIn('const ASSISTANT_DELTA_KIND_THINKING_PLACEHOLDER = "thinking_placeholder";', app_js)
        self.assertIn("payload.metadata?.assistant_delta_kind", app_js)
        self.assertIn("hasThinkingPlaceholder: false", app_js)
        self.assertIn("activeAssistantStream.hasThinkingPlaceholder = true;", app_js)
        self.assertIn("activeAssistantStream.content = \"\";", app_js)
        self.assertIn("activeAssistantStream.hasThinkingPlaceholder ? \"\" : activeAssistantStream.content", app_js)

    def test_webchat_supports_structured_ppt_confirmation_controls(self) -> None:
        app_js = Path("src/webchat/static/app.js").read_text(encoding="utf-8")
        styles_css = Path("src/webchat/static/styles.css").read_text(encoding="utf-8")

        self.assertIn("payload.metadata?.ppt_confirmation_request", app_js)
        self.assertIn("function renderPptConfirmationControls", app_js)
        self.assertIn('type: "ppt_confirmation"', app_js)
        self.assertIn("confirmationId: request.confirmation_id", app_js)
        self.assertIn(".cc-ppt-confirmation", styles_css)
        self.assertIn(".cc-ppt-confirmation-textarea", styles_css)

    def test_markdown_resource_urls_use_workspace_route(self) -> None:
        app_js = Path("src/webchat/static/app.js").read_text(encoding="utf-8")

        self.assertIn("function normalizeMarkdownResourceUrl", app_js)
        self.assertIn("function isWorkspaceRelativePath", app_js)
        self.assertIn("return `/workspace/${value}`;", app_js)
        self.assertIn("normalizeMarkdownResourceUrl(rawUrl, { allowMailto: false })", app_js)
        self.assertIn("normalizeMarkdownResourceUrl(rawUrl, { allowMailto: true })", app_js)

    def test_tldraw_add_to_chat_reuses_single_artifact_without_success_toast(self) -> None:
        tldraw_source = Path("src/webchat/tldraw_app/main.jsx").read_text(encoding="utf-8")
        app_js = Path("src/webchat/static/app.js").read_text(encoding="utf-8")
        tldraw_bundle = Path("src/webchat/static/tldraw-assets/creative-claw-tldraw.js").read_text(encoding="utf-8")

        self.assertIn("selectedArtifactReference(editor)", tldraw_source)
        self.assertIn("artifact: referencedArtifact", tldraw_source)
        self.assertIn("creativeClawArtifactPath", tldraw_source)
        self.assertNotIn("creative-claw-selection-attached", tldraw_source)
        self.assertNotIn("已添加到对话", tldraw_source)
        self.assertNotIn("已添加到对话", tldraw_bundle)
        self.assertIn("creative-claw-selection-attach-failed", tldraw_source)

        self.assertIn("payload?.artifact", app_js)
        self.assertIn("function attachExistingWorkspaceArtifact", app_js)
        self.assertIn("Referenced tldraw image artifact.", app_js)
        self.assertIn("Selected tldraw canvas export.", app_js)

    def test_web_channel_accepts_workspace_relative_tldraw_attachments(self) -> None:
        async def _handler(_message: InboundMessage) -> None:
            return None

        channel = WebChannel(
            config=WebChannelConfig(host="127.0.0.1", port=0, open_browser=False),
            inbound_handler=_handler,
        )
        generated_file = generated_root() / f"web_channel_tldraw_{uuid.uuid4().hex[:8]}.png"
        generated_file.write_bytes(b"fake tldraw image")

        try:
            relative_path = workspace_relative_path(generated_file)
            attachments = channel._attachments_from_chat_payload(
                {
                    "attachments": [
                        {
                            "name": generated_file.name,
                            "path": relative_path,
                            "mimeType": "image/png",
                            "description": "Referenced tldraw image artifact.",
                        }
                    ]
                }
            )

            self.assertEqual(len(attachments), 1)
            self.assertEqual(Path(attachments[0].path).resolve(), generated_file.resolve())
            self.assertEqual(attachments[0].name, generated_file.name)
            self.assertEqual(attachments[0].mime_type, "image/png")
            self.assertEqual(attachments[0].description, "Referenced tldraw image artifact.")

            with self.assertRaisesRegex(ValueError, "not a valid uploaded file"):
                channel._attachments_from_chat_payload(
                    {"attachments": [{"name": "remote.png", "path": "https://example.com/remote.png"}]}
                )
        finally:
            with contextlib.suppress(FileNotFoundError):
                generated_file.unlink()


class PptxPreviewRenderingTests(unittest.TestCase):
    def test_pptx_preview_prefers_libreoffice_pdf_conversion(self) -> None:
        pptx_path = Path("deck.pptx")
        pdf_path = Path("deck.pdf")

        with (
            patch.object(web_channel_module, "_convert_pptx_to_pdf_preview", return_value=pdf_path) as convert,
            patch.object(web_channel_module, "_render_pdf_preview_html", return_value="<html>pdf preview</html>") as render_pdf,
            patch.object(web_channel_module, "_render_pptx_browser_preview_html") as render_browser_fallback,
        ):
            html = web_channel_module._render_pptx_preview_html(pptx_path)

        self.assertEqual(html, "<html>pdf preview</html>")
        convert.assert_called_once_with(pptx_path)
        render_pdf.assert_called_once_with(pdf_path, title="deck.pptx", page_label="Slide")
        render_browser_fallback.assert_not_called()

    def test_pptx_preview_uses_browser_fallback_when_pdf_conversion_is_unavailable(self) -> None:
        pptx_path = Path("deck.pptx")

        with (
            patch.object(
                web_channel_module,
                "_convert_pptx_to_pdf_preview",
                side_effect=RuntimeError("LibreOffice missing"),
            ) as convert,
            patch.object(
                web_channel_module,
                "_render_pptx_browser_preview_html",
                return_value="<html>browser fallback</html>",
            ) as render_browser_fallback,
        ):
            html = web_channel_module._render_pptx_preview_html(pptx_path)

        self.assertEqual(html, "<html>browser fallback</html>")
        convert.assert_called_once_with(pptx_path)
        render_browser_fallback.assert_called_once_with(pptx_path)

    def test_pptx_browser_fallback_shell_loads_workspace_file_and_bundle(self) -> None:
        generated_file = generated_root() / f"web_channel_{uuid.uuid4().hex[:8]}.pptx"
        generated_file.write_bytes(b"fake pptx")
        try:
            relative_path = workspace_relative_path(generated_file)
            html = web_channel_module._render_pptx_browser_preview_html(generated_file)
        finally:
            with contextlib.suppress(FileNotFoundError):
                generated_file.unlink()

        self.assertIn("pptx-preview-root", html)
        self.assertIn("/pptx-preview-assets/creative-claw-pptx-preview.js", html)
        self.assertIn(f"/workspace/{quote(relative_path)}", html)

    def test_pptx_to_pdf_conversion_uses_cached_output_path(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir_name:
            temp_dir = Path(temp_dir_name)
            pptx_path = temp_dir / "deck.pptx"
            cached_pdf_path = temp_dir / "cache" / "deck.pdf"
            pptx_path.write_bytes(b"pptx bytes")

            def fake_run(command, **_kwargs):
                output_dir = Path(command[command.index("--outdir") + 1])
                (output_dir / "deck.pdf").write_bytes(b"%PDF-1.4\n")
                return subprocess.CompletedProcess(command, 0, stdout="converted", stderr="")

            with (
                patch.object(web_channel_module, "_find_libreoffice_executable", return_value="/usr/bin/libreoffice"),
                patch.object(web_channel_module, "_pptx_preview_cache_pdf_path", return_value=cached_pdf_path),
                patch.object(web_channel_module.subprocess, "run", side_effect=fake_run) as run,
            ):
                converted_path = web_channel_module._convert_pptx_to_pdf_preview(pptx_path)
                cached_path = web_channel_module._convert_pptx_to_pdf_preview(pptx_path)

            self.assertEqual(converted_path, cached_pdf_path)
            self.assertEqual(cached_path, cached_pdf_path)
            self.assertEqual(run.call_count, 1)
            self.assertEqual(cached_pdf_path.read_bytes(), b"%PDF-1.4\n")


class WebChannelTests(unittest.IsolatedAsyncioTestCase):
    async def asyncSetUp(self) -> None:
        self.inbound_messages: list[InboundMessage] = []

        async def _handler(message: InboundMessage) -> None:
            self.inbound_messages.append(message)

        self.channel = WebChannel(
            config=WebChannelConfig(
                host="127.0.0.1",
                port=0,
                open_browser=False,
                title="CreativeClaw Web Chat",
            ),
            inbound_handler=_handler,
        )
        await self.channel.start()

    async def asyncTearDown(self) -> None:
        await self.channel.stop()

    async def test_web_channel_serves_index_page(self) -> None:
        def fetch_index():
            with urlopen(self.channel.url) as response:  # noqa: S310 - local test server
                return response.status, response.read().decode("utf-8")

        status, body = await asyncio.to_thread(fetch_index)
        self.assertEqual(status, 200)
        self.assertIn("CreativeClaw Web Chat", body)
        self.assertIn("/app.js", body)
        self.assertNotIn("Creative flow in one surface", body)
        self.assertNotIn("A local browser chat surface", body)
        self.assertNotIn("Recent Sessions", body)
        self.assertIn('data-preview-tab="tldraw"', body)
        self.assertIn('data-preview-tab="html"', body)
        self.assertIn('data-preview-tab="ppt"', body)
        self.assertIn('data-preview-tab="model3d"', body)
        self.assertIn("Visual Board", body)
        self.assertIn("Design", body)
        self.assertIn("3D", body)
        self.assertIn("No Design preview", body)
        self.assertIn("No 3D preview", body)
        self.assertIn("/model3d-assets/creative-claw-model3d.js", body)
        self.assertNotIn("No HTML preview", body)
        self.assertIn('aria-label="Send message"', body)
        self.assertIn('aria-label="Attach files"', body)
        self.assertIn('id="session-history"', body)
        self.assertIn('id="session-popover"', body)
        self.assertIn('id="file-input"', body)
        self.assertNotIn("Describe the image, prompt", body)

    async def test_web_channel_serves_pptx_preview_page(self) -> None:
        generated_file = generated_root() / f"web_channel_{uuid.uuid4().hex[:8]}.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        textbox.text = "Quarterly roadmap preview"
        presentation.save(generated_file)

        try:
            relative_path = workspace_relative_path(generated_file)

            def fetch_preview():
                preview_url = f"{self.channel.url}/workspace-preview/{quote(relative_path)}"
                with urlopen(preview_url) as response:  # noqa: S310 - local test server
                    return response.status, response.headers.get("Content-Type", ""), response.read().decode("utf-8")

            with patch.object(
                web_channel_module,
                "_convert_pptx_to_pdf_preview",
                side_effect=RuntimeError("LibreOffice unavailable in test"),
            ):
                status, content_type, body = await asyncio.to_thread(fetch_preview)
            self.assertEqual(status, 200)
            self.assertIn("text/html", content_type)
            self.assertIn("pptx-preview-root", body)
            self.assertIn("/pptx-preview-assets/creative-claw-pptx-preview.js", body)
            self.assertIn(quote(relative_path), body)
            self.assertNotIn("Quarterly roadmap preview", body)
        finally:
            with contextlib.suppress(FileNotFoundError):
                generated_file.unlink()

    async def test_web_channel_marks_private_ppt_html_deck_artifact(self) -> None:
        session_dir = generated_root() / f"art_session_web_channel_{uuid.uuid4().hex[:8]}"
        private_dir = session_dir / "turn_1" / "ppt_private_skill_step_1"
        private_html = private_dir / "index.html"
        ordinary_html = generated_root() / f"web_channel_{uuid.uuid4().hex[:8]}.html"
        private_dir.mkdir(parents=True, exist_ok=True)
        private_html.write_text("<!doctype html><html><body><main id='deck'></main></body></html>", encoding="utf-8")
        ordinary_html.write_text("<!doctype html><html><body>Poster page</body></html>", encoding="utf-8")

        try:
            async with websockets.connect(f"ws://127.0.0.1:{self.channel._port}/ws?session_id=ppt-html-session") as websocket:
                ready = json.loads(await asyncio.wait_for(websocket.recv(), timeout=2))
                self.assertEqual(ready["type"], "ready")

                await self.channel.send(
                    OutboundMessage(
                        channel="web",
                        chat_id="ppt-html-session",
                        text="deck ready",
                        artifact_paths=[str(private_html), str(ordinary_html)],
                        metadata={"display_style": "final"},
                    )
                )
                final_message = await self._recv_until(websocket, "assistant_message")
                artifacts_by_path = {artifact["path"]: artifact for artifact in final_message["artifacts"]}
                private_artifact = artifacts_by_path[workspace_relative_path(private_html)]
                ordinary_artifact = artifacts_by_path[workspace_relative_path(ordinary_html)]

                self.assertEqual(private_artifact["artifactKind"], "interactive_ppt_html")
                self.assertEqual(private_artifact["mimeType"], "text/html")
                self.assertNotIn("artifactKind", ordinary_artifact)
                self.assertEqual(ordinary_artifact["mimeType"], "text/html")
        finally:
            shutil.rmtree(session_dir, ignore_errors=True)
            with contextlib.suppress(FileNotFoundError):
                ordinary_html.unlink()

    def test_webchat_routes_interactive_html_decks_to_ppt_preview(self) -> None:
        app_js = Path("src/webchat/static/app.js").read_text(encoding="utf-8")

        preview_router = app_js[
            app_js.index("function previewTabForArtifact") : app_js.index("function isHtmlArtifact")
        ]
        self.assertIn("function isInteractiveHtmlDeckArtifact", app_js)
        self.assertIn('const INTERACTIVE_PPT_HTML_KIND = "interactive_ppt_html";', app_js)
        self.assertLess(
            preview_router.index("isInteractiveHtmlDeckArtifact"),
            preview_router.index("isHtmlArtifact"),
        )
        self.assertIn('iframe.className = "ppt-preview-frame ppt-html-deck-frame";', app_js)
        self.assertIn("iframe.src = artifact.url;", app_js)

    async def test_web_channel_serves_3d_artifact_metadata_and_asset(self) -> None:
        cases = [
            (".fbx", "application/octet-stream", b"Kaydara FBX Binary  \x00"),
            (".glb", "model/gltf-binary", b"fake-glb"),
            (".usd", "model/vnd.usd", b"#usda 1.0\n"),
            (".obj", "model/obj", b"o cube\nv 0 0 0\n"),
            (".stl", "model/stl", b"solid cube\nendsolid cube\n"),
            (".usdz", "model/vnd.usdz+zip", b"fake-usdz"),
        ]
        generated_files = []
        for extension, _mime_type, body in cases:
            generated_file = generated_root() / f"web_channel_{uuid.uuid4().hex[:8]}{extension}"
            generated_file.write_bytes(body)
            generated_files.append(generated_file)

        try:
            async with websockets.connect(f"ws://127.0.0.1:{self.channel._port}/ws?session_id=model-session") as websocket:
                ready = json.loads(await asyncio.wait_for(websocket.recv(), timeout=2))
                self.assertEqual(ready["type"], "ready")

                await self.channel.send(
                    OutboundMessage(
                        channel="web",
                        chat_id="model-session",
                        text="model ready",
                        artifact_paths=[str(path) for path in generated_files],
                        metadata={"display_style": "final"},
                    )
                )
                final_message = await self._recv_until(websocket, "assistant_message")
                self.assertEqual(len(final_message["artifacts"]), len(cases))
                artifacts_by_name = {artifact["name"]: artifact for artifact in final_message["artifacts"]}

                for generated_file, (_extension, expected_mime_type, expected_body) in zip(generated_files, cases):
                    artifact = artifacts_by_name[generated_file.name]
                    self.assertEqual(artifact["name"], generated_file.name)
                    self.assertFalse(artifact["isImage"])
                    self.assertTrue(artifact["is3D"])
                    self.assertEqual(artifact["mimeType"], expected_mime_type)
                    self.assertEqual(artifact["sizeBytes"], len(expected_body))

                    def fetch_artifact():
                        with urlopen(f"{self.channel.url}{artifact['url']}") as response:  # noqa: S310 - local test server
                            return response.status, response.headers.get("Content-Type", ""), response.read()

                    status, content_type, body = await asyncio.to_thread(fetch_artifact)
                    self.assertEqual(status, 200)
                    self.assertIn(expected_mime_type, content_type)
                    self.assertEqual(body, expected_body)
        finally:
            for generated_file in generated_files:
                with contextlib.suppress(FileNotFoundError):
                    generated_file.unlink()

    async def test_web_channel_selects_fbx_usd_from_zip_model_package(self) -> None:
        generated_file = generated_root() / f"web_channel_hy3d_model_{uuid.uuid4().hex[:8]}.zip"
        fbx_body = b"Kaydara FBX Binary  \x00"
        usdz_body = b"fake-usdz"
        with zipfile.ZipFile(generated_file, "w") as archive:
            archive.writestr("models/preview.usdz", usdz_body)
            archive.writestr("models/preview.fbx", fbx_body)
            archive.writestr("notes/readme.txt", b"not a model")

        try:
            relative_path = workspace_relative_path(generated_file)

            def fetch_json(path: str):
                with urlopen(f"{self.channel.url}{path}") as response:  # noqa: S310 - local test server
                    return response.status, response.headers.get("Content-Type", ""), json.loads(
                        response.read().decode("utf-8")
                    )

            def fetch_bytes(path: str):
                with urlopen(f"{self.channel.url}{path}") as response:  # noqa: S310 - local test server
                    return response.status, response.headers.get("Content-Type", ""), response.read()

            manifest_path = f"/workspace-3d-package/manifest/{quote(relative_path)}"
            status, content_type, manifest = await asyncio.to_thread(fetch_json, manifest_path)
            self.assertEqual(status, 200)
            self.assertIn("application/json", content_type)
            self.assertEqual(manifest["modelEntry"], "models/preview.fbx")
            self.assertEqual(manifest["modelDirectory"], "models")
            self.assertEqual(manifest["modelSizeBytes"], len(fbx_body))

            status, content_type, body = await asyncio.to_thread(fetch_bytes, manifest["modelUrl"])
            self.assertEqual(status, 200)
            self.assertIn("application/octet-stream", content_type)
            self.assertEqual(body, fbx_body)
        finally:
            with contextlib.suppress(FileNotFoundError):
                generated_file.unlink()

    async def test_web_channel_serves_zip_model_package_manifest_and_entries(self) -> None:
        generated_file = generated_root() / f"web_channel_model_{uuid.uuid4().hex[:8]}.zip"
        gltf_body = json.dumps(
            {
                "asset": {"version": "2.0"},
                "buffers": [{"uri": "model.bin", "byteLength": 0}],
                "scenes": [{"nodes": []}],
                "scene": 0,
            }
        ).encode("utf-8")
        bin_body = b""
        with zipfile.ZipFile(generated_file, "w") as archive:
            archive.writestr("models/model.gltf", gltf_body)
            archive.writestr("models/model.bin", bin_body)
            archive.writestr("notes/readme.txt", b"not a model")

        try:
            relative_path = workspace_relative_path(generated_file)

            def fetch_json(path: str):
                with urlopen(f"{self.channel.url}{path}") as response:  # noqa: S310 - local test server
                    return response.status, response.headers.get("Content-Type", ""), json.loads(
                        response.read().decode("utf-8")
                    )

            def fetch_bytes(path: str):
                with urlopen(f"{self.channel.url}{path}") as response:  # noqa: S310 - local test server
                    return response.status, response.headers.get("Content-Type", ""), response.read()

            manifest_path = f"/workspace-3d-package/manifest/{quote(relative_path)}"
            status, content_type, manifest = await asyncio.to_thread(fetch_json, manifest_path)
            self.assertEqual(status, 200)
            self.assertIn("application/json", content_type)
            self.assertEqual(manifest["modelEntry"], "models/model.gltf")
            self.assertEqual(manifest["modelDirectory"], "models")
            self.assertEqual(manifest["modelSizeBytes"], len(gltf_body))
            self.assertIn("/workspace-3d-package/file/", manifest["modelUrl"])
            self.assertEqual(manifest["entryCount"], 3)

            status, content_type, body = await asyncio.to_thread(fetch_bytes, manifest["modelUrl"])
            self.assertEqual(status, 200)
            self.assertIn("model/gltf+json", content_type)
            self.assertEqual(body, gltf_body)

            bin_path = f"/workspace-3d-package/file/{quote(relative_path)}?entry={quote('models/model.bin', safe='')}"
            status, content_type, body = await asyncio.to_thread(fetch_bytes, bin_path)
            self.assertEqual(status, 200)
            self.assertIn("application/octet-stream", content_type)
            self.assertEqual(body, bin_body)

            async with websockets.connect(f"ws://127.0.0.1:{self.channel._port}/ws?session_id=zip-model-session") as websocket:
                ready = json.loads(await asyncio.wait_for(websocket.recv(), timeout=2))
                self.assertEqual(ready["type"], "ready")
                await self.channel.send(
                    OutboundMessage(
                        channel="web",
                        chat_id="zip-model-session",
                        text="zip model ready",
                        artifact_paths=[str(generated_file)],
                        metadata={"display_style": "final"},
                    )
                )
                final_message = await self._recv_until(websocket, "assistant_message")
                self.assertEqual(len(final_message["artifacts"]), 1)
                artifact = final_message["artifacts"][0]
                self.assertTrue(artifact["is3D"])
                self.assertEqual(artifact["mimeType"], "application/zip")
                self.assertEqual(artifact["sizeBytes"], generated_file.stat().st_size)
        finally:
            with contextlib.suppress(FileNotFoundError):
                generated_file.unlink()

    async def test_web_channel_serves_design_system_catalog_and_preview(self) -> None:
        def fetch_json(path: str):
            with urlopen(f"{self.channel.url}{path}") as response:  # noqa: S310 - local test server
                return response.status, response.headers.get("Content-Type", ""), json.loads(
                    response.read().decode("utf-8")
                )

        def fetch_html(path: str):
            with urlopen(f"{self.channel.url}{path}") as response:  # noqa: S310 - local test server
                return response.status, response.headers.get("Content-Type", ""), response.read().decode("utf-8")

        status, content_type, payload = await asyncio.to_thread(fetch_json, "/api/design-systems")
        self.assertEqual(status, 200)
        self.assertIn("application/json", content_type)
        systems = {item["id"]: item for item in payload["designSystems"]}
        self.assertIn("claude", systems)
        self.assertIn("/api/design-systems/claude/preview", systems["claude"]["previewUrl"])

        status, content_type, body = await asyncio.to_thread(fetch_html, "/api/design-systems/claude/showcase")
        self.assertEqual(status, 200)
        self.assertIn("text/html", content_type)
        self.assertIn("Claude Design System Preview", body)

    async def test_web_channel_bridges_websocket_messages_and_artifacts(self) -> None:
        generated_file = generated_root() / f"web_channel_{uuid.uuid4().hex[:8]}.png"
        generated_file.write_bytes(b"fake-image")

        try:
            async with websockets.connect(f"ws://127.0.0.1:{self.channel._port}/ws?session_id=test-session") as websocket:
                ready = json.loads(await asyncio.wait_for(websocket.recv(), timeout=2))
                self.assertEqual(ready["type"], "ready")
                self.assertEqual(ready["sessionId"], "test-session")

                await websocket.send(json.dumps({"type": "chat", "content": "hello web"}))
                started = await self._recv_until(websocket, "task_started")
                self.assertTrue(started["runId"])
                inbound = await asyncio.wait_for(self._consume_inbound(), timeout=2)
                self.assertEqual(inbound.channel, "web")
                self.assertEqual(inbound.chat_id, "test-session")
                self.assertEqual(inbound.text, "hello web")
                self.assertEqual(inbound.metadata["run_id"], started["runId"])
                await self._recv_until(websocket, "task_finished")

                await self.channel.send(
                    OutboundMessage(
                        channel="web",
                        chat_id="test-session",
                        text="working on it",
                        metadata={
                            "display_style": "progress",
                            "stage_title": "Planning",
                            "activity_group_id": "runtime-session:turn:1",
                        },
                    )
                )
                progress = await self._recv_until(websocket, "progress")
                self.assertEqual(progress["type"], "progress")
                self.assertEqual(progress["metadata"]["stage_title"], "Planning")
                self.assertEqual(progress["metadata"]["activity_group_id"], "runtime-session:turn:1")

                await self.channel.send(
                    OutboundMessage(
                        channel="web",
                        chat_id="test-session",
                        text="final answer",
                        artifact_paths=[str(generated_file)],
                        metadata={"display_style": "final"},
                    )
                )
                delta = await self._recv_until(websocket, "assistant_delta")
                self.assertEqual(delta["delta"], "final answer")
                final_message = await self._recv_until(websocket, "assistant_message")
                self.assertEqual(final_message["type"], "assistant_message")
                self.assertEqual(final_message["content"], "final answer")
                self.assertTrue(final_message["streamComplete"])
                self.assertEqual(len(final_message["artifacts"]), 1)
                artifact = final_message["artifacts"][0]
                self.assertTrue(artifact["isImage"])

                def fetch_artifact():
                    with urlopen(f"{self.channel.url}{artifact['url']}") as response:  # noqa: S310 - local test server
                        return response.status, response.read()

                status, body = await asyncio.to_thread(fetch_artifact)
                self.assertEqual(status, 200)
                self.assertEqual(body, b"fake-image")
        finally:
            with contextlib.suppress(FileNotFoundError):
                generated_file.unlink()

    async def test_web_channel_accepts_structured_ppt_confirmation_payloads(self) -> None:
        async with websockets.connect(f"ws://127.0.0.1:{self.channel._port}/ws?session_id=ppt-hitl-input") as websocket:
            ready = json.loads(await asyncio.wait_for(websocket.recv(), timeout=2))
            self.assertEqual(ready["type"], "ready")

            await websocket.send(
                json.dumps(
                    {
                        "type": "ppt_confirmation",
                        "action": "confirm",
                        "confirmationId": "requirement-1",
                        "stage": "awaiting_requirement_confirmation",
                        "runId": "ppt-confirm-1",
                    }
                )
            )
            await self._recv_until(websocket, "task_started")
            inbound = await asyncio.wait_for(self._consume_inbound(), timeout=2)
            self.assertEqual(inbound.text, "确认")
            self.assertEqual(inbound.metadata["ppt_confirmation_response"]["action"], "confirm")
            self.assertEqual(inbound.metadata["ppt_confirmation_response"]["confirmation_id"], "requirement-1")
            await self._recv_until(websocket, "task_finished")

            await websocket.send(
                json.dumps(
                    {
                        "type": "chat",
                        "content": "",
                        "pptConfirmation": {
                            "action": "revise",
                            "message": "改成 5 页，并强化结论页。",
                            "stage": "awaiting_content_plan_confirmation",
                        },
                        "runId": "ppt-confirm-2",
                    }
                )
            )
            await self._recv_until(websocket, "task_started")
            inbound = await asyncio.wait_for(self._consume_inbound(), timeout=2)
            self.assertEqual(inbound.text, "改成 5 页，并强化结论页。")
            self.assertEqual(inbound.metadata["ppt_confirmation_response"]["action"], "revise")
            self.assertEqual(
                inbound.metadata["ppt_confirmation_response"]["stage"],
                "awaiting_content_plan_confirmation",
            )
            await self._recv_until(websocket, "task_finished")

    async def test_web_channel_streams_question_form_messages(self) -> None:
        form_message = (
            '<cc-question-form id="design-brief" version="design-brief-form-v1">\n'
            '{"questions":[{"id":"style","label":"Style","type":"single_choice","options":[]}]}'
            "\n</cc-question-form>"
        )

        async with websockets.connect(f"ws://127.0.0.1:{self.channel._port}/ws?session_id=form-session") as websocket:
            ready = json.loads(await asyncio.wait_for(websocket.recv(), timeout=2))
            self.assertEqual(ready["type"], "ready")

            await self.channel.send(
                OutboundMessage(
                    channel="web",
                    chat_id="form-session",
                    text=form_message,
                    metadata={"display_style": "final"},
                )
            )
            delta = await self._recv_until(websocket, "assistant_delta")
            self.assertIn("<cc-question-form", delta["delta"])
            final_message = await self._recv_until(websocket, "assistant_message")
            self.assertEqual(final_message["content"], form_message)
            self.assertTrue(final_message["streamComplete"])

    async def test_web_channel_forwards_realtime_assistant_delta(self) -> None:
        async with websockets.connect(f"ws://127.0.0.1:{self.channel._port}/ws?session_id=delta-session") as websocket:
            ready = json.loads(await asyncio.wait_for(websocket.recv(), timeout=2))
            self.assertEqual(ready["type"], "ready")

            await self.channel.send(
                OutboundMessage(
                    channel="web",
                    chat_id="delta-session",
                    text="Hel",
                    metadata={"display_style": "assistant_delta", "session_id": "runtime-session"},
                )
            )
            delta = await self._recv_until(websocket, "assistant_delta")
            self.assertEqual(delta["delta"], "Hel")
            self.assertEqual(delta["content"], "Hel")
            self.assertEqual(delta["metadata"]["session_id"], "runtime-session")

    async def test_web_channel_accepts_chunked_file_uploads(self) -> None:
        upload_id = f"upload-{uuid.uuid4().hex[:8]}"
        body = b"hello uploaded file"
        uploaded_path: Path | None = None

        async with websockets.connect(f"ws://127.0.0.1:{self.channel._port}/ws?session_id=upload-session") as websocket:
            ready = json.loads(await asyncio.wait_for(websocket.recv(), timeout=2))
            self.assertEqual(ready["type"], "ready")

            await websocket.send(
                json.dumps(
                    {
                        "type": "upload_start",
                        "uploadId": upload_id,
                        "name": "note.txt",
                        "size": len(body),
                        "mimeType": "text/plain",
                    }
                )
            )
            started = json.loads(await asyncio.wait_for(websocket.recv(), timeout=2))
            self.assertEqual(started["type"], "upload_started")
            self.assertEqual(started["uploadId"], upload_id)

            await websocket.send(
                json.dumps(
                    {
                        "type": "upload_chunk",
                        "uploadId": upload_id,
                        "data": base64.b64encode(body).decode("ascii"),
                    }
                )
            )
            chunk = json.loads(await asyncio.wait_for(websocket.recv(), timeout=2))
            self.assertEqual(chunk["type"], "upload_chunk_received")
            self.assertEqual(chunk["received"], len(body))

            await websocket.send(json.dumps({"type": "upload_finish", "uploadId": upload_id}))
            complete = json.loads(await asyncio.wait_for(websocket.recv(), timeout=2))
            self.assertEqual(complete["type"], "upload_complete")
            self.assertEqual(complete["name"], "note.txt")
            uploaded_path = Path(complete["path"])
            self.assertEqual(uploaded_path.read_bytes(), body)

            await websocket.send(
                json.dumps(
                    {
                        "type": "chat",
                        "content": "use this uploaded note",
                        "attachments": [
                            {
                                "name": "note.txt",
                                "path": complete["path"],
                                "mimeType": "text/plain",
                                "description": "uploaded test note",
                            }
                        ],
                    }
                )
            )
            await self._recv_until(websocket, "task_started")
            inbound = await asyncio.wait_for(self._consume_inbound(), timeout=2)
            self.assertEqual(inbound.text, "use this uploaded note")
            self.assertEqual(len(inbound.attachments), 1)
            self.assertEqual(Path(inbound.attachments[0].path).resolve(), uploaded_path.resolve())
            self.assertEqual(inbound.attachments[0].name, "note.txt")
            self.assertEqual(inbound.attachments[0].mime_type, "text/plain")
            self.assertEqual(inbound.attachments[0].description, "uploaded test note")
            await self._recv_until(websocket, "task_finished")

        if uploaded_path is not None:
            with contextlib.suppress(FileNotFoundError):
                uploaded_path.unlink()
            with contextlib.suppress(OSError):
                uploaded_path.parent.rmdir()

    async def test_web_channel_stop_cancels_active_run_without_blocking_read_loop(self) -> None:
        started = asyncio.Event()
        released = asyncio.Event()

        async def _blocking_handler(message: InboundMessage) -> None:
            self.inbound_messages.append(message)
            started.set()
            await released.wait()

        self.channel.inbound_handler = _blocking_handler

        async with websockets.connect(f"ws://127.0.0.1:{self.channel._port}/ws?session_id=stop-session") as websocket:
            ready = json.loads(await asyncio.wait_for(websocket.recv(), timeout=2))
            self.assertEqual(ready["type"], "ready")

            run_id = f"run-{uuid.uuid4().hex}"
            await websocket.send(json.dumps({"type": "chat", "content": "long task", "runId": run_id}))
            started_payload = await self._recv_until(websocket, "task_started")
            self.assertEqual(started_payload["runId"], run_id)
            await asyncio.wait_for(started.wait(), timeout=2)

            await websocket.send(json.dumps({"type": "stop", "runId": run_id}))
            stopping = await self._recv_until(websocket, "task_stopping")
            self.assertEqual(stopping["runId"], run_id)
            finished = await self._recv_until(websocket, "task_finished")
            self.assertEqual(finished["runId"], run_id)
            self.assertEqual(finished["reason"], "cancelled")

            next_run_id = f"run-{uuid.uuid4().hex}"
            await websocket.send(json.dumps({"type": "chat", "content": "next task", "runId": next_run_id}))
            next_started = await self._recv_until(websocket, "task_started")
            self.assertEqual(next_started["runId"], next_run_id)

        released.set()

    async def _consume_inbound(self) -> InboundMessage:
        for _ in range(20):
            if self.inbound_messages:
                return self.inbound_messages.pop(0)
            await asyncio.sleep(0.01)
        raise AssertionError("Inbound message was not received.")

    async def _recv_until(self, websocket, expected_type: str) -> dict[str, object]:
        for _ in range(10):
            payload = json.loads(await asyncio.wait_for(websocket.recv(), timeout=2))
            if payload.get("type") == expected_type:
                return payload
        raise AssertionError(f"Did not receive payload type {expected_type!r}.")


class WebChannelRuntimePptSmokeTests(unittest.IsolatedAsyncioTestCase):
    async def asyncSetUp(self) -> None:
        self.runtime = CreativeClawRuntime()
        self.manager = ChannelManager(self.runtime)
        self.channel = WebChannel(
            config=WebChannelConfig(
                host="127.0.0.1",
                port=0,
                open_browser=False,
                title="CreativeClaw Web Chat",
            ),
            inbound_handler=self.manager.handle_inbound,
        )
        self.manager.register(self.channel)
        await self.channel.start()

    async def asyncTearDown(self) -> None:
        await self.channel.stop()

    async def test_web_channel_ppt_adk_hitl_smoke_reaches_final_delivery(self) -> None:
        task = "做一个 3 页 PPTX，用于产品发布，受众为管理层。"
        session_key = "ppt-web-smoke"

        with RuntimePptSmokePatch(task=task).install() as smoke:
            async with websockets.connect(f"ws://127.0.0.1:{self.channel._port}/ws?session_id={session_key}") as websocket:
                ready = json.loads(await asyncio.wait_for(websocket.recv(), timeout=5))
                self.assertEqual(ready["type"], "ready")

                first_final = await self._send_chat_and_wait_for_final(websocket, task, run_id="ppt-web-smoke-1")
                self.assertIn("请确认 PPT 需求参数", first_final["content"])
                self.assertEqual(
                    first_final["metadata"]["ppt_confirmation_request"]["confirmation_type"],
                    "requirement",
                )

                second_final = await self._send_chat_and_wait_for_final(websocket, "确认", run_id="ppt-web-smoke-2")
                self.assertIn("请确认 PPT 内容规划", second_final["content"])
                self.assertEqual(
                    second_final["metadata"]["ppt_confirmation_request"]["confirmation_type"],
                    "content_plan",
                )

                third_final = await self._send_chat_and_wait_for_final(websocket, "确认", run_id="ppt-web-smoke-3")

        self.assertIn(
            "HTML route generated the PPTX after requirement and content-plan confirmation.",
            third_final["content"],
        )
        self.assertEqual(len(third_final["artifacts"]), 1)
        artifact = third_final["artifacts"][0]
        self.assertTrue(artifact["path"].endswith(".pptx"))
        self.assertEqual(artifact["mimeType"], "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        self.assertEqual(len(smoke.fake_llms), 3)
        self.assertEqual(len(smoke.fake_llms[0].requests), 1)
        self.assertEqual(len(smoke.fake_llms[2].requests), 1)

        runtime_session_id = self.runtime._session_keys[f"web:{session_key}"]
        session = await self.runtime.session_service.get_session(
            app_name=SYS_CONFIG.app_name,
            user_id=ready["clientId"],
            session_id=runtime_session_id,
        )
        self.assertEqual(session.state["ppt_product_result"]["status"], "success")
        self.assertEqual(session.state["ppt_workflow_state"]["stage"], "completed")
        self.assertIsNone(session.state.get("ppt_adk_pending_confirmation"))
        self.assertEqual(session.state["final_file_paths"], [artifact["path"]])
        self.assertTrue(resolve_workspace_path(artifact["path"]).is_file())

    async def test_web_channel_ppt_structured_hitl_smoke_reaches_final_delivery(self) -> None:
        task = "做一个 3 页 PPTX，用于产品发布，受众为管理层。"
        session_key = "ppt-web-structured-smoke"

        with RuntimePptSmokePatch(task=task).install():
            async with websockets.connect(f"ws://127.0.0.1:{self.channel._port}/ws?session_id={session_key}") as websocket:
                ready = json.loads(await asyncio.wait_for(websocket.recv(), timeout=5))
                self.assertEqual(ready["type"], "ready")

                first_final = await self._send_chat_and_wait_for_final(
                    websocket,
                    task,
                    run_id="ppt-web-structured-1",
                )
                self.assertIn("请确认 PPT 需求参数", first_final["content"])

                second_final = await self._send_payload_and_wait_for_final(
                    websocket,
                    {
                        "type": "ppt_confirmation",
                        "action": "confirm",
                        "confirmationId": first_final["metadata"]["ppt_confirmation_request"]["confirmation_id"],
                        "stage": first_final["metadata"]["ppt_confirmation_request"]["stage"],
                        "runId": "ppt-web-structured-2",
                    },
                    run_id="ppt-web-structured-2",
                )
                self.assertIn("请确认 PPT 内容规划", second_final["content"])

                third_final = await self._send_payload_and_wait_for_final(
                    websocket,
                    {
                        "type": "chat",
                        "content": "",
                        "pptConfirmation": {
                            "action": "confirm",
                            "confirmationId": second_final["metadata"]["ppt_confirmation_request"]["confirmation_id"],
                            "stage": second_final["metadata"]["ppt_confirmation_request"]["stage"],
                        },
                        "runId": "ppt-web-structured-3",
                    },
                    run_id="ppt-web-structured-3",
                )

        self.assertIn(
            "HTML route generated the PPTX after requirement and content-plan confirmation.",
            third_final["content"],
        )
        self.assertEqual(len(third_final["artifacts"]), 1)
        artifact = third_final["artifacts"][0]
        self.assertTrue(artifact["path"].endswith(".pptx"))

    async def test_web_channel_ppt_structured_revision_smoke_reaches_final_delivery(self) -> None:
        task = "做一个 3 页 PPTX，用于产品发布，受众为管理层。"
        revision = "改成 4 页，受众: 研发负责人。"
        session_key = "ppt-web-structured-revision-smoke"

        with RuntimePptSmokePatch(task=task).install():
            async with websockets.connect(f"ws://127.0.0.1:{self.channel._port}/ws?session_id={session_key}") as websocket:
                ready = json.loads(await asyncio.wait_for(websocket.recv(), timeout=5))
                self.assertEqual(ready["type"], "ready")

                first_final = await self._send_chat_and_wait_for_final(
                    websocket,
                    task,
                    run_id="ppt-web-structured-revision-1",
                )
                first_request = first_final["metadata"]["ppt_confirmation_request"]
                self.assertEqual(first_request["confirmation_type"], "requirement")

                second_final = await self._send_payload_and_wait_for_final(
                    websocket,
                    {
                        "type": "ppt_confirmation",
                        "action": "revise",
                        "message": revision,
                        "confirmationId": first_request["confirmation_id"],
                        "stage": first_request["stage"],
                        "runId": "ppt-web-structured-revision-2",
                    },
                    run_id="ppt-web-structured-revision-2",
                )
                self.assertIn("请确认 PPT 需求参数", second_final["content"])
                second_request = second_final["metadata"]["ppt_confirmation_request"]
                self.assertEqual(second_request["confirmation_type"], "requirement")

                third_final = await self._send_payload_and_wait_for_final(
                    websocket,
                    {
                        "type": "ppt_confirmation",
                        "action": "confirm",
                        "confirmationId": second_request["confirmation_id"],
                        "stage": second_request["stage"],
                        "runId": "ppt-web-structured-revision-3",
                    },
                    run_id="ppt-web-structured-revision-3",
                )
                self.assertIn("请确认 PPT 内容规划", third_final["content"])
                third_request = third_final["metadata"]["ppt_confirmation_request"]
                self.assertEqual(third_request["confirmation_type"], "content_plan")

                fourth_final = await self._send_payload_and_wait_for_final(
                    websocket,
                    {
                        "type": "chat",
                        "content": "",
                        "pptConfirmation": {
                            "action": "confirm",
                            "confirmationId": third_request["confirmation_id"],
                            "stage": third_request["stage"],
                        },
                        "runId": "ppt-web-structured-revision-4",
                    },
                    run_id="ppt-web-structured-revision-4",
                )

        self.assertIn(
            "HTML route generated the PPTX after requirement and content-plan confirmation.",
            fourth_final["content"],
        )
        self.assertEqual(len(fourth_final["artifacts"]), 1)
        artifact = fourth_final["artifacts"][0]
        self.assertTrue(artifact["path"].endswith(".pptx"))

        runtime_session_id = self.runtime._session_keys[f"web:{session_key}"]
        session = await self.runtime.session_service.get_session(
            app_name=SYS_CONFIG.app_name,
            user_id=ready["clientId"],
            session_id=runtime_session_id,
        )
        self.assertEqual(session.state["ppt_product_result"]["status"], "success")
        self.assertEqual(session.state["ppt_workflow_state"]["stage"], "completed")
        self.assertEqual(session.state["ppt_confirmed_requirement"]["slide_count_policy"]["target"], 4)
        self.assertIn("研发负责人", session.state["ppt_confirmed_requirement"]["audience"])

    async def _send_chat_and_wait_for_final(
        self,
        websocket,
        content: str,
        *,
        run_id: str,
    ) -> dict[str, object]:
        return await self._send_payload_and_wait_for_final(
            websocket,
            {"type": "chat", "content": content, "runId": run_id},
            run_id=run_id,
        )

    async def _send_payload_and_wait_for_final(
        self,
        websocket,
        payload: dict[str, object],
        *,
        run_id: str,
    ) -> dict[str, object]:
        await websocket.send(json.dumps(payload, ensure_ascii=False))
        final_message: dict[str, object] | None = None
        while True:
            payload = json.loads(await asyncio.wait_for(websocket.recv(), timeout=10))
            payload_type = payload.get("type")
            if payload_type == "error":
                self.fail(f"Unexpected WebChannel error payload: {payload}")
            if payload_type == "assistant_message":
                final_message = payload
            if payload_type == "task_finished" and payload.get("runId") == run_id:
                if final_message is None:
                    self.fail(f"Task finished without assistant message: {payload}")
                return final_message


if __name__ == "__main__":
    unittest.main()
