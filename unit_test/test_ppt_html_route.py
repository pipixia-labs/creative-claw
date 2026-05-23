import json
import tempfile
import unittest
from pathlib import Path
from types import SimpleNamespace

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.productions.ppt.ppt_product_manager import PptProductManager
from src.productions.ppt.schemas import DeckPageAsset
from src.productions.ppt.routes.html import (
    HTML_DELIVERY_STAGE,
    HTML_PAGE_GENERATION_CONTENT_PLAN_KEY,
    HTML_PAGE_GENERATION_PAGES_KEY,
    HTML_ROUTE_STAGE_SEQUENCE,
    PPT_HTML_PAGE_GENERATION_EXPERT_NAME,
    HtmlPageGenerationResult,
    build_html_page_generation_agent,
    build_ppt_html_page_generation_expert,
    build_html_route,
    deliver_html_route_quality,
    export_html_pptx,
    generate_html_pages,
    prepare_html_route_paths,
    prepare_html_template,
    save_html_route_pages,
)
from src.productions.ppt.routes.html import route as html_route_module
from src.productions.ppt.routes.html.html_to_pptx import (
    convert_html_pages_to_pptx,
    extract_html_slide_models,
    preflight_html_slide_models,
)
from src.productions.ppt.templates.html_registry import list_html_templates, load_html_template_package
from src.runtime.workspace import workspace_relative_path, workspace_root


def _write_route_test_image(tmpdir: Path) -> str:
    image_path = tmpdir / "slide_asset.png"
    Image.new("RGB", (640, 360), "#2457D6").save(image_path)
    return str(image_path)


class PptHtmlRouteTests(unittest.TestCase):
    def test_html_template_registry_loads_default_template(self) -> None:
        templates = list_html_templates()
        template = load_html_template_package("clean_business", aspect_ratio="16:9")

        self.assertGreaterEqual(len(templates), 1)
        self.assertEqual(template.template_id, "clean_business")
        self.assertEqual(template.viewport_width, 1280)
        self.assertIn("cover", template.page_types)
        self.assertEqual(template.pptx_strategy, "native_editable")
        self.assertIn("editable", template.editability_notes)

    def test_html_template_preparation_uses_free_design_without_template(self) -> None:
        template_stage = prepare_html_template(template_id="", aspect_ratio="16:9")

        self.assertEqual(template_stage.template.template_id, "free_design")
        self.assertEqual(template_stage.template.label, "Free Design")
        self.assertIn("No system HTML template", template_stage.template.editability_notes)

    def test_html_page_generation_agent_uses_free_html_contract(self) -> None:
        agent = build_ppt_html_page_generation_expert()
        compatibility_agent = build_html_page_generation_agent()

        self.assertEqual(agent.name, PPT_HTML_PAGE_GENERATION_EXPERT_NAME)
        self.assertEqual(compatibility_agent.name, PPT_HTML_PAGE_GENERATION_EXPERT_NAME)
        self.assertEqual(agent.output_key, "ppt_html_page_generation_agent_message")
        self.assertEqual(agent.include_contents, "none")
        self.assertIn("editable PPT-friendly HTML fragments", agent.description)
        self.assertIn("Do not use a fixed template", agent.instruction)
        self.assertIn("one HTML fragment per slide", agent.instruction)
        self.assertIn("PPTX conversion compatibility", agent.instruction)
        self.assertEqual({tool.__name__ for tool in agent.tools}, {"save_html_route_pages"})

    def test_save_html_route_pages_accepts_per_slide_html_fragments(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="给我做一个ppt，用来给幼儿园小朋友讲英语单词。图文并茂。小于10页。",
            inputs=[],
            output={"format": "pptx"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        tool_context = SimpleNamespace(
            state={
                HTML_PAGE_GENERATION_CONTENT_PLAN_KEY: content_plan.model_dump(mode="json"),
            }
        )
        pages = [
            {
                "slide_number": page.slide_number,
                "html": (
                    f"<section><div style='font-size:72px'>{page.title}</div>"
                    f"<p>{page.key_takeaway}</p></section>"
                ),
            }
            for page in content_plan.pages
        ]

        result = save_html_route_pages(pages, tool_context)

        self.assertEqual(result["status"], "success")
        self.assertEqual(result["page_count"], len(content_plan.pages))
        saved_pages = tool_context.state[HTML_PAGE_GENERATION_PAGES_KEY]
        self.assertEqual(len(saved_pages), len(content_plan.pages))
        self.assertIn('class="slide generated-slide', saved_pages[0]["html"])
        self.assertIn('data-slide-number="01"', saved_pages[0]["html"])

    def test_html_page_generation_message_exposes_route_local_asset_html_src(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 5 页 PPTX，用于产品发布。",
            inputs=[],
            output={"format": "pptx"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)

        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            image_path = Path(tmpdir) / "slide_asset.png"
            Image.new("RGB", (640, 360), "#2457D6").save(image_path)
            relative_image_path = workspace_relative_path(image_path)
            content_plan.pages[0].assets = [
                DeckPageAsset(
                    asset_id="material_figure_1",
                    source_kind="material_figure",
                    status="ready",
                    path=relative_image_path,
                    alt="Existing image",
                )
            ]
            content_plan.pages[1].assets = [
                DeckPageAsset(
                    asset_id="material_figure_2",
                    source_kind="material_figure",
                    status="ready",
                    path="source_01_input_1_files/page_003_image_01.png",
                    alt="Missing image",
                )
            ]
            route_output_dir = Path(tmpdir) / "route_output"
            paths = prepare_html_route_paths(route_output_dir)
            template = prepare_html_template(template_id="", aspect_ratio="16:9").template

            message = html_route_module._build_html_page_generation_user_message(
                content_plan=content_plan,
                template=template,
                paths=paths,
            )

            self.assertIn(f'"path": "{relative_image_path}"', message)
            self.assertIn('"html_src": "assets/slide_001_asset_001.png"', message)
            self.assertNotIn('"html_src": "file://', message)
            self.assertNotIn(str(workspace_root()), message)
            self.assertIn('"missing_asset": true', message)
            self.assertNotIn("/workspace/source_01_input_1_files/page_003_image_01.png", message)
            self.assertTrue((route_output_dir / "assets" / "slide_001_asset_001.png").exists())

    def test_html_route_builds_html_previews_and_pptx(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 5 页 PPTX，用于产品发布。",
            inputs=[],
            output={"format": "pptx"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)

        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            route_build = build_html_route(
                content_plan=content_plan,
                output_dir=Path(tmpdir),
                aspect_ratio="16:9",
            )

            pptx_path = workspace_root() / route_build.pptx_path
            html_path = workspace_root() / route_build.html_deck_path
            quality_path = workspace_root() / route_build.quality_report_path
            build_log_path = workspace_root() / route_build.build_log_path

            self.assertTrue(pptx_path.exists())
            self.assertTrue(html_path.exists())
            self.assertTrue(quality_path.exists())
            self.assertTrue(build_log_path.exists())
            self.assertEqual(len(route_build.preview_paths), len(content_plan.pages))
            self.assertEqual(route_build.warnings, [])

            prs = Presentation(str(pptx_path))
            pptx_text = "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
            )
            picture_count = sum(
                1
                for slide in prs.slides
                for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            )
            quality_report = json.loads(quality_path.read_text(encoding="utf-8"))
            build_log = json.loads(build_log_path.read_text(encoding="utf-8"))

            self.assertEqual(len(prs.slides), len(content_plan.pages))
            self.assertIn(content_plan.pages[0].title, pptx_text)
            self.assertEqual(picture_count, 0)
            self.assertTrue(quality_report["checks"]["pptx_contains_editable_text"])
            self.assertTrue(quality_report["checks"]["pptx_titles_present"])
            self.assertEqual(quality_report["route_stages"], list(HTML_ROUTE_STAGE_SEQUENCE))
            self.assertEqual(quality_report["delivery_stage"], HTML_DELIVERY_STAGE)
            self.assertEqual(quality_report["pptx_conversion"]["final_strategy"], "native_editable")
            self.assertEqual(len(quality_report["pptx_conversion"]["pages"]), len(content_plan.pages))
            self.assertEqual(build_log["workflow_name"], "HtmlRouteSequentialAgent")
            self.assertEqual(build_log["route_stages"], list(HTML_ROUTE_STAGE_SEQUENCE))
            self.assertEqual(build_log["delivery_stage"], HTML_DELIVERY_STAGE)
            self.assertEqual(build_log["pptx_conversion"]["editable_level"], "high")
            self.assertEqual(build_log["template_id"], "free_design")
            self.assertEqual(build_log["template_source"], "none")
            self.assertEqual(build_log["page_generation_mode"], "free_design")
            self.assertIn('data-deck-template="free_design"', html_path.read_text(encoding="utf-8"))

    def test_html_route_renders_ready_assets_into_html_previews_and_pptx(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 5 页 PPTX，用于产品发布。",
            inputs=[],
            output={"format": "pptx"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)

        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            asset_path = _write_route_test_image(Path(tmpdir))
            content_plan.pages[0].assets = [
                DeckPageAsset(
                    asset_id="cover_asset",
                    source_kind="user_upload",
                    status="ready",
                    path=asset_path,
                    alt="Cover visual",
                )
            ]
            route_build = build_html_route(
                content_plan=content_plan,
                output_dir=Path(tmpdir) / "route_output",
                aspect_ratio="16:9",
            )

            pptx_path = workspace_root() / route_build.pptx_path
            html_path = workspace_root() / route_build.html_deck_path
            quality_path = workspace_root() / route_build.quality_report_path
            prs = Presentation(str(pptx_path))
            picture_count = sum(
                1
                for slide in prs.slides
                for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            )
            quality_report = json.loads(quality_path.read_text(encoding="utf-8"))

            html_text = html_path.read_text(encoding="utf-8")
            self.assertIn("<img", html_text)
            self.assertIn('src="assets/slide_001_asset_001.png"', html_text)
            self.assertNotIn("file://", html_text)
            self.assertTrue((Path(tmpdir) / "route_output" / "assets" / "slide_001_asset_001.png").exists())
            self.assertGreaterEqual(picture_count, 1)
            self.assertEqual(quality_report["ready_asset_count"], 1)
            self.assertTrue(quality_report["checks"]["pptx_ready_assets_rendered"])

    def test_html_route_stage_functions_run_in_order(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 5 页 PPTX，用于产品发布。",
            inputs=[],
            output={"format": "pptx"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)

        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            paths = prepare_html_route_paths(Path(tmpdir))
            template_stage = prepare_html_template(template_id="clean_business", aspect_ratio="16:9")
            page_stage = generate_html_pages(
                content_plan=content_plan,
                template=template_stage.template,
                paths=paths,
            )
            pptx_stage = export_html_pptx(
                content_plan=content_plan,
                template=template_stage.template,
                page_generation=page_stage,
                paths=paths,
            )
            quality_stage = deliver_html_route_quality(
                content_plan=content_plan,
                template=template_stage.template,
                page_generation=page_stage,
                pptx_output=pptx_stage,
                paths=paths,
            )

            self.assertEqual(
                [
                    template_stage.stage,
                    page_stage.stage,
                    pptx_stage.stage,
                ],
                list(HTML_ROUTE_STAGE_SEQUENCE),
            )
            self.assertEqual(quality_stage.stage, HTML_DELIVERY_STAGE)
            self.assertTrue(page_stage.html_path.exists())
            self.assertTrue(pptx_stage.pptx_path.exists())
            self.assertEqual(quality_stage.quality_report["status"], "pass")

    def test_html_pptx_output_converts_agent_html_to_editable_pptx(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="给幼儿园小朋友讲 2 个英语单词：猫和狗。",
            inputs=[],
            output={"format": "pptx", "slide_count": 2},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        content_plan.pages = content_plan.pages[:2]

        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            paths = prepare_html_route_paths(Path(tmpdir))
            template = prepare_html_template(template_id="", aspect_ratio="16:9").template
            html_pages = [
                (
                    "<section>"
                    f"<h1 style='position:absolute;left:80px;top:64px;width:940px;height:90px;font-size:44px;color:#172033;'>{page.title}</h1>"
                    f"<p style='position:absolute;left:86px;top:174px;width:760px;height:52px;font-size:24px;color:#5f6472;'>{page.key_takeaway}</p>"
                    "<div style='position:absolute;left:86px;top:280px;width:260px;height:140px;background:#dff8f4;border:1px solid #0f9f8f;'></div>"
                    "</section>"
                )
                for page in content_plan.pages
            ]
            page_stage = HtmlPageGenerationResult(
                html_path=paths.html_path,
                preview_paths=[],
                html_pages=html_pages,
                generation_mode="llm_agent_html",
            )

            pptx_stage = export_html_pptx(
                content_plan=content_plan,
                template=template,
                page_generation=page_stage,
                paths=paths,
            )

            prs = Presentation(str(pptx_stage.pptx_path))
            pptx_text = "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
            )
            picture_count = sum(
                1
                for slide in prs.slides
                for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            )

            self.assertEqual(pptx_stage.pptx_strategy, "html_to_pptx")
            self.assertTrue(
                pptx_stage.warnings == []
                or any("Browser HTML-to-PPTX converter was unavailable" in warning for warning in pptx_stage.warnings)
            )
            self.assertEqual(pptx_stage.conversion_report["final_strategy"], "html_to_pptx")
            self.assertFalse(pptx_stage.conversion_report["fallback_used"])
            self.assertEqual(pptx_stage.conversion_report["pages"][0]["status"], "html_to_pptx")
            self.assertEqual(pptx_stage.conversion_report["pages"][0]["editable_level"], "high")
            self.assertEqual(len(prs.slides), len(content_plan.pages))
            self.assertIn(content_plan.pages[0].title, pptx_text)
            self.assertEqual(picture_count, 0)

    def test_html_pptx_output_fails_without_html_pages_instead_of_screenshot_fallback(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 3 页 PPTX，用于产品发布。",
            inputs=[],
            output={"format": "pptx"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)

        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            paths = prepare_html_route_paths(Path(tmpdir))
            template = prepare_html_template(template_id="", aspect_ratio="16:9").template
            fallback_pages = generate_html_pages(
                content_plan=content_plan,
                template=template,
                paths=paths,
            )
            page_stage = HtmlPageGenerationResult(
                html_path=fallback_pages.html_path,
                preview_paths=fallback_pages.preview_paths,
                html_pages=[],
                generation_mode="llm_agent_html",
            )

            pptx_stage = export_html_pptx(
                content_plan=content_plan,
                template=template,
                page_generation=page_stage,
                paths=paths,
            )

            self.assertEqual(pptx_stage.pptx_strategy, "html_to_pptx_failed")
            self.assertTrue(any("HTML-to-PPTX conversion failed" in warning for warning in pptx_stage.warnings))
            self.assertFalse(pptx_stage.pptx_path.exists())
            self.assertFalse(pptx_stage.conversion_report["fallback_used"])
            self.assertEqual(pptx_stage.conversion_report["final_strategy"], "html_to_pptx_failed")
            self.assertEqual(pptx_stage.conversion_report["pages"][0]["status"], "html_to_pptx_failed")
            self.assertEqual(pptx_stage.conversion_report["pages"][0]["editable_level"], "none")

    def test_html_to_pptx_converts_text_shapes_lines_and_images(self) -> None:
        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            tmp_path = Path(tmpdir)
            image_path = _write_route_test_image(tmp_path)
            template = prepare_html_template(template_id="", aspect_ratio="16:9").template
            pptx_path = tmp_path / "html_to_pptx.pptx"
            html_pages = [
                (
                    "<section style='background:#fff8e7;'>"
                    "<h1 style='position:absolute;left:64px;top:44px;width:720px;height:72px;font-size:38px;color:#172033;'>Cat 猫</h1>"
                    "<p style='position:absolute;left:68px;top:128px;width:680px;height:46px;font-size:22px;color:#5f6472;'>Cat means 猫. Listen and repeat.</p>"
                    "<div style='position:absolute;left:64px;top:208px;width:420px;height:2px;background:#ff6b57;'></div>"
                    "<div style='position:absolute;right:80px;bottom:72px;width:320px;height:150px;background:#dff8f4;border:2px solid #0f9f8f;border-radius:18px;'></div>"
                    f"<img src='{Path(image_path).as_uri()}' style='position:absolute;left:820px;top:220px;width:320px;height:190px;' />"
                    "</section>"
                )
            ]

            result = convert_html_pages_to_pptx(
                html_pages=html_pages,
                pptx_path=pptx_path,
                template=template,
            )

            prs = Presentation(str(pptx_path))
            pptx_text = "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
            )
            picture_count = sum(
                1
                for slide in prs.slides
                for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            )

            self.assertTrue(result.ok)
            self.assertEqual(result.errors, [])
            self.assertIn("Cat 猫", pptx_text)
            self.assertIn("Cat means 猫", pptx_text)
            self.assertEqual(picture_count, 1)

    def test_html_to_pptx_resolves_route_local_image_src(self) -> None:
        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            tmp_path = Path(tmpdir)
            assets_dir = tmp_path / "assets"
            assets_dir.mkdir(parents=True, exist_ok=True)
            image_path = assets_dir / "slide_001_asset_001.png"
            Image.new("RGB", (640, 360), "#2457D6").save(image_path)
            template = prepare_html_template(template_id="", aspect_ratio="16:9").template
            pptx_path = tmp_path / "route_local_asset.pptx"
            html_pages = [
                (
                    "<section>"
                    "<h1 style='position:absolute;left:64px;top:44px;width:720px;height:72px;font-size:38px;color:#172033;'>Cat 猫</h1>"
                    "<p style='position:absolute;left:68px;top:128px;width:680px;height:46px;font-size:22px;color:#5f6472;'>Cat means 猫.</p>"
                    "<img src='assets/slide_001_asset_001.png' style='position:absolute;left:820px;top:220px;width:320px;height:190px;' />"
                    "</section>"
                )
            ]

            result = convert_html_pages_to_pptx(
                html_pages=html_pages,
                pptx_path=pptx_path,
                template=template,
                asset_base_dir=tmp_path,
            )

            prs = Presentation(str(pptx_path))
            picture_count = sum(
                1
                for slide in prs.slides
                for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            )

            self.assertTrue(result.ok)
            self.assertEqual(result.errors, [])
            self.assertEqual(picture_count, 1)

    def test_html_to_pptx_preflight_reports_text_overflow(self) -> None:
        template = prepare_html_template(template_id="", aspect_ratio="16:9").template
        long_text = "这是一个非常长的文本 " * 80
        models = extract_html_slide_models(
            html_pages=[
                (
                    "<section>"
                    f"<p style='position:absolute;left:80px;top:80px;width:220px;height:28px;font-size:24px;'>{long_text}</p>"
                    "</section>"
                )
            ],
            template=template,
        )

        report = preflight_html_slide_models(models=models, template=template)

        self.assertFalse(report.ok)
        self.assertTrue(any(finding.code == "text_overflow_risk" for finding in report.findings))

    def test_html_to_pptx_expands_positioned_flow_text_to_editable_boxes(self) -> None:
        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            tmp_path = Path(tmpdir)
            template = prepare_html_template(template_id="", aspect_ratio="16:9").template
            pptx_path = tmp_path / "flow_text.pptx"
            html_pages = [
                (
                    "<section style='position:relative;width:1280px;height:720px;background:#f7f5ee;'>"
                    "<div style='position:absolute;left:56px;top:42px;width:610px;z-index:2;'>"
                    "<p style='margin:0 0 12px 0;font-size:18px;color:#2a8c93;font-weight:700;'>Thinking with Visual Primitives</p>"
                    "<h1 style='margin:0;font-size:48px;line-height:1.12;font-weight:800;color:#172033;'>这篇论文在解决什么问题</h1>"
                    "<p style='margin:24px 0 0 0;font-size:25px;line-height:1.42;font-weight:700;color:#c65a3b;'>核心不是让模型看得更多，而是让模型指得更准。</p>"
                    "</div>"
                    "<div style='position:absolute;left:57px;top:285px;width:500px;z-index:3;'>"
                    "<div style='display:flex;align-items:flex-start;margin-bottom:18px;'>"
                    "<div style='flex:0 0 32px;height:32px;font-size:18px;font-weight:700;'>1</div>"
                    "<p style='margin:2px 0 0 0;font-size:21px;line-height:1.35;color:#172033;'>多模态大模型已经能看图问答，但复杂空间推理仍容易崩溃</p>"
                    "</div>"
                    "<div style='display:flex;align-items:flex-start;margin-bottom:18px;'>"
                    "<div style='flex:0 0 32px;height:32px;font-size:18px;font-weight:700;'>2</div>"
                    "<p style='margin:2px 0 0 0;font-size:21px;line-height:1.35;color:#172033;'>传统 CoT 主要发生在语言空间，难以精确绑定图像中的对象位置</p>"
                    "</div>"
                    "</div>"
                    "</section>"
                )
            ]

            result = convert_html_pages_to_pptx(
                html_pages=html_pages,
                pptx_path=pptx_path,
                template=template,
            )

            prs = Presentation(str(pptx_path))
            pptx_text = "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
            )
            text_shape_count = sum(
                1
                for slide in prs.slides
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and str(shape.text or "").strip()
            )

            self.assertTrue(result.ok)
            self.assertEqual(result.errors, [])
            self.assertIn("这篇论文在解决什么问题", pptx_text)
            self.assertIn("多模态大模型", pptx_text)
            self.assertGreaterEqual(text_shape_count, 5)

    def test_html_pptx_output_falls_back_for_high_risk_html(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 2 页 PPTX，用于产品发布。",
            inputs=[],
            output={"format": "pptx", "slide_count": 2},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        content_plan.pages = content_plan.pages[:2]

        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            paths = prepare_html_route_paths(Path(tmpdir))
            template = prepare_html_template(template_id="", aspect_ratio="16:9").template
            fallback_pages = generate_html_pages(
                content_plan=content_plan,
                template=template,
                paths=paths,
            )
            page_stage = HtmlPageGenerationResult(
                html_path=fallback_pages.html_path,
                preview_paths=fallback_pages.preview_paths,
                html_pages=[
                    (
                        "<section>"
                        "<h1 style='position:absolute;left:80px;top:60px;width:780px;height:90px;font-size:42px;'>Valid title</h1>"
                        "<div style='position:absolute;left:120px;top:200px;width:620px;height:240px;background:linear-gradient(90deg,#fff,#000);'></div>"
                        "</section>"
                    )
                ],
                generation_mode="llm_agent_html",
            )

            pptx_stage = export_html_pptx(
                content_plan=content_plan,
                template=template,
                page_generation=page_stage,
                paths=paths,
            )

            self.assertEqual(pptx_stage.pptx_strategy, "html_to_pptx_failed")
            self.assertTrue(any("unsupported_css" in warning or "Gradient" in warning for warning in pptx_stage.warnings))
            self.assertEqual(pptx_stage.conversion_report["final_strategy"], "html_to_pptx_failed")
            self.assertFalse(pptx_stage.conversion_report["fallback_used"])
            self.assertTrue(pptx_stage.conversion_report["pages"][0]["errors"])
            self.assertFalse(pptx_stage.pptx_path.exists())


if __name__ == "__main__":
    unittest.main()
