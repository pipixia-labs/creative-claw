import tempfile
import unittest
import zipfile
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch

from openpyxl import Workbook
from pptx import Presentation

from src.agents.experts.anything_to_md.anything_to_md_expert import AnythingToMDExpert
from src.agents.experts.anything_to_md.tool import convert_anything_to_markdown
from src.runtime.expert_dispatcher import normalize_invoke_agent_parameters
from src.runtime.workspace import workspace_relative_path, workspace_root


def _build_ctx(state: dict) -> SimpleNamespace:
    return SimpleNamespace(
        session=SimpleNamespace(
            state=state,
            app_name="test_app",
            user_id="user_1",
            id="session_1",
        ),
    )


class AnythingToMDExpertTests(unittest.IsolatedAsyncioTestCase):
    def test_normalize_requires_structured_payload(self) -> None:
        with self.assertRaisesRegex(ValueError, "requires structured invoke_agent parameters"):
            normalize_invoke_agent_parameters(
                agent_name="AnythingToMD",
                prompt="convert this file",
                state={},
            )

    async def test_plain_text_file_converts_to_markdown_output_file(self) -> None:
        agent = AnythingToMDExpert(name="AnythingToMD")
        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmp_dir:
            source = Path(tmp_dir) / "notes.txt"
            source.write_text("hello world", encoding="utf-8")
            relative_source = workspace_relative_path(source)

            ctx = _build_ctx(
                {
                    "current_parameters": {"input_path": relative_source},
                    "turn_index": 2,
                    "step": 3,
                    "expert_step": 1,
                }
            )
            events = [event async for event in agent._run_async_impl(ctx)]

        current_output = events[0].actions.state_delta["current_output"]
        self.assertEqual(current_output["status"], "success")
        self.assertEqual(current_output["results"]["method"], "primary:plain_text")
        self.assertIn("hello world", current_output["output_text"])
        self.assertTrue(current_output["output_files"][0]["path"].endswith(".md"))
        self.assertEqual(
            events[0].actions.state_delta["anything_to_md_results"]["input_paths"],
            [relative_source],
        )

    async def test_pptx_file_uses_primary_powerpoint_converter(self) -> None:
        agent = AnythingToMDExpert(name="AnythingToMD")
        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmp_dir:
            source = Path(tmp_dir) / "deck.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = "Quarterly Plan"
            textbox = slide.shapes.add_textbox(100000, 1000000, 5000000, 1000000)
            textbox.text_frame.text = "Launch campaign"
            presentation.save(source)
            relative_source = workspace_relative_path(source)

            ctx = _build_ctx({"current_parameters": {"input_path": relative_source}})
            events = [event async for event in agent._run_async_impl(ctx)]

        current_output = events[0].actions.state_delta["current_output"]
        self.assertEqual(current_output["status"], "success")
        self.assertEqual(current_output["results"]["method"], "primary:pptx")
        self.assertIn("## Slide 1", current_output["output_text"])
        self.assertIn("Quarterly Plan", current_output["output_text"])
        self.assertIn("Launch campaign", current_output["output_text"])

    async def test_xlsx_file_uses_primary_excel_converter_with_limits(self) -> None:
        agent = AnythingToMDExpert(name="AnythingToMD")
        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmp_dir:
            source = Path(tmp_dir) / "sheet.xlsx"
            workbook = Workbook()
            worksheet = workbook.active
            worksheet.title = "Plan"
            worksheet.append(["Name", "Score", "Hidden"])
            worksheet.append(["Alpha", 10, "ignored"])
            worksheet.append(["Beta", 20, "ignored"])
            workbook.save(source)
            relative_source = workspace_relative_path(source)

            ctx = _build_ctx(
                {
                    "current_parameters": {
                        "input_path": relative_source,
                        "max_rows": 2,
                        "max_cols": 2,
                    }
                }
            )
            events = [event async for event in agent._run_async_impl(ctx)]

        current_output = events[0].actions.state_delta["current_output"]
        self.assertEqual(current_output["status"], "success")
        self.assertEqual(current_output["results"]["method"], "primary:excel")
        self.assertIn("## Sheet: Plan", current_output["output_text"])
        self.assertIn("| Name | Score |", current_output["output_text"])
        self.assertIn("| Alpha | 10 |", current_output["output_text"])
        self.assertNotIn("Beta", current_output["output_text"])
        self.assertNotIn("Hidden", current_output["output_text"])

    async def test_pdf_file_prefers_doc2x_v3_when_key_is_available(self) -> None:
        captured_call = {}

        class _FakeDoc2X:
            def __init__(self, **kwargs):
                captured_call["init"] = kwargs

            def pdf2file(self, **kwargs):
                captured_call["pdf2file"] = kwargs
                output_path = Path(kwargs["output_path"])
                output_path.mkdir(parents=True, exist_ok=True)
                zip_path = output_path / "doc2x_md.zip"
                with zipfile.ZipFile(zip_path, "w") as archive:
                    archive.writestr("combined.md", "# Doc2X Markdown\n\nMerged PDF content")
                return [str(zip_path)], [{"error": "", "path": ""}], False

        agent = AnythingToMDExpert(name="AnythingToMD")
        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmp_dir:
            source = Path(tmp_dir) / "source.pdf"
            source.write_bytes(b"%PDF-1.4\n")
            relative_source = workspace_relative_path(source)
            ctx = _build_ctx(
                {
                    "current_parameters": {
                        "input_path": relative_source,
                        "formula_level": 1,
                    }
                }
            )
            with (
                patch.dict("os.environ", {"DOC2X_API_KEY": "test-key"}, clear=True),
                patch.dict("sys.modules", {"pdfdeal": SimpleNamespace(Doc2X=_FakeDoc2X)}),
            ):
                events = [event async for event in agent._run_async_impl(ctx)]

        current_output = events[0].actions.state_delta["current_output"]
        self.assertEqual(current_output["status"], "success")
        self.assertEqual(current_output["results"]["method"], "primary:doc2x_v3")
        self.assertIn("Merged PDF content", current_output["output_text"])
        self.assertEqual(captured_call["init"]["apikey"], "test-key")
        self.assertEqual(captured_call["pdf2file"]["model"], "v3-2026")
        self.assertEqual(captured_call["pdf2file"]["output_format"], "md")
        self.assertEqual(captured_call["pdf2file"]["formula_level"], 1)

    def test_pdf_file_skips_doc2x_without_key_and_uses_local_pdf_converter(self) -> None:
        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmp_dir:
            source = Path(tmp_dir) / "source.pdf"
            source.write_bytes(b"%PDF-1.4\n")
            relative_source = workspace_relative_path(source)
            with (
                patch.dict("os.environ", {}, clear=True),
                patch(
                    "src.agents.experts.anything_to_md.tool._convert_pdf",
                    return_value="# Local PDF Markdown",
                ) as local_pdf,
            ):
                result = convert_anything_to_markdown({"input_path": relative_source})

        self.assertEqual(result["status"], "success")
        self.assertEqual(result["results"]["method"], "primary:pdf")
        self.assertIn("Local PDF Markdown", result["output_text"])
        local_pdf.assert_called_once()

    def test_unsupported_file_can_fallback_to_markitdown(self) -> None:
        class _FakeResult:
            text_content = "# Converted by fallback"

        class _FakeMarkItDown:
            def convert(self, source):
                return _FakeResult()

        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmp_dir:
            source = Path(tmp_dir) / "source.bin"
            source.write_bytes(b"binary")
            relative_source = workspace_relative_path(source)
            with patch.dict("sys.modules", {"markitdown": SimpleNamespace(MarkItDown=_FakeMarkItDown)}):
                result = convert_anything_to_markdown({"input_path": relative_source})

        self.assertEqual(result["status"], "success")
        self.assertEqual(result["results"]["method"], "fallback:markitdown")
        self.assertIn("Converted by fallback", result["output_text"])

    def test_primary_failure_can_fallback_to_markitdown(self) -> None:
        class _FakeResult:
            text_content = "# Converted after primary failure"

        class _FakeMarkItDown:
            def convert(self, source):
                return _FakeResult()

        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmp_dir:
            source = Path(tmp_dir) / "source.pdf"
            source.write_bytes(b"placeholder")
            relative_source = workspace_relative_path(source)
            with (
                patch.dict("sys.modules", {"markitdown": SimpleNamespace(MarkItDown=_FakeMarkItDown)}),
                patch(
                    "src.agents.experts.anything_to_md.tool._convert_file_primary",
                    side_effect=RuntimeError("primary failed"),
                ),
            ):
                result = convert_anything_to_markdown({"input_path": relative_source})

        self.assertEqual(result["status"], "success")
        self.assertEqual(result["results"]["method"], "fallback:markitdown")
        self.assertIn("Converted after primary failure", result["output_text"])

    def test_missing_input_returns_error(self) -> None:
        result = convert_anything_to_markdown({})

        self.assertEqual(result["status"], "error")
        self.assertIn("input_path is required", result["message"])

    def test_url_input_returns_error(self) -> None:
        result = convert_anything_to_markdown({"url": "https://example.com/source.pdf"})

        self.assertEqual(result["status"], "error")
        self.assertIn("URL inputs are not supported", result["message"])
