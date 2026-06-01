import tempfile
import unittest
import warnings
from pathlib import Path
from types import SimpleNamespace
from typing import Any, AsyncGenerator
from unittest.mock import patch

from google.adk import Context, Workflow
from google.adk.apps import App, ResumabilityConfig
from google.adk.agents import LlmAgent
from google.adk.artifacts import InMemoryArtifactService
from google.adk.models import BaseLlm
from google.adk.models.llm_request import LlmRequest
from google.adk.models.llm_response import LlmResponse
from google.adk.runners import Runner
from google.adk.sessions import InMemorySessionService
from google.adk.tools.tool_context import ToolContext
from google.adk.workflow import node
from google.genai.types import Content, FunctionCall, FunctionResponse, Part
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pydantic import PrivateAttr

from src.productions.ppt.planning.content_planner import (
    _build_content_planning_user_message,
    _select_content_page_type,
)
from src.productions.ppt.ppt_product_manager import (
    PptProductManager,
    PptProductRequest,
    ProductPptSkillRegistry,
)
from src.productions.ppt.ppt_product_manager.ppt_product_manager import (
    PPT_ADK_CONFIRMATION_REQUEST_STATE_KEY,
    PPT_CONFIRMED_REQUIREMENT_STATE_KEY,
    PPT_ASSET_RESOLUTION_WORKFLOW_OUTPUT_KEY,
    PPT_AUTO_CONFIRM_WORKFLOW_OUTPUT_KEY,
    PPT_CONTENT_PLAN_REVISION_OUTPUT_STATE_KEY,
    PPT_CONTENT_PLAN_REVISION_RESULT_STATE_KEY,
    PPT_CONTENT_PLAN_REVISION_WORKFLOW_OUTPUT_KEY,
    PPT_CONTENT_PLANNING_WORKFLOW_OUTPUT_KEY,
    PPT_CONTENT_PLAN_CONFIRMATION_WORKFLOW_OUTPUT_KEY,
    PPT_FINAL_DELIVERY_WORKFLOW_OUTPUT_KEY,
    PPT_INITIAL_REQUEST_WORKFLOW_OUTPUT_KEY,
    PPT_PRIVATE_SKILL_EXECUTION_RESULT_STATE_KEY,
    PPT_PRIVATE_SKILL_EXECUTION_WORKFLOW_OUTPUT_KEY,
    PPT_PRIVATE_SKILL_DELIVERY_WORKFLOW_OUTPUT_KEY,
    PPT_PRODUCT_REQUEST_STATE_KEY,
    PPT_REQUIREMENT_ANALYSIS_AGENT_MESSAGE_KEY,
    PPT_REQUIREMENT_ANALYSIS_RESULT_STATE_KEY,
    PPT_REQUIREMENT_ANALYSIS_WORKFLOW_OUTPUT_KEY,
    PPT_REQUIREMENT_REVISION_OUTPUT_STATE_KEY,
    PPT_REQUIREMENT_REVISION_RESULT_STATE_KEY,
    PPT_REQUIREMENT_REVISION_WORKFLOW_OUTPUT_KEY,
    PPT_SYSTEM_SELECTION_OUTPUT_STATE_KEY,
    PPT_SYSTEM_SELECTION_RESULT_STATE_KEY,
    PPT_SYSTEM_SELECTION_WORKFLOW_OUTPUT_KEY,
    PPT_REQUIREMENT_CONFIRMATION_WORKFLOW_OUTPUT_KEY,
    PPT_ROUTE_EXECUTION_WORKFLOW_OUTPUT_KEY,
    PPT_SOURCE_PREPARATION_WORKFLOW_OUTPUT_KEY,
    PPT_SYSTEM_SELECTION_AGENT_MESSAGE_KEY,
    PPT_SYSTEM_SELECTION_STATE_KEY,
    _build_product_manager_skill_run_user_message,
    _snapshot_workspace_pptx_files,
)
from src.productions.ppt.schemas import (
    DeckContentPlan,
    DeckPageAsset,
    DeckPagePlan,
    HtmlRouteBuildPackage,
    HtmlTemplatePackage,
    PptAssetResolutionResult,
    PptAdkConfirmationResponse,
    PptContentPlanRevisionResult,
    PptContentPlanningResult,
    PptFinalDeliveryResult,
    PptPrivateSkillBuild,
    PptPrivateSkillDeliveryResult,
    PptPrivateSkillExecutionResult,
    PptRequirementAnalysisResult,
    PptRequirementRevisionResult,
    PptRouteExecutionResult,
    PptSystemSelectionResult,
    PptWorkflowState,
    SourceUnderstanding,
)
from src.productions.ppt.routes.html import PPT_HTML_PAGE_GENERATION_EXPERT_NAME
from src.productions.ppt.routes.svg import (
    PPT_DESIGN_STRATEGY_EXPERT_NAME,
    PPT_SVG_DECK_EXECUTOR_EXPERT_NAME,
    PPT_SVG_EXECUTION_PLAN_STATE_KEY,
)
from src.runtime.workspace import (
    build_workspace_file_record,
    resolve_workspace_path,
    workspace_relative_path,
    workspace_root,
)
from src.skills.registry import SkillRegistry


ADK_REQUEST_CONFIRMATION_FUNCTION_NAME = "adk_request_confirmation"


async def _collect_events(events: AsyncGenerator[Any, None]) -> list[Any]:
    """Drain an ADK event stream into a list for assertions."""
    return [event async for event in events]


def _write_markdown_source(name: str, text: str) -> str:
    source_dir = workspace_root() / "inbox" / "ppt_product_manager_tests"
    source_dir.mkdir(parents=True, exist_ok=True)
    source_path = source_dir / name
    source_path.write_text(text, encoding="utf-8")
    return workspace_relative_path(source_path)


def _write_test_image(name: str) -> str:
    image_dir = workspace_root() / "inbox" / "ppt_product_manager_tests"
    image_dir.mkdir(parents=True, exist_ok=True)
    image_path = image_dir / name
    Image.new("RGB", (640, 360), "#2457D6").save(image_path)
    return workspace_relative_path(image_path)


def _page(slide_number: int, page_type: str) -> DeckPagePlan:
    return DeckPagePlan(
        slide_number=slide_number,
        page_type=page_type,
        title=f"Slide {slide_number}",
        purpose="Explain the planned message.",
        key_takeaway="Audience remembers the core point.",
        asset_intent="Use a simple supporting visual.",
    )


class _DictState(dict):
    def to_dict(self) -> dict:
        return dict(self)


class _PptContentPlannerToolCallingFakeLlm(BaseLlm):
    """Fake planner model that drives the real content-planning tools."""

    _markdown_plan: str = PrivateAttr()
    _requests: list[LlmRequest] = PrivateAttr(default_factory=list)

    def __init__(self, *, markdown_plan: str) -> None:
        super().__init__(model="fake-ppt-content-planner")
        self._markdown_plan = markdown_plan

    @property
    def requests(self) -> list[LlmRequest]:
        return self._requests

    async def generate_content_async(
        self,
        llm_request: LlmRequest,
        stream: bool = False,
    ) -> AsyncGenerator[LlmResponse, None]:
        self._requests.append(llm_request)
        request_index = len(self._requests)
        if request_index == 1:
            yield LlmResponse(
                content=Content(
                    role="model",
                    parts=[
                        Part(
                            function_call=FunctionCall(
                                name="read_ppt_markdown_sources",
                                args={},
                            )
                        )
                    ],
                )
            )
            return
        if request_index == 2:
            yield LlmResponse(
                content=Content(
                    role="model",
                    parts=[
                        Part(
                            function_call=FunctionCall(
                                name="save_ppt_deck_content_plan_markdown",
                                args={"markdown": self._markdown_plan},
                            )
                        )
                    ],
                )
            )
            return
        yield LlmResponse(
            content=Content(
                role="model",
                parts=[Part(text="PptContentPlanningAgent produced DeckContentPlan.")],
            )
        )


class _PptProductManagerToolCallingFakeLlm(BaseLlm):
    """Fake PPT product-manager model that drives a planned tool sequence."""

    _function_calls: list[FunctionCall] = PrivateAttr()
    _final_text: str = PrivateAttr()
    _requests: list[LlmRequest] = PrivateAttr(default_factory=list)

    def __init__(
        self,
        *,
        function_calls: list[FunctionCall],
        final_text: str,
    ) -> None:
        super().__init__(model="fake-ppt-product-manager-internal")
        self._function_calls = function_calls
        self._final_text = final_text

    @property
    def requests(self) -> list[LlmRequest]:
        return self._requests

    async def generate_content_async(
        self,
        llm_request: LlmRequest,
        stream: bool = False,
    ) -> AsyncGenerator[LlmResponse, None]:
        self._requests.append(llm_request)
        call_index = len(self._requests) - 1
        if call_index < len(self._function_calls):
            yield LlmResponse(
                content=Content(
                    role="model",
                    parts=[Part(function_call=self._function_calls[call_index])],
                )
            )
            return
        yield LlmResponse(
            content=Content(role="model", parts=[Part(text=self._final_text)])
        )


def _function_declaration_names(llm_request: LlmRequest) -> list[str]:
    names: list[str] = []
    for tool in llm_request.config.tools or []:
        for declaration in tool.function_declarations or []:
            if declaration.name:
                names.append(declaration.name)
    return names


class _FakeRemoteResponse:
    def __init__(self, data: bytes, headers: dict[str, str]):
        self._data = data
        self._offset = 0
        self.headers = headers

    def __enter__(self):
        return self

    def __exit__(self, exc_type, exc, traceback):
        return False

    def read(self, size: int = -1) -> bytes:
        if self._offset >= len(self._data):
            return b""
        if size is None or size < 0:
            size = len(self._data) - self._offset
        chunk = self._data[self._offset : self._offset + size]
        self._offset += len(chunk)
        return chunk


async def _fake_source_converter(source_input, parameters: dict) -> dict:
    output_path = str(parameters["output_path"])
    output_file = resolve_workspace_path(output_path)
    output_file.parent.mkdir(parents=True, exist_ok=True)
    asset_dir = output_file.parent / "figures"
    asset_dir.mkdir(parents=True, exist_ok=True)
    chart_path = asset_dir / "activation.png"
    Image.new("RGB", (640, 360), "#43A6FF").save(chart_path)
    markdown = "# Growth Launch\n\n![Activation chart](figures/activation.png)\n"
    output_file.write_text(markdown, encoding="utf-8")
    return {
        "status": "success",
        "message": "converted",
        "output_text": markdown,
        "results": {
            "method": "test:markdown",
            "output_path": output_path,
        },
        "output_files": [
            build_workspace_file_record(
                output_file,
                description="Converted Markdown source.",
                source="expert",
                name=output_file.name,
            )
        ],
    }


class PptProductManagerTests(unittest.IsolatedAsyncioTestCase):
    def test_instruction_prioritizes_pptx_and_adk_workflow(self) -> None:
        manager = PptProductManager()

        instruction = manager.build_instruction()

        self.assertIsInstance(manager, LlmAgent)
        self.assertIs(manager.build_agent(), manager)
        self.assertEqual(manager.include_contents, "none")
        self.assertEqual(
            {tool.__name__ for tool in manager.tools},
            {
                "list_product_ppt_skills",
                "read_product_ppt_skill",
                "read_product_ppt_skill_file",
                "list_session_files",
                "list_dir",
                "glob",
                "grep",
                "read_file",
                "write_file",
                "edit_file",
                "exec_command",
                "process_session",
                "list_ppt_experts",
                "invoke_ppt_expert",
                "save_ppt_system_selection",
                "save_ppt_private_skill_html",
                "save_ppt_private_skill_pptx",
                "save_ppt_design_strategy",
                "save_ppt_svg_execution_plan",
                "read_ppt_svg_execution_plan",
                "save_ppt_svg_page",
                "check_ppt_svg_quality",
                "export_ppt_svg_to_pptx",
                "dispatch_ppt_route",
            },
        )
        self.assertIn("PPT and PowerPoint production", instruction)
        self.assertIn("ADK workflow", instruction)
        self.assertIn("currently implemented built-in route", instruction)
        self.assertIn("PPT system selection", instruction)
        self.assertIn("skills/product-ppt-skills", instruction)
        self.assertIn("built-in HTML route", instruction)
        self.assertIn("uploaded input includes PPTX/PPTM/POTX/POTM", instruction)
        self.assertIn("choose between the built-in HTML route and the built-in SVG route", instruction)
        self.assertIn("hard-coded keyword-to-skill rules", instruction)
        self.assertIn("you run that skill workflow directly as PptProductManager", instruction)
        self.assertIn("PptHtmlPageGenerationExpert", instruction)
        self.assertIn("PptDesignStrategyExpert", instruction)
        self.assertIn("PptSvgDeckExecutorExpert", instruction)
        self.assertIn("invoke_ppt_expert", instruction)
        self.assertIn("export_ppt_svg_to_pptx", instruction)
        self.assertIn("template-based PPTX workflow", instruction)
        self.assertIn("immediately after the file is generated and verified", instruction)
        self.assertIn("must not block delivery", instruction)
        self.assertIn("Template analysis artifacts", instruction)
        self.assertIn("Do not stop after `thumbnail.py`", instruction)
        self.assertIn("return a concrete blocker", instruction)
        self.assertIn("Do not claim PPTX generation succeeded", instruction)

    def test_ppt_product_request_schema_normalizes_public_contract(self) -> None:
        request = PptProductRequest.model_validate(
            {
                "task": "  生成一个产品发布 PPT。 ",
                "inputs": {"brief": "inbox/demo/brief.md"},
                "output": None,
            }
        )

        self.assertEqual(request.task, "生成一个产品发布 PPT。")
        self.assertEqual(request.inputs, {"brief": "inbox/demo/brief.md"})
        self.assertEqual(request.output, {})
        self.assertEqual(
            request.to_state_dict(),
            {
                "task": "生成一个产品发布 PPT。",
                "inputs": {"brief": "inbox/demo/brief.md"},
                "output": {},
            },
        )

    def test_ppt_workflow_state_schema_preserves_dict_contract(self) -> None:
        workflow_state = PptWorkflowState.model_validate(
            {
                "workflow_id": " workflow-1 ",
                "stage": " awaiting_requirement_confirmation ",
                "revision": "2",
                "waiting_since_turn_index": "3",
                "confirmed_requirement": {"topic": "产品发布"},
                "raw_inputs": None,
                "debug_marker": "kept",
            }
        ).to_state_dict()

        self.assertEqual(workflow_state["workflow_id"], "workflow-1")
        self.assertEqual(workflow_state["stage"], "awaiting_requirement_confirmation")
        self.assertEqual(workflow_state["revision"], 2)
        self.assertEqual(workflow_state["waiting_since_turn_index"], 3)
        self.assertEqual(workflow_state["confirmed_requirement"]["topic"], "产品发布")
        self.assertEqual(workflow_state["debug_marker"], "kept")
        self.assertNotIn("raw_inputs", workflow_state)

    def test_private_skill_build_schema_preserves_dict_result_contract(self) -> None:
        private_build = PptPrivateSkillBuild.model_validate(
            {
                "status": " success ",
                "source": " save_ppt_private_skill_html ",
                "output_path": " generated/session/private.html ",
                "output_files": [{"path": "generated/session/private.html"}],
                "extra_private_field": "kept",
            }
        ).to_state_dict()

        execution = PptPrivateSkillExecutionResult.model_validate(
            {
                "skill_name": " magazine ",
                "output_format": " html ",
                "input_signature": " sig ",
                "private_build": private_build,
            }
        )
        delivery = PptPrivateSkillDeliveryResult.model_validate(
            {
                "product_result": {
                    "status": "success",
                    "phase": "private_skill_delivery",
                    "message": "done",
                    "selected_route": "html",
                },
                "private_build": private_build,
            }
        )

        self.assertIsInstance(execution.private_build, dict)
        self.assertEqual(execution.skill_name, "magazine")
        self.assertEqual(execution.private_build["source"], "save_ppt_private_skill_html")
        self.assertEqual(execution.private_build["extra_private_field"], "kept")
        self.assertEqual(delivery.private_build["output_path"], "generated/session/private.html")

    def test_ppt_adk_confirmation_response_maps_to_text_protocol(self) -> None:
        confirm_response = PptAdkConfirmationResponse.model_validate(
            {"action": "approve", "message": "ignored for confirmation"}
        )
        revise_response = PptAdkConfirmationResponse.model_validate(
            {"action": "revise", "message": "改成 5 页，并面向投资人。"}
        )
        raw_text_response = PptAdkConfirmationResponse.model_validate("改成 4 页。")

        self.assertEqual(confirm_response.to_user_response(), "确认")
        self.assertEqual(revise_response.to_user_response(), "改成 5 页，并面向投资人。")
        self.assertEqual(raw_text_response.action, "revise")
        self.assertEqual(raw_text_response.to_user_response(), "改成 4 页。")

    async def test_run_product_request_rejects_non_object_output(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={})

        result = await manager.run_product_request(
            task="生成一个产品发布 PPT。",
            output=["pptx"],
            tool_context=tool_context,
        )

        self.assertEqual(result["status"], "error")
        self.assertEqual(result["product_line"], "ppt")
        self.assertEqual(result["result_schema_version"], "ppt-product-result-v1")
        self.assertIn("output to be an object", result["message"])

    async def test_adk_tool_confirmation_bridge_advances_to_content_plan_gate(self) -> None:
        manager = PptProductManager()
        task = "做一个 3 页 PPTX，用于产品发布，受众为管理层。"
        output = {"format": "pptx", "slide_count": 3, "confirmation_mode": "adk_hitl"}

        async def _prepare_initial_requirement_phase_stub(**kwargs: Any) -> PptRequirementAnalysisResult:
            requirement = manager.prepare_confirmed_requirement(
                task=kwargs["task"],
                inputs=kwargs["raw_inputs"],
                output=kwargs["output"],
                source_understanding=kwargs["source_understanding"],
            )
            return PptRequirementAnalysisResult(
                confirmed_requirement=requirement,
                analysis_output={"source": "stub"},
                agent_message="stubbed requirement analysis",
            )

        async def _select_ppt_system_phase_stub(**kwargs: Any) -> PptSystemSelectionResult:
            requirement = kwargs["requirement"]
            return PptSystemSelectionResult(
                system_selection=manager._build_default_system_selection(requirement),
                selection_output={"source": "stub"},
                agent_message="stubbed system selection",
            )

        async def run_ppt_product_for_adk_hitl(task: str, tool_context: ToolContext) -> dict[str, Any]:
            return await manager.run_product_request(
                task=task,
                inputs=[],
                output=output,
                tool_context=tool_context,
                content_plan_builder=manager.build_initial_deck_content_plan,
            )

        fake_llm = _PptProductManagerToolCallingFakeLlm(
            function_calls=[
                FunctionCall(
                    name="run_ppt_product_for_adk_hitl",
                    id="run_ppt_product_for_adk_hitl_call",
                    args={"task": task},
                )
            ],
            final_text="PPT HITL bridge completed.",
        )
        with warnings.catch_warnings():
            warnings.filterwarnings("ignore", message=r"\[EXPERIMENTAL\].*", category=UserWarning)
            agent = LlmAgent(
                name="PptAdkToolConfirmationBridgeRoot",
                model=fake_llm,
                instruction="Call the PPT product tool.",
                tools=[run_ppt_product_for_adk_hitl],
            )
            app = App(
                name="PptAdkToolConfirmationBridgeApp",
                root_agent=agent,
                resumability_config=ResumabilityConfig(is_resumable=True),
            )
        session_service = InMemorySessionService()
        artifact_service = InMemoryArtifactService()
        runner = Runner(
            app=app,
            session_service=session_service,
            artifact_service=artifact_service,
        )
        user_id = "user-ppt-adk-tool-confirmation"
        session_id = "session-ppt-adk-tool-confirmation"
        try:
            await session_service.create_session(
                app_name=app.name,
                user_id=user_id,
                session_id=session_id,
                state={
                    "sid": "ppt-adk-hitl-test",
                    "turn_index": 1,
                    "step": 1,
                },
            )
            with (
                patch.object(
                    manager,
                    "_prepare_initial_requirement_phase",
                    side_effect=_prepare_initial_requirement_phase_stub,
                ),
                patch.object(
                    manager,
                    "_select_ppt_system_phase",
                    side_effect=_select_ppt_system_phase_stub,
                ),
            ):
                first_events = await _collect_events(
                    runner.run_async(
                        user_id=user_id,
                        session_id=session_id,
                        new_message=Content(role="user", parts=[Part(text="Start PPT product.")]),
                    )
                )
                first_invocation_id = first_events[0].invocation_id
                first_confirmation_calls = [
                    call
                    for event in first_events
                    for call in event.get_function_calls()
                    if call.name == ADK_REQUEST_CONFIRMATION_FUNCTION_NAME
                ]
                self.assertEqual(len(first_confirmation_calls), 1)
                first_request_payload = first_confirmation_calls[0].args["toolConfirmation"]["payload"]
                self.assertEqual(first_request_payload["stage"], "awaiting_requirement_confirmation")
                self.assertEqual(first_request_payload["confirmation_type"], "requirement")
                self.assertEqual(first_request_payload["allowed_actions"], ["confirm", "revise"])

                first_session = await session_service.get_session(
                    app_name=app.name,
                    user_id=user_id,
                    session_id=session_id,
                )
                self.assertEqual(
                    first_session.state["ppt_product_result"]["status"],
                    "awaiting_requirement_confirmation",
                )
                self.assertEqual(
                    first_session.state[PPT_ADK_CONFIRMATION_REQUEST_STATE_KEY]["confirmation_id"],
                    first_request_payload["confirmation_id"],
                )

                resume_payload = {
                    "schema_version": "ppt-adk-confirmation-response-v1",
                    "action": "confirm",
                    "confirmation_id": first_request_payload["confirmation_id"],
                    "stage": first_request_payload["stage"],
                }
                second_events = await _collect_events(
                    runner.run_async(
                        user_id=user_id,
                        session_id=session_id,
                        invocation_id=first_invocation_id,
                        new_message=Content(
                            role="user",
                            parts=[
                                Part(
                                    function_response=FunctionResponse(
                                        id=first_confirmation_calls[0].id,
                                        name=ADK_REQUEST_CONFIRMATION_FUNCTION_NAME,
                                        response={"confirmed": True, "payload": resume_payload},
                                    )
                                )
                            ],
                        ),
                    )
                )
        finally:
            await runner.close()

        second_confirmation_calls = [
            call
            for event in second_events
            for call in event.get_function_calls()
            if call.name == ADK_REQUEST_CONFIRMATION_FUNCTION_NAME
        ]
        self.assertEqual(second_confirmation_calls, [])

        second_session = await session_service.get_session(
            app_name=app.name,
            user_id=user_id,
            session_id=session_id,
        )
        self.assertEqual(
            second_session.state["ppt_product_result"]["status"],
            "awaiting_content_plan_confirmation",
        )
        self.assertEqual(
            second_session.state[PPT_ADK_CONFIRMATION_REQUEST_STATE_KEY],
            {},
        )

    async def test_adk_tool_confirmation_bridge_requests_content_plan_gate(self) -> None:
        manager = PptProductManager()
        task = "做一个 3 页 PPTX，用于产品发布，受众为管理层。"
        output = {"format": "pptx", "slide_count": 3, "confirmation_mode": "adk_hitl"}
        requirement = manager.prepare_confirmed_requirement(task=task, inputs=[], output=output)
        content_plan = manager.build_initial_deck_content_plan(requirement)
        workflow_state = {
            "workflow_id": "ppt-adk-hitl-content-test:ppt:1",
            "stage": "awaiting_content_plan_confirmation",
            "revision": 2,
            "task": task,
            "raw_inputs": [],
            "output": output,
            "confirmed_requirement": requirement.model_dump(mode="json"),
            "deck_content_plan": content_plan.model_dump(mode="json"),
            "system_selection": manager._build_default_system_selection(requirement),
            "waiting_since_turn_index": 1,
            "confirmation_id": "ppt-adk-hitl-content-test:ppt:1:awaiting_content_plan_confirmation:2:1",
        }

        async def run_ppt_product_for_adk_hitl(task: str, tool_context: ToolContext) -> dict[str, Any]:
            return await manager.run_product_request(
                task=task,
                inputs=[],
                output=output,
                tool_context=tool_context,
            )

        fake_llm = _PptProductManagerToolCallingFakeLlm(
            function_calls=[
                FunctionCall(
                    name="run_ppt_product_for_adk_hitl",
                    id="run_ppt_product_for_adk_hitl_content_call",
                    args={"task": task},
                )
            ],
            final_text="PPT HITL content gate completed.",
        )
        with warnings.catch_warnings():
            warnings.filterwarnings("ignore", message=r"\[EXPERIMENTAL\].*", category=UserWarning)
            agent = LlmAgent(
                name="PptAdkToolConfirmationContentGateRoot",
                model=fake_llm,
                instruction="Call the PPT product tool.",
                tools=[run_ppt_product_for_adk_hitl],
            )
            app = App(
                name="PptAdkToolConfirmationContentGateApp",
                root_agent=agent,
                resumability_config=ResumabilityConfig(is_resumable=True),
            )
        session_service = InMemorySessionService()
        runner = Runner(
            app=app,
            session_service=session_service,
            artifact_service=InMemoryArtifactService(),
        )
        user_id = "user-ppt-adk-tool-confirmation-content"
        session_id = "session-ppt-adk-tool-confirmation-content"
        try:
            await session_service.create_session(
                app_name=app.name,
                user_id=user_id,
                session_id=session_id,
                state={
                    "sid": "ppt-adk-hitl-content-test",
                    "turn_index": 1,
                    "step": 1,
                    "ppt_workflow_state": workflow_state,
                },
            )
            events = await _collect_events(
                runner.run_async(
                    user_id=user_id,
                    session_id=session_id,
                    new_message=Content(role="user", parts=[Part(text="Request current content gate.")]),
                )
            )
        finally:
            await runner.close()

        confirmation_calls = [
            call
            for event in events
            for call in event.get_function_calls()
            if call.name == ADK_REQUEST_CONFIRMATION_FUNCTION_NAME
        ]
        self.assertEqual(len(confirmation_calls), 1)
        request_payload = confirmation_calls[0].args["toolConfirmation"]["payload"]
        self.assertEqual(request_payload["stage"], "awaiting_content_plan_confirmation")
        self.assertEqual(request_payload["confirmation_type"], "content_plan")
        self.assertEqual(request_payload["confirmation_id"], workflow_state["confirmation_id"])

        session = await session_service.get_session(
            app_name=app.name,
            user_id=user_id,
            session_id=session_id,
        )
        self.assertEqual(
            session.state["ppt_product_result"]["status"],
            "awaiting_content_plan_confirmation",
        )
        self.assertEqual(
            session.state[PPT_ADK_CONFIRMATION_REQUEST_STATE_KEY]["confirmation_id"],
            workflow_state["confirmation_id"],
        )

    def test_private_pptx_template_skill_prompt_enforces_delivery_checklist(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="用这个模板做一个毕业答辩 PPT。",
            inputs=[{"name": "template.pptx", "path": "inbox/demo/template.pptx"}],
            output={"format": "pptx", "route": "xml"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)

        user_message = _build_product_manager_skill_run_user_message(
            requirement=requirement,
            content_plan=content_plan,
            system_selection={
                "system_type": "private_skill",
                "route": "xml",
                "skill_name": "pptx",
                "output_format": "pptx",
                "reason": "Use pptx skill for uploaded template.",
            },
            skill_content="# PPTX Skill\nTemplate-Based Workflow",
            available_experts=[],
        )

        self.assertIn("user_template_pptx_workflow_checklist", user_message)
        self.assertIn("confirmed_requirement_json.source_inputs", user_message)
        self.assertIn("template_requirement.template_path", user_message)
        self.assertIn("list_session_files(section='uploaded')", user_message)
        self.assertIn("uploaded_history", user_message)
        self.assertIn("ppt_private_skill_output_dir", user_message)
        self.assertIn("Template thumbnails", user_message)
        self.assertIn("not final deliverables", user_message)
        self.assertIn("Do not stop after thumbnail.py", user_message)
        self.assertIn("generate a real .pptx", user_message)
        self.assertIn("immediately register it with save_ppt_private_skill_pptx", user_message)
        self.assertIn("optional QA or expert failures", user_message)
        self.assertIn("save_ppt_private_skill_pptx", user_message)
        self.assertIn("concrete blocker", user_message)

    def test_list_session_files_uploaded_falls_back_to_latest_history(self) -> None:
        manager = PptProductManager()
        source_dir = workspace_root() / "inbox" / "ppt_product_manager_tests"
        source_dir.mkdir(parents=True, exist_ok=True)
        template_path = source_dir / "history-template.pptx"
        template_path.write_bytes(b"pptx")
        file_record = build_workspace_file_record(
            template_path,
            description="Uploaded template.",
            source="channel",
            name="history-template.pptx",
            turn=1,
        )
        tool_context = SimpleNamespace(
            state={
                "uploaded": [],
                "input_files": [],
                "uploaded_history": [
                    {"turn": 0, "files": []},
                    {"turn": 1, "files": [file_record]},
                ],
            }
        )

        uploaded_result = manager.list_session_files(section="uploaded", tool_context=tool_context)
        input_result = manager.list_session_files(section="input", tool_context=tool_context)

        self.assertEqual(uploaded_result["uploaded"], [file_record])
        self.assertEqual(input_result["input_files"], [file_record])

    def test_private_skill_source_inputs_are_visible_as_uploaded_files(self) -> None:
        manager = PptProductManager()
        source_dir = workspace_root() / "inbox" / "ppt_product_manager_tests"
        source_dir.mkdir(parents=True, exist_ok=True)
        template_path = source_dir / "confirmed-template.pptx"
        template_path.write_bytes(b"pptx")
        relative_template_path = workspace_relative_path(template_path)
        requirement = manager.prepare_confirmed_requirement(
            task="用这个模板做一个 PPT。",
            inputs=[{"name": "confirmed-template.pptx", "path": relative_template_path}],
            output={"format": "pptx", "route": "xml"},
        )
        state = {
            "uploaded": [],
            "input_files": [],
            "uploaded_history": [],
            "turn_index": 2,
            "step": 3,
        }

        appended = manager._ensure_private_skill_source_files_visible(state, requirement)

        self.assertEqual(len(appended), 1)
        self.assertEqual(appended[0]["path"], relative_template_path)
        self.assertEqual(state["uploaded"][0]["path"], relative_template_path)
        self.assertEqual(state["input_files"][0]["path"], relative_template_path)
        self.assertEqual(state["uploaded"][0]["source"], "confirmed_requirement")

    def test_product_manager_registers_html_page_generation_expert(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={})

        experts = manager.product_expert_agents
        listed = manager.list_ppt_experts(tool_context)

        self.assertIn(PPT_HTML_PAGE_GENERATION_EXPERT_NAME, experts)
        self.assertIn(PPT_DESIGN_STRATEGY_EXPERT_NAME, experts)
        self.assertIn(PPT_SVG_DECK_EXECUTOR_EXPERT_NAME, experts)
        self.assertEqual(experts[PPT_HTML_PAGE_GENERATION_EXPERT_NAME].name, PPT_HTML_PAGE_GENERATION_EXPERT_NAME)
        self.assertEqual(
            {tool.__name__ for tool in experts[PPT_HTML_PAGE_GENERATION_EXPERT_NAME].tools},
            {"save_html_route_pages"},
        )
        self.assertEqual(
            {tool.__name__ for tool in experts[PPT_DESIGN_STRATEGY_EXPERT_NAME].tools},
            {"save_ppt_design_strategy", "save_ppt_svg_execution_plan"},
        )
        self.assertEqual(
            {tool.__name__ for tool in experts[PPT_SVG_DECK_EXECUTOR_EXPERT_NAME].tools},
            {"read_ppt_svg_execution_plan", "save_ppt_svg_page"},
        )
        self.assertIn(PPT_HTML_PAGE_GENERATION_EXPERT_NAME, listed["experts"])
        self.assertIn(PPT_DESIGN_STRATEGY_EXPERT_NAME, listed["experts"])
        self.assertIn(PPT_SVG_DECK_EXECUTOR_EXPERT_NAME, listed["experts"])
        self.assertEqual(tool_context.state["ppt_skill_available_experts"], listed["experts"])

    async def test_invoke_ppt_html_page_generation_expert_uses_ppt_state(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 3 页 PPTX 产品介绍。",
            inputs=[],
            output={"format": "pptx", "slide_count": 3},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        captured: dict[str, object] = {}

        async def _fake_run_html_page_generation_expert(**kwargs):
            captured.update(kwargs)
            return [{"slide_number": 1, "html": "<section><h1>Slide</h1></section>"}]

        tool_context = SimpleNamespace(
            state={
                "ppt_confirmed_requirement": requirement.model_dump(mode="json"),
                "ppt_deck_content_plan": content_plan.model_dump(mode="json"),
            },
            _invocation_context=SimpleNamespace(app_name="creative_claw", user_id="test-user"),
        )

        with patch(
            "src.productions.ppt.ppt_product_manager.ppt_product_manager.run_html_page_generation_expert",
            _fake_run_html_page_generation_expert,
        ):
            result = await manager.invoke_ppt_expert(
                PPT_HTML_PAGE_GENERATION_EXPERT_NAME,
                "Generate editable HTML slide fragments.",
                tool_context,
            )

        self.assertEqual(result["status"], "success")
        self.assertEqual(result["agent_name"], PPT_HTML_PAGE_GENERATION_EXPERT_NAME)
        self.assertEqual(result["current_output"]["html_pages"][0]["slide_number"], 1)
        self.assertIs(
            captured["page_generation_agent"],
            manager.product_expert_agents[PPT_HTML_PAGE_GENERATION_EXPERT_NAME],
        )
        self.assertEqual(captured["template"].template_id, "free_design")
        self.assertEqual(tool_context.state["current_output"]["status"], "success")
        self.assertEqual(
            tool_context.state["ppt_skill_last_expert_result"]["agent_name"],
            PPT_HTML_PAGE_GENERATION_EXPERT_NAME,
        )

    async def test_invoke_ppt_expert_returns_error_payload_for_bad_parameters(self) -> None:
        manager = PptProductManager()
        manager._skill_runtime_expert_agents = {
            "ImageUnderstandingAgent": SimpleNamespace(name="ImageUnderstandingAgent")
        }
        tool_context = SimpleNamespace(
            state=_DictState({"sid": "ppt-expert-error-test", "turn_index": 1, "step": 1}),
            _invocation_context=SimpleNamespace(app_name="creative_claw", user_id="test-user"),
        )

        result = await manager.invoke_ppt_expert(
            "ImageUnderstandingAgent",
            "Visually inspect these rendered PPT slides.",
            tool_context,
        )

        self.assertEqual(result["status"], "error")
        self.assertEqual(result["agent_name"], "ImageUnderstandingAgent")
        self.assertIn("requires structured invoke_agent parameters", result["message"])
        self.assertEqual(tool_context.state["ppt_skill_last_expert_result"], result)

    def test_private_product_ppt_skill_registry_lists_complete_workflow(self) -> None:
        registry = ProductPptSkillRegistry()

        skills = registry.list_skills()
        skill_names = {skill.name for skill in skills}
        content = registry.read_skill("ppt-complete-workflow")
        skill_file_content = registry.read_skill_file("ppt-complete-workflow", "SKILL.md")

        self.assertIn("ppt-complete-workflow", skill_names)
        self.assertIn("PPT Complete Workflow", content)
        self.assertEqual(content, skill_file_content)
        self.assertIn("Built-in HTML route", content)
        self.assertIn("If the user explicitly specifies", content)
        self.assertIn("Do not use local absolute paths", content)

    def test_private_product_ppt_skill_registry_lists_pptx_template_skill(self) -> None:
        registry = ProductPptSkillRegistry()

        skills = registry.list_skills()
        skill_names = {skill.name for skill in skills}
        content = registry.read_skill("pptx")
        editing = registry.read_skill_file("pptx", "editing.md")

        self.assertIn("pptx", skill_names)
        self.assertIn("PPTX Skill", content)
        self.assertIn("uploaded PPTX/POTX template workflows", content)
        self.assertIn("Edit or create from template", content)
        self.assertIn("Template-Based Workflow", editing)

    def test_svg_route_accepts_explicit_system_layout_template_requirement(self) -> None:
        manager = PptProductManager()

        requirement = manager.prepare_confirmed_requirement(
            task="做一个战略咨询汇报 PPT，走 svg route，用 mckinsey 模板。",
            inputs=[],
            output={"format": "pptx", "route": "svg", "template_id": "mckinsey"},
        )

        self.assertEqual(requirement.route, "svg")
        self.assertTrue(requirement.template_requirement.use_template)
        self.assertEqual(requirement.template_requirement.template_source, "system")
        self.assertEqual(requirement.template_requirement.template_id, "mckinsey")

    def test_global_skill_registry_does_not_expose_private_product_ppt_skills(self) -> None:
        global_registry = SkillRegistry()

        skill_names = {skill.name for skill in global_registry.list_skills()}

        self.assertNotIn("ppt-complete-workflow", skill_names)
        self.assertNotIn("product-ppt-skills", skill_names)

    def test_private_ppt_skill_tools_list_and_read_skills(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={})

        listed = manager.list_product_ppt_skills(tool_context)
        read = manager.read_product_ppt_skill("ppt-complete-workflow", tool_context)

        self.assertEqual(listed["status"], "success")
        self.assertGreaterEqual(listed["count"], 1)
        self.assertEqual(read["status"], "success")
        self.assertEqual(read["name"], "ppt-complete-workflow")
        self.assertIn("PPT Complete Workflow", read["content"])
        self.assertEqual(tool_context.state["active_product_ppt_skill"]["name"], "ppt-complete-workflow")

    def test_private_ppt_skill_tool_reads_skill_relative_files(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={})

        result = manager.read_product_ppt_skill_file(
            name="ppt-complete-workflow",
            relative_path="SKILL.md",
            tool_context=tool_context,
        )

        self.assertEqual(result["status"], "success")
        self.assertEqual(result["name"], "ppt-complete-workflow")
        self.assertEqual(result["relative_path"], "SKILL.md")
        self.assertIn("PPT Complete Workflow", result["content"])

    def test_private_ppt_skill_workspace_tools_use_runtime_workspace(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={"sid": "ppt-private-tools-test"})
        path = "generated/ppt-private-tools-test/notes.txt"

        write_result = manager.write_file(path, "alpha\nbeta\n", tool_context)
        read_result = manager.read_file(path, tool_context)
        grep_result = manager.grep(
            "beta",
            path="generated/ppt-private-tools-test",
            output_mode="content",
            tool_context=tool_context,
        )
        glob_result = manager.glob(
            "*.txt",
            path="generated/ppt-private-tools-test",
            tool_context=tool_context,
        )
        exec_result = manager.exec_command("printf ppt-tool-ok", timeout=10, tool_context=tool_context)

        self.assertIn("Successfully wrote", write_result)
        self.assertEqual(read_result, "alpha\nbeta\n")
        self.assertIn("beta", grep_result)
        self.assertIn("notes.txt", glob_result)
        self.assertEqual(exec_result, "ppt-tool-ok")

    def test_save_ppt_private_skill_pptx_registers_final_artifact(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={"sid": "ppt-private-pptx-test", "turn_index": 1, "step": 2})
        pptx_path = resolve_workspace_path("generated/ppt-private-pptx-test/final.pptx")
        pptx_path.parent.mkdir(parents=True, exist_ok=True)
        deck = Presentation()
        deck.slides.add_slide(deck.slide_layouts[0])
        deck.save(pptx_path)

        result = manager.save_ppt_private_skill_pptx(
            pptx_path=workspace_relative_path(pptx_path),
            description="Private PPTX generated by test.",
            tool_context=tool_context,
        )

        self.assertEqual(result["status"], "success")
        self.assertEqual(result["artifact_type"], "pptx")
        self.assertEqual(result["source"], "save_ppt_private_skill_pptx")
        self.assertEqual(result["pptx_path"], workspace_relative_path(pptx_path))
        self.assertEqual(tool_context.state["final_file_paths"], [workspace_relative_path(pptx_path)])
        self.assertEqual(tool_context.state["ppt_private_skill_build"]["pptx_path"], workspace_relative_path(pptx_path))

    def test_record_output_files_normalizes_absolute_final_paths(self) -> None:
        state = {"sid": "ppt-final-path-test", "turn_index": 1, "step": 2}
        pptx_path = resolve_workspace_path("generated/ppt-final-path-test/final.pptx")
        pptx_path.parent.mkdir(parents=True, exist_ok=True)
        deck = Presentation()
        deck.slides.add_slide(deck.slide_layouts[0])
        deck.save(pptx_path)

        PptProductManager._record_output_files(
            state,
            [workspace_relative_path(pptx_path)],
            final_file_paths=[str(pptx_path)],
        )

        self.assertEqual(state["final_file_paths"], [workspace_relative_path(pptx_path)])
        self.assertNotIn(str(workspace_root()), state["final_file_paths"][0])

    def test_recover_unregistered_private_skill_pptx_registers_final_artifact(self) -> None:
        manager = PptProductManager()
        state = {"sid": "ppt-private-recovery-test", "turn_index": 1, "step": 3}
        before_snapshot = _snapshot_workspace_pptx_files()
        source_dir = Path(tempfile.mkdtemp(prefix="work_ai_pptx_recovery_", dir=workspace_root()))
        source_path = source_dir / "recovered_deck.pptx"
        deck = Presentation()
        deck.slides.add_slide(deck.slide_layouts[0])
        deck.save(source_path)

        result = manager._recover_unregistered_private_skill_pptx(
            state,
            skill_name="pptx",
            before_snapshot=before_snapshot,
        )

        self.assertEqual(result["status"], "success")
        self.assertEqual(result["artifact_type"], "pptx")
        self.assertEqual(result["source"], "private_skill_pptx_recovery")
        self.assertEqual(result["recovered_from_path"], workspace_relative_path(source_path))
        self.assertNotEqual(result["pptx_path"], workspace_relative_path(source_path))
        self.assertTrue(
            result["pptx_path"].startswith(
                "generated/ppt-private-recovery-test/turn_1/ppt_private_skill_step_3/"
            )
        )
        self.assertTrue(resolve_workspace_path(result["pptx_path"]).exists())
        self.assertEqual(state["final_file_paths"], [result["pptx_path"]])
        self.assertEqual(state["ppt_private_skill_build"]["pptx_path"], result["pptx_path"])

    def test_private_skill_delivery_accepts_svg_pptx_export(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="用 pptx skill 做 2 页 PPTX。",
            inputs=[],
            output={"format": "pptx", "route": "svg", "slide_count": 2},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        pptx_path = "generated/session/turn_1/ppt_svg_route/deck.pptx"
        file_record = build_workspace_file_record(
            pptx_path,
            description="PPT product SVG route PPTX artifact.",
            source="ppt_product_manager",
        )

        private_build = manager._resolve_private_skill_build_from_state(
            {
                "ppt_svg_pptx_export": {
                    "status": "success",
                    "message": "Exported PPT SVG pages.",
                    "pptx_path": pptx_path,
                    "conversion_report": {"ok": True},
                    "output_files": [file_record],
                }
            },
            skill_name="pptx",
        )
        result = manager._build_private_skill_delivery_result(
            requirement=requirement,
            content_plan=content_plan,
            system_selection={
                "system_type": "private_skill",
                "skill_name": "pptx",
                "output_format": "pptx",
            },
            private_build=private_build,
        )

        self.assertEqual(private_build["artifact_type"], "pptx")
        self.assertEqual(private_build["source"], "export_ppt_svg_to_pptx")
        self.assertEqual(result.status, "success")
        self.assertEqual(result.delivery_manifest.final_pptx, pptx_path)
        self.assertEqual(result.delivery_manifest.intermediate_artifacts, [])
        self.assertIn("editable PPTX", " ".join(result.warnings))

    def test_private_skill_delivery_phase_records_typed_result(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="用 private skill 做 2 页 HTML PPT。",
            inputs=[],
            output={"format": "pptx", "route": "html", "slide_count": 2},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        output_path = "generated/session/turn_1/ppt_private_skill_step_1/index.html"
        file_record = build_workspace_file_record(
            output_path,
            description="Private PPT skill HTML deck artifact.",
            source="ppt_product_manager",
        )
        private_build = {
            "skill_name": "ppt-complete-workflow",
            "output_path": output_path,
            "artifact_type": "html",
            "output_format": "html",
            "output_files": [file_record],
            "source": "save_ppt_private_skill_html",
        }
        tool_context = SimpleNamespace(state={})

        result = manager._finalize_private_skill_delivery_phase(
            requirement=requirement,
            content_plan=content_plan,
            system_selection={
                "system_type": "private_skill",
                "skill_name": "ppt-complete-workflow",
                "output_format": "html",
            },
            private_build=private_build,
            tool_context=tool_context,
        )

        self.assertIsInstance(result, PptPrivateSkillDeliveryResult)
        self.assertEqual(result.product_result.status, "success")
        self.assertEqual(result.delivery_manifest.intermediate_artifacts, [output_path])
        self.assertEqual(result.output_files, [file_record])
        self.assertEqual(
            tool_context.state["ppt_private_skill_delivery_result"]["private_build"]["output_path"],
            output_path,
        )

    def test_system_selection_confirmation_uses_table_without_narrow_system_selection_column(self) -> None:
        summary = PptProductManager._format_system_selection_confirmation(
            {
                "system_type": "private_skill",
                "skill_name": "pptx",
                "output_format": "pptx",
                "reason": "使用 PPTX 模板 skill 生成可编辑 PPTX。",
            }
        )

        self.assertIn("### 系统选择", summary)
        self.assertIn("| 项目 | 当前值 |", summary)
        self.assertIn("| 制作系统 | 私有 PPT skill `pptx` |", summary)
        self.assertIn("| 输出方式 | pptx |", summary)
        self.assertIn("| 选择理由 | 使用 PPTX 模板 skill 生成可编辑 PPTX。 |", summary)
        self.assertNotIn("| 系统选择 |", summary)

    def test_route_registry_registers_all_routes(self) -> None:
        manager = PptProductManager()

        routes = manager.list_registered_routes()

        self.assertEqual(set(routes), {"html", "svg", "xml"})
        self.assertTrue(routes["html"]["implemented"])
        self.assertTrue(routes["svg"]["implemented"])
        self.assertFalse(routes["xml"]["implemented"])

    def test_content_planning_agent_exposes_material_tools(self) -> None:
        manager = PptProductManager()

        agent = manager.content_planner.build_agent()

        self.assertIsInstance(agent, LlmAgent)
        self.assertEqual(agent.name, "PptContentPlanningAgent")
        self.assertEqual(agent.output_key, "ppt_content_planning_agent_message")
        self.assertEqual(agent.include_contents, "none")
        self.assertEqual(
            {tool.__name__ for tool in agent.tools},
            {"read_ppt_markdown_sources", "save_ppt_deck_content_plan_markdown"},
        )
        self.assertIn("do not force cover, toc, chapter_start", agent.instruction)
        self.assertIn("template requirements only", agent.instruction)
        self.assertIn("Do not overuse `content`", agent.instruction)
        self.assertIn("`comparison` for tradeoffs", agent.instruction)

    async def test_content_planning_agenttool_main_path_saves_plan_without_fallback(self) -> None:
        manager = PptProductManager()
        source_path = _write_markdown_source(
            "agenttool_planning_brief.md",
            "# Planning Brief\n\nActivation rose after onboarding.\n",
        )
        source_understanding = SourceUnderstanding(
            document_type="markdown",
            markdown_sources=[
                {
                    "name": "agenttool_planning_brief.md",
                    "source_path": source_path,
                    "method": "test",
                    "output_path": source_path,
                }
            ],
        )
        requirement = manager.prepare_confirmed_requirement(
            task="基于材料生成 5 页 PPTX。",
            inputs=[{"name": "agenttool_planning_brief.md", "path": source_path}],
            output={"format": "pptx"},
            source_understanding=source_understanding,
        )
        markdown_plan = """# Deck: Planning Brief
Audience: Internal team
Language: en
SlideCount: 5
Narrative: Explain activation changes.

## Slide 1 | cover | Planning Brief
Purpose: Introduce the planning brief.
Takeaway: Activation rose after onboarding.
Content:
- Audience: Growth team
Visual:
- placeholder | role=hero | description=clean title area

## Slide 2 | toc | Agenda
Purpose: Preview the deck.
Takeaway: The deck covers evidence and next steps.
Content:
- Activation
- Evidence
- Next steps
Visual:
- placeholder | role=list | description=agenda list

## Slide 3 | chapter_start | Activation
Purpose: Start the activation chapter.
Takeaway: Activation rose after onboarding.
Content:
- Activation rose after onboarding.
Visual:
- search | role=reference | query=activation onboarding chart | description=visual reference for activation onboarding

## Slide 4 | chapter_content | Evidence
Purpose: Explain the evidence.
Takeaway: Guided onboarding improved activation.
Content:
- Activation rose after onboarding.
- Enterprise teams need proof.
Visual:
- ai | role=supporting_visual | description=friendly product onboarding illustration

## Slide 5 | ending | Next Steps
Purpose: Close with next steps.
Takeaway: Use the activation proof in the story.
Content:
- Review the evidence
- Prepare the launch story
Visual:
- placeholder | role=summary | description=closing icon area
"""
        fake_llm = _PptContentPlannerToolCallingFakeLlm(markdown_plan=markdown_plan)
        session_service = InMemorySessionService()
        artifact_service = InMemoryArtifactService()

        @node(name="PptContentPlannerHarnessNode", rerun_on_resume=True)
        async def content_planner_harness(ctx: Context, node_input: str) -> dict:
            plan = await manager.content_planner.build_plan_with_agent(
                requirement,
                tool_context=ctx,
                app_name="creative_claw",
                artifact_service=artifact_service,
            )
            return plan.model_dump(mode="json")

        workflow = Workflow(
            name="PptContentPlannerHarness",
            edges=[("START", content_planner_harness)],
        )
        runner = Runner(
            node=workflow,
            session_service=session_service,
            artifact_service=artifact_service,
        )
        user_id = "user-ppt-content-planner"
        session_id = "session-ppt-content-planner"

        try:
            await session_service.create_session(
                app_name=workflow.name,
                user_id=user_id,
                session_id=session_id,
                state={
                    "sid": "ppt-content-planner-agenttool-test",
                    "turn_index": 1,
                    "step": 1,
                },
            )
            with patch(
                "src.productions.ppt.planning.content_planner.build_llm",
                return_value=fake_llm,
            ):
                async for _ in runner.run_async(
                    user_id=user_id,
                    session_id=session_id,
                    new_message=Content(role="user", parts=[Part(text="Plan deck")]),
                ):
                    pass
            session = await session_service.get_session(
                app_name=workflow.name,
                user_id=user_id,
                session_id=session_id,
            )
        finally:
            await runner.close()

        self.assertIsNotNone(session)
        self.assertEqual(session.state["ppt_content_planning_output"]["source"], "llm_agent")
        self.assertEqual(session.state["ppt_content_planning_output"]["status"], "success")
        self.assertEqual(session.state["ppt_deck_content_plan"]["title"], "Planning Brief")
        self.assertEqual(len(session.state["ppt_deck_content_plan"]["pages"]), 5)
        self.assertEqual(session.state["ppt_deck_content_plan"]["pages"][2]["asset_source_preference"], "search")
        self.assertEqual(session.state["ppt_deck_content_plan"]["pages"][3]["asset_source_preference"], "ai")
        self.assertIn("Activation rose", session.state["ppt_markdown_source_texts"][0]["text"])
        self.assertIn("Slide 1 | cover | Planning Brief", session.state["ppt_deck_content_plan_markdown"])
        self.assertGreaterEqual(len(fake_llm.requests), 3)
        first_request_tools = _function_declaration_names(fake_llm.requests[0])
        self.assertIn("read_ppt_markdown_sources", first_request_tools)
        self.assertIn("save_ppt_deck_content_plan_markdown", first_request_tools)

    async def test_content_planning_phase_persists_typed_result_without_resolving_assets(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="给小学生做一个 AI 科普 PPTX。",
            inputs=[],
            output={"format": "pptx", "slide_count": 5},
        )
        tool_context = SimpleNamespace(state={"sid": "ppt-content-phase-test", "turn_index": 1, "step": 1})

        def _content_plan_builder(_requirement):
            plan = DeckContentPlan(
                title="AI for Kids",
                core_narrative="Explain AI through concrete classroom examples.",
                pages=[_page(index, "content") for index in range(1, 6)],
            )
            plan.pages[2].assets = [
                DeckPageAsset(
                    asset_id="slide_03_pending_visual",
                    source_kind="image_generation",
                    status="pending",
                    description="Classroom AI illustration.",
                    prompt="Classroom AI illustration.",
                )
            ]
            return plan

        result = await manager._build_deck_content_plan_phase(
            requirement,
            tool_context=tool_context,
            app_name="creative_claw",
            artifact_service=None,
            expert_agents={},
            content_plan_builder=_content_plan_builder,
        )

        self.assertIsInstance(result, PptContentPlanningResult)
        self.assertEqual(result.content_plan.title, "AI for Kids")
        self.assertEqual(result.planning_output["source"], "injected")
        self.assertEqual(tool_context.state["ppt_deck_content_plan"]["title"], "AI for Kids")
        self.assertEqual(tool_context.state["ppt_content_planning_output"]["source"], "injected")
        self.assertNotIn("ppt_resolved_asset_manifest", tool_context.state)
        self.assertEqual(result.content_plan.pages[2].assets[0].status, "pending")

    def test_requirement_analysis_agent_saves_confirmed_requirement_json(self) -> None:
        manager = PptProductManager()

        agent = manager.build_requirement_analysis_agent()

        self.assertIsInstance(agent, LlmAgent)
        self.assertEqual(agent.name, "PptRequirementAnalysisAgent")
        self.assertEqual(agent.output_key, "ppt_requirement_analysis_agent_message")
        self.assertEqual(agent.include_contents, "none")
        self.assertEqual(
            {tool.__name__ for tool in agent.tools},
            {"save_ppt_confirmed_requirement_json"},
        )
        self.assertIn("ConfirmedRequirement JSON", agent.instruction)
        self.assertIn("multiple PPT systems", agent.instruction)
        self.assertIn("system-selection agent", agent.instruction)
        self.assertIn("source_inputs include PPTX/PPTM/POTX/POTM", agent.instruction)
        self.assertIn("Do not infer routes from keyword matching", agent.instruction)
        self.assertIn("受众为", agent.instruction)

    async def test_requirement_analysis_agenttool_main_path_saves_requirement(self) -> None:
        manager = PptProductManager()
        task = "做一个 3 页 PPTX，用于产品发布，受众为管理层。"
        output = {"format": "pptx", "slide_count": 3}
        expected_requirement = manager.prepare_confirmed_requirement(
            task=task,
            inputs=[],
            output=output,
        )
        requirement_payload = expected_requirement.model_dump(mode="json")
        requirement_payload["audience"] = "管理层"
        fake_llm = _PptProductManagerToolCallingFakeLlm(
            function_calls=[
                FunctionCall(
                    name="save_ppt_confirmed_requirement_json",
                    args={"requirement_json": requirement_payload},
                )
            ],
            final_text="PptRequirementAnalysisAgent saved ConfirmedRequirement.",
        )
        session_service = InMemorySessionService()
        artifact_service = InMemoryArtifactService()

        @node(name="PptRequirementAnalysisAgentToolHarnessNode", rerun_on_resume=True)
        async def requirement_harness(ctx: Context, node_input: str) -> dict:
            requirement = await manager.prepare_confirmed_requirement_with_agent(
                task=task,
                inputs=[],
                output=output,
                source_understanding=None,
                tool_context=ctx,
                app_name="creative_claw",
                artifact_service=artifact_service,
            )
            return requirement.model_dump(mode="json")

        workflow = Workflow(
            name="PptRequirementAnalysisAgentToolHarness",
            edges=[("START", requirement_harness)],
        )
        runner = Runner(
            node=workflow,
            session_service=session_service,
            artifact_service=artifact_service,
        )
        user_id = "user-ppt-requirement-agenttool"
        session_id = "session-ppt-requirement-agenttool"

        try:
            await session_service.create_session(
                app_name=workflow.name,
                user_id=user_id,
                session_id=session_id,
                state={
                    "sid": "ppt-requirement-agenttool-test",
                    "turn_index": 1,
                    "step": 1,
                },
            )
            with patch(
                "src.productions.ppt.ppt_product_manager.ppt_product_manager.build_llm",
                return_value=fake_llm,
            ):
                async for _ in runner.run_async(
                    user_id=user_id,
                    session_id=session_id,
                    new_message=Content(role="user", parts=[Part(text="Analyze PPT requirement")]),
                ):
                    pass
            session = await session_service.get_session(
                app_name=workflow.name,
                user_id=user_id,
                session_id=session_id,
            )
        finally:
            await runner.close()

        self.assertIsNotNone(session)
        self.assertEqual(
            session.state["ppt_requirement_analysis_output"]["source"],
            "llm_agent",
        )
        self.assertEqual(session.state[PPT_CONFIRMED_REQUIREMENT_STATE_KEY]["audience"], "管理层")
        self.assertEqual(session.state[PPT_CONFIRMED_REQUIREMENT_STATE_KEY]["output_format"], "pptx")
        self.assertIn(
            "PptRequirementAnalysisAgent saved ConfirmedRequirement.",
            session.state[PPT_REQUIREMENT_ANALYSIS_AGENT_MESSAGE_KEY],
        )
        first_request_tools = _function_declaration_names(fake_llm.requests[0])
        self.assertIn("save_ppt_confirmed_requirement_json", first_request_tools)
        self.assertGreaterEqual(len(fake_llm.requests), 2)
        self.assertNotEqual(
            session.state["ppt_requirement_analysis_output"]["source"],
            "deterministic_fallback",
        )

    async def test_requirement_analysis_phase_persists_typed_result(self) -> None:
        source_path = _write_markdown_source("initial_requirement_phase_source.md", "# Source")
        manager = PptProductManager()
        source_inputs = manager._normalize_source_inputs([source_path])
        tool_context = SimpleNamespace(
            state={"sid": "ppt-requirement-analysis-phase-test", "turn_index": 1, "step": 1}
        )

        result = await manager._prepare_initial_requirement_phase(
            task="基于素材做一个 4 页项目汇报 PPT。",
            raw_inputs=[source_path],
            output={"format": "pptx", "slide_count": 4},
            source_understanding=SourceUnderstanding(document_type="markdown"),
            source_inputs=source_inputs,
            tool_context=tool_context,
            app_name="creative_claw",
            artifact_service=None,
        )

        self.assertIsInstance(result, PptRequirementAnalysisResult)
        self.assertEqual(result.confirmed_requirement.slide_count_policy.target, 4)
        self.assertEqual(result.analysis_output["source"], "deterministic_fallback")
        self.assertEqual(
            tool_context.state[PPT_REQUIREMENT_ANALYSIS_RESULT_STATE_KEY]["confirmed_requirement"]["source_inputs"][0]["path"],
            source_path,
        )
        self.assertEqual(
            tool_context.state[PPT_REQUIREMENT_ANALYSIS_RESULT_STATE_KEY]["analysis_output"]["source"],
            "deterministic_fallback",
        )
        self.assertEqual(tool_context.state[PPT_CONFIRMED_REQUIREMENT_STATE_KEY]["source_inputs"][0]["path"], source_path)

    async def test_requirement_analysis_phase_uses_adk_workflow_context(self) -> None:
        manager = PptProductManager()
        task = "做一个 3 页 PPTX，用于产品发布，受众为管理层。"
        output = {"format": "pptx", "slide_count": 3}
        expected_requirement = manager.prepare_confirmed_requirement(
            task=task,
            inputs=[],
            output=output,
            source_understanding=SourceUnderstanding(document_type="none"),
        )
        requirement_payload = expected_requirement.model_dump(mode="json")
        requirement_payload["audience"] = "管理层"
        fake_llm = _PptProductManagerToolCallingFakeLlm(
            function_calls=[
                FunctionCall(
                    name="save_ppt_confirmed_requirement_json",
                    args={"requirement_json": requirement_payload},
                )
            ],
            final_text="PptRequirementAnalysisAgent saved initial requirement.",
        )

        async def _run_requirement_harness():
            @node(name="PptRequirementAnalysisWorkflowHarnessNode", rerun_on_resume=True)
            async def requirement_harness(ctx: Context, node_input: str) -> dict[str, Any]:
                result = await manager._prepare_initial_requirement_phase(
                    task=task,
                    raw_inputs=[],
                    output=output,
                    source_understanding=SourceUnderstanding(document_type="none"),
                    source_inputs=None,
                    tool_context=ctx,
                    app_name="creative_claw",
                    artifact_service=artifact_service,
                )
                return result.model_dump(mode="json")

            workflow = Workflow(
                name="PptRequirementAnalysisWorkflowHarness",
                edges=[("START", requirement_harness)],
            )
            session_service = InMemorySessionService()
            runner = Runner(
                node=workflow,
                session_service=session_service,
                artifact_service=artifact_service,
            )
            user_id = "user-ppt-requirement-analysis-workflow"
            session_id = "session-ppt-requirement-analysis-workflow"
            try:
                await session_service.create_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                    state={
                        "sid": "ppt-requirement-analysis-workflow-test",
                        "turn_index": 1,
                        "step": 1,
                    },
                )
                with patch(
                    "src.productions.ppt.ppt_product_manager.ppt_product_manager.build_llm",
                    return_value=fake_llm,
                ):
                    async for _ in runner.run_async(
                        user_id=user_id,
                        session_id=session_id,
                        new_message=Content(role="user", parts=[Part(text="Analyze initial PPT requirement")]),
                    ):
                        pass
                return await session_service.get_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                )
            finally:
                await runner.close()

        artifact_service = InMemoryArtifactService()
        session = await _run_requirement_harness()

        self.assertIsNotNone(session)
        self.assertEqual(session.state[PPT_REQUIREMENT_ANALYSIS_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_REQUIREMENT_ANALYSIS_WORKFLOW_OUTPUT_KEY]["analysis_source"], "llm_agent")
        self.assertEqual(session.state[PPT_REQUIREMENT_ANALYSIS_RESULT_STATE_KEY]["analysis_output"]["source"], "llm_agent")
        self.assertEqual(session.state[PPT_CONFIRMED_REQUIREMENT_STATE_KEY]["audience"], "管理层")
        first_request_tools = _function_declaration_names(fake_llm.requests[0])
        self.assertIn("save_ppt_confirmed_requirement_json", first_request_tools)
        self.assertGreaterEqual(len(fake_llm.requests), 2)

    async def test_initial_request_uses_parent_adk_workflow_context(self) -> None:
        manager = PptProductManager()
        task = "做一个 3 页 PPTX，用于产品发布，受众为管理层。"
        output = {"format": "pptx", "slide_count": 3}
        expected_requirement = manager.prepare_confirmed_requirement(
            task=task,
            inputs=[],
            output=output,
            source_understanding=SourceUnderstanding(document_type="none"),
        )
        requirement_payload = expected_requirement.model_dump(mode="json")
        requirement_payload["audience"] = "管理层"
        fake_llm = _PptProductManagerToolCallingFakeLlm(
            function_calls=[
                FunctionCall(
                    name="save_ppt_confirmed_requirement_json",
                    args={"requirement_json": requirement_payload},
                )
            ],
            final_text="PptInitialRequestWorkflow saved initial requirement.",
        )

        def _system_selection_builder(**_kwargs):
            return {
                "system_type": "built_in_route",
                "route": "html",
                "skill_name": "",
                "output_format": "pptx",
                "reason": "Use the built-in HTML route for the initial request Workflow test.",
            }

        async def _run_initial_request_harness():
            @node(name="PptInitialRequestWorkflowHarnessNode", rerun_on_resume=True)
            async def initial_request_harness(ctx: Context, node_input: str) -> dict[str, Any]:
                with patch.object(
                    manager,
                    "_start_interactive_product_request_direct",
                    side_effect=AssertionError("parent workflow must own initial request"),
                ):
                    return await manager.run_product_request(
                        task=node_input,
                        inputs=[],
                        output=output,
                        tool_context=ctx,
                        system_selection_builder=_system_selection_builder,
                    )

            workflow = Workflow(
                name="PptInitialRequestWorkflowHarness",
                edges=[("START", initial_request_harness)],
            )
            session_service = InMemorySessionService()
            runner = Runner(
                node=workflow,
                session_service=session_service,
                artifact_service=artifact_service,
            )
            user_id = "user-ppt-initial-request-workflow"
            session_id = "session-ppt-initial-request-workflow"
            try:
                await session_service.create_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                    state={
                        "sid": "ppt-initial-request-workflow-test",
                        "turn_index": 1,
                        "step": 1,
                    },
                )
                with patch(
                    "src.productions.ppt.ppt_product_manager.ppt_product_manager.build_llm",
                    return_value=fake_llm,
                ):
                    async for _ in runner.run_async(
                        user_id=user_id,
                        session_id=session_id,
                        new_message=Content(role="user", parts=[Part(text=task)]),
                    ):
                        pass
                return await session_service.get_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                )
            finally:
                await runner.close()

        artifact_service = InMemoryArtifactService()
        session = await _run_initial_request_harness()

        self.assertIsNotNone(session)
        self.assertEqual(session.state["ppt_product_result"]["status"], "awaiting_requirement_confirmation")
        self.assertEqual(session.state["ppt_workflow_state"]["stage"], "awaiting_requirement_confirmation")
        self.assertEqual(session.state[PPT_INITIAL_REQUEST_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(
            session.state[PPT_INITIAL_REQUEST_WORKFLOW_OUTPUT_KEY]["branch"],
            "requirement_confirmation",
        )
        self.assertEqual(session.state[PPT_REQUIREMENT_ANALYSIS_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_REQUIREMENT_ANALYSIS_WORKFLOW_OUTPUT_KEY]["analysis_source"], "llm_agent")
        self.assertEqual(session.state[PPT_SYSTEM_SELECTION_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_SYSTEM_SELECTION_WORKFLOW_OUTPUT_KEY]["selection_source"], "injected")
        self.assertEqual(session.state[PPT_CONFIRMED_REQUIREMENT_STATE_KEY]["audience"], "管理层")
        self.assertNotIn(PPT_SOURCE_PREPARATION_WORKFLOW_OUTPUT_KEY, session.state)
        self.assertNotIn(PPT_CONTENT_PLANNING_WORKFLOW_OUTPUT_KEY, session.state)
        self.assertNotIn("final_file_paths", session.state)

    async def test_auto_confirm_uses_parent_adk_workflow_context(self) -> None:
        manager = PptProductManager()
        task = "做一个 3 页 PPTX，用于产品发布，受众为管理层。"
        output = {"format": "pptx", "slide_count": 3, "auto_confirm": True}
        expected_requirement = manager.prepare_confirmed_requirement(
            task=task,
            inputs=[],
            output=output,
            source_understanding=SourceUnderstanding(document_type="none"),
        )
        requirement_payload = expected_requirement.model_dump(mode="json")
        requirement_payload["audience"] = "管理层"
        fake_llm = _PptProductManagerToolCallingFakeLlm(
            function_calls=[
                FunctionCall(
                    name="save_ppt_confirmed_requirement_json",
                    args={"requirement_json": requirement_payload},
                )
            ],
            final_text="PptAutoConfirmWorkflow saved initial requirement.",
        )
        route_calls: list[dict[str, Any]] = []
        final_pptx = "generated/ppt-auto-confirm-workflow-test/turn_1/final.pptx"

        def _system_selection_builder(**_kwargs):
            return {
                "system_type": "built_in_route",
                "route": "html",
                "skill_name": "",
                "output_format": "pptx",
                "reason": "Use the built-in HTML route for the auto-confirm Workflow test.",
            }

        def _content_plan_builder(_requirement):
            return DeckContentPlan(
                title="Auto Confirm Launch",
                core_narrative="Present the product launch clearly.",
                pages=[_page(1, "cover"), _page(2, "content"), _page(3, "ending")],
            )

        async def _asset_resolver(asset, _page, _requirement):
            return {
                "asset_id": asset.asset_id,
                "status": "ready",
                "path": "",
                "provider": "test_resolver",
            }

        async def _fake_dispatch_ppt_route(**kwargs):
            route_calls.append(kwargs)
            return HtmlRouteBuildPackage(
                template=HtmlTemplatePackage(template_id="auto-confirm-test", label="Auto Confirm Test"),
                html_deck_path="generated/ppt-auto-confirm-workflow-test/turn_1/deck.html",
                preview_paths=["generated/ppt-auto-confirm-workflow-test/turn_1/preview.png"],
                pptx_path=final_pptx,
                quality_report_path="generated/ppt-auto-confirm-workflow-test/turn_1/quality.json",
                build_log_path="generated/ppt-auto-confirm-workflow-test/turn_1/build.log",
                warnings=[],
            )

        async def _run_auto_confirm_harness():
            @node(name="PptAutoConfirmWorkflowHarnessNode", rerun_on_resume=True)
            async def auto_confirm_harness(ctx: Context, node_input: str) -> dict[str, Any]:
                with patch.object(
                    manager,
                    "_run_auto_confirm_product_request_direct",
                    side_effect=AssertionError("parent workflow must own auto-confirm delivery"),
                ):
                    return await manager.run_product_request(
                        task=node_input,
                        inputs=[],
                        output=output,
                        tool_context=ctx,
                        content_plan_builder=_content_plan_builder,
                        asset_resolver=_asset_resolver,
                        system_selection_builder=_system_selection_builder,
                    )

            workflow = Workflow(
                name="PptAutoConfirmWorkflowHarness",
                edges=[("START", auto_confirm_harness)],
            )
            session_service = InMemorySessionService()
            runner = Runner(
                node=workflow,
                session_service=session_service,
                artifact_service=artifact_service,
            )
            user_id = "user-ppt-auto-confirm-workflow"
            session_id = "session-ppt-auto-confirm-workflow"
            try:
                await session_service.create_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                    state={
                        "sid": "ppt-auto-confirm-workflow-test",
                        "turn_index": 1,
                        "step": 1,
                    },
                )
                with (
                    patch(
                        "src.productions.ppt.ppt_product_manager.ppt_product_manager.build_llm",
                        return_value=fake_llm,
                    ),
                    patch.object(manager, "_dispatch_ppt_route", _fake_dispatch_ppt_route),
                ):
                    async for _ in runner.run_async(
                        user_id=user_id,
                        session_id=session_id,
                        new_message=Content(role="user", parts=[Part(text=task)]),
                    ):
                        pass
                return await session_service.get_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                )
            finally:
                await runner.close()

        artifact_service = InMemoryArtifactService()
        session = await _run_auto_confirm_harness()

        self.assertIsNotNone(session)
        self.assertEqual(session.state["ppt_product_result"]["status"], "success")
        self.assertEqual(session.state["ppt_product_result"]["phase"], "html_route_delivery")
        self.assertEqual(session.state[PPT_AUTO_CONFIRM_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_AUTO_CONFIRM_WORKFLOW_OUTPUT_KEY]["branch"], "built_in_route")
        self.assertEqual(session.state[PPT_SOURCE_PREPARATION_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_REQUIREMENT_ANALYSIS_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_SYSTEM_SELECTION_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_CONTENT_PLANNING_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_ASSET_RESOLUTION_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_ROUTE_EXECUTION_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_FINAL_DELIVERY_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_CONFIRMED_REQUIREMENT_STATE_KEY]["audience"], "管理层")
        self.assertEqual(session.state["final_file_paths"], [final_pptx])
        self.assertEqual(len(route_calls), 1)
        self.assertNotIn("ppt_workflow_state", session.state)

    async def test_auto_confirm_workflow_reentry_reuses_side_effect_phases(self) -> None:
        manager = PptProductManager()
        task = "针对远程报告做一个 3 页 PPTX，用于产品复盘，受众为管理层。"
        source_url = "https://example.com/ppt-reentry-source.pdf"
        raw_inputs = [{"name": "ppt-reentry-source.pdf", "url": source_url}]
        output = {"format": "pptx", "slide_count": 3, "auto_confirm": True}
        expected_requirement = manager.prepare_confirmed_requirement(
            task=task,
            inputs=raw_inputs,
            output=output,
            source_understanding=SourceUnderstanding(document_type="markdown"),
        )
        requirement_payload = expected_requirement.model_dump(mode="json")
        requirement_payload["audience"] = "管理层"
        source_image_path = _write_test_image("auto_confirm_reentry_asset.png")
        urlopen_calls = 0
        converter_calls = 0
        asset_resolver_calls = 0
        route_calls: list[Path] = []
        fake_llms: list[_PptProductManagerToolCallingFakeLlm] = []

        def _fake_urlopen(*_args, **_kwargs):
            nonlocal urlopen_calls
            urlopen_calls += 1
            return _FakeRemoteResponse(
                b"%PDF-1.4\nremote pdf fixture\n",
                {"content-type": "application/pdf", "content-length": "28"},
            )

        async def _counting_source_converter(source_input, parameters: dict) -> dict:
            nonlocal converter_calls
            converter_calls += 1
            return await _fake_source_converter(source_input, parameters)

        def _build_requirement_fake_llm(*_args, **_kwargs):
            fake_llm = _PptProductManagerToolCallingFakeLlm(
                function_calls=[
                    FunctionCall(
                        name="save_ppt_confirmed_requirement_json",
                        args={"requirement_json": requirement_payload},
                    )
                ],
                final_text="PptAutoConfirmWorkflow re-entry saved requirement.",
            )
            fake_llms.append(fake_llm)
            return fake_llm

        def _system_selection_builder(**_kwargs):
            return {
                "system_type": "built_in_route",
                "route": "html",
                "skill_name": "",
                "output_format": "pptx",
                "reason": "Use the built-in HTML route for the auto-confirm re-entry test.",
            }

        def _content_plan_builder(_requirement):
            plan = DeckContentPlan(
                title="Re-entry Launch",
                core_narrative="Present the source material clearly.",
                pages=[_page(1, "cover"), _page(2, "content"), _page(3, "ending")],
            )
            plan.pages[1].assets = [
                DeckPageAsset(
                    asset_id="slide_02_reentry_visual",
                    source_kind="image_generation",
                    status="pending",
                    description="A management review visual.",
                    prompt="A management review visual.",
                )
            ]
            return plan

        async def _asset_resolver(asset, _page, _requirement):
            nonlocal asset_resolver_calls
            asset_resolver_calls += 1
            return {
                "asset_id": asset.asset_id,
                "status": "ready",
                "path": source_image_path,
                "provider": "test_resolver",
            }

        async def _fake_dispatch_ppt_route(**kwargs):
            output_dir = Path(kwargs["output_dir"])
            route_calls.append(output_dir)
            final_pptx = output_dir / "final.pptx"
            final_pptx.write_bytes(b"fake-pptx")
            html_path = output_dir / "deck.html"
            html_path.write_text("<html></html>", encoding="utf-8")
            preview_path = output_dir / "preview.png"
            preview_path.write_bytes(b"fake-preview")
            quality_path = output_dir / "quality.json"
            quality_path.write_text("{}", encoding="utf-8")
            build_log_path = output_dir / "build.log"
            build_log_path.write_text("ok", encoding="utf-8")
            return HtmlRouteBuildPackage(
                template=HtmlTemplatePackage(template_id="auto-confirm-reentry-test", label="Auto Confirm Re-entry"),
                html_deck_path=workspace_relative_path(html_path),
                preview_paths=[workspace_relative_path(preview_path)],
                pptx_path=workspace_relative_path(final_pptx),
                quality_report_path=workspace_relative_path(quality_path),
                build_log_path=workspace_relative_path(build_log_path),
                warnings=[],
            )

        async def _run_auto_confirm_reentry_harness():
            @node(name="PptAutoConfirmReentryWorkflowHarnessNode", rerun_on_resume=True)
            async def auto_confirm_reentry_harness(ctx: Context, node_input: str) -> dict[str, Any]:
                with patch.object(
                    manager,
                    "_run_auto_confirm_product_request_direct",
                    side_effect=AssertionError("parent workflow must own auto-confirm delivery"),
                ):
                    return await manager.run_product_request(
                        task=node_input,
                        inputs=raw_inputs,
                        output=output,
                        tool_context=ctx,
                        source_converter=_counting_source_converter,
                        content_plan_builder=_content_plan_builder,
                        asset_resolver=_asset_resolver,
                        system_selection_builder=_system_selection_builder,
                    )

            workflow = Workflow(
                name="PptAutoConfirmReentryWorkflowHarness",
                edges=[("START", auto_confirm_reentry_harness)],
            )
            session_service = InMemorySessionService()
            runner = Runner(
                node=workflow,
                session_service=session_service,
                artifact_service=artifact_service,
            )
            user_id = "user-ppt-auto-confirm-reentry-workflow"
            session_id = "session-ppt-auto-confirm-reentry-workflow"
            try:
                await session_service.create_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                    state={
                        "sid": "ppt-auto-confirm-reentry-workflow-test",
                        "turn_index": 1,
                        "step": 1,
                    },
                )
                with (
                    patch(
                        "src.productions.ppt.ppt_product_manager.ppt_product_manager.urlopen",
                        side_effect=_fake_urlopen,
                    ),
                    patch(
                        "src.productions.ppt.ppt_product_manager.ppt_product_manager.build_llm",
                        side_effect=_build_requirement_fake_llm,
                    ),
                    patch.object(manager, "_dispatch_ppt_route", _fake_dispatch_ppt_route),
                ):
                    for _ in range(2):
                        async for _event in runner.run_async(
                            user_id=user_id,
                            session_id=session_id,
                            new_message=Content(role="user", parts=[Part(text=task)]),
                        ):
                            pass
                return await session_service.get_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                )
            finally:
                await runner.close()

        artifact_service = InMemoryArtifactService()
        session = await _run_auto_confirm_reentry_harness()

        self.assertIsNotNone(session)
        self.assertEqual(session.state["ppt_product_result"]["status"], "success")
        self.assertEqual(urlopen_calls, 1)
        self.assertEqual(converter_calls, 1)
        self.assertEqual(asset_resolver_calls, 1)
        self.assertEqual(len(route_calls), 1)
        self.assertEqual(len(fake_llms), 2)
        self.assertTrue(session.state[PPT_SOURCE_PREPARATION_WORKFLOW_OUTPUT_KEY]["reused_existing_preparation"])
        self.assertTrue(session.state["ppt_source_preparation_result"]["reused_existing_preparation"])
        self.assertTrue(session.state[PPT_ASSET_RESOLUTION_WORKFLOW_OUTPUT_KEY]["reused_existing_resolution"])
        self.assertTrue(session.state["ppt_asset_resolution_result"]["reused_existing_resolution"])
        self.assertTrue(session.state[PPT_ROUTE_EXECUTION_WORKFLOW_OUTPUT_KEY]["reused_existing_build"])
        self.assertEqual(session.state["final_file_paths"], [session.state["ppt_route_build"]["pptx_path"]])
        output_files = session.state["ppt_product_result"]["output_files"]
        self.assertEqual(session.state["new_files"], output_files)
        self.assertEqual(session.state["generated"], output_files)
        self.assertEqual(session.state["files_history"], [output_files])

    async def test_auto_confirm_private_skill_branch_uses_parent_adk_workflow_context(self) -> None:
        manager = PptProductManager()
        task = "给小学生做一个 2 页 AI 科普 PPT，用 private skill 完成。"
        output = {"format": "pptx", "slide_count": 2, "auto_confirm": True}
        expected_requirement = manager.prepare_confirmed_requirement(
            task=task,
            inputs=[],
            output=output,
            source_understanding=SourceUnderstanding(document_type="none"),
        )
        requirement_payload = expected_requirement.model_dump(mode="json")
        requirement_payload["audience"] = "小学生"
        fake_llm = _PptProductManagerToolCallingFakeLlm(
            function_calls=[
                FunctionCall(
                    name="save_ppt_confirmed_requirement_json",
                    args={"requirement_json": requirement_payload},
                )
            ],
            final_text="PptAutoConfirmWorkflow saved private-skill requirement.",
        )
        private_execution_calls: list[dict[str, Any]] = []

        def _system_selection_builder(**_kwargs):
            return {
                "system_type": "private_skill",
                "route": "html",
                "skill_name": "ppt-complete-workflow",
                "output_format": "html",
                "reason": "Use a private PPT skill for the auto-confirm private branch test.",
            }

        def _content_plan_builder(_requirement):
            return DeckContentPlan(
                title="Private Skill AI",
                core_narrative="Explain AI with classroom examples.",
                pages=[_page(1, "cover"), _page(2, "ending")],
            )

        async def _fake_private_execution(**kwargs):
            private_execution_calls.append(kwargs)
            private_build = manager.save_ppt_private_skill_html(
                file_name="index.html",
                html_content="<!doctype html><html><body><h1>Private Skill AI</h1></body></html>",
                description="Private-skill auto-confirm HTML deck.",
                tool_context=kwargs["tool_context"],
            )
            return PptPrivateSkillExecutionResult(
                skill_name="ppt-complete-workflow",
                output_format="html",
                private_build={**private_build, "skill_name": "ppt-complete-workflow"},
                execution_output={
                    "status": "success",
                    "source": "test_private_skill_execution",
                    "message": "Private skill execution was faked at the phase boundary.",
                },
            )

        async def _run_auto_confirm_private_skill_harness():
            @node(name="PptAutoConfirmPrivateSkillWorkflowHarnessNode", rerun_on_resume=True)
            async def auto_confirm_private_skill_harness(ctx: Context, node_input: str) -> dict[str, Any]:
                with patch.object(
                    manager,
                    "_run_auto_confirm_product_request_direct",
                    side_effect=AssertionError("parent workflow must own auto-confirm delivery"),
                ):
                    return await manager.run_product_request(
                        task=node_input,
                        inputs=[],
                        output=output,
                        tool_context=ctx,
                        content_plan_builder=_content_plan_builder,
                        system_selection_builder=_system_selection_builder,
                    )

            workflow = Workflow(
                name="PptAutoConfirmPrivateSkillWorkflowHarness",
                edges=[("START", auto_confirm_private_skill_harness)],
            )
            session_service = InMemorySessionService()
            runner = Runner(
                node=workflow,
                session_service=session_service,
                artifact_service=artifact_service,
            )
            user_id = "user-ppt-auto-confirm-private-skill-workflow"
            session_id = "session-ppt-auto-confirm-private-skill-workflow"
            try:
                await session_service.create_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                    state={
                        "sid": "ppt-auto-confirm-private-skill-workflow-test",
                        "turn_index": 1,
                        "step": 1,
                    },
                )
                with (
                    patch(
                        "src.productions.ppt.ppt_product_manager.ppt_product_manager.build_llm",
                        return_value=fake_llm,
                    ),
                    patch.object(manager, "_execute_private_ppt_skill_phase", _fake_private_execution),
                    patch.object(
                        manager,
                        "_dispatch_ppt_route",
                        side_effect=AssertionError("private-skill branch must not dispatch a built-in route"),
                    ),
                ):
                    async for _ in runner.run_async(
                        user_id=user_id,
                        session_id=session_id,
                        new_message=Content(role="user", parts=[Part(text=task)]),
                    ):
                        pass
                return await session_service.get_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                )
            finally:
                await runner.close()

        artifact_service = InMemoryArtifactService()
        session = await _run_auto_confirm_private_skill_harness()

        self.assertIsNotNone(session)
        self.assertEqual(session.state["ppt_product_result"]["status"], "success")
        self.assertEqual(session.state["ppt_product_result"]["phase"], "private_skill_delivery")
        self.assertEqual(session.state[PPT_AUTO_CONFIRM_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_AUTO_CONFIRM_WORKFLOW_OUTPUT_KEY]["branch"], "private_skill")
        self.assertEqual(session.state[PPT_SOURCE_PREPARATION_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_REQUIREMENT_ANALYSIS_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_SYSTEM_SELECTION_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_CONTENT_PLANNING_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_PRIVATE_SKILL_EXECUTION_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_PRIVATE_SKILL_DELIVERY_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_CONFIRMED_REQUIREMENT_STATE_KEY]["audience"], "小学生")
        private_output_path = session.state["ppt_private_skill_build"]["output_path"]
        self.assertEqual(session.state["final_file_paths"], [private_output_path])
        self.assertEqual(session.state["ppt_system_selection"]["system_type"], "private_skill")
        self.assertEqual(session.state["ppt_private_skill_build"]["output_path"], private_output_path)
        self.assertEqual(len(private_execution_calls), 1)
        self.assertNotIn(PPT_ASSET_RESOLUTION_WORKFLOW_OUTPUT_KEY, session.state)
        self.assertNotIn(PPT_ROUTE_EXECUTION_WORKFLOW_OUTPUT_KEY, session.state)
        self.assertNotIn(PPT_FINAL_DELIVERY_WORKFLOW_OUTPUT_KEY, session.state)
        self.assertNotIn("ppt_workflow_state", session.state)

    def test_system_selection_agent_chooses_without_keyword_rules(self) -> None:
        manager = PptProductManager()

        agent = manager.build_system_selection_agent()

        self.assertIsInstance(agent, LlmAgent)
        self.assertEqual(agent.name, "PptSystemSelectionAgent")
        self.assertEqual(agent.output_key, "ppt_system_selection_agent_message")
        self.assertEqual(
            {tool.__name__ for tool in agent.tools},
            {
                "list_product_ppt_skills",
                "read_product_ppt_skill",
                "save_ppt_system_selection",
            },
        )
        self.assertIn("Do not use hard-coded keyword rules", agent.instruction)
        self.assertIn("source_inputs include PPTX/PPTM/POTX/POTM", agent.instruction)
        self.assertIn("built-in `html` or `svg`", agent.instruction)

    async def test_system_selection_agenttool_main_path_saves_selection(self) -> None:
        manager = PptProductManager()
        task = "做一个 5 页 SVG route 产品介绍 PPT。"
        output = {"format": "pptx", "route": "svg", "slide_count": 5}
        requirement = manager.prepare_confirmed_requirement(
            task=task,
            inputs=[],
            output=output,
        )
        selection_payload = {
            "system_type": "built_in_route",
            "route": "svg",
            "skill_name": "",
            "output_format": "pptx",
            "reason": "The built-in SVG route best fits an editable vector deck.",
        }
        fake_llm = _PptProductManagerToolCallingFakeLlm(
            function_calls=[
                FunctionCall(name="list_product_ppt_skills", args={}),
                FunctionCall(
                    name="save_ppt_system_selection",
                    args={"selection_json": selection_payload},
                ),
            ],
            final_text="PptSystemSelectionAgent saved the PPT system selection.",
        )
        session_service = InMemorySessionService()
        artifact_service = InMemoryArtifactService()

        @node(name="PptSystemSelectionAgentToolHarnessNode", rerun_on_resume=True)
        async def selection_harness(ctx: Context, node_input: str) -> dict:
            selection = await manager.select_ppt_system_with_agent(
                task=task,
                output=output,
                requirement=requirement,
                tool_context=ctx,
                app_name="creative_claw",
                artifact_service=artifact_service,
            )
            return dict(selection)

        workflow = Workflow(
            name="PptSystemSelectionAgentToolHarness",
            edges=[("START", selection_harness)],
        )
        runner = Runner(
            node=workflow,
            session_service=session_service,
            artifact_service=artifact_service,
        )
        user_id = "user-ppt-system-selection-agenttool"
        session_id = "session-ppt-system-selection-agenttool"

        try:
            await session_service.create_session(
                app_name=workflow.name,
                user_id=user_id,
                session_id=session_id,
                state={
                    "sid": "ppt-system-selection-agenttool-test",
                    "turn_index": 1,
                    "step": 1,
                },
            )
            with patch(
                "src.productions.ppt.ppt_product_manager.ppt_product_manager.build_llm",
                return_value=fake_llm,
            ):
                async for _ in runner.run_async(
                    user_id=user_id,
                    session_id=session_id,
                    new_message=Content(role="user", parts=[Part(text="Select PPT system")]),
                ):
                    pass
            session = await session_service.get_session(
                app_name=workflow.name,
                user_id=user_id,
                session_id=session_id,
            )
        finally:
            await runner.close()

        self.assertIsNotNone(session)
        self.assertEqual(session.state[PPT_SYSTEM_SELECTION_STATE_KEY]["route"], "svg")
        self.assertEqual(
            session.state[PPT_SYSTEM_SELECTION_STATE_KEY]["reason"],
            selection_payload["reason"],
        )
        self.assertIn(PPT_SYSTEM_SELECTION_AGENT_MESSAGE_KEY, session.state)
        self.assertEqual(len(session.state["product_ppt_skills"]), len(manager.skill_registry.list_skills()))
        first_request_tools = _function_declaration_names(fake_llm.requests[0])
        self.assertIn("list_product_ppt_skills", first_request_tools)
        self.assertIn("read_product_ppt_skill", first_request_tools)
        self.assertIn("save_ppt_system_selection", first_request_tools)
        self.assertGreaterEqual(len(fake_llm.requests), 3)

    async def test_system_selection_phase_persists_injected_built_in_result(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 SVG route 产品介绍 PPT。",
            inputs=[],
            output={"format": "pptx", "route": "svg"},
        )
        tool_context = SimpleNamespace(
            state={"sid": "ppt-system-selection-phase-injected-test", "turn_index": 1, "step": 1}
        )

        def _selection_builder(**_kwargs):
            return {
                "system_type": "built_in_route",
                "route": "svg",
                "skill_name": "",
                "output_format": "pptx",
                "reason": "Injected selector chose SVG for editable vector output.",
            }

        result = await manager._select_ppt_system_phase(
            task="做一个 SVG route 产品介绍 PPT。",
            output={"format": "pptx", "route": "svg"},
            requirement=requirement,
            tool_context=tool_context,
            app_name="creative_claw",
            artifact_service=None,
            system_selection_builder=_selection_builder,
        )

        self.assertIsInstance(result, PptSystemSelectionResult)
        self.assertEqual(result.system_selection.system_type, "built_in_route")
        self.assertEqual(result.system_selection.route, "svg")
        self.assertEqual(result.selection_output["source"], "injected")
        self.assertEqual(tool_context.state[PPT_SYSTEM_SELECTION_STATE_KEY]["route"], "svg")
        self.assertEqual(tool_context.state[PPT_SYSTEM_SELECTION_OUTPUT_STATE_KEY]["source"], "injected")
        self.assertEqual(
            tool_context.state[PPT_SYSTEM_SELECTION_RESULT_STATE_KEY]["system_selection"]["reason"],
            "Injected selector chose SVG for editable vector output.",
        )

    async def test_system_selection_phase_persists_private_template_fallback_result(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="套用用户上传 PPTX 模板生成汇报。",
            inputs=[{"name": "template.pptx", "path": "inbox/demo/template.pptx"}],
            output={"format": "pptx", "route": "xml"},
        )
        tool_context = SimpleNamespace(
            state={"sid": "ppt-system-selection-phase-fallback-test", "turn_index": 1, "step": 1}
        )

        result = await manager._select_ppt_system_phase(
            task="套用用户上传 PPTX 模板生成汇报。",
            output={"format": "pptx", "route": "xml"},
            requirement=requirement,
            tool_context=tool_context,
            app_name="creative_claw",
            artifact_service=None,
        )

        self.assertIsInstance(result, PptSystemSelectionResult)
        self.assertEqual(result.system_selection.system_type, "private_skill")
        self.assertEqual(result.system_selection.route, "xml")
        self.assertEqual(result.system_selection.skill_name, "pptx")
        self.assertEqual(result.system_selection.output_format, "pptx")
        self.assertEqual(result.selection_output["source"], "deterministic_fallback")
        self.assertEqual(tool_context.state[PPT_SYSTEM_SELECTION_STATE_KEY]["skill_name"], "pptx")
        self.assertEqual(
            tool_context.state[PPT_SYSTEM_SELECTION_RESULT_STATE_KEY]["selection_output"]["source"],
            "deterministic_fallback",
        )

    async def test_system_selection_phase_uses_adk_workflow_context(self) -> None:
        manager = PptProductManager()
        task = "做一个 5 页 SVG route 产品介绍 PPT。"
        output = {"format": "pptx", "route": "svg", "slide_count": 5}
        requirement = manager.prepare_confirmed_requirement(
            task=task,
            inputs=[],
            output=output,
        )
        selection_payload = {
            "system_type": "built_in_route",
            "route": "svg",
            "skill_name": "",
            "output_format": "pptx",
            "reason": "The built-in SVG route best fits an editable vector deck.",
        }
        fake_llm = _PptProductManagerToolCallingFakeLlm(
            function_calls=[
                FunctionCall(name="list_product_ppt_skills", args={}),
                FunctionCall(
                    name="save_ppt_system_selection",
                    args={"selection_json": selection_payload},
                ),
            ],
            final_text="PptSystemSelectionAgent saved the PPT system selection.",
        )

        async def _run_selection_harness():
            @node(name="PptSystemSelectionWorkflowHarnessNode", rerun_on_resume=True)
            async def selection_harness(ctx: Context, node_input: str) -> dict[str, Any]:
                result = await manager._select_ppt_system_phase(
                    task=task,
                    output=output,
                    requirement=requirement,
                    tool_context=ctx,
                    app_name="creative_claw",
                    artifact_service=artifact_service,
                )
                return result.model_dump(mode="json")

            workflow = Workflow(
                name="PptSystemSelectionWorkflowHarness",
                edges=[("START", selection_harness)],
            )
            session_service = InMemorySessionService()
            runner = Runner(
                node=workflow,
                session_service=session_service,
                artifact_service=artifact_service,
            )
            user_id = "user-ppt-system-selection-workflow"
            session_id = "session-ppt-system-selection-workflow"
            try:
                await session_service.create_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                    state={
                        "sid": "ppt-system-selection-workflow-test",
                        "turn_index": 1,
                        "step": 1,
                    },
                )
                with patch(
                    "src.productions.ppt.ppt_product_manager.ppt_product_manager.build_llm",
                    return_value=fake_llm,
                ):
                    async for _ in runner.run_async(
                        user_id=user_id,
                        session_id=session_id,
                        new_message=Content(role="user", parts=[Part(text="Select PPT system")]),
                    ):
                        pass
                return await session_service.get_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                )
            finally:
                await runner.close()

        artifact_service = InMemoryArtifactService()
        session = await _run_selection_harness()

        self.assertIsNotNone(session)
        self.assertEqual(
            session.state[PPT_SYSTEM_SELECTION_WORKFLOW_OUTPUT_KEY]["source"],
            "adk_workflow",
        )
        self.assertEqual(
            session.state[PPT_SYSTEM_SELECTION_WORKFLOW_OUTPUT_KEY]["selection_source"],
            "llm_agent",
        )
        self.assertEqual(session.state[PPT_SYSTEM_SELECTION_STATE_KEY]["route"], "svg")
        self.assertEqual(session.state[PPT_SYSTEM_SELECTION_OUTPUT_STATE_KEY]["source"], "llm_agent")
        self.assertEqual(
            session.state[PPT_SYSTEM_SELECTION_RESULT_STATE_KEY]["system_selection"]["reason"],
            selection_payload["reason"],
        )
        self.assertIn(PPT_SYSTEM_SELECTION_AGENT_MESSAGE_KEY, session.state)
        first_request_tools = _function_declaration_names(fake_llm.requests[0])
        self.assertIn("list_product_ppt_skills", first_request_tools)
        self.assertIn("save_ppt_system_selection", first_request_tools)

    def test_ppt_svg_strategy_tools_save_and_read_execution_plan(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={})

        strategy_result = manager.save_ppt_design_strategy(
            {
                "style_name": "clean_svg",
                "design_direction": "Use editable SVG primitives.",
                "palette": ["#FFFFFF", "#172033", "#2457D6"],
            },
            {
                "summary": "Use clean SVG route styling.",
                "decisions": ["Use 16:9"],
            },
            tool_context,
        )
        plan_result = manager.save_ppt_svg_execution_plan(
            {
                "aspect_ratio": "16:9",
                "canvas_width": 1280,
                "canvas_height": 720,
                "accent_color": "#2457D6",
            },
            tool_context,
        )
        read_result = manager.read_ppt_svg_execution_plan(tool_context)

        self.assertEqual(strategy_result["status"], "success")
        self.assertEqual(plan_result["status"], "success")
        self.assertEqual(read_result["status"], "success")
        self.assertEqual(tool_context.state["ppt_design_strategy"]["style_name"], "clean_svg")
        self.assertEqual(tool_context.state[PPT_SVG_EXECUTION_PLAN_STATE_KEY]["canvas_width"], 1280)

    def test_save_ppt_svg_execution_plan_preserves_route_guidance_fields(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(
            state={
                PPT_SVG_EXECUTION_PLAN_STATE_KEY: {
                    "aspect_ratio": "16:9",
                    "canvas_width": 1280,
                    "canvas_height": 720,
                    "page_rhythm_by_slide": {"P01": "anchor", "P02": "dense"},
                    "typography_ramp": {"page_title": 40, "body": 22},
                    "page_type_layout_guidance": {"comparison": "Use matched comparison lanes."},
                    "template_adherence_rules": {"content": "Preserve template header and footer."},
                    "quality_constraints": ["Respect page rhythm."],
                }
            }
        )

        result = manager.save_ppt_svg_execution_plan(
            {
                "aspect_ratio": "16:9",
                "canvas_width": 1280,
                "canvas_height": 720,
                "accent_color": "#FF0000",
                "page_rhythm_by_slide": {},
            },
            tool_context,
        )
        saved_plan = tool_context.state[PPT_SVG_EXECUTION_PLAN_STATE_KEY]

        self.assertEqual(result["status"], "success")
        self.assertEqual(saved_plan["accent_color"], "#FF0000")
        self.assertEqual(saved_plan["page_rhythm_by_slide"]["P02"], "dense")
        self.assertEqual(saved_plan["typography_ramp"]["body"], 22)
        self.assertIn("comparison", saved_plan["page_type_layout_guidance"])
        self.assertEqual(saved_plan["template_adherence_rules"]["content"], "Preserve template header and footer.")

    def test_content_planning_page_type_heuristics_use_richer_svg_types(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 SVG route PPT，对比两个方案并说明实施流程和关键指标。",
            inputs=[],
            output={"format": "pptx", "route": "svg"},
        )

        self.assertEqual(
            _select_content_page_type(
                requirement=requirement,
                point="方案 A 和方案 B 的成本收益对比",
                support="资源投入和风险差异明显",
                index=0,
                default="content",
            ),
            "comparison",
        )
        self.assertEqual(
            _select_content_page_type(
                requirement=requirement,
                point="实施流程分为需求确认、试点、推广三个步骤",
                support="每一步都有明确交付物",
                index=1,
                default="content",
            ),
            "process",
        )
        self.assertEqual(
            _select_content_page_type(
                requirement=requirement,
                point="核心指标提升 35%",
                support="效果已经超过目标",
                index=2,
                default="content",
            ),
            "stat",
        )

    def test_product_manager_skill_runner_masks_saved_long_html_arguments(self) -> None:
        manager = PptProductManager()
        long_html = "<!DOCTYPE html>" + ("x" * 9000)
        output_path = "generated/session/turn_3/ppt_private_skill_step_3/index.html"
        callback_context = SimpleNamespace(
            state={
                "ppt_private_skill_build": {
                    "output_path": output_path,
                }
            }
        )
        llm_request = SimpleNamespace(
            contents=[
                Content(
                    role="model",
                    parts=[
                        Part.from_function_call(
                            name="save_ppt_private_skill_html",
                            args={
                                "file_name": "index.html",
                                "html_content": long_html,
                                "description": "HTML deck",
                            },
                        )
                    ],
                )
            ]
        )

        self.assertEqual(manager.include_contents, "none")
        manager.before_model_callback(callback_context, llm_request)

        args = llm_request.contents[0].parts[0].function_call.args
        self.assertEqual(args["file_name"], "index.html")
        self.assertEqual(args["description"], "HTML deck")
        self.assertNotIn("<!DOCTYPE html>", args["html_content"])
        self.assertIn(output_path, args["html_content"])
        self.assertIn("<tool_output_masked>", args["html_content"])
        self.assertEqual(callback_context.state["ppt_private_skill_masked_html_content_count"], 1)

    def test_product_manager_skill_runner_keeps_short_html_arguments(self) -> None:
        manager = PptProductManager()
        short_html = "<!DOCTYPE html><html><body>Short</body></html>"
        callback_context = SimpleNamespace(state={})
        llm_request = SimpleNamespace(
            contents=[
                Content(
                    role="model",
                    parts=[
                        Part.from_function_call(
                            name="save_ppt_private_skill_html",
                            args={
                                "file_name": "index.html",
                                "html_content": short_html,
                            },
                        )
                    ],
                )
            ]
        )

        manager.before_model_callback(callback_context, llm_request)

        args = llm_request.contents[0].parts[0].function_call.args
        self.assertEqual(args["html_content"], short_html)
        self.assertNotIn("ppt_private_skill_masked_html_content_count", callback_context.state)

    def test_private_skill_execution_agent_is_not_exposed(self) -> None:
        manager = PptProductManager()

        self.assertFalse(hasattr(manager, "build_private_skill_execution_agent"))

    def test_requirement_analysis_save_tool_preserves_source_inputs(self) -> None:
        source_path = _write_markdown_source("requirement_source.pdf", "%PDF test fixture")
        manager = PptProductManager()
        fallback_requirement = manager.prepare_confirmed_requirement(
            task="针对这个素材，给我做一个ppt。",
            inputs={"files": [source_path]},
            output={"format": "pptx"},
        )
        tool_context = SimpleNamespace(
            state={
                "ppt_requirement_analysis_base": {
                    "fallback_requirement": fallback_requirement.model_dump(mode="json"),
                }
            }
        )

        result = manager.save_ppt_confirmed_requirement_json(
            {
                "route": "html",
                "topic": "视觉原语推理",
                "audience": "同组的同学",
                "scenario": "组会",
                "slide_count_policy": {
                    "minimum": 14,
                    "maximum": 16,
                    "target": 15,
                    "source": "user",
                },
                "source_inputs": [{"name": "invented.pdf", "path": "missing.pdf"}],
            },
            tool_context,
        )

        self.assertEqual(result["status"], "success")
        saved_requirement = tool_context.state["ppt_confirmed_requirement"]
        self.assertEqual(saved_requirement["topic"], "视觉原语推理")
        self.assertEqual(saved_requirement["audience"], "同组的同学")
        self.assertEqual(saved_requirement["scenario"], "组会")
        self.assertEqual(saved_requirement["slide_count_policy"]["target"], 15)
        self.assertEqual(saved_requirement["source_inputs"][0]["path"], source_path)

    def test_content_planning_user_message_includes_requirement_json(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="给我做一个ppt，用来给幼儿园小朋友讲英语单词。图文并茂。小于10页。",
            inputs=[],
            output={
                "format": "pptx",
                "language": "zh-CN",
                "slide_count": "小于10页",
                "style": "图文并茂、活泼可爱、适合儿童英语启蒙",
            },
        )

        user_message = _build_content_planning_user_message(requirement)

        self.assertIn("ConfirmedRequirement JSON", user_message)
        self.assertIn('"request_brief": "给我做一个ppt，用来给幼儿园小朋友讲英语单词。图文并茂。小于10页。"', user_message)
        self.assertIn('"topic": "英语单词"', user_message)
        self.assertIn('"audience": "幼儿园小朋友"', user_message)
        self.assertIn("Do not invent a generic business communication deck", user_message)

    def test_content_planning_tools_read_and_save_plan(self) -> None:
        source_path = _write_markdown_source(
            "planning_brief.md",
            "# Planning Brief\n\n- Activation rose after onboarding.\n",
        )
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="基于材料生成 5 页 PPTX。",
            inputs=[{"name": "planning_brief.md", "path": source_path}],
            output={"format": "pptx"},
            source_understanding=SourceUnderstanding(
                document_type="markdown",
                markdown_sources=[
                    {
                        "name": "planning_brief.md",
                        "source_path": source_path,
                        "method": "test",
                        "output_path": source_path,
                    }
                ],
            ),
        )
        tool_context = SimpleNamespace(
            state={
                "ppt_confirmed_requirement": requirement.model_dump(mode="json"),
            }
        )

        source_result = manager.content_planner.read_ppt_markdown_sources(tool_context)
        markdown_plan = """# Deck: Planning Brief
Audience: Internal team
Language: en
SlideCount: 5
Narrative: Explain activation changes.

## Slide 1 | cover | Planning Brief
Purpose: Introduce the planning brief.
Takeaway: Activation rose after onboarding.
Content:
- Audience: Growth team
Visual:
- placeholder | role=hero | description=clean title area

## Slide 2 | toc | Agenda
Purpose: Preview the deck.
Takeaway: The deck covers evidence and next steps.
Content:
- Activation
- Evidence
- Next steps
Visual:
- placeholder | role=list | description=agenda list

## Slide 3 | chapter_start | Activation
Purpose: Start the activation chapter.
Takeaway: Activation rose after onboarding.
Content:
- Activation rose after onboarding.
Visual:
- search | role=reference | query=activation onboarding chart | description=visual reference for activation onboarding

## Slide 4 | chapter_content | Evidence
Purpose: Explain the evidence.
Takeaway: Guided onboarding improved activation.
Content:
- Activation rose after onboarding.
- Enterprise teams need proof.
Visual:
- ai | role=supporting_visual | description=friendly product onboarding illustration

## Slide 5 | ending | Next Steps
Purpose: Close with next steps.
Takeaway: Use the activation proof in the story.
Content:
- Review the evidence
- Prepare the launch story
Visual:
- placeholder | role=summary | description=closing icon area
"""
        save_result = manager.content_planner.save_ppt_deck_content_plan_markdown(
            markdown_plan,
            tool_context,
        )

        self.assertEqual(source_result["status"], "success")
        self.assertIn("Activation rose", source_result["source_texts"][0]["text"])
        self.assertEqual(save_result["status"], "success")
        self.assertEqual(tool_context.state["ppt_deck_content_plan"]["title"], "Planning Brief")
        self.assertIn("ppt_deck_content_plan_markdown", tool_context.state)
        self.assertEqual(tool_context.state["ppt_deck_content_plan"]["pages"][2]["asset_source_preference"], "search")
        self.assertEqual(tool_context.state["ppt_deck_content_plan"]["pages"][3]["asset_source_preference"], "ai")

    def test_content_planning_rejects_off_task_kindergarten_business_plan(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="给我做一个ppt，用来给幼儿园小朋友讲英语单词。图文并茂。小于10页。",
            inputs=[],
            output={"format": "pptx"},
        )
        tool_context = SimpleNamespace(
            state={
                "ppt_confirmed_requirement": requirement.model_dump(mode="json"),
            }
        )
        bad_markdown_plan = """# Deck: 目标对齐沟通稿
Audience: 团队成员
Language: zh-CN
SlideCount: 6
Narrative: 通过清晰的背景、目标、关键信息和协作安排，帮助团队快速形成一致理解。

## Slide 1 | cover | 目标对齐沟通稿
Purpose: 建立主题氛围，说明本次沟通聚焦于统一理解与推进协作。
Takeaway: 团队需要先对目标与重点形成共同认知。
Content:
- 聚焦共同目标
- 明确核心信息
Visual:
- ai | role=hero | description=现代团队围绕简洁白板讨论目标，明亮办公空间，专业、清爽、无文字

## Slide 2 | toc | 内容一览
Purpose: 展示整体结构。
Takeaway: 本次内容将从背景、目标、重点和协作安排展开。
Content:
- 背景与目标
- 关键信息梳理
Visual:
- placeholder | role=grid | description=四段式目录布局

## Slide 3 | chapter_start | 第一部分：背景与目标
Purpose: 开启背景说明章节。
Takeaway: 明确背景是判断重点与行动方向的前提。
Content:
- 先看现状
- 再定方向
Visual:
- ai | role=hero | description=抽象路线图从起点延伸到目标旗帜，简洁商务插画风，无文字

## Slide 4 | chapter_content | 明确沟通对象
Purpose: 梳理受众关注点。
Takeaway: 面向不同对象时，信息重点与表达深度需要有所侧重。
Content:
- 识别主要听众与决策角色
- 提炼听众最关心的问题
Visual:
- placeholder | role=grid | description=人物角色卡片与关注点列表

## Slide 5 | ending | 后续协作
Purpose: 收束内容，推动会后形成明确协作节奏。
Takeaway: 共识需要转化为责任、时间和交付物。
Content:
- 确认负责人和参与方
- 明确近期交付物
Visual:
- ai | role=hero | description=团队成员把任务卡片贴到看板上，现代扁平插画，无文字
"""

        with self.assertRaisesRegex(ValueError, "kindergarten English-word task"):
            manager.content_planner.save_ppt_deck_content_plan_markdown(
                bad_markdown_plan,
                tool_context,
            )

    async def test_dispatch_ppt_route_tool_uses_state_registry(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 5 页 PPTX 产品介绍。",
            inputs=[],
            output={"format": "pptx"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        tool_context = SimpleNamespace(
            state={
                "sid": "ppt-dispatch-tool-test",
                "turn_index": 1,
                "step": 1,
                "ppt_confirmed_requirement": requirement.model_dump(mode="json"),
                "ppt_deck_content_plan": content_plan.model_dump(mode="json"),
            }
        )

        result = await manager.dispatch_ppt_route(route="html", tool_context=tool_context)

        self.assertEqual(result["status"], "success")
        self.assertEqual(result["selected_route"], "html")
        self.assertEqual(tool_context.state["ppt_route_build"]["template"]["template_id"], "free_design")
        self.assertTrue(result["output_files"])

    async def test_route_execution_phase_reuses_existing_successful_build(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 5 页 PPTX 产品介绍。",
            inputs=[],
            output={"format": "pptx"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        tool_context = SimpleNamespace(state={"sid": "ppt-route-phase-test", "turn_index": 1, "step": 1})
        route_calls: list[Path] = []

        async def _fake_dispatch_ppt_route(**kwargs):
            output_dir = Path(kwargs["output_dir"])
            route_calls.append(output_dir)
            final_pptx = output_dir / "deck.pptx"
            final_pptx.write_bytes(b"fake-pptx")
            html_path = output_dir / "deck.html"
            html_path.write_text("<html></html>", encoding="utf-8")
            return HtmlRouteBuildPackage(
                template=HtmlTemplatePackage(template_id="workflow-test", label="Workflow Test"),
                html_deck_path=workspace_relative_path(html_path),
                preview_paths=[],
                pptx_path=workspace_relative_path(final_pptx),
                quality_report_path="",
                build_log_path="",
                warnings=[],
            )

        with patch.object(manager, "_dispatch_ppt_route", _fake_dispatch_ppt_route):
            first_result = await manager._execute_ppt_route_phase(
                requirement=requirement,
                content_plan=content_plan,
                tool_context=tool_context,
                expert_agents={},
            )
            second_result = await manager._execute_ppt_route_phase(
                requirement=requirement,
                content_plan=content_plan,
                tool_context=tool_context,
                expert_agents={},
            )

        self.assertIsInstance(first_result, PptRouteExecutionResult)
        self.assertFalse(first_result.reused_existing_build)
        self.assertTrue(second_result.reused_existing_build)
        self.assertEqual(len(route_calls), 1)
        self.assertEqual(first_result.route_build.pptx_path, second_result.route_build.pptx_path)
        self.assertEqual(
            tool_context.state["ppt_route_execution_result"]["route_build"]["pptx_path"],
            first_result.route_build.pptx_path,
        )
        self.assertTrue(tool_context.state["ppt_route_execution_result"]["input_signature"])
        self.assertEqual(tool_context.state["ppt_route_build"]["pptx_path"], first_result.route_build.pptx_path)

    async def test_route_execution_phase_reruns_when_inputs_change(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 5 页 PPTX 产品介绍。",
            inputs=[],
            output={"format": "pptx"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        changed_plan = content_plan.model_copy(update={"title": "Changed Product Deck"}, deep=True)
        tool_context = SimpleNamespace(state={"sid": "ppt-route-phase-input-test", "turn_index": 1, "step": 1})
        route_calls: list[Path] = []

        async def _fake_dispatch_ppt_route(**kwargs):
            output_dir = Path(kwargs["output_dir"])
            route_calls.append(output_dir)
            final_pptx = output_dir / "deck.pptx"
            final_pptx.write_bytes(f"fake-pptx-{len(route_calls)}".encode("utf-8"))
            html_path = output_dir / "deck.html"
            html_path.write_text("<html></html>", encoding="utf-8")
            return HtmlRouteBuildPackage(
                template=HtmlTemplatePackage(template_id="workflow-test", label="Workflow Test"),
                html_deck_path=workspace_relative_path(html_path),
                preview_paths=[],
                pptx_path=workspace_relative_path(final_pptx),
                quality_report_path="",
                build_log_path="",
                warnings=[],
            )

        with patch.object(manager, "_dispatch_ppt_route", _fake_dispatch_ppt_route):
            first_result = await manager._execute_ppt_route_phase(
                requirement=requirement,
                content_plan=content_plan,
                tool_context=tool_context,
                expert_agents={},
            )
            second_result = await manager._execute_ppt_route_phase(
                requirement=requirement,
                content_plan=changed_plan,
                tool_context=tool_context,
                expert_agents={},
            )

        self.assertFalse(first_result.reused_existing_build)
        self.assertFalse(second_result.reused_existing_build)
        self.assertEqual(len(route_calls), 2)
        self.assertNotEqual(first_result.input_signature, second_result.input_signature)
        self.assertEqual(tool_context.state["ppt_route_execution_result"]["input_signature"], second_result.input_signature)

    async def test_route_final_delivery_phase_records_typed_result(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 5 页 PPTX 产品介绍。",
            inputs=[],
            output={"format": "pptx"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        tool_context = SimpleNamespace(state={"sid": "ppt-final-delivery-phase-test", "turn_index": 1, "step": 1})
        output_dir = manager._build_route_output_dir(tool_context.state, route=requirement.route)
        final_pptx = output_dir / "deck.pptx"
        final_pptx.write_bytes(b"fake-pptx")
        html_path = output_dir / "deck.html"
        html_path.write_text("<html></html>", encoding="utf-8")
        route_build = HtmlRouteBuildPackage(
            template=HtmlTemplatePackage(template_id="workflow-test", label="Workflow Test"),
            html_deck_path=workspace_relative_path(html_path),
            preview_paths=[],
            pptx_path=workspace_relative_path(final_pptx),
            quality_report_path="",
            build_log_path="",
            warnings=["route warning"],
        )
        route_execution = PptRouteExecutionResult(
            route=requirement.route,
            output_dir=workspace_relative_path(output_dir),
            input_signature=manager._route_execution_input_signature(requirement, content_plan),
            route_build=route_build,
        )

        result = manager._finalize_route_delivery_phase(
            requirement=requirement,
            content_plan=content_plan,
            route_execution=route_execution,
            tool_context=tool_context,
        )

        self.assertIsInstance(result, PptFinalDeliveryResult)
        self.assertEqual(result.product_result.status, "success")
        self.assertEqual(result.delivery_manifest.final_pptx, route_build.pptx_path)
        self.assertEqual(tool_context.state["final_file_paths"], [route_build.pptx_path])
        self.assertEqual(
            tool_context.state["ppt_final_delivery_result"]["delivery_manifest"]["final_pptx"],
            route_build.pptx_path,
        )
        self.assertIn("route warning", result.product_result.warnings)

    async def test_route_final_delivery_phase_failure_does_not_register_final_pptx(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 5 页 PPTX 产品介绍。",
            inputs=[],
            output={"format": "pptx"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        tool_context = SimpleNamespace(state={"sid": "ppt-final-delivery-failure-test", "turn_index": 1, "step": 1})
        output_dir = manager._build_route_output_dir(tool_context.state, route=requirement.route)
        html_path = output_dir / "deck.html"
        html_path.write_text("<html></html>", encoding="utf-8")
        route_build = HtmlRouteBuildPackage(
            template=HtmlTemplatePackage(template_id="workflow-test", label="Workflow Test"),
            html_deck_path=workspace_relative_path(html_path),
            preview_paths=[],
            pptx_path="",
            quality_report_path="",
            build_log_path="",
            warnings=["html to pptx conversion failed"],
        )
        route_execution = PptRouteExecutionResult(
            route=requirement.route,
            output_dir=workspace_relative_path(output_dir),
            input_signature=manager._route_execution_input_signature(requirement, content_plan),
            route_build=route_build,
        )

        result = manager._finalize_route_delivery_phase(
            requirement=requirement,
            content_plan=content_plan,
            route_execution=route_execution,
            tool_context=tool_context,
        )

        self.assertEqual(result.product_result.status, "generation_failed")
        self.assertEqual(result.delivery_manifest.final_pptx, "")
        self.assertNotIn("final_file_paths", tool_context.state)
        self.assertIn(workspace_relative_path(html_path), [record["path"] for record in result.output_files])
        self.assertEqual(
            tool_context.state["ppt_final_delivery_result"]["product_result"]["status"],
            "generation_failed",
        )

    async def test_route_final_delivery_phase_is_idempotent_for_same_outputs(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 5 页 PPTX 产品介绍。",
            inputs=[],
            output={"format": "pptx"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        tool_context = SimpleNamespace(state={"sid": "ppt-final-delivery-idempotent-test", "turn_index": 1, "step": 1})
        output_dir = manager._build_route_output_dir(tool_context.state, route=requirement.route)
        final_pptx = output_dir / "deck.pptx"
        final_pptx.write_bytes(b"fake-pptx")
        html_path = output_dir / "deck.html"
        html_path.write_text("<html></html>", encoding="utf-8")
        route_build = HtmlRouteBuildPackage(
            template=HtmlTemplatePackage(template_id="workflow-test", label="Workflow Test"),
            html_deck_path=workspace_relative_path(html_path),
            preview_paths=[],
            pptx_path=workspace_relative_path(final_pptx),
            quality_report_path="",
            build_log_path="",
            warnings=[],
        )
        route_execution = PptRouteExecutionResult(
            route=requirement.route,
            output_dir=workspace_relative_path(output_dir),
            input_signature=manager._route_execution_input_signature(requirement, content_plan),
            route_build=route_build,
        )

        first_result = manager._finalize_route_delivery_phase(
            requirement=requirement,
            content_plan=content_plan,
            route_execution=route_execution,
            tool_context=tool_context,
        )
        second_result = manager._finalize_route_delivery_phase(
            requirement=requirement,
            content_plan=content_plan,
            route_execution=route_execution,
            tool_context=tool_context,
        )

        self.assertEqual(first_result.output_files, second_result.output_files)
        self.assertEqual(tool_context.state["new_files"], first_result.output_files)
        self.assertEqual(tool_context.state["generated"], first_result.output_files)
        self.assertEqual(tool_context.state["files_history"], [first_result.output_files])
        self.assertEqual(tool_context.state["final_file_paths"], [route_build.pptx_path])

    async def test_dispatch_ppt_route_injects_product_html_page_expert(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 5 页 PPTX 产品介绍。",
            inputs=[],
            output={"format": "pptx"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        captured: dict[str, object] = {}

        async def _fake_build_html_route_with_agent(**kwargs):
            captured.update(kwargs)
            return SimpleNamespace(pptx_path="generated/test/deck.pptx")

        with tempfile.TemporaryDirectory(dir=workspace_root()) as tmpdir:
            with patch(
                "src.productions.ppt.ppt_product_manager.ppt_product_manager.build_html_route_with_agent",
                _fake_build_html_route_with_agent,
            ):
                result = await manager._dispatch_ppt_route(
                    requirement=requirement,
                    content_plan=content_plan,
                    output_dir=Path(tmpdir),
                    tool_context=SimpleNamespace(state={}),
                    expert_agents=manager.product_expert_agents,
                )

        self.assertEqual(result.pptx_path, "generated/test/deck.pptx")
        self.assertIs(
            captured["page_generation_agent"],
            manager.product_expert_agents[PPT_HTML_PAGE_GENERATION_EXPERT_NAME],
        )

    def test_prepare_confirmed_requirement_defaults_to_html_mvp_for_pptx(self) -> None:
        manager = PptProductManager()

        requirement = manager.prepare_confirmed_requirement(
            task="做一个 6 页 PPTX，用于产品发布会。",
            inputs=[{"name": "brief.md", "path": "inbox/demo/brief.md"}],
            output={"format": "pptx"},
        )

        self.assertEqual(requirement.route, "html")
        self.assertEqual(requirement.output_format, "pptx")
        self.assertEqual(requirement.slide_count_policy.target, 6)
        self.assertEqual(requirement.slide_count_policy.source, "user")
        self.assertEqual(requirement.language, "zh-CN")
        self.assertEqual(requirement.source_understanding.document_type, "markdown")
        self.assertFalse(requirement.template_requirement.use_template)
        self.assertEqual(requirement.template_requirement.template_source, "none")
        self.assertEqual(requirement.editability_requirement.level, "high")
        self.assertFalse(requirement.confirmed_by_user)

    def test_prepare_confirmed_requirement_keeps_task_brief_without_documents(self) -> None:
        manager = PptProductManager()
        task = "给我做一个ppt，用来给幼儿园小朋友讲英语单词。图文并茂。小于10页。"

        requirement = manager.prepare_confirmed_requirement(
            task=task,
            inputs=[],
            output={"format": "pptx"},
        )

        self.assertEqual(requirement.request_brief, task)
        self.assertEqual(requirement.source_inputs, [])
        self.assertEqual(requirement.source_understanding.document_type, "brief")

    def test_prepare_confirmed_requirement_accepts_file_path_strings(self) -> None:
        source_path = _write_markdown_source("creative_agent_NeurIPS_2026_10_.pdf", "%PDF test fixture")
        manager = PptProductManager()

        requirement = manager.prepare_confirmed_requirement(
            task="针对这个素材，给我做一个ppt，用来组会上给团队的同学讲解。",
            inputs={"files": [source_path]},
            output={"format": "pptx", "language": "zh-CN", "purpose": "组会讲解"},
        )

        self.assertEqual(len(requirement.source_inputs), 1)
        self.assertEqual(requirement.source_inputs[0].name, "creative_agent_NeurIPS_2026_10_.pdf")
        self.assertEqual(requirement.source_inputs[0].path, source_path)
        self.assertEqual(requirement.source_understanding.document_type, "pdf")
        self.assertIn("| 输入材料 | 1 个 |", manager._format_requirement_confirmation(requirement))

    def test_prepare_confirmed_requirement_separates_task_documents_and_ignored_outline(self) -> None:
        source_path = _write_markdown_source("kid_words.md", "# Words\n\n- Apple\n")
        manager = PptProductManager()
        task = (
            "重新制作PPT：主题必须是“给幼儿园小朋友讲英语单词”，不是商务汇报。"
            "必须包含 Apple、Cat、Dog、Sun、Ball。"
        )

        requirement = manager.prepare_confirmed_requirement(
            task=task,
            inputs={
                "outline": [{"slide": 1, "title": "不要把这个当文档"}],
                "documents": [{"name": "kid_words.md", "path": source_path, "mime_type": "text/markdown"}],
            },
            output={
                "format": "pptx",
                "language": "zh-CN",
                "slide_count": 8,
                "style": "儿童友好、卡通、图文并茂、明亮柔和配色",
                "must_not_include": "商务、目标共识、推进路径、团队协作、行动计划",
            },
        )

        self.assertEqual(requirement.request_brief, task)
        self.assertEqual(requirement.topic, "英语单词")
        self.assertEqual(requirement.slide_count_policy.target, 8)
        self.assertEqual(requirement.language, "zh-CN")
        self.assertEqual(len(requirement.source_inputs), 1)
        self.assertEqual(requirement.source_inputs[0].name, "kid_words.md")
        self.assertEqual(requirement.source_understanding.document_type, "markdown")
        self.assertNotIn("business", requirement.style_requirement.style_keywords)
        self.assertIn("playful", requirement.style_requirement.style_keywords)
        self.assertIn("kid_friendly", requirement.style_requirement.style_keywords)
        self.assertIn("illustrated", requirement.style_requirement.style_keywords)

    def test_prepare_confirmed_requirement_extracts_public_topic_and_audience(self) -> None:
        manager = PptProductManager()

        requirement = manager.prepare_confirmed_requirement(
            task="给我做一个pptx，用于向大学文科学生科普ai",
            inputs=[],
            output={"format": "pptx"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)

        self.assertEqual(requirement.topic, "AI科普")
        self.assertEqual(requirement.audience, "大学文科学生")
        self.assertNotIn("给我", requirement.topic)
        self.assertNotIn("pptx", requirement.topic.lower())
        self.assertNotIn("用于", content_plan.pages[0].title)
        self.assertNotIn("给我做", content_plan.pages[0].title)
        self.assertNotIn("pptx", content_plan.pages[0].title.lower())

    def test_prepare_confirmed_requirement_detects_illustrated_kid_word_deck(self) -> None:
        manager = PptProductManager()

        requirement = manager.prepare_confirmed_requirement(
            task="给我做一个ppt，用来给幼儿园小朋友讲英语单词。图文并茂。小于10页。",
            inputs=[],
            output={"format": "pptx"},
        )

        self.assertEqual(requirement.topic, "英语单词")
        self.assertEqual(requirement.audience, "幼儿园小朋友")
        self.assertEqual(requirement.slide_count_policy.maximum, 9)
        self.assertLessEqual(requirement.slide_count_policy.target, 9)
        self.assertIn("illustrated", requirement.style_requirement.style_keywords)
        self.assertIn("kid_friendly", requirement.style_requirement.style_keywords)
        self.assertIn("playful", requirement.style_requirement.style_keywords)

    def test_content_plan_honors_exact_kindergarten_word_pages(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="给我做一个ppt，用来给幼儿园小朋友讲英语单词。3页，分别讲 猫、狗、鸭子。",
            inputs=[],
            output={"format": "pptx"},
        )

        plan = manager.build_initial_deck_content_plan(requirement)

        self.assertEqual(requirement.slide_count_policy.target, 3)
        self.assertEqual(len(plan.pages), 3)
        self.assertEqual([page.page_type for page in plan.pages], ["content", "content", "content"])
        self.assertEqual([page.title for page in plan.pages], ["Cat 猫", "Dog 狗", "Duck 鸭子"])
        self.assertNotIn("cover", {page.page_type for page in plan.pages})
        self.assertNotIn("toc", {page.page_type for page in plan.pages})
        self.assertNotIn("chapter_start", {page.page_type for page in plan.pages})

    def test_prepare_confirmed_requirement_cleans_orchestrator_style_task(self) -> None:
        manager = PptProductManager()

        requirement = manager.prepare_confirmed_requirement(
            task=(
                "制作一个面向大学文科学生的AI科普PPTX，语言为中文，风格清晰现代、适合课堂/讲座使用。"
                "内容需帮助非理工背景学生理解AI：AI是什么、发展简史、核心概念。"
            ),
            inputs=[],
            output={"format": "pptx"},
        )

        self.assertEqual(requirement.topic, "AI科普")
        self.assertEqual(requirement.audience, "大学文科学生")
        self.assertEqual(requirement.scenario, "课堂/讲座")

    def test_prepare_confirmed_requirement_extracts_topic_from_given_audience_phrase(self) -> None:
        manager = PptProductManager()

        requirement = manager.prepare_confirmed_requirement(
            task="基于上传材料做一个给大学文科学生的AI科普PPTX",
            inputs={
                "outline": [{"slide": 1, "title": "not a document"}],
                "documents": [{"name": "brief.md", "path": "input/brief.md", "mime_type": "text/markdown"}],
            },
            output={"format": "pptx", "slide_count": 8},
        )

        self.assertEqual(requirement.topic, "AI科普")
        self.assertEqual(requirement.audience, "大学文科学生")
        self.assertEqual(len(requirement.source_inputs), 1)
        self.assertEqual(requirement.slide_count_policy.target, 8)

    def test_prepare_confirmed_requirement_honors_explicit_route(self) -> None:
        manager = PptProductManager()

        requirement = manager.prepare_confirmed_requirement(
            task="套用用户上传 PPTX 模板生成汇报。",
            inputs=[{"name": "template.pptx", "path": "inbox/demo/template.pptx"}],
            output={"route": "xml"},
        )

        self.assertEqual(requirement.route, "xml")
        self.assertTrue(requirement.confirmed_by_user)
        self.assertTrue(requirement.template_requirement.use_template)
        self.assertEqual(requirement.template_requirement.template_source, "user")
        self.assertEqual(requirement.template_requirement.template_path, "inbox/demo/template.pptx")
        self.assertEqual(requirement.editability_requirement.level, "native")

    def test_prepare_confirmed_requirement_defaults_xml_for_powerpoint_input(self) -> None:
        manager = PptProductManager()

        requirement = manager.prepare_confirmed_requirement(
            task="用这个模版给我做一个ppt，用于宣传我的花店。",
            inputs=[{"name": "flower-template.pptx", "path": "inbox/demo/flower-template.pptx"}],
            output={"format": "pptx"},
        )

        self.assertEqual(requirement.route, "xml")
        self.assertFalse(requirement.confirmed_by_user)
        self.assertTrue(requirement.template_requirement.use_template)
        self.assertEqual(requirement.template_requirement.template_source, "user")
        self.assertEqual(requirement.template_requirement.template_path, "inbox/demo/flower-template.pptx")
        self.assertEqual(requirement.editability_requirement.level, "native")

        selection = manager._build_default_system_selection(requirement)

        self.assertEqual(selection["system_type"], "private_skill")
        self.assertEqual(selection["route"], "xml")
        self.assertEqual(selection["skill_name"], "pptx")
        self.assertEqual(selection["output_format"], "pptx")

    def test_prepare_confirmed_requirement_explicit_route_overrides_powerpoint_input(self) -> None:
        manager = PptProductManager()

        requirement = manager.prepare_confirmed_requirement(
            task="用这个模版给我做一个ppt，用于宣传我的花店。",
            inputs=[{"name": "flower-template.pptx", "path": "inbox/demo/flower-template.pptx"}],
            output={"format": "pptx", "route": "html"},
        )

        self.assertEqual(requirement.route, "html")
        self.assertTrue(requirement.confirmed_by_user)
        self.assertFalse(requirement.template_requirement.use_template)

        selection = manager._build_default_system_selection(requirement)

        self.assertEqual(selection["system_type"], "built_in_route")
        self.assertEqual(selection["route"], "html")
        self.assertEqual(selection["skill_name"], "")

    def test_default_system_selection_uses_pptx_skill_for_user_templates(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="套用用户上传 PPTX 模板生成汇报。",
            inputs=[{"name": "template.pptx", "path": "inbox/demo/template.pptx"}],
            output={"route": "xml"},
        )

        selection = manager._build_default_system_selection(requirement)

        self.assertEqual(selection["system_type"], "private_skill")
        self.assertEqual(selection["skill_name"], "pptx")
        self.assertEqual(selection["output_format"], "pptx")

    async def test_run_generates_html_route_outputs_and_writes_state(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={"sid": "ppt-manager-test", "turn_index": 1, "step": 1})

        result = await manager.run_product_request(
            task="生成一个 PPTX 产品介绍。",
            inputs=[],
            output={"format": "pptx", "auto_confirm": True},
            tool_context=tool_context,
        )

        self.assertEqual(result["status"], "success")
        self.assertEqual(result["result_schema_version"], "ppt-product-result-v1")
        self.assertEqual(result["product_line"], "ppt")
        self.assertEqual(result["selected_route"], "html")
        self.assertIn("ppt_confirmed_requirement", tool_context.state)
        self.assertIn("ppt_deck_content_plan", tool_context.state)
        self.assertIn("ppt_route_build", tool_context.state)
        self.assertEqual(
            tool_context.state[PPT_PRODUCT_REQUEST_STATE_KEY],
            {
                "task": "生成一个 PPTX 产品介绍。",
                "inputs": [],
                "output": {"format": "pptx", "auto_confirm": True},
            },
        )
        self.assertEqual(tool_context.state["product_line"], "ppt")
        self.assertEqual(tool_context.state["ppt_product_result"]["status"], "success")
        self.assertEqual(len(result["output_files"]), len(result["delivery_manifest"]["output_files"]))
        self.assertTrue(result["delivery_manifest"]["final_pptx"].endswith(".pptx"))
        self.assertEqual(tool_context.state["final_file_paths"], [result["delivery_manifest"]["final_pptx"]])

        pptx_path = resolve_workspace_path(result["delivery_manifest"]["final_pptx"])
        html_path = resolve_workspace_path(result["delivery_manifest"]["intermediate_artifacts"][0])
        self.assertTrue(pptx_path.exists())
        self.assertTrue(html_path.exists())
        self.assertGreater(len(result["delivery_manifest"]["previews"]), 0)
        self.assertEqual(len(Presentation(str(pptx_path)).slides), len(result["deck_content_plan"]["pages"]))

    async def test_run_generates_svg_route_outputs_and_writes_state(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={"sid": "ppt-manager-svg-test", "turn_index": 1, "step": 1})

        result = await manager.run_product_request(
            task="生成一个 3 页 PPTX 产品介绍，使用 SVG route。",
            inputs=[],
            output={"format": "pptx", "route": "svg", "slide_count": 3, "auto_confirm": True},
            tool_context=tool_context,
        )

        self.assertEqual(result["status"], "success")
        self.assertEqual(result["selected_route"], "svg")
        self.assertEqual(result["phase"], "svg_route_delivery")
        self.assertEqual(len(result["route_build"]["svg_page_paths"]), len(result["deck_content_plan"]["pages"]))
        self.assertEqual(result["delivery_manifest"]["intermediate_artifacts"], result["route_build"]["svg_page_paths"])
        self.assertTrue(result["delivery_manifest"]["final_pptx"].endswith(".pptx"))
        self.assertEqual(tool_context.state["final_file_paths"], [result["delivery_manifest"]["final_pptx"]])
        self.assertIn("ppt_design_strategy", tool_context.state)
        self.assertIn("ppt_svg_execution_plan", tool_context.state)

        pptx_path = resolve_workspace_path(result["delivery_manifest"]["final_pptx"])
        svg_path = resolve_workspace_path(result["route_build"]["svg_page_paths"][0])
        quality_path = resolve_workspace_path(result["delivery_manifest"]["quality_report"])
        self.assertTrue(pptx_path.exists())
        self.assertTrue(svg_path.exists())
        self.assertTrue(quality_path.exists())
        self.assertEqual(len(Presentation(str(pptx_path)).slides), len(result["deck_content_plan"]["pages"]))

    async def test_run_can_deliver_private_skill_html_when_selector_chooses_skill(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={"sid": "ppt-manager-private-skill-test", "turn_index": 1, "step": 1})

        def _content_plan_builder(_requirement):
            return DeckContentPlan(
                title="AI for Kids",
                core_narrative="Explain AI through familiar classroom examples.",
                pages=[
                    DeckPagePlan(
                        slide_number=1,
                        page_type="cover",
                        title="AI 是什么",
                        purpose="Introduce AI in simple language.",
                        key_takeaway="AI can help computers learn patterns.",
                        content_blocks=[{"items": ["AI 像一个会观察和练习的小助手。"]}],
                    ),
                    DeckPagePlan(
                        slide_number=2,
                        page_type="content",
                        title="AI 怎么学习",
                        purpose="Explain learning from examples.",
                        key_takeaway="Examples help AI get better.",
                        content_blocks=[{"items": ["看到很多图片后，AI 能学会分类。"]}],
                    ),
                ],
            )

        async def _selector(**_kwargs):
            return {
                "system_type": "private_skill",
                "route": "html",
                "skill_name": "ppt-complete-workflow",
                "output_format": "html",
                "reason": "Test selector chose the private PPT skill.",
            }

        result = await manager.run_product_request(
            task="给我做一个ppt，用来和小学生科普AI，不超过8页。用最合适的 skill 完成。",
            inputs=[],
            output={"format": "pptx", "auto_confirm": True},
            tool_context=tool_context,
            content_plan_builder=_content_plan_builder,
            system_selection_builder=_selector,
        )

        self.assertEqual(result["status"], "success")
        self.assertEqual(result["phase"], "private_skill_delivery")
        self.assertEqual(result["selected_route"], "html")
        self.assertEqual(result["delivery_manifest"]["final_pptx"], "")
        self.assertEqual(tool_context.state["ppt_system_selection"]["system_type"], "private_skill")
        self.assertEqual(tool_context.state["ppt_system_selection"]["skill_name"], "ppt-complete-workflow")
        self.assertEqual(tool_context.state["active_product_ppt_skill"]["name"], "ppt-complete-workflow")
        self.assertIn("PptProductManager skill runner", tool_context.state["ppt_private_skill_execution_output"]["message"])
        self.assertIn("final_file_paths", tool_context.state)
        self.assertEqual(tool_context.state["final_file_paths"], result["delivery_manifest"]["intermediate_artifacts"])

        html_path = resolve_workspace_path(result["delivery_manifest"]["intermediate_artifacts"][0])
        self.assertTrue(html_path.exists())
        html_text = html_path.read_text(encoding="utf-8")
        self.assertIn("AI 是什么", html_text)
        self.assertIn("ppt-complete-workflow", html_text)

    async def test_pptx_private_skill_without_runner_does_not_fallback_to_html(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(
            state={"sid": "pptx-private-skill-failure-test", "turn_index": 1, "step": 1}
        )
        requirement = manager.prepare_confirmed_requirement(
            task="用这个模板做一个花店宣传 PPT。",
            inputs=[{"name": "template.pptx", "path": "inbox/demo/template.pptx"}],
            output={"format": "pptx", "route": "xml"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)

        result = await manager.execute_private_ppt_skill(
            requirement=requirement,
            content_plan=content_plan,
            system_selection={
                "system_type": "private_skill",
                "route": "xml",
                "skill_name": "pptx",
                "output_format": "pptx",
                "reason": "Use pptx skill for uploaded template.",
            },
            tool_context=tool_context,
            expert_agents={},
            app_name="creative_claw",
            artifact_service=None,
        )

        self.assertEqual(result["status"], "error")
        self.assertEqual(result["output_path"], "")
        self.assertEqual(result["output_format"], "pptx")
        self.assertNotIn("final_file_paths", tool_context.state)
        self.assertIn("ppt_private_skill_runtime_path", tool_context.state)
        self.assertTrue(resolve_workspace_path(tool_context.state["ppt_private_skill_runtime_path"]).exists())

    async def test_private_skill_execution_phase_records_typed_result(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(
            state={"sid": "private-skill-execution-phase-test", "turn_index": 1, "step": 1}
        )
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 HTML private skill PPT。",
            inputs=[],
            output={"format": "html", "route": "html"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)

        result = await manager._execute_private_ppt_skill_phase(
            requirement=requirement,
            content_plan=content_plan,
            system_selection={
                "system_type": "private_skill",
                "route": "html",
                "skill_name": "ppt-complete-workflow",
                "output_format": "html",
            },
            tool_context=tool_context,
            expert_agents={},
            app_name="creative_claw",
            artifact_service=None,
        )

        self.assertIsInstance(result, PptPrivateSkillExecutionResult)
        self.assertEqual(result.skill_name, "ppt-complete-workflow")
        self.assertEqual(result.output_format, "html")
        self.assertEqual(result.execution_output["source"], "deterministic_fallback")
        self.assertEqual(result.private_build["source"], "save_ppt_private_skill_html")
        self.assertEqual(
            tool_context.state[PPT_PRIVATE_SKILL_EXECUTION_RESULT_STATE_KEY]["private_build"]["output_path"],
            result.private_build["output_path"],
        )

    async def test_private_skill_execution_phase_reuses_same_input_successful_build(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(
            state={"sid": "private-skill-reuse-test", "turn_index": 1, "step": 1}
        )
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 HTML private skill PPT。",
            inputs=[],
            output={"format": "html", "route": "html"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        selection = {
            "system_type": "private_skill",
            "route": "html",
            "skill_name": "ppt-complete-workflow",
            "output_format": "html",
        }
        execution_calls: list[dict[str, Any]] = []

        async def _execute_private_skill(self, **kwargs):
            execution_calls.append(kwargs)
            return self.save_ppt_private_skill_html(
                file_name="reused-private-skill.html",
                html_content="<!doctype html><html><body><h1>Reusable private skill</h1></body></html>",
                description="Reusable private-skill HTML deck.",
                tool_context=kwargs["tool_context"],
            )

        async def _unexpected_execute_private_skill(*_args, **_kwargs):
            raise AssertionError("same private-skill input should be reused")

        with patch.object(PptProductManager, "execute_private_ppt_skill", _execute_private_skill):
            first_result = await manager._execute_private_ppt_skill_phase(
                requirement=requirement,
                content_plan=content_plan,
                system_selection=selection,
                tool_context=tool_context,
                expert_agents={},
                app_name="creative_claw",
                artifact_service=None,
            )

        with patch.object(PptProductManager, "execute_private_ppt_skill", _unexpected_execute_private_skill):
            second_result = await manager._execute_private_ppt_skill_phase(
                requirement=requirement,
                content_plan=content_plan,
                system_selection=selection,
                tool_context=tool_context,
                expert_agents={},
                app_name="creative_claw",
                artifact_service=None,
            )

        self.assertEqual(len(execution_calls), 1)
        self.assertFalse(first_result.reused_existing_build)
        self.assertTrue(second_result.reused_existing_build)
        self.assertEqual(first_result.input_signature, second_result.input_signature)
        self.assertEqual(first_result.private_build["output_path"], second_result.private_build["output_path"])
        self.assertTrue(resolve_workspace_path(second_result.private_build["output_path"]).is_file())
        self.assertEqual(tool_context.state["final_file_paths"], [second_result.private_build["output_path"]])
        self.assertTrue(
            tool_context.state[PPT_PRIVATE_SKILL_EXECUTION_RESULT_STATE_KEY]["reused_existing_build"]
        )

    async def test_private_skill_execution_phase_reruns_when_inputs_change(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(
            state={"sid": "private-skill-rerun-test", "turn_index": 1, "step": 1}
        )
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 HTML private skill PPT。",
            inputs=[],
            output={"format": "html", "route": "html"},
        )
        first_plan = manager.build_initial_deck_content_plan(requirement)
        second_plan = DeckContentPlan(
            title="Changed Private Skill Plan",
            core_narrative="Changed plan should force private skill execution.",
            pages=[_page(1, "cover"), _page(2, "ending")],
        )
        selection = {
            "system_type": "private_skill",
            "route": "html",
            "skill_name": "ppt-complete-workflow",
            "output_format": "html",
        }
        execution_calls: list[str] = []

        async def _execute_private_skill(self, **kwargs):
            file_name = f"private-skill-rerun-{len(execution_calls) + 1}.html"
            execution_calls.append(file_name)
            return self.save_ppt_private_skill_html(
                file_name=file_name,
                html_content=f"<!doctype html><html><body><h1>{file_name}</h1></body></html>",
                description="Rerun private-skill HTML deck.",
                tool_context=kwargs["tool_context"],
            )

        with patch.object(PptProductManager, "execute_private_ppt_skill", _execute_private_skill):
            first_result = await manager._execute_private_ppt_skill_phase(
                requirement=requirement,
                content_plan=first_plan,
                system_selection=selection,
                tool_context=tool_context,
                expert_agents={},
                app_name="creative_claw",
                artifact_service=None,
            )
            second_result = await manager._execute_private_ppt_skill_phase(
                requirement=requirement,
                content_plan=second_plan,
                system_selection=selection,
                tool_context=tool_context,
                expert_agents={},
                app_name="creative_claw",
                artifact_service=None,
            )

        self.assertEqual(execution_calls, ["private-skill-rerun-1.html", "private-skill-rerun-2.html"])
        self.assertFalse(first_result.reused_existing_build)
        self.assertFalse(second_result.reused_existing_build)
        self.assertNotEqual(first_result.input_signature, second_result.input_signature)
        self.assertTrue(second_result.private_build["output_path"].endswith("private-skill-rerun-2.html"))
        self.assertEqual(
            tool_context.state[PPT_PRIVATE_SKILL_EXECUTION_RESULT_STATE_KEY]["input_signature"],
            second_result.input_signature,
        )

    async def test_private_skill_execution_phase_uses_adk_workflow_context(self) -> None:
        fake_llm = _PptProductManagerToolCallingFakeLlm(
            function_calls=[
                FunctionCall(
                    name="save_ppt_private_skill_html",
                    args={
                        "file_name": "agenttool-private-skill-workflow.html",
                        "html_content": "<!doctype html><html><body><h1>Workflow Skill Deck</h1></body></html>",
                        "description": "Workflow private skill HTML deck.",
                    },
                )
            ],
            final_text="PptProductManager saved the workflow private skill HTML deck.",
        )
        manager = PptProductManager(model=fake_llm)
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 HTML private skill PPT。",
            inputs=[],
            output={"format": "html", "route": "html"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        selection = {
            "system_type": "private_skill",
            "route": "html",
            "skill_name": "ppt-complete-workflow",
            "output_format": "html",
            "reason": "Exercise the private-skill execution workflow node.",
        }

        async def _run_execution_harness():
            @node(name="PptPrivateSkillExecutionWorkflowHarnessNode", rerun_on_resume=True)
            async def execution_harness(ctx: Context, node_input: str) -> dict[str, Any]:
                result = await manager._run_private_skill_execution_phase(
                    requirement=requirement,
                    content_plan=content_plan,
                    system_selection=selection,
                    tool_context=ctx,
                    expert_agents={},
                    app_name="creative_claw",
                    artifact_service=artifact_service,
                )
                return result.model_dump(mode="json")

            workflow = Workflow(
                name="PptPrivateSkillExecutionWorkflowHarness",
                edges=[("START", execution_harness)],
            )
            session_service = InMemorySessionService()
            runner = Runner(
                node=workflow,
                session_service=session_service,
                artifact_service=artifact_service,
            )
            user_id = "user-ppt-private-skill-execution-workflow"
            session_id = "session-ppt-private-skill-execution-workflow"
            try:
                await session_service.create_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                    state={
                        "sid": "ppt-private-skill-execution-workflow-test",
                        "turn_index": 1,
                        "step": 3,
                    },
                )
                async for _ in runner.run_async(
                    user_id=user_id,
                    session_id=session_id,
                    new_message=Content(role="user", parts=[Part(text="Execute private skill")]),
                ):
                    pass
                return await session_service.get_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                )
            finally:
                await runner.close()

        artifact_service = InMemoryArtifactService()
        session = await _run_execution_harness()

        self.assertIsNotNone(session)
        self.assertEqual(
            session.state[PPT_PRIVATE_SKILL_EXECUTION_WORKFLOW_OUTPUT_KEY]["source"],
            "adk_workflow",
        )
        self.assertEqual(
            session.state[PPT_PRIVATE_SKILL_EXECUTION_RESULT_STATE_KEY]["execution_output"]["source"],
            "ppt_product_manager",
        )
        self.assertEqual(
            session.state[PPT_PRIVATE_SKILL_EXECUTION_RESULT_STATE_KEY]["private_build"]["source"],
            "save_ppt_private_skill_html",
        )
        self.assertNotEqual(
            session.state["ppt_private_skill_execution_output"]["source"],
            "deterministic_fallback",
        )
        output_path = session.state[PPT_PRIVATE_SKILL_EXECUTION_RESULT_STATE_KEY]["private_build"]["output_path"]
        self.assertEqual(session.state["final_file_paths"], [output_path])
        html_path = resolve_workspace_path(output_path)
        self.assertTrue(html_path.exists())
        self.assertIn("Workflow Skill Deck", html_path.read_text(encoding="utf-8"))
        self.assertIn("save_ppt_private_skill_html", _function_declaration_names(fake_llm.requests[0]))

    async def test_private_skill_agenttool_main_path_saves_html(self) -> None:
        fake_llm = _PptProductManagerToolCallingFakeLlm(
            function_calls=[
                FunctionCall(
                    name="save_ppt_private_skill_html",
                    args={
                        "file_name": "agenttool-private-skill.html",
                        "html_content": "<!doctype html><html><body><h1>AgentTool Skill Deck</h1></body></html>",
                        "description": "AgentTool private skill HTML deck.",
                    },
                )
            ],
            final_text="PptProductManager saved the private skill HTML deck.",
        )
        manager = PptProductManager(model=fake_llm)
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 HTML private skill PPT。",
            inputs=[],
            output={"format": "html", "route": "html"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        selection = {
            "system_type": "private_skill",
            "route": "html",
            "skill_name": "ppt-complete-workflow",
            "output_format": "html",
            "reason": "Exercise the AgentTool private skill transport.",
        }
        session_service = InMemorySessionService()
        artifact_service = InMemoryArtifactService()

        @node(name="PptPrivateSkillAgentToolHarnessNode", rerun_on_resume=True)
        async def private_skill_harness(ctx: Context, node_input: str) -> dict:
            return await manager.execute_private_ppt_skill(
                requirement=requirement,
                content_plan=content_plan,
                system_selection=selection,
                tool_context=ctx,
                expert_agents={},
                app_name="creative_claw",
                artifact_service=artifact_service,
            )

        workflow = Workflow(
            name="PptPrivateSkillAgentToolHarness",
            edges=[("START", private_skill_harness)],
        )
        runner = Runner(
            node=workflow,
            session_service=session_service,
            artifact_service=artifact_service,
        )
        user_id = "user-ppt-private-skill-agenttool"
        session_id = "session-ppt-private-skill-agenttool"

        try:
            await session_service.create_session(
                app_name=workflow.name,
                user_id=user_id,
                session_id=session_id,
                state={
                    "sid": "ppt-private-skill-agenttool-test",
                    "turn_index": 1,
                    "step": 3,
                },
            )
            async for _ in runner.run_async(
                user_id=user_id,
                session_id=session_id,
                new_message=Content(role="user", parts=[Part(text="Run private PPT skill")]),
            ):
                pass
            session = await session_service.get_session(
                app_name=workflow.name,
                user_id=user_id,
                session_id=session_id,
            )
        finally:
            await runner.close()

        self.assertIsNotNone(session)
        private_build = session.state["ppt_private_skill_build"]
        self.assertEqual(private_build["status"], "success")
        self.assertEqual(private_build["source"], "save_ppt_private_skill_html")
        self.assertEqual(
            session.state["ppt_private_skill_execution_output"]["source"],
            "ppt_product_manager",
        )
        self.assertNotEqual(
            session.state["ppt_private_skill_execution_output"]["source"],
            "deterministic_fallback",
        )
        self.assertEqual(session.state["final_file_paths"], [private_build["output_path"]])
        html_path = resolve_workspace_path(private_build["output_path"])
        self.assertTrue(html_path.exists())
        self.assertIn("AgentTool Skill Deck", html_path.read_text(encoding="utf-8"))
        first_request_tools = _function_declaration_names(fake_llm.requests[0])
        self.assertIn("save_ppt_private_skill_html", first_request_tools)
        self.assertIn("invoke_ppt_expert", first_request_tools)
        self.assertGreaterEqual(len(fake_llm.requests), 2)

    async def test_private_skill_delivery_phase_uses_adk_workflow_context(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="做一个 HTML private skill PPT。",
            inputs=[],
            output={"format": "html", "route": "html"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)
        output_path = "generated/ppt-private-skill-delivery-workflow-test/turn_1/index.html"
        file_record = build_workspace_file_record(
            output_path,
            description="Private PPT skill HTML deck artifact.",
            source="ppt_product_manager",
        )
        private_build = {
            "skill_name": "ppt-complete-workflow",
            "output_path": output_path,
            "artifact_type": "html",
            "output_format": "html",
            "output_files": [file_record],
            "source": "save_ppt_private_skill_html",
        }
        selection = {
            "system_type": "private_skill",
            "route": "html",
            "skill_name": "ppt-complete-workflow",
            "output_format": "html",
        }

        async def _run_delivery_harness():
            @node(name="PptPrivateSkillDeliveryWorkflowHarnessNode", rerun_on_resume=True)
            async def delivery_harness(ctx: Context, node_input: str) -> dict[str, Any]:
                result = await manager._run_private_skill_delivery_phase(
                    requirement=requirement,
                    content_plan=content_plan,
                    system_selection=selection,
                    private_build=private_build,
                    tool_context=ctx,
                )
                return result.model_dump(mode="json")

            workflow = Workflow(
                name="PptPrivateSkillDeliveryWorkflowHarness",
                edges=[("START", delivery_harness)],
            )
            session_service = InMemorySessionService()
            artifact_service = InMemoryArtifactService()
            runner = Runner(
                node=workflow,
                session_service=session_service,
                artifact_service=artifact_service,
            )
            user_id = "user-ppt-private-skill-delivery-workflow"
            session_id = "session-ppt-private-skill-delivery-workflow"
            try:
                await session_service.create_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                    state={
                        "sid": "ppt-private-skill-delivery-workflow-test",
                        "turn_index": 1,
                        "step": 1,
                    },
                )
                async for _ in runner.run_async(
                    user_id=user_id,
                    session_id=session_id,
                    new_message=Content(role="user", parts=[Part(text="Deliver private skill")]),
                ):
                    pass
                return await session_service.get_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                )
            finally:
                await runner.close()

        session = await _run_delivery_harness()

        self.assertIsNotNone(session)
        self.assertEqual(
            session.state[PPT_PRIVATE_SKILL_DELIVERY_WORKFLOW_OUTPUT_KEY]["source"],
            "adk_workflow",
        )
        self.assertEqual(
            session.state["ppt_private_skill_delivery_result"]["product_result"]["status"],
            "success",
        )
        self.assertEqual(
            session.state["ppt_private_skill_delivery_result"]["private_build"]["output_path"],
            output_path,
        )

    async def test_private_skill_recovers_generated_pptx_after_child_runner_error(self) -> None:
        manager = PptProductManager()
        captured: dict[str, object] = {}
        tool_context = SimpleNamespace(
            state=_DictState({
                "sid": "pptx-private-skill-child-error-test",
                "turn_index": 1,
                "step": 3,
            }),
            _invocation_context=SimpleNamespace(
                app_name="creative_claw",
                user_id="test-user",
                plugin_manager=SimpleNamespace(plugins=[]),
                credential_service=None,
            ),
        )
        requirement = manager.prepare_confirmed_requirement(
            task="用这个模板做一个品牌故事 PPT。",
            inputs=[{"name": "template.pptx", "path": "inbox/demo/template.pptx"}],
            output={"format": "pptx", "route": "xml"},
        )
        content_plan = manager.build_initial_deck_content_plan(requirement)

        async def _failing_agent_tool_transport(**kwargs):
            captured.update(kwargs)
            output_dir = resolve_workspace_path(tool_context.state["ppt_private_skill_output_dir"]) / "output"
            output_dir.mkdir(parents=True, exist_ok=True)
            with tempfile.NamedTemporaryFile(
                prefix="deck_after_error_",
                suffix=".pptx",
                dir=output_dir,
                delete=False,
            ) as handle:
                pptx_name = handle.name
            Path(pptx_name).unlink()
            deck = Presentation()
            deck.slides.add_slide(deck.slide_layouts[0])
            deck.save(pptx_name)
            raise ValueError("optional QA failed after PPTX generation")

        with patch(
            "src.productions.ppt.ppt_product_manager.ppt_product_manager._run_ppt_internal_agent_tool",
            _failing_agent_tool_transport,
        ):
            result = await manager.execute_private_ppt_skill(
                requirement=requirement,
                content_plan=content_plan,
                system_selection={
                    "system_type": "private_skill",
                    "route": "xml",
                    "skill_name": "pptx",
                    "output_format": "pptx",
                    "reason": "Use pptx skill for uploaded template.",
                },
                tool_context=tool_context,
                expert_agents={},
                app_name="creative_claw",
                artifact_service=None,
            )

        self.assertEqual(result["status"], "success")
        self.assertEqual(result["artifact_type"], "pptx")
        self.assertEqual(result["source"], "private_skill_pptx_recovery")
        self.assertIn("optional QA failed", result["execution_warning"])
        self.assertEqual(tool_context.state["final_file_paths"], [result["pptx_path"]])
        self.assertEqual(
            tool_context.state["ppt_private_skill_execution_output"]["status"],
            "success_with_warning",
        )
        self.assertTrue(resolve_workspace_path(result["pptx_path"]).exists())
        self.assertIn("ppt_product_manager_skill_run_base", captured["initial_state"])

    async def test_interactive_workflow_pauses_for_two_confirmations(self) -> None:
        image_path = _write_test_image("interactive_kid_word_asset.png")
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={"sid": "ppt-manager-interactive-test", "turn_index": 1, "step": 1})
        resolved_assets: list[str] = []

        async def _asset_resolver(asset, _page, _requirement):
            resolved_assets.append(asset.asset_id)
            return {
                "asset_id": asset.asset_id,
                "status": "ready",
                "path": image_path,
                "provider": "test_resolver",
            }

        requirement_result = await manager.run_product_request(
            task="给我做一个ppt，用来给幼儿园小朋友讲英语单词。3页，分别讲 猫、狗、鸭子。",
            inputs=[],
            output={"format": "pptx"},
            tool_context=tool_context,
            asset_resolver=_asset_resolver,
        )

        self.assertEqual(requirement_result["status"], "awaiting_requirement_confirmation")
        self.assertEqual(tool_context.state["ppt_workflow_state"]["stage"], "awaiting_requirement_confirmation")
        self.assertNotIn("final_file_paths", tool_context.state)
        self.assertIn("summary_markdown", requirement_result["confirmation_request"])
        self.assertIn("### 系统选择", requirement_result["confirmation_request"]["summary_markdown"])
        self.assertEqual(resolved_assets, [])
        self.assertEqual(tool_context.state["ppt_workflow_state"]["waiting_since_turn_index"], 1)

        same_turn_requirement_result = await manager.continue_product_request(
            user_response="确认",
            tool_context=tool_context,
            asset_resolver=_asset_resolver,
        )

        self.assertEqual(same_turn_requirement_result["status"], "awaiting_requirement_confirmation")
        self.assertEqual(tool_context.state["ppt_workflow_state"]["stage"], "awaiting_requirement_confirmation")
        self.assertEqual(resolved_assets, [])

        tool_context.state["turn_index"] = 2
        plan_result = await manager.continue_product_request(
            user_response="确认",
            tool_context=tool_context,
            asset_resolver=_asset_resolver,
        )

        self.assertEqual(plan_result["status"], "awaiting_content_plan_confirmation")
        self.assertEqual(tool_context.state["ppt_workflow_state"]["stage"], "awaiting_content_plan_confirmation")
        self.assertEqual([page["title"] for page in plan_result["deck_content_plan"]["pages"]], ["Cat 猫", "Dog 狗", "Duck 鸭子"])
        self.assertNotIn("### 系统选择", plan_result["confirmation_request"]["summary_markdown"])
        self.assertEqual(resolved_assets, [])
        self.assertEqual(tool_context.state["ppt_workflow_state"]["waiting_since_turn_index"], 2)

        same_turn_plan_result = await manager.continue_product_request(
            user_response="确认",
            tool_context=tool_context,
            asset_resolver=_asset_resolver,
        )

        self.assertEqual(same_turn_plan_result["status"], "awaiting_content_plan_confirmation")
        self.assertEqual(tool_context.state["ppt_workflow_state"]["stage"], "awaiting_content_plan_confirmation")
        self.assertEqual(resolved_assets, [])
        self.assertNotIn("final_file_paths", tool_context.state)

        tool_context.state["turn_index"] = 3
        final_result = await manager.continue_product_request(
            user_response="确认",
            tool_context=tool_context,
            asset_resolver=_asset_resolver,
        )

        self.assertEqual(final_result["status"], "success")
        self.assertEqual(tool_context.state["ppt_workflow_state"]["stage"], "completed")
        self.assertGreaterEqual(len(resolved_assets), 1)
        self.assertTrue(final_result["delivery_manifest"]["final_pptx"].endswith(".pptx"))
        self.assertEqual(tool_context.state["final_file_paths"], [final_result["delivery_manifest"]["final_pptx"]])

    async def test_requirement_confirmation_uses_adk_workflow_context(self) -> None:
        manager = PptProductManager()
        task = "给幼儿园小朋友做一个 3 页英语单词 PPT。"
        output = {"format": "pptx", "slide_count": 3}
        requirement = manager.prepare_confirmed_requirement(task=task, inputs=[], output=output)
        system_selection = manager._build_default_system_selection(requirement)

        async def _run_ppt_harness():
            @node(name="PptRequirementConfirmationWorkflowHarnessNode", rerun_on_resume=True)
            async def ppt_harness(ctx: Context, node_input: str) -> dict[str, Any]:
                return await manager.continue_product_request(
                    user_response=node_input,
                    tool_context=ctx,
                    content_plan_builder=manager.build_initial_deck_content_plan,
                )

            workflow = Workflow(
                name="PptRequirementConfirmationWorkflowHarness",
                edges=[("START", ppt_harness)],
            )
            session_service = InMemorySessionService()
            artifact_service = InMemoryArtifactService()
            runner = Runner(
                node=workflow,
                session_service=session_service,
                artifact_service=artifact_service,
            )
            user_id = "user-ppt-requirement-confirmation-workflow"
            session_id = "session-ppt-requirement-confirmation-workflow"
            try:
                await session_service.create_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                    state={
                        "sid": "ppt-requirement-confirmation-workflow-test",
                        "turn_index": 2,
                        "step": 1,
                        "ppt_workflow_state": {
                            "workflow_id": "ppt-requirement-confirmation-workflow-test:ppt:1",
                            "stage": "awaiting_requirement_confirmation",
                            "revision": 1,
                            "task": task,
                            "raw_inputs": [],
                            "output": output,
                            "confirmed_requirement": requirement.model_dump(mode="json"),
                            "system_selection": system_selection,
                            "waiting_since_turn_index": 1,
                        },
                    },
                )
                async for _ in runner.run_async(
                    user_id=user_id,
                    session_id=session_id,
                    new_message=Content(role="user", parts=[Part(text="确认")]),
                ):
                    pass
                return await session_service.get_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                )
            finally:
                await runner.close()

        session = await _run_ppt_harness()

        self.assertIsNotNone(session)
        self.assertEqual(session.state["ppt_product_result"]["status"], "awaiting_content_plan_confirmation")
        self.assertEqual(session.state["ppt_workflow_state"]["stage"], "awaiting_content_plan_confirmation")
        self.assertEqual(
            session.state[PPT_REQUIREMENT_CONFIRMATION_WORKFLOW_OUTPUT_KEY]["source"],
            "adk_workflow",
        )
        self.assertEqual(
            session.state[PPT_REQUIREMENT_CONFIRMATION_WORKFLOW_OUTPUT_KEY]["status"],
            "awaiting_content_plan_confirmation",
        )
        self.assertEqual(session.state[PPT_SOURCE_PREPARATION_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_CONTENT_PLANNING_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertIn("deck_content_plan", session.state["ppt_workflow_state"])
        self.assertNotIn("final_file_paths", session.state)

    async def test_content_plan_confirmation_uses_adk_workflow_context(self) -> None:
        manager = PptProductManager()
        task = "给幼儿园小朋友做一个 3 页英语单词 PPT。"
        output = {"format": "pptx", "slide_count": 3}
        requirement = manager.prepare_confirmed_requirement(task=task, inputs=[], output=output)
        content_plan = manager.build_initial_deck_content_plan(requirement)
        system_selection = manager._build_default_system_selection(requirement)
        image_path = _write_test_image("content_plan_workflow_asset.png")
        final_pptx = "generated/ppt-content-plan-workflow-test/turn_3/final.pptx"
        route_calls: list[dict[str, Any]] = []

        async def _asset_resolver(asset, _page, _requirement):
            return {
                "asset_id": asset.asset_id,
                "status": "ready",
                "path": image_path,
                "provider": "test_resolver",
            }

        async def _fake_dispatch_ppt_route(**kwargs):
            route_calls.append(kwargs)
            return HtmlRouteBuildPackage(
                template=HtmlTemplatePackage(template_id="workflow-test", label="Workflow Test"),
                html_deck_path="generated/ppt-content-plan-workflow-test/turn_3/deck.html",
                preview_paths=["generated/ppt-content-plan-workflow-test/turn_3/preview.png"],
                pptx_path=final_pptx,
                quality_report_path="generated/ppt-content-plan-workflow-test/turn_3/quality.json",
                build_log_path="generated/ppt-content-plan-workflow-test/turn_3/build.log",
                warnings=[],
            )

        async def _run_ppt_harness():
            @node(name="PptContentPlanConfirmationWorkflowHarnessNode", rerun_on_resume=True)
            async def ppt_harness(ctx: Context, node_input: str) -> dict[str, Any]:
                return await manager.continue_product_request(
                    user_response=node_input,
                    tool_context=ctx,
                    asset_resolver=_asset_resolver,
                )

            workflow = Workflow(
                name="PptContentPlanConfirmationWorkflowHarness",
                edges=[("START", ppt_harness)],
            )
            session_service = InMemorySessionService()
            artifact_service = InMemoryArtifactService()
            runner = Runner(
                node=workflow,
                session_service=session_service,
                artifact_service=artifact_service,
            )
            user_id = "user-ppt-content-plan-confirmation-workflow"
            session_id = "session-ppt-content-plan-confirmation-workflow"
            try:
                await session_service.create_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                    state={
                        "sid": "ppt-content-plan-workflow-test",
                        "turn_index": 3,
                        "step": 1,
                        "ppt_workflow_state": {
                            "workflow_id": "ppt-content-plan-workflow-test:ppt:1",
                            "stage": "awaiting_content_plan_confirmation",
                            "revision": 2,
                            "task": task,
                            "raw_inputs": [],
                            "output": output,
                            "confirmed_requirement": requirement.model_dump(mode="json"),
                            "deck_content_plan": content_plan.model_dump(mode="json"),
                            "deck_content_plan_markdown": "",
                            "system_selection": system_selection,
                            "waiting_since_turn_index": 2,
                        },
                    },
                )
                with patch.object(manager, "_dispatch_ppt_route", _fake_dispatch_ppt_route):
                    async for _ in runner.run_async(
                        user_id=user_id,
                        session_id=session_id,
                        new_message=Content(role="user", parts=[Part(text="确认")]),
                    ):
                        pass
                return await session_service.get_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                )
            finally:
                await runner.close()

        session = await _run_ppt_harness()

        self.assertIsNotNone(session)
        self.assertEqual(session.state["ppt_product_result"]["status"], "success")
        self.assertEqual(session.state["ppt_workflow_state"]["stage"], "completed")
        self.assertEqual(
            session.state[PPT_CONTENT_PLAN_CONFIRMATION_WORKFLOW_OUTPUT_KEY]["source"],
            "adk_workflow",
        )
        self.assertEqual(
            session.state[PPT_CONTENT_PLAN_CONFIRMATION_WORKFLOW_OUTPUT_KEY]["status"],
            "success",
        )
        self.assertEqual(session.state[PPT_ASSET_RESOLUTION_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_ROUTE_EXECUTION_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state[PPT_FINAL_DELIVERY_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(session.state["final_file_paths"], [final_pptx])
        self.assertEqual(len(route_calls), 1)
        self.assertEqual(route_calls[0]["content_plan"].pages[0].assets[0].path, image_path)

    async def test_requirement_revision_branch_uses_parent_adk_workflow_context(self) -> None:
        manager = PptProductManager()
        task = "给团队做一个 5 页产品复盘 PPT。"
        output = {"format": "pptx", "slide_count": 5}
        requirement = manager.prepare_confirmed_requirement(task=task, inputs=[], output=output)
        system_selection = manager._build_default_system_selection(requirement)
        revised_payload = requirement.model_dump(mode="json")
        revised_payload.update(
            {
                "audience": "技术负责人",
                "scenario": "季度复盘",
                "slide_count_policy": {
                    **revised_payload["slide_count_policy"],
                    "target": 12,
                },
            }
        )
        selection_payload = {
            "system_type": "built_in_route",
            "route": "html",
            "skill_name": "",
            "output_format": "pptx",
            "reason": "Keep the built-in HTML route after requirement revision.",
        }
        fake_llm = _PptProductManagerToolCallingFakeLlm(
            function_calls=[
                FunctionCall(
                    name="save_ppt_confirmed_requirement_json",
                    args={"requirement_json": revised_payload},
                ),
                FunctionCall(name="list_product_ppt_skills", args={}),
                FunctionCall(
                    name="save_ppt_system_selection",
                    args={"selection_json": selection_payload},
                ),
            ],
            final_text="PPT parent workflow revision branch completed.",
        )

        async def _run_ppt_harness():
            @node(name="PptRequirementRevisionParentWorkflowHarnessNode", rerun_on_resume=True)
            async def ppt_harness(ctx: Context, node_input: str) -> dict[str, Any]:
                with patch.object(
                    manager,
                    "_continue_after_requirement_confirmation_direct",
                    side_effect=AssertionError("parent workflow must own requirement revision"),
                ):
                    return await manager.continue_product_request(
                        user_response=node_input,
                        tool_context=ctx,
                    )

            workflow = Workflow(
                name="PptRequirementRevisionParentWorkflowHarness",
                edges=[("START", ppt_harness)],
            )
            session_service = InMemorySessionService()
            artifact_service = InMemoryArtifactService()
            runner = Runner(
                node=workflow,
                session_service=session_service,
                artifact_service=artifact_service,
            )
            user_id = "user-ppt-requirement-revision-parent-workflow"
            session_id = "session-ppt-requirement-revision-parent-workflow"
            try:
                await session_service.create_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                    state={
                        "sid": "ppt-requirement-revision-parent-workflow-test",
                        "turn_index": 2,
                        "step": 1,
                        "ppt_workflow_state": {
                            "workflow_id": "ppt-requirement-revision-parent-workflow-test:ppt:1",
                            "stage": "awaiting_requirement_confirmation",
                            "revision": 1,
                            "task": task,
                            "raw_inputs": [],
                            "output": output,
                            "confirmed_requirement": requirement.model_dump(mode="json"),
                            "system_selection": system_selection,
                            "waiting_since_turn_index": 1,
                        },
                    },
                )
                with patch(
                    "src.productions.ppt.ppt_product_manager.ppt_product_manager.build_llm",
                    return_value=fake_llm,
                ):
                    async for _ in runner.run_async(
                        user_id=user_id,
                        session_id=session_id,
                        new_message=Content(
                            role="user",
                            parts=[Part(text="受众改为技术负责人，场景改成季度复盘，页数改为 12 页。")],
                        ),
                    ):
                        pass
                return await session_service.get_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                )
            finally:
                await runner.close()

        session = await _run_ppt_harness()

        self.assertIsNotNone(session)
        self.assertEqual(session.state["ppt_product_result"]["status"], "awaiting_requirement_confirmation")
        self.assertEqual(session.state["ppt_workflow_state"]["stage"], "awaiting_requirement_confirmation")
        self.assertEqual(
            session.state[PPT_REQUIREMENT_CONFIRMATION_WORKFLOW_OUTPUT_KEY]["branch"],
            "revision",
        )
        self.assertEqual(
            session.state[PPT_REQUIREMENT_REVISION_WORKFLOW_OUTPUT_KEY]["source"],
            "adk_workflow",
        )
        self.assertEqual(
            session.state[PPT_SYSTEM_SELECTION_WORKFLOW_OUTPUT_KEY]["source"],
            "adk_workflow",
        )
        self.assertEqual(session.state[PPT_CONFIRMED_REQUIREMENT_STATE_KEY]["audience"], "技术负责人")
        self.assertEqual(session.state[PPT_CONFIRMED_REQUIREMENT_STATE_KEY]["scenario"], "季度复盘")
        self.assertEqual(session.state[PPT_CONFIRMED_REQUIREMENT_STATE_KEY]["slide_count_policy"]["target"], 12)
        self.assertNotIn(PPT_SOURCE_PREPARATION_WORKFLOW_OUTPUT_KEY, session.state)
        first_request_tools = _function_declaration_names(fake_llm.requests[0])
        self.assertIn("save_ppt_confirmed_requirement_json", first_request_tools)

    async def test_content_plan_revision_branch_uses_parent_adk_workflow_context(self) -> None:
        manager = PptProductManager()
        task = "给幼儿园小朋友做一个 3 页英语单词 PPT。"
        output = {"format": "pptx", "slide_count": 3}
        requirement = manager.prepare_confirmed_requirement(task=task, inputs=[], output=output)
        content_plan = manager.build_initial_deck_content_plan(requirement)
        system_selection = manager._build_default_system_selection(requirement)
        captured: dict[str, Any] = {}

        def _content_plan_builder(revised_requirement):
            captured["request_brief"] = revised_requirement.request_brief
            return DeckContentPlan(
                title="Parent Workflow Revised Plan",
                core_narrative="Use the parent Workflow revision branch.",
                pages=[_page(1, "cover"), _page(2, "content"), _page(3, "ending")],
            )

        async def _run_ppt_harness():
            @node(name="PptContentPlanRevisionParentWorkflowHarnessNode", rerun_on_resume=True)
            async def ppt_harness(ctx: Context, node_input: str) -> dict[str, Any]:
                with patch.object(
                    manager,
                    "_continue_after_content_plan_confirmation_direct",
                    side_effect=AssertionError("parent workflow must own content-plan revision"),
                ):
                    return await manager.continue_product_request(
                        user_response=node_input,
                        tool_context=ctx,
                        content_plan_builder=_content_plan_builder,
                    )

            workflow = Workflow(
                name="PptContentPlanRevisionParentWorkflowHarness",
                edges=[("START", ppt_harness)],
            )
            session_service = InMemorySessionService()
            artifact_service = InMemoryArtifactService()
            runner = Runner(
                node=workflow,
                session_service=session_service,
                artifact_service=artifact_service,
            )
            user_id = "user-ppt-content-plan-revision-parent-workflow"
            session_id = "session-ppt-content-plan-revision-parent-workflow"
            try:
                await session_service.create_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                    state={
                        "sid": "ppt-content-plan-revision-parent-workflow-test",
                        "turn_index": 3,
                        "step": 1,
                        "ppt_workflow_state": {
                            "workflow_id": "ppt-content-plan-revision-parent-workflow-test:ppt:1",
                            "stage": "awaiting_content_plan_confirmation",
                            "revision": 2,
                            "task": task,
                            "raw_inputs": [],
                            "output": output,
                            "confirmed_requirement": requirement.model_dump(mode="json"),
                            "deck_content_plan": content_plan.model_dump(mode="json"),
                            "deck_content_plan_markdown": "",
                            "system_selection": system_selection,
                            "waiting_since_turn_index": 2,
                        },
                    },
                )
                async for _ in runner.run_async(
                    user_id=user_id,
                    session_id=session_id,
                    new_message=Content(
                        role="user",
                        parts=[Part(text="把第 2 页改成兔子，并加一个结束页。")],
                    ),
                ):
                    pass
                return await session_service.get_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                )
            finally:
                await runner.close()

        session = await _run_ppt_harness()

        self.assertIsNotNone(session)
        self.assertEqual(session.state["ppt_product_result"]["status"], "awaiting_content_plan_confirmation")
        self.assertEqual(session.state["ppt_workflow_state"]["stage"], "awaiting_content_plan_confirmation")
        self.assertIn("Content plan revision", captured["request_brief"])
        self.assertEqual(
            session.state[PPT_CONTENT_PLAN_CONFIRMATION_WORKFLOW_OUTPUT_KEY]["branch"],
            "revision",
        )
        self.assertEqual(
            session.state[PPT_CONTENT_PLAN_REVISION_WORKFLOW_OUTPUT_KEY]["source"],
            "adk_workflow",
        )
        self.assertEqual(
            session.state[PPT_CONTENT_PLANNING_WORKFLOW_OUTPUT_KEY]["source"],
            "adk_workflow",
        )
        self.assertEqual(session.state["ppt_deck_content_plan"]["title"], "Parent Workflow Revised Plan")
        self.assertNotIn(PPT_ASSET_RESOLUTION_WORKFLOW_OUTPUT_KEY, session.state)
        self.assertNotIn(PPT_ROUTE_EXECUTION_WORKFLOW_OUTPUT_KEY, session.state)
        self.assertNotIn("final_file_paths", session.state)

    async def test_multi_turn_ppt_workflow_runs_through_revisions_and_delivery(self) -> None:
        manager = PptProductManager()
        task = "给幼儿园小朋友做一个 3 页英语单词 PPT。"
        output = {"format": "pptx", "slide_count": 3}
        initial_requirement = manager.prepare_confirmed_requirement(task=task, inputs=[], output=output)
        revised_requirement_payload = initial_requirement.model_dump(mode="json")
        revised_requirement_payload.update(
            {
                "audience": "一年级学生",
                "slide_count_policy": {
                    **revised_requirement_payload["slide_count_policy"],
                    "target": 4,
                },
            }
        )
        revision_selection_payload = {
            "system_type": "built_in_route",
            "route": "html",
            "skill_name": "",
            "output_format": "pptx",
            "reason": "Use the built-in HTML route after requirement revision.",
        }
        initial_requirement_llm = _PptProductManagerToolCallingFakeLlm(
            function_calls=[
                FunctionCall(
                    name="save_ppt_confirmed_requirement_json",
                    args={"requirement_json": initial_requirement.model_dump(mode="json")},
                )
            ],
            final_text="Initial requirement saved.",
        )
        revised_requirement_llm = _PptProductManagerToolCallingFakeLlm(
            function_calls=[
                FunctionCall(
                    name="save_ppt_confirmed_requirement_json",
                    args={"requirement_json": revised_requirement_payload},
                )
            ],
            final_text="Revised requirement saved.",
        )
        revision_selection_llm = _PptProductManagerToolCallingFakeLlm(
            function_calls=[
                FunctionCall(name="list_product_ppt_skills", args={}),
                FunctionCall(
                    name="save_ppt_system_selection",
                    args={"selection_json": revision_selection_payload},
                ),
            ],
            final_text="Revised system selection saved.",
        )
        llm_models = iter([initial_requirement_llm, revised_requirement_llm, revision_selection_llm])
        final_pptx = "generated/ppt-multi-turn-workflow-test/turn_5/final.pptx"
        route_calls: list[dict[str, Any]] = []
        planning_briefs: list[str] = []

        def _system_selection_builder(**_kwargs):
            return {
                "system_type": "built_in_route",
                "route": "html",
                "skill_name": "",
                "output_format": "pptx",
                "reason": "Use the built-in HTML route for the interactive multi-turn test.",
            }

        def _content_plan_builder(requirement):
            planning_briefs.append(requirement.request_brief)
            if "Content plan revision" in requirement.request_brief:
                return DeckContentPlan(
                    title="Revised Rabbit Words",
                    core_narrative="Teach the revised animal words with a rabbit slide.",
                    pages=[
                        _page(1, "cover"),
                        _page(2, "content"),
                        _page(3, "content"),
                        _page(4, "ending"),
                    ],
                )
            return DeckContentPlan(
                title="Initial Animal Words",
                core_narrative="Teach simple animal words.",
                pages=[_page(1, "cover"), _page(2, "content"), _page(3, "ending")],
            )

        async def _asset_resolver(asset, _page, _requirement):
            return {
                "asset_id": asset.asset_id,
                "status": "ready",
                "path": "",
                "provider": "test_resolver",
            }

        async def _fake_dispatch_ppt_route(**kwargs):
            route_calls.append(kwargs)
            return HtmlRouteBuildPackage(
                template=HtmlTemplatePackage(template_id="multi-turn-test", label="Multi Turn Test"),
                html_deck_path="generated/ppt-multi-turn-workflow-test/turn_5/deck.html",
                preview_paths=["generated/ppt-multi-turn-workflow-test/turn_5/preview.png"],
                pptx_path=final_pptx,
                quality_report_path="generated/ppt-multi-turn-workflow-test/turn_5/quality.json",
                build_log_path="generated/ppt-multi-turn-workflow-test/turn_5/build.log",
                warnings=[],
            )

        @node(name="PptMultiTurnWorkflowHarnessNode", rerun_on_resume=True)
        async def ppt_harness(ctx: Context, node_input: str) -> dict[str, Any]:
            ctx.state["turn_index"] = int(ctx.state.get("turn_index", 0) or 0) + 1
            if "ppt_workflow_state" not in ctx.state:
                return await manager.run_product_request(
                    task=node_input,
                    inputs=[],
                    output=output,
                    tool_context=ctx,
                    content_plan_builder=_content_plan_builder,
                    asset_resolver=_asset_resolver,
                    system_selection_builder=_system_selection_builder,
                )
            return await manager.continue_product_request(
                user_response=node_input,
                tool_context=ctx,
                content_plan_builder=_content_plan_builder,
                asset_resolver=_asset_resolver,
            )

        workflow = Workflow(
            name="PptMultiTurnWorkflowHarness",
            edges=[("START", ppt_harness)],
        )
        session_service = InMemorySessionService()
        artifact_service = InMemoryArtifactService()
        runner = Runner(
            node=workflow,
            session_service=session_service,
            artifact_service=artifact_service,
        )
        user_id = "user-ppt-multi-turn-workflow"
        session_id = "session-ppt-multi-turn-workflow"
        turn_sessions = []
        try:
            await session_service.create_session(
                app_name=workflow.name,
                user_id=user_id,
                session_id=session_id,
                state={"sid": "ppt-multi-turn-workflow-test", "turn_index": 0, "step": 1},
            )
            with (
                patch(
                    "src.productions.ppt.ppt_product_manager.ppt_product_manager.build_llm",
                    side_effect=lambda *_args, **_kwargs: next(llm_models),
                ),
                patch.object(manager, "_dispatch_ppt_route", _fake_dispatch_ppt_route),
            ):
                for message in [
                    task,
                    "受众改成一年级学生，页数改成 4 页。",
                    "确认",
                    "把第 2 页改成兔子，并保留 4 页结构。",
                    "确认",
                ]:
                    async for _ in runner.run_async(
                        user_id=user_id,
                        session_id=session_id,
                        new_message=Content(role="user", parts=[Part(text=message)]),
                    ):
                        pass
                    turn_sessions.append(
                        await session_service.get_session(
                            app_name=workflow.name,
                            user_id=user_id,
                            session_id=session_id,
                        )
                    )
        finally:
            await runner.close()

        self.assertEqual(turn_sessions[0].state["ppt_product_result"]["status"], "awaiting_requirement_confirmation")
        self.assertEqual(turn_sessions[0].state[PPT_INITIAL_REQUEST_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(
            turn_sessions[0].state[PPT_INITIAL_REQUEST_WORKFLOW_OUTPUT_KEY]["branch"],
            "requirement_confirmation",
        )
        self.assertEqual(turn_sessions[0].state[PPT_REQUIREMENT_ANALYSIS_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(
            turn_sessions[0].state[PPT_REQUIREMENT_ANALYSIS_RESULT_STATE_KEY]["analysis_output"]["source"],
            "llm_agent",
        )
        self.assertEqual(turn_sessions[1].state["ppt_product_result"]["status"], "awaiting_requirement_confirmation")
        self.assertEqual(turn_sessions[1].state[PPT_REQUIREMENT_CONFIRMATION_WORKFLOW_OUTPUT_KEY]["branch"], "revision")
        self.assertEqual(turn_sessions[1].state[PPT_REQUIREMENT_REVISION_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(turn_sessions[1].state[PPT_SYSTEM_SELECTION_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")

        self.assertEqual(turn_sessions[2].state["ppt_product_result"]["status"], "awaiting_content_plan_confirmation")
        self.assertEqual(turn_sessions[2].state[PPT_SOURCE_PREPARATION_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(turn_sessions[2].state[PPT_CONTENT_PLANNING_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")

        self.assertEqual(turn_sessions[3].state["ppt_product_result"]["status"], "awaiting_content_plan_confirmation")
        self.assertEqual(turn_sessions[3].state[PPT_CONTENT_PLAN_CONFIRMATION_WORKFLOW_OUTPUT_KEY]["branch"], "revision")
        self.assertEqual(turn_sessions[3].state[PPT_CONTENT_PLAN_REVISION_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertNotIn(PPT_ASSET_RESOLUTION_WORKFLOW_OUTPUT_KEY, turn_sessions[3].state)
        self.assertNotIn("final_file_paths", turn_sessions[3].state)

        final_state = turn_sessions[4].state
        self.assertEqual(final_state["ppt_product_result"]["status"], "success")
        self.assertEqual(final_state["ppt_workflow_state"]["stage"], "completed")
        self.assertEqual(final_state[PPT_ASSET_RESOLUTION_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(final_state[PPT_ROUTE_EXECUTION_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(final_state[PPT_FINAL_DELIVERY_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertEqual(final_state["final_file_paths"], [final_pptx])
        self.assertEqual(final_state["ppt_deck_content_plan"]["title"], "Revised Rabbit Words")
        self.assertEqual(len(final_state["ppt_deck_content_plan"]["pages"]), 4)
        self.assertEqual(len(route_calls), 1)
        self.assertEqual(route_calls[0]["content_plan"].title, "Revised Rabbit Words")
        self.assertTrue(any("Content plan revision" in brief for brief in planning_briefs))

    async def test_interactive_workflow_allows_revision_on_later_turn(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={"sid": "ppt-manager-revision-test", "turn_index": 1, "step": 1})

        await manager.run_product_request(
            task="给我做一个ppt，用来给幼儿园小朋友讲英语单词。3页，分别讲 猫、狗、鸭子。",
            inputs=[],
            output={"format": "pptx"},
            tool_context=tool_context,
        )

        tool_context.state["turn_index"] = 2
        plan_result = await manager.continue_product_request(
            user_response="确认",
            tool_context=tool_context,
        )
        self.assertEqual(plan_result["status"], "awaiting_content_plan_confirmation")

        tool_context.state["turn_index"] = 3
        revised_result = await manager.continue_product_request(
            user_response="把第 2 页改成兔子。",
            tool_context=tool_context,
        )

        self.assertEqual(revised_result["status"], "awaiting_content_plan_confirmation")
        self.assertEqual(tool_context.state["ppt_workflow_state"]["stage"], "awaiting_content_plan_confirmation")
        self.assertEqual(tool_context.state["ppt_workflow_state"]["waiting_since_turn_index"], 3)
        self.assertIn("Content plan revision", tool_context.state["ppt_workflow_state"]["confirmed_requirement"]["request_brief"])
        self.assertNotIn("final_file_paths", tool_context.state)

    async def test_content_plan_revision_phase_persists_typed_result(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(
            state={"sid": "ppt-content-plan-revision-phase-test", "turn_index": 3, "step": 1}
        )
        requirement = manager.prepare_confirmed_requirement(
            task="给幼儿园小朋友做一个 3 页英语单词 PPT。",
            inputs=[],
            output={"format": "pptx", "slide_count": 3},
        )
        captured: dict[str, Any] = {}

        def _content_plan_builder(revised_requirement):
            captured["request_brief"] = revised_requirement.request_brief
            return DeckContentPlan(
                title="Revised Animal Words",
                core_narrative="Use the revised animal word sequence.",
                pages=[_page(1, "cover"), _page(2, "content")],
            )

        result = await manager._revise_content_plan_phase(
            requirement=requirement,
            user_response="把第 2 页改成兔子。",
            tool_context=tool_context,
            app_name="creative_claw",
            artifact_service=None,
            expert_agents={},
            content_plan_builder=_content_plan_builder,
        )

        self.assertIsInstance(result, PptContentPlanRevisionResult)
        self.assertEqual(result.content_plan.title, "Revised Animal Words")
        self.assertIn("Content plan revision", result.confirmed_requirement.request_brief)
        self.assertIn("把第 2 页改成兔子", captured["request_brief"])
        self.assertEqual(result.revision_output["source"], "injected")
        self.assertEqual(result.revision_output["page_count"], 2)
        self.assertEqual(
            tool_context.state[PPT_CONTENT_PLAN_REVISION_OUTPUT_STATE_KEY]["source"],
            "injected",
        )
        self.assertEqual(
            tool_context.state[PPT_CONTENT_PLAN_REVISION_RESULT_STATE_KEY]["content_plan"]["title"],
            "Revised Animal Words",
        )
        self.assertEqual(tool_context.state["ppt_deck_content_plan"]["title"], "Revised Animal Words")
        self.assertNotIn("final_file_paths", tool_context.state)

    async def test_content_plan_revision_phase_uses_adk_workflow_context(self) -> None:
        manager = PptProductManager()
        requirement = manager.prepare_confirmed_requirement(
            task="给幼儿园小朋友做一个 3 页英语单词 PPT。",
            inputs=[],
            output={"format": "pptx", "slide_count": 3},
        )
        captured: dict[str, Any] = {}

        def _content_plan_builder(revised_requirement):
            captured["request_brief"] = revised_requirement.request_brief
            return DeckContentPlan(
                title="Workflow Revised Animal Words",
                core_narrative="Use the revised animal word sequence from Workflow.",
                pages=[_page(1, "cover"), _page(2, "content"), _page(3, "ending")],
            )

        async def _run_revision_harness():
            @node(name="PptContentPlanRevisionWorkflowHarnessNode", rerun_on_resume=True)
            async def revision_harness(ctx: Context, node_input: str) -> dict[str, Any]:
                result = await manager._revise_content_plan_phase(
                    requirement=requirement,
                    user_response="把第 2 页改成兔子。",
                    tool_context=ctx,
                    app_name="creative_claw",
                    artifact_service=artifact_service,
                    expert_agents={},
                    content_plan_builder=_content_plan_builder,
                )
                return result.model_dump(mode="json")

            workflow = Workflow(
                name="PptContentPlanRevisionWorkflowHarness",
                edges=[("START", revision_harness)],
            )
            session_service = InMemorySessionService()
            runner = Runner(
                node=workflow,
                session_service=session_service,
                artifact_service=artifact_service,
            )
            user_id = "user-ppt-content-plan-revision-workflow"
            session_id = "session-ppt-content-plan-revision-workflow"
            try:
                await session_service.create_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                    state={
                        "sid": "ppt-content-plan-revision-workflow-test",
                        "turn_index": 3,
                        "step": 1,
                    },
                )
                async for _ in runner.run_async(
                    user_id=user_id,
                    session_id=session_id,
                    new_message=Content(role="user", parts=[Part(text="Revise PPT content plan")]),
                ):
                    pass
                return await session_service.get_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                )
            finally:
                await runner.close()

        artifact_service = InMemoryArtifactService()
        session = await _run_revision_harness()

        self.assertIsNotNone(session)
        self.assertIn("Content plan revision", captured["request_brief"])
        self.assertIn("把第 2 页改成兔子", captured["request_brief"])
        self.assertEqual(
            session.state[PPT_CONTENT_PLAN_REVISION_WORKFLOW_OUTPUT_KEY]["source"],
            "adk_workflow",
        )
        self.assertEqual(
            session.state[PPT_CONTENT_PLAN_REVISION_WORKFLOW_OUTPUT_KEY]["revision_source"],
            "injected",
        )
        self.assertEqual(
            session.state[PPT_CONTENT_PLAN_REVISION_WORKFLOW_OUTPUT_KEY]["page_count"],
            3,
        )
        self.assertEqual(session.state[PPT_CONTENT_PLAN_REVISION_OUTPUT_STATE_KEY]["source"], "injected")
        self.assertEqual(
            session.state[PPT_CONTENT_PLAN_REVISION_RESULT_STATE_KEY]["content_plan"]["title"],
            "Workflow Revised Animal Words",
        )
        self.assertEqual(session.state["ppt_deck_content_plan"]["title"], "Workflow Revised Animal Words")
        self.assertEqual(session.state[PPT_CONTENT_PLANNING_WORKFLOW_OUTPUT_KEY]["source"], "adk_workflow")
        self.assertNotIn("final_file_paths", session.state)
        self.assertNotIn(PPT_ASSET_RESOLUTION_WORKFLOW_OUTPUT_KEY, session.state)

    async def test_requirement_revision_phase_persists_typed_result(self) -> None:
        source_path = _write_markdown_source("revision_phase_source.md", "# Revision source")
        manager = PptProductManager()
        tool_context = SimpleNamespace(
            state={"sid": "ppt-requirement-revision-phase-test", "turn_index": 2, "step": 1}
        )
        existing_requirement = manager.prepare_confirmed_requirement(
            task=(
                "针对这个素材，做一个用于组会和同学讲解的 PPT。"
                "素材为论文/文档 revision_phase_source.md，需要提炼内容。"
            ),
            inputs={"files": [source_path]},
            output={"format": "pptx", "language": "zh-CN"},
        )

        result = await manager._revise_requirement_phase(
            existing_requirement=existing_requirement,
            user_response="受众为同组的同学，场景为组会，页数改成15页左右。",
            task=existing_requirement.request_brief,
            raw_inputs=[source_path],
            output={"format": "pptx", "language": "zh-CN"},
            source_understanding=existing_requirement.source_understanding,
            tool_context=tool_context,
            app_name="creative_claw",
            artifact_service=None,
        )

        self.assertIsInstance(result, PptRequirementRevisionResult)
        self.assertEqual(result.confirmed_requirement.audience, "同组的同学")
        self.assertEqual(result.confirmed_requirement.scenario, "组会")
        self.assertEqual(result.confirmed_requirement.slide_count_policy.target, 15)
        self.assertEqual(result.revision_output["source"], "deterministic_fallback")
        self.assertEqual(
            tool_context.state[PPT_REQUIREMENT_REVISION_OUTPUT_STATE_KEY]["source"],
            "deterministic_fallback",
        )
        self.assertEqual(
            tool_context.state[PPT_REQUIREMENT_REVISION_RESULT_STATE_KEY]["confirmed_requirement"]["audience"],
            "同组的同学",
        )
        self.assertEqual(tool_context.state[PPT_CONFIRMED_REQUIREMENT_STATE_KEY]["source_inputs"][0]["path"], source_path)
        self.assertNotIn("User revision", result.confirmed_requirement.topic)

    async def test_requirement_revision_phase_uses_adk_workflow_context(self) -> None:
        source_path = _write_markdown_source("revision_workflow_source.md", "# Revision workflow source")
        manager = PptProductManager()
        existing_requirement = manager.prepare_confirmed_requirement(
            task=(
                "针对这个素材，做一个用于产品复盘的 PPT。"
                "素材为论文/文档 revision_workflow_source.md，需要提炼内容。"
            ),
            inputs={"files": [source_path]},
            output={"format": "pptx", "language": "zh-CN"},
        )
        revised_payload = existing_requirement.model_dump(mode="json")
        revised_payload.update(
            {
                "audience": "技术负责人",
                "scenario": "季度复盘",
                "slide_count_policy": {
                    **revised_payload["slide_count_policy"],
                    "target": 12,
                },
            }
        )
        fake_llm = _PptProductManagerToolCallingFakeLlm(
            function_calls=[
                FunctionCall(
                    name="save_ppt_confirmed_requirement_json",
                    args={"requirement_json": revised_payload},
                )
            ],
            final_text="PptRequirementAnalysisAgent saved revised requirement.",
        )

        async def _run_revision_harness():
            @node(name="PptRequirementRevisionWorkflowHarnessNode", rerun_on_resume=True)
            async def revision_harness(ctx: Context, node_input: str) -> dict[str, Any]:
                result = await manager._revise_requirement_phase(
                    existing_requirement=existing_requirement,
                    user_response="受众改为技术负责人，场景改为季度复盘，页数改成 12 页。",
                    task=existing_requirement.request_brief,
                    raw_inputs=[source_path],
                    output={"format": "pptx", "language": "zh-CN"},
                    source_understanding=existing_requirement.source_understanding,
                    tool_context=ctx,
                    app_name="creative_claw",
                    artifact_service=artifact_service,
                )
                return result.model_dump(mode="json")

            workflow = Workflow(
                name="PptRequirementRevisionWorkflowHarness",
                edges=[("START", revision_harness)],
            )
            session_service = InMemorySessionService()
            runner = Runner(
                node=workflow,
                session_service=session_service,
                artifact_service=artifact_service,
            )
            user_id = "user-ppt-requirement-revision-workflow"
            session_id = "session-ppt-requirement-revision-workflow"
            try:
                await session_service.create_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                    state={
                        "sid": "ppt-requirement-revision-workflow-test",
                        "turn_index": 2,
                        "step": 1,
                    },
                )
                with patch(
                    "src.productions.ppt.ppt_product_manager.ppt_product_manager.build_llm",
                    return_value=fake_llm,
                ):
                    async for _ in runner.run_async(
                        user_id=user_id,
                        session_id=session_id,
                        new_message=Content(role="user", parts=[Part(text="Revise PPT requirement")]),
                    ):
                        pass
                return await session_service.get_session(
                    app_name=workflow.name,
                    user_id=user_id,
                    session_id=session_id,
                )
            finally:
                await runner.close()

        artifact_service = InMemoryArtifactService()
        session = await _run_revision_harness()

        self.assertIsNotNone(session)
        self.assertEqual(
            session.state[PPT_REQUIREMENT_REVISION_WORKFLOW_OUTPUT_KEY]["source"],
            "adk_workflow",
        )
        self.assertEqual(
            session.state[PPT_REQUIREMENT_REVISION_WORKFLOW_OUTPUT_KEY]["revision_source"],
            "llm_agent",
        )
        self.assertEqual(session.state[PPT_REQUIREMENT_REVISION_OUTPUT_STATE_KEY]["source"], "llm_agent")
        self.assertEqual(
            session.state[PPT_REQUIREMENT_REVISION_RESULT_STATE_KEY]["confirmed_requirement"]["audience"],
            "技术负责人",
        )
        self.assertEqual(session.state[PPT_CONFIRMED_REQUIREMENT_STATE_KEY]["scenario"], "季度复盘")
        self.assertEqual(session.state[PPT_CONFIRMED_REQUIREMENT_STATE_KEY]["slide_count_policy"]["target"], 12)
        first_request_tools = _function_declaration_names(fake_llm.requests[0])
        self.assertIn("save_ppt_confirmed_requirement_json", first_request_tools)
        self.assertGreaterEqual(len(fake_llm.requests), 2)

    async def test_requirement_confirmation_revision_updates_structured_fields(self) -> None:
        source_path = _write_markdown_source("Thinking_with_Visual_Primitives.pdf", "%PDF test fixture")
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={"sid": "ppt-manager-requirement-revision-test", "turn_index": 1, "step": 1})

        initial_result = await manager.run_product_request(
            task=(
                "针对这个素材，做一个用于组会和同学讲解的 PPT。"
                "素材为论文/文档 Thinking_with_Visual_Primitives.pdf，需要提炼内容。"
            ),
            inputs={"files": [source_path]},
            output={"format": "pptx", "language": "zh-CN", "use_case": "组会讲解"},
            tool_context=tool_context,
        )

        self.assertEqual(initial_result["status"], "awaiting_requirement_confirmation")
        self.assertEqual(initial_result["confirmed_requirement"]["scenario"], "组会")
        self.assertNotIn("User revision", initial_result["confirmed_requirement"]["topic"])
        self.assertEqual(len(initial_result["confirmed_requirement"]["source_inputs"]), 1)

        tool_context.state["turn_index"] = 2
        page_count_result = await manager.continue_product_request(
            user_response="页数改成15页左右。",
            tool_context=tool_context,
        )

        self.assertEqual(page_count_result["status"], "awaiting_requirement_confirmation")
        requirement_after_page_count = page_count_result["confirmed_requirement"]
        self.assertEqual(requirement_after_page_count["slide_count_policy"]["target"], 15)
        self.assertEqual(requirement_after_page_count["source_inputs"][0]["path"], source_path)
        self.assertNotIn("User revision", requirement_after_page_count["topic"])
        self.assertNotIn("User revision", requirement_after_page_count["request_brief"])

        tool_context.state["turn_index"] = 3
        audience_result = await manager.continue_product_request(
            user_response="受众为同组的同学，场景为组会。",
            tool_context=tool_context,
        )

        self.assertEqual(audience_result["status"], "awaiting_requirement_confirmation")
        revised_requirement = audience_result["confirmed_requirement"]
        self.assertEqual(revised_requirement["audience"], "同组的同学")
        self.assertEqual(revised_requirement["scenario"], "组会")
        self.assertEqual(revised_requirement["slide_count_policy"]["target"], 15)
        self.assertEqual(revised_requirement["source_inputs"][0]["path"], source_path)
        self.assertNotIn("User revision", revised_requirement["topic"])
        self.assertIn("| 受众 | 同组的同学 |", audience_result["confirmation_request"]["summary_markdown"])
        self.assertIn("| 场景 | 组会 |", audience_result["confirmation_request"]["summary_markdown"])

    async def test_requirement_confirmation_stages_external_upload_before_source_conversion(self) -> None:
        with tempfile.TemporaryDirectory(dir="/private/tmp") as temp_dir:
            external_pdf = Path(temp_dir) / "Thinking_with_Visual_Primitives.pdf"
            external_pdf.write_bytes(b"%PDF external upload fixture")

            manager = PptProductManager()
            tool_context = SimpleNamespace(
                state={
                    "sid": "ppt-manager-external-upload-test",
                    "channel": "web",
                    "turn_index": 1,
                    "step": 1,
                }
            )
            captured: dict[str, str] = {}

            async def _capturing_source_converter(source_input, parameters: dict) -> dict:
                captured["source_input_path"] = source_input.path
                captured["input_path"] = parameters["input_path"]
                return await _fake_source_converter(source_input, parameters)

            await manager.run_product_request(
                task="针对这个素材，给我做一个ppt，用于组会和同学讲解。",
                inputs=[str(external_pdf)],
                output={"format": "pptx", "language": "zh-CN", "usage": "组会讲解"},
                tool_context=tool_context,
            )

            tool_context.state["turn_index"] = 2
            result = await manager.continue_product_request(
                user_response="确认",
                tool_context=tool_context,
                source_converter=_capturing_source_converter,
            )

        self.assertEqual(result["status"], "awaiting_content_plan_confirmation")
        self.assertNotEqual(captured["input_path"], str(external_pdf))
        self.assertTrue(captured["input_path"].startswith("inbox/web/ppt-manager-external-upload-test/turn_2/"))
        self.assertTrue(resolve_workspace_path(captured["input_path"]).exists())
        self.assertEqual(result["confirmed_requirement"]["source_inputs"][0]["path"], captured["source_input_path"])
        self.assertEqual(result["confirmed_requirement"]["source_inputs"][0]["path"], captured["input_path"])

    async def test_requirement_confirmation_downloads_remote_url_before_source_conversion(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(
            state={
                "sid": "ppt-manager-remote-source-test",
                "channel": "web",
                "turn_index": 1,
                "step": 1,
            }
        )
        source_url = "https://arxiv.org/pdf/1706.03762"
        captured: dict[str, str] = {}

        async def _capturing_source_converter(source_input, parameters: dict) -> dict:
            captured["source_input_path"] = source_input.path
            captured["input_path"] = parameters["input_path"]
            captured["has_url_parameter"] = str("url" in parameters)
            return await _fake_source_converter(source_input, parameters)

        await manager.run_product_request(
            task="针对这个论文，给我做一个ppt，用于组会和同学讲解。",
            inputs=[{"name": "1706.03762", "url": source_url}],
            output={"format": "pptx", "language": "zh-CN", "usage": "组会讲解"},
            tool_context=tool_context,
        )

        tool_context.state["turn_index"] = 2
        fake_response = _FakeRemoteResponse(
            b"%PDF-1.4\nremote pdf fixture\n",
            {"content-type": "application/pdf", "content-length": "28"},
        )
        with patch(
            "src.productions.ppt.ppt_product_manager.ppt_product_manager.urlopen",
            return_value=fake_response,
        ):
            result = await manager.continue_product_request(
                user_response="确认",
                tool_context=tool_context,
                source_converter=_capturing_source_converter,
            )

        self.assertEqual(result["status"], "awaiting_content_plan_confirmation")
        self.assertEqual(captured["has_url_parameter"], "False")
        self.assertTrue(captured["input_path"].startswith("generated/ppt-manager-remote-source-test/turn_2/"))
        self.assertTrue(captured["input_path"].endswith(".pdf"))
        self.assertTrue(resolve_workspace_path(captured["input_path"]).exists())
        self.assertEqual(result["confirmed_requirement"]["source_inputs"][0]["path"], captured["source_input_path"])
        self.assertEqual(result["confirmed_requirement"]["source_inputs"][0]["path"], captured["input_path"])
        self.assertEqual(tool_context.state["ppt_remote_source_downloads"][0]["source_url"], source_url)

    async def test_run_returns_deferred_status_for_unimplemented_xml_route(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={"sid": "ppt-manager-test", "turn_index": 1, "step": 1})

        result = await manager.run_product_request(
            task="使用 xml route 生成一份增长策略汇报。",
            inputs=[],
            output={"route": "xml", "auto_confirm": True},
            tool_context=tool_context,
        )

        self.assertEqual(result["status"], "route_not_implemented")
        self.assertEqual(result["selected_route"], "xml")
        self.assertEqual(result["output_files"], [])
        self.assertNotIn("final_file_paths", tool_context.state)

    async def test_run_records_source_materials_and_resets_current_output(self) -> None:
        source_path = _write_markdown_source(
            "launch_brief.md",
            """# Growth Launch

## Customer Proof
- Activation rose after guided onboarding.
- Enterprise pipeline needs proof-led messaging.
""",
        )
        manager = PptProductManager()
        tool_context = SimpleNamespace(
            state={
                "sid": "ppt-manager-source-test",
                "turn_index": 2,
                "step": 1,
                "current_output": {"status": "success", "message": "stale expert output"},
            }
        )

        result = await manager.run_product_request(
            task="基于材料生成 6 页 PPTX，用于增长发布会。",
            inputs=[{"name": "launch_brief.md", "path": source_path}],
            output={"format": "pptx", "auto_confirm": True},
            tool_context=tool_context,
            source_converter=_fake_source_converter,
        )

        self.assertEqual(result["status"], "success")
        source_materials = result["confirmed_requirement"]["source_understanding"]
        self.assertEqual(source_materials["markdown_sources"][0]["name"], "launch_brief.md")
        self.assertEqual(source_materials["figures"][0]["alt"], "Activation chart")
        ready_assets = [
            asset
            for page in result["deck_content_plan"]["pages"]
            for asset in page.get("assets", [])
            if asset.get("status") == "ready"
        ]
        self.assertEqual(ready_assets[0]["source_kind"], "material_figure")
        self.assertTrue(ready_assets[0]["path"].endswith("activation.png"))
        self.assertEqual(tool_context.state["ppt_resolved_asset_manifest"]["ready_asset_count"], 1)
        plan_text = str(result["deck_content_plan"])
        self.assertIn("prepared source materials", plan_text)
        self.assertIn("ppt_source_markdown_sources", tool_context.state)
        self.assertIn("ppt_source_figures", tool_context.state)
        self.assertTrue(tool_context.state["ppt_source_output_files"])
        self.assertEqual(tool_context.state["current_output"]["product_line"], "ppt")
        self.assertEqual(tool_context.state["current_output"]["status"], "success")

        html_path = resolve_workspace_path(result["delivery_manifest"]["intermediate_artifacts"][0])
        html_text = html_path.read_text(encoding="utf-8")
        self.assertIn("Growth Launch", html_text)
        self.assertIn("Activation chart", html_text)

        pptx_path = resolve_workspace_path(result["delivery_manifest"]["final_pptx"])
        pptx_text = "\n".join(
            shape.text
            for slide in Presentation(str(pptx_path)).slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        self.assertIn("Growth Launch", pptx_text)
        self.assertIn("Use the provided figures", pptx_text)
        picture_count = sum(
            1
            for slide in Presentation(str(pptx_path)).slides
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        )
        self.assertGreaterEqual(picture_count, 1)

    async def test_run_converts_web_uploaded_file_path_strings(self) -> None:
        source_path = _write_markdown_source("creative_agent_NeurIPS_2026_10_.pdf", "%PDF test fixture")
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={"sid": "ppt-manager-web-upload-test", "turn_index": 2, "step": 1})

        result = await manager.run_product_request(
            task="针对这个素材，给我做一个ppt，用来组会上给团队的同学讲解。",
            inputs={"files": [source_path]},
            output={"format": "pptx", "language": "zh-CN", "purpose": "组会讲解", "auto_confirm": True},
            tool_context=tool_context,
            source_converter=_fake_source_converter,
        )

        self.assertEqual(result["status"], "success")
        self.assertEqual(result["confirmed_requirement"]["source_inputs"][0]["path"], source_path)
        source_materials = result["confirmed_requirement"]["source_understanding"]
        self.assertEqual(source_materials["document_type"], "pdf")
        self.assertEqual(source_materials["markdown_sources"][0]["name"], "creative_agent_NeurIPS_2026_10_.pdf")
        self.assertTrue(tool_context.state["ppt_source_markdown_sources"])

    async def test_run_uses_existing_markdown_source_without_anything_to_md(self) -> None:
        source_path = _write_markdown_source(
            "local_markdown_brief.md",
            """# Local Markdown Brief

- Retention improved after guided onboarding.
- Sales teams need a simple proof-led deck.
""",
        )
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={"sid": "ppt-manager-local-md-test", "turn_index": 1, "step": 1})

        result = await manager.run_product_request(
            task="基于本地 Markdown 生成 6 页 PPTX。",
            inputs=[{"name": "local_markdown_brief.md", "path": source_path}],
            output={"format": "pptx", "auto_confirm": True},
            tool_context=tool_context,
        )

        self.assertEqual(result["status"], "success")
        source_materials = result["confirmed_requirement"]["source_understanding"]
        self.assertEqual(source_materials["markdown_sources"][0]["method"], "local:markdown_passthrough")
        self.assertEqual(source_materials["markdown_sources"][0]["output_path"], source_path)
        self.assertIn("ppt_markdown_source_texts", tool_context.state)
        self.assertIn("Retention improved", str(result["deck_content_plan"]))

        pptx_path = resolve_workspace_path(result["delivery_manifest"]["final_pptx"])
        pptx_text = "\n".join(
            shape.text
            for slide in Presentation(str(pptx_path)).slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        self.assertIn("Retention improved", pptx_text)

    async def test_content_planning_resolves_pending_generated_asset_before_route(self) -> None:
        image_path = _write_test_image("generated_asset_fixture.png")
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={"sid": "ppt-manager-asset-test", "turn_index": 1, "step": 1})

        def _content_plan_builder(_requirement):
            plan = DeckContentPlan(
                title="AI for Kids",
                core_narrative="Explain AI through concrete classroom examples.",
                pages=[
                    _page(1, "cover"),
                    _page(2, "toc"),
                    _page(3, "chapter_start"),
                    _page(4, "chapter_content"),
                    _page(5, "ending"),
                ],
            )
            plan.pages[3].asset_source_preference = "ai"
            plan.pages[3].assets = [
                DeckPageAsset(
                    asset_id="slide_04_ai_visual",
                    source_kind="image_generation",
                    status="pending",
                    description="A friendly classroom illustration showing students learning AI.",
                    prompt="A friendly classroom illustration showing students learning AI.",
                )
            ]
            return plan

        async def _asset_resolver(asset, _page, _requirement):
            return {
                "asset_id": asset.asset_id,
                "status": "ready",
                "path": image_path,
                "provider": "test_resolver",
            }

        result = await manager.run_product_request(
            task="给小学生做一个 AI 科普 PPTX。",
            inputs=[],
            output={"format": "pptx", "auto_confirm": True},
            tool_context=tool_context,
            content_plan_builder=_content_plan_builder,
            asset_resolver=_asset_resolver,
        )

        self.assertEqual(result["status"], "success")
        resolved_asset = result["deck_content_plan"]["pages"][3]["assets"][0]
        self.assertEqual(resolved_asset["status"], "ready")
        self.assertEqual(resolved_asset["path"], image_path)
        self.assertEqual(tool_context.state["ppt_resolved_asset_manifest"]["ready_asset_count"], 1)
        progress_events = list(tool_context.state.get("orchestration_events") or [])
        image_generation_events = [
            event for event in progress_events if event.get("title") == "PPT Image Generation"
        ]
        self.assertEqual(len(image_generation_events), 2)
        self.assertIn("Status: started", image_generation_events[0]["detail"])
        self.assertIn("Status: success", image_generation_events[1]["detail"])
        self.assertIn("slide_04_ai_visual", image_generation_events[1]["detail"])

        pptx_path = resolve_workspace_path(result["delivery_manifest"]["final_pptx"])
        picture_count = sum(
            1
            for slide in Presentation(str(pptx_path)).slides
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        )
        self.assertGreaterEqual(picture_count, 1)

    async def test_asset_resolution_phase_persists_typed_result(self) -> None:
        image_path = _write_test_image("asset_resolution_phase_fixture.png")
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={"sid": "ppt-asset-phase-test", "turn_index": 1, "step": 1})
        requirement = manager.prepare_confirmed_requirement(
            task="给小学生做一个 AI 科普 PPTX。",
            inputs=[],
            output={"format": "pptx", "slide_count": 5},
        )
        content_plan = DeckContentPlan(
            title="AI for Kids",
            core_narrative="Explain AI through concrete classroom examples.",
            pages=[
                _page(1, "cover"),
                _page(2, "toc"),
                _page(3, "chapter_start"),
                _page(4, "chapter_content"),
                _page(5, "ending"),
            ],
        )
        content_plan.pages[3].assets = [
            DeckPageAsset(
                asset_id="slide_04_ai_visual",
                source_kind="image_generation",
                status="pending",
                description="A friendly classroom illustration showing students learning AI.",
                prompt="A friendly classroom illustration showing students learning AI.",
            )
        ]

        async def _asset_resolver(asset, _page, _requirement):
            return {
                "asset_id": asset.asset_id,
                "status": "ready",
                "path": image_path,
                "provider": "test_resolver",
            }

        result = await manager._resolve_deck_assets_phase(
            content_plan,
            requirement,
            tool_context=tool_context,
            expert_agents={},
            app_name="creative_claw",
            artifact_service=None,
            asset_resolver=_asset_resolver,
        )

        self.assertIsInstance(result, PptAssetResolutionResult)
        resolved_asset = result.content_plan.pages[3].assets[0]
        self.assertEqual(resolved_asset.status, "ready")
        self.assertEqual(resolved_asset.path, image_path)
        self.assertEqual(result.resolved_asset_manifest["ready_asset_count"], 1)
        self.assertEqual(tool_context.state["ppt_resolved_asset_manifest"], result.resolved_asset_manifest)
        self.assertEqual(tool_context.state["ppt_deck_content_plan"]["pages"][3]["assets"][0]["path"], image_path)
        self.assertEqual(
            tool_context.state["ppt_asset_resolution_result"]["resolved_asset_manifest"],
            result.resolved_asset_manifest,
        )

    async def test_asset_resolution_phase_reuses_same_input_ready_assets(self) -> None:
        image_path = _write_test_image("asset_resolution_reuse_fixture.png")
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={"sid": "ppt-asset-reuse-test", "turn_index": 1, "step": 1})
        requirement = manager.prepare_confirmed_requirement(
            task="给小学生做一个 AI 科普 PPTX。",
            inputs=[],
            output={"format": "pptx", "slide_count": 5},
        )
        content_plan = DeckContentPlan(
            title="AI for Kids",
            core_narrative="Explain AI through concrete classroom examples.",
            pages=[
                _page(1, "cover"),
                _page(2, "toc"),
                _page(3, "chapter_start"),
                _page(4, "chapter_content"),
                _page(5, "ending"),
            ],
        )
        content_plan.pages[3].assets = [
            DeckPageAsset(
                asset_id="slide_04_ai_visual",
                source_kind="image_generation",
                status="pending",
                description="A friendly classroom illustration showing students learning AI.",
                prompt="A friendly classroom illustration showing students learning AI.",
            )
        ]
        resolver_calls = 0

        async def _asset_resolver(asset, _page, _requirement):
            nonlocal resolver_calls
            resolver_calls += 1
            return {
                "asset_id": asset.asset_id,
                "status": "ready",
                "path": image_path,
                "provider": "test_resolver",
            }

        first_result = await manager._resolve_deck_assets_phase(
            content_plan,
            requirement,
            tool_context=tool_context,
            expert_agents={},
            app_name="creative_claw",
            artifact_service=None,
            asset_resolver=_asset_resolver,
        )
        second_result = await manager._resolve_deck_assets_phase(
            content_plan,
            requirement,
            tool_context=tool_context,
            expert_agents={},
            app_name="creative_claw",
            artifact_service=None,
            asset_resolver=_asset_resolver,
        )

        self.assertFalse(first_result.reused_existing_resolution)
        self.assertTrue(second_result.reused_existing_resolution)
        self.assertEqual(resolver_calls, 1)
        self.assertEqual(first_result.input_signature, second_result.input_signature)
        self.assertEqual(second_result.content_plan.pages[3].assets[0].path, image_path)
        self.assertTrue(tool_context.state["ppt_asset_resolution_result"]["reused_existing_resolution"])

    async def test_asset_resolution_phase_reruns_when_ready_asset_is_missing(self) -> None:
        first_image_path = _write_test_image("asset_resolution_missing_first.png")
        second_image_path = _write_test_image("asset_resolution_missing_second.png")
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={"sid": "ppt-asset-missing-test", "turn_index": 1, "step": 1})
        requirement = manager.prepare_confirmed_requirement(
            task="给小学生做一个 AI 科普 PPTX。",
            inputs=[],
            output={"format": "pptx", "slide_count": 5},
        )
        content_plan = DeckContentPlan(
            title="AI for Kids",
            core_narrative="Explain AI through concrete classroom examples.",
            pages=[
                _page(1, "cover"),
                _page(2, "toc"),
                _page(3, "chapter_start"),
                _page(4, "chapter_content"),
                _page(5, "ending"),
            ],
        )
        content_plan.pages[3].assets = [
            DeckPageAsset(
                asset_id="slide_04_ai_visual",
                source_kind="image_generation",
                status="pending",
                description="A friendly classroom illustration showing students learning AI.",
                prompt="A friendly classroom illustration showing students learning AI.",
            )
        ]
        resolver_paths = [first_image_path, second_image_path]
        resolver_calls = 0

        async def _asset_resolver(asset, _page, _requirement):
            nonlocal resolver_calls
            path = resolver_paths[resolver_calls]
            resolver_calls += 1
            return {
                "asset_id": asset.asset_id,
                "status": "ready",
                "path": path,
                "provider": "test_resolver",
            }

        first_result = await manager._resolve_deck_assets_phase(
            content_plan,
            requirement,
            tool_context=tool_context,
            expert_agents={},
            app_name="creative_claw",
            artifact_service=None,
            asset_resolver=_asset_resolver,
        )
        resolve_workspace_path(first_image_path).unlink()
        second_result = await manager._resolve_deck_assets_phase(
            content_plan,
            requirement,
            tool_context=tool_context,
            expert_agents={},
            app_name="creative_claw",
            artifact_service=None,
            asset_resolver=_asset_resolver,
        )

        self.assertFalse(first_result.reused_existing_resolution)
        self.assertFalse(second_result.reused_existing_resolution)
        self.assertEqual(resolver_calls, 2)
        self.assertEqual(second_result.content_plan.pages[3].assets[0].path, second_image_path)
        self.assertFalse(tool_context.state["ppt_asset_resolution_result"]["reused_existing_resolution"])

    async def test_illustrated_kid_word_deck_generates_plan_assets(self) -> None:
        image_path = _write_test_image("kid_word_generated_asset.png")
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={"sid": "ppt-manager-kid-word-test", "turn_index": 1, "step": 1})

        async def _asset_resolver(asset, _page, _requirement):
            return {
                "asset_id": asset.asset_id,
                "status": "ready",
                "path": image_path,
                "provider": "test_resolver",
            }

        result = await manager.run_product_request(
            task="给我做一个ppt，用来给幼儿园小朋友讲英语单词。图文并茂。小于10页。",
            inputs=[],
            output={"format": "pptx", "auto_confirm": True},
            tool_context=tool_context,
            asset_resolver=_asset_resolver,
        )

        self.assertEqual(result["status"], "success")
        self.assertEqual(result["confirmed_requirement"]["topic"], "英语单词")
        self.assertEqual(result["confirmed_requirement"]["audience"], "幼儿园小朋友")
        self.assertLessEqual(result["confirmed_requirement"]["slide_count_policy"]["target"], 9)

        pages = result["deck_content_plan"]["pages"]
        page_titles = [page["title"] for page in pages]
        self.assertIn("Apple 苹果", page_titles)
        self.assertIn("Cat 猫", page_titles)
        self.assertIn("Dog 狗", page_titles)
        self.assertNotIn("Context", page_titles)
        self.assertNotIn("Insight", page_titles)
        self.assertNotIn("Next Steps", page_titles)
        self.assertNotIn("No source file", str(pages))
        self.assertNotIn("ContentPlanningAgent", str(pages))

        ai_pages = [page for page in pages if page["asset_source_preference"] == "ai"]
        self.assertGreaterEqual(len(ai_pages), 1)
        self.assertTrue(all(page["page_type"] == "content" for page in pages))
        self.assertTrue(all(page["asset_source_preference"] == "ai" for page in pages))

        ready_assets = [
            asset
            for page in pages
            for asset in page.get("assets", [])
            if asset.get("status") == "ready"
        ]
        self.assertGreaterEqual(len(ready_assets), 1)
        self.assertTrue(all(asset["source_kind"] == "image_generation" for asset in ready_assets))
        self.assertGreaterEqual(tool_context.state["ppt_resolved_asset_manifest"]["ready_asset_count"], 1)

        pptx_path = resolve_workspace_path(result["delivery_manifest"]["final_pptx"])
        picture_count = sum(
            1
            for slide in Presentation(str(pptx_path)).slides
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        )
        self.assertGreaterEqual(picture_count, 1)

    async def test_run_returns_needs_clarification_for_too_thin_request(self) -> None:
        manager = PptProductManager()
        tool_context = SimpleNamespace(state={"sid": "ppt-manager-test", "turn_index": 1, "step": 1})

        result = await manager.run_product_request(
            task="做个 PPT",
            inputs=[],
            output={"format": "pptx"},
            tool_context=tool_context,
        )

        self.assertEqual(result["status"], "needs_clarification")
        self.assertEqual(result["selected_route"], "html")
        self.assertIn("补充 PPT 的主题", result["next_actions"][0])
        self.assertNotIn("final_file_paths", tool_context.state)

    def test_deck_content_plan_allows_no_template_page_types(self) -> None:
        plan = DeckContentPlan(
            title="Demo deck",
            core_narrative="A concise direct narrative.",
            pages=[
                _page(1, "content"),
                _page(2, "quote"),
                _page(3, "activity"),
            ],
        )

        self.assertEqual(len(plan.pages), 3)
        self.assertEqual({page.page_type for page in plan.pages}, {"content", "quote", "activity"})

        with self.assertRaisesRegex(ValueError, "duplicate slide numbers"):
            DeckContentPlan(
                title="Broken deck",
                core_narrative="Duplicate slide numbers.",
                pages=[
                    _page(1, "content"),
                    _page(1, "content"),
                ],
            )


if __name__ == "__main__":
    unittest.main()
