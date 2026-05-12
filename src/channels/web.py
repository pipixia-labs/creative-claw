"""Local browser-based web chat channel for CreativeClaw."""

from __future__ import annotations

import asyncio
import base64
import contextlib
import html
import json
import mimetypes
import re
import tempfile
import uuid
import webbrowser
import zipfile
from collections.abc import Awaitable, Callable
from dataclasses import dataclass
from http import HTTPStatus
from importlib import resources
from pathlib import Path, PurePosixPath
from typing import Any
from urllib.parse import parse_qs, quote, unquote, urlparse

from websockets.exceptions import ConnectionClosed
from websockets.datastructures import Headers
from websockets.legacy.server import WebSocketServerProtocol, serve

from conf.channel import WebChannelConfig
from src.logger import logger
from src.productions.design.design_systems import (
    list_design_systems,
    read_design_system,
    resolve_design_system_preview,
)
from src.runtime.cancellation import get_cancellation_manager
from src.runtime.process_sessions import ProcessKillSummary
from src.runtime import InboundMessage, MessageAttachment
from src.runtime.workspace import looks_like_image, resolve_workspace_path, workspace_relative_path

from .base import BaseChannel
from .events import OutboundMessage


STATIC_PACKAGE = "src.webchat.static"
INDEX_FILE = "index.html"
PPTX_PREVIEW_ERROR_TITLE = "PPTX preview unavailable"
PDF_PREVIEW_ERROR_TITLE = "PDF preview unavailable"
UPLOAD_SIZE_LIMIT = 100 * 1024 * 1024
UPLOAD_ROOT = Path(tempfile.gettempdir()) / "creative-claw-web-uploads"
INTERACTIVE_PPT_HTML_KIND = "interactive_ppt_html"
MODEL_MIME_TYPES = {
    ".fbx": "application/octet-stream",
    ".glb": "model/gltf-binary",
    ".gltf": "model/gltf+json",
    ".obj": "model/obj",
    ".stl": "model/stl",
    ".usd": "model/vnd.usd",
    ".usda": "model/vnd.usd",
    ".usdc": "model/vnd.usd",
    ".usdz": "model/vnd.usdz+zip",
}
ASSISTANT_STREAM_CHUNK_SIZE = 96
MODEL_PACKAGE_EXTENSIONS = {".fbx", ".glb", ".gltf", ".obj", ".stl", ".usd", ".usda", ".usdc", ".usdz"}
MODEL_PACKAGE_EXTENSION_PRIORITY = {
    ".glb": 0,
    ".gltf": 1,
    ".obj": 2,
    ".fbx": 3,
    ".usdz": 4,
    ".usd": 5,
    ".usda": 6,
    ".usdc": 7,
    ".stl": 8,
}
MODEL_PACKAGE_MAX_ENTRIES = 2000
MODEL_PACKAGE_MAX_ENTRY_BYTES = 300 * 1024 * 1024
MODEL_PACKAGE_NAME_PATTERN = re.compile(
    r"(^|[._/\-])(3d|hy3d|seed3d|hyper3d|hitem3d|model|mesh)([._/\-]|$)",
    re.IGNORECASE,
)


@dataclass(slots=True)
class _ClientConnection:
    """Connected browser client information."""

    websocket: WebSocketServerProtocol
    session_id: str
    client_id: str


@dataclass(slots=True)
class _PendingUpload:
    """One file upload currently being streamed from a browser client."""

    path: Path
    original_name: str
    mime_type: str
    expected_size: int
    received_size: int = 0


@dataclass(slots=True)
class _ActiveRun:
    """One run currently executing for a Web Chat session."""

    task: asyncio.Task[None]
    run_id: str
    session_id: str
    runtime_session_id: str | None = None
    stopping: bool = False
    cleanup_task: asyncio.Task[ProcessKillSummary | None] | None = None


def _guess_content_type(filename: str) -> str:
    """Guess one response content type from a file name."""
    model_mime_type = MODEL_MIME_TYPES.get(Path(filename).suffix.lower())
    if model_mime_type:
        return model_mime_type
    guessed, _ = mimetypes.guess_type(filename)
    return guessed or "application/octet-stream"


def _artifact_kind_for_workspace_file(relative_path: str, resolved: Path) -> str:
    """Return a browser-facing artifact kind for special workspace files."""
    if resolved.suffix.lower() not in {".html", ".htm"}:
        return ""
    parts = PurePosixPath(relative_path).parts
    if any(part.startswith("ppt_private_skill_step_") for part in parts):
        return INTERACTIVE_PPT_HTML_KIND
    return ""


def _looks_like_3d_model(filename: str) -> bool:
    """Return whether one file name appears to be a 3D model artifact."""
    suffix = Path(filename).suffix.lower()
    return suffix in MODEL_MIME_TYPES or (suffix == ".zip" and _looks_like_3d_package(filename))


def _looks_like_3d_package(filename: str) -> bool:
    """Return whether one zip name looks like a 3D model package."""
    return bool(MODEL_PACKAGE_NAME_PATTERN.search(str(filename or "")))


def _normalize_zip_entry(raw_entry: str) -> str | None:
    """Normalize one zip entry name and reject path traversal."""
    decoded = unquote(str(raw_entry or "").strip()).replace("\\", "/")
    if not decoded:
        return None
    pure = PurePosixPath(decoded)
    if pure.is_absolute() or ".." in pure.parts:
        return None
    normalized = pure.as_posix().strip("/")
    return normalized or None


def _zip_entry_sort_key(entry: dict[str, Any]) -> tuple[int, int, str]:
    """Return model package entry priority for automatic preview selection."""
    extension = Path(str(entry.get("name") or "")).suffix.lower()
    depth = str(entry.get("name") or "").count("/")
    return (MODEL_PACKAGE_EXTENSION_PRIORITY.get(extension, 99), depth, str(entry.get("name") or ""))


def _should_stream_assistant_payload(payload: dict[str, Any]) -> bool:
    """Return whether one browser payload should be sent as assistant text deltas."""
    if payload.get("type") != "assistant_message":
        return False
    text = str(payload.get("content") or "")
    if not text.strip():
        return False
    if payload.get("metadata", {}).get("disable_stream"):
        return False
    return True


def _chunk_text(text: str, size: int) -> list[str]:
    """Split text into stable non-empty chunks for browser streaming."""
    chunk_size = max(1, int(size or ASSISTANT_STREAM_CHUNK_SIZE))
    value = str(text or "")
    return [value[index : index + chunk_size] for index in range(0, len(value), chunk_size)]


def _normalize_static_path(raw_path: str) -> str:
    """Normalize one static asset path and reject traversal."""
    path = raw_path.split("?", 1)[0].split("#", 1)[0] or "/"
    if path == "/":
        return INDEX_FILE
    pure = PurePosixPath(path.lstrip("/"))
    if pure.is_absolute() or ".." in pure.parts:
        return INDEX_FILE
    normalized = pure.as_posix()
    return normalized or INDEX_FILE


def _normalize_workspace_relative_path(raw_path: str) -> str | None:
    """Normalize one workspace asset URL path and reject traversal."""
    decoded = unquote(raw_path or "").strip()
    pure = PurePosixPath(decoded.lstrip("/"))
    if pure.is_absolute() or ".." in pure.parts:
        return None
    normalized = pure.as_posix().strip()
    return normalized or None


def _safe_upload_segment(value: str, *, fallback: str) -> str:
    """Return one filesystem-safe upload path segment."""
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", value.strip())
    cleaned = cleaned.strip("._")
    return cleaned[:120] or fallback


def _html_response(body: str, *, status: HTTPStatus = HTTPStatus.OK) -> tuple[HTTPStatus, list[tuple[str, str]], bytes]:
    """Build one no-cache HTML response."""
    return (
        status,
        [("Content-Type", "text/html; charset=utf-8"), ("Cache-Control", "no-cache")],
        body.encode("utf-8"),
    )


def _json_response(data: Any, *, status: HTTPStatus = HTTPStatus.OK) -> tuple[HTTPStatus, list[tuple[str, str]], bytes]:
    """Build one no-cache JSON response."""
    return (
        status,
        [("Content-Type", "application/json; charset=utf-8"), ("Cache-Control", "no-cache")],
        json.dumps(data, ensure_ascii=False).encode("utf-8"),
    )


def _simple_preview_error(title: str, message: str) -> str:
    """Render a minimal browser-readable preview error page."""
    return f"""<!doctype html>
<html lang="en">
  <head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1">
    <title>{html.escape(title)}</title>
    <style>
      body {{
        margin: 0;
        min-height: 100vh;
        display: grid;
        place-items: center;
        background: #f3f5ef;
        color: #18211d;
        font-family: Avenir Next, Segoe UI, sans-serif;
      }}
      main {{
        width: min(560px, calc(100vw - 48px));
        padding: 28px;
        border: 1px solid rgba(29, 39, 34, 0.14);
        border-radius: 18px;
        background: #fffef9;
        box-shadow: 0 18px 48px rgba(35, 42, 36, 0.12);
      }}
      h1 {{ margin: 0 0 10px; font-size: 22px; }}
      p {{ margin: 0; color: #5b675f; line-height: 1.6; }}
    </style>
  </head>
  <body>
    <main>
      <h1>{html.escape(title)}</h1>
      <p>{html.escape(message)}</p>
    </main>
  </body>
</html>"""


def _pct(value: Any, total: int) -> float:
    """Convert one EMU coordinate into a slide-relative percentage."""
    try:
        numeric = float(value)
    except (TypeError, ValueError):
        numeric = 0.0
    return 0.0 if total <= 0 else max(0.0, numeric / total * 100.0)


def _shape_style(shape: Any, slide_width: int, slide_height: int) -> str:
    """Build absolute-position CSS for one PPTX shape."""
    left = _pct(getattr(shape, "left", 0), slide_width)
    top = _pct(getattr(shape, "top", 0), slide_height)
    width = _pct(getattr(shape, "width", 0), slide_width)
    height = _pct(getattr(shape, "height", 0), slide_height)
    return f"left:{left:.4f}%;top:{top:.4f}%;width:{width:.4f}%;height:{height:.4f}%;"


def _pptx_text_html(shape: Any) -> str:
    """Render text content from one PPTX text shape."""
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return ""
    paragraphs: list[str] = []
    for paragraph in text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs) or paragraph.text
        stripped = text.strip()
        if stripped:
            paragraphs.append(f"<p>{html.escape(stripped)}</p>")
    return "".join(paragraphs)


def _pptx_table_html(shape: Any) -> str:
    """Render one PPTX table shape into HTML."""
    table = getattr(shape, "table", None)
    if table is None:
        return ""
    rows: list[str] = []
    for row in table.rows:
        cells = "".join(f"<td>{html.escape(cell.text.strip())}</td>" for cell in row.cells)
        rows.append(f"<tr>{cells}</tr>")
    return f"<table>{''.join(rows)}</table>" if rows else ""


def _pptx_image_html(shape: Any) -> str:
    """Render one PPTX picture shape as an embedded image."""
    image = getattr(shape, "image", None)
    if image is None:
        return ""
    mime_type = getattr(image, "content_type", None) or _guess_content_type(getattr(image, "filename", "image"))
    encoded = base64.b64encode(image.blob).decode("ascii")
    alt = html.escape(str(getattr(shape, "name", "") or "slide image"))
    return f'<img src="data:{html.escape(mime_type)};base64,{encoded}" alt="{alt}">'


def _render_pptx_shape(shape: Any, slide_width: int, slide_height: int) -> str:
    """Render one PPTX shape into positioned HTML."""
    content = ""
    if getattr(shape, "has_table", False):
        content = _pptx_table_html(shape)
    elif getattr(shape, "has_text_frame", False):
        content = _pptx_text_html(shape)
    if not content:
        content = _pptx_image_html(shape)
    if not content:
        return ""

    style = _shape_style(shape, slide_width, slide_height)
    return f'<div class="shape" style="{style}">{content}</div>'


def _render_pptx_preview_html(pptx_path: Path) -> str:
    """Render a PPTX file into one standalone HTML preview."""
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - dependency is declared by the project.
        raise RuntimeError("python-pptx is required to preview PPTX files.") from exc

    presentation = Presentation(str(pptx_path))
    slide_width = int(presentation.slide_width) or 12192000
    slide_height = int(presentation.slide_height) or 6858000
    aspect_ratio = f"{slide_width} / {slide_height}"
    slide_items: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        shapes = "".join(
            rendered
            for shape in slide.shapes
            if (rendered := _render_pptx_shape(shape, slide_width, slide_height))
        )
        empty = '<div class="empty-slide">No visible text or images on this slide.</div>' if not shapes else ""
        slide_items.append(
            f"""<section class="slide-wrap">
  <div class="slide-label">Slide {index}</div>
  <div class="slide" style="aspect-ratio:{aspect_ratio};">{shapes}{empty}</div>
</section>"""
        )

    title = html.escape(pptx_path.name)
    slides = "\n".join(slide_items) or '<div class="empty-deck">No slides found.</div>'
    return f"""<!doctype html>
<html lang="en">
  <head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1">
    <title>{title}</title>
    <style>
      :root {{
        --bg: #eef1ea;
        --paper: #fffef9;
        --ink: #18211d;
        --muted: #68746d;
        --border: rgba(29, 39, 34, 0.16);
      }}
      * {{ box-sizing: border-box; }}
      body {{
        margin: 0;
        padding: 24px;
        background: var(--bg);
        color: var(--ink);
        font-family: Avenir Next, Segoe UI, sans-serif;
      }}
      .deck {{
        max-width: 1180px;
        margin: 0 auto;
        display: grid;
        gap: 22px;
      }}
      .slide-label {{
        margin-bottom: 8px;
        color: var(--muted);
        font-size: 12px;
        font-weight: 700;
        letter-spacing: 0.08em;
        text-transform: uppercase;
      }}
      .slide {{
        position: relative;
        overflow: hidden;
        width: 100%;
        border: 1px solid var(--border);
        border-radius: 16px;
        background: var(--paper);
        box-shadow: 0 18px 46px rgba(35, 42, 36, 0.13);
      }}
      .shape {{
        position: absolute;
        overflow: hidden;
        color: var(--ink);
      }}
      .shape p {{
        margin: 0 0 0.35em;
        font-size: clamp(10px, 1.8vw, 30px);
        line-height: 1.24;
        white-space: pre-wrap;
      }}
      .shape img {{
        width: 100%;
        height: 100%;
        object-fit: contain;
        display: block;
      }}
      .shape table {{
        width: 100%;
        height: 100%;
        border-collapse: collapse;
        font-size: clamp(8px, 1.2vw, 18px);
      }}
      .shape td {{
        border: 1px solid rgba(29, 39, 34, 0.16);
        padding: 0.35em;
        vertical-align: top;
      }}
      .empty-slide,
      .empty-deck {{
        height: 100%;
        display: grid;
        place-items: center;
        color: var(--muted);
        font-size: 14px;
      }}
    </style>
  </head>
  <body>
    <main class="deck">{slides}</main>
  </body>
</html>"""


def _render_pdf_preview_html(pdf_path: Path) -> str:
    """Render a PDF file into one standalone HTML preview."""
    try:
        import fitz
    except ImportError as exc:  # pragma: no cover - dependency is declared by the project environment.
        raise RuntimeError("PyMuPDF is required to preview PDF files.") from exc

    pages: list[str] = []
    with fitz.open(pdf_path) as document:
        for index, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            encoded = base64.b64encode(pixmap.tobytes("png")).decode("ascii")
            pages.append(
                f"""<section class="page-wrap">
  <div class="page-label">Page {index}</div>
  <img class="page-image" src="data:image/png;base64,{encoded}" alt="Page {index}">
</section>"""
            )

    title = html.escape(pdf_path.name)
    body = "\n".join(pages) or '<div class="empty-document">No pages found.</div>'
    return f"""<!doctype html>
<html lang="en">
  <head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1">
    <title>{title}</title>
    <style>
      :root {{
        --bg: #eef1ea;
        --paper: #fffef9;
        --ink: #18211d;
        --muted: #68746d;
        --border: rgba(29, 39, 34, 0.16);
      }}
      * {{ box-sizing: border-box; }}
      body {{
        margin: 0;
        padding: 24px;
        background: var(--bg);
        color: var(--ink);
        font-family: Avenir Next, Segoe UI, sans-serif;
      }}
      .document {{
        max-width: 980px;
        margin: 0 auto;
        display: grid;
        gap: 22px;
      }}
      .page-label {{
        margin-bottom: 8px;
        color: var(--muted);
        font-size: 12px;
        font-weight: 700;
        letter-spacing: 0.08em;
        text-transform: uppercase;
      }}
      .page-image {{
        display: block;
        width: 100%;
        height: auto;
        border: 1px solid var(--border);
        border-radius: 16px;
        background: var(--paper);
        box-shadow: 0 18px 46px rgba(35, 42, 36, 0.13);
      }}
      .empty-document {{
        min-height: 100vh;
        display: grid;
        place-items: center;
        color: var(--muted);
        font-size: 14px;
      }}
    </style>
  </head>
  <body>
    <main class="document">{body}</main>
  </body>
</html>"""


class WebChannel(BaseChannel):
    """Serve one local browser chat surface over HTTP and WebSocket."""

    name = "web"

    def __init__(
        self,
        *,
        config: WebChannelConfig,
        inbound_handler: Callable[[InboundMessage], Awaitable[None]],
    ) -> None:
        super().__init__()
        self.config = config
        self.inbound_handler = inbound_handler
        self._server = None
        self._server_task: asyncio.Task[None] | None = None
        self._started = asyncio.Event()
        self._sessions: dict[str, list[_ClientConnection]] = {}
        self._client_seq = 0
        self._host = config.host
        self._port = config.port
        self._pending_uploads: dict[str, _PendingUpload] = {}
        self._active_runs: dict[str, _ActiveRun] = {}

    @property
    def url(self) -> str:
        """Return the local browser URL."""
        return f"http://{self._host}:{self._port}"

    async def wait_until_started(self) -> None:
        """Wait until the HTTP/WebSocket server is listening."""
        await self._started.wait()

    async def start(self) -> None:
        """Start the Web chat service in the background."""
        if self._running:
            return
        self._running = True
        self._started = asyncio.Event()
        self._server_task = asyncio.create_task(self._run_server(), name="creative-claw-webchat")
        started_waiter = asyncio.create_task(self._started.wait(), name="creative-claw-webchat-started")
        done, pending = await asyncio.wait(
            {started_waiter, self._server_task},
            return_when=asyncio.FIRST_COMPLETED,
        )
        for pending_task in pending:
            pending_task.cancel()
        if self._server_task in done and not self._started.is_set():
            self._running = False
            await self._server_task
            raise RuntimeError("Web chat server stopped before reporting readiness.")
        logger.info("CreativeClaw Web chat listening on {}", self.url)
        if self.config.open_browser:
            with contextlib.suppress(Exception):
                webbrowser.open(self.url)

    async def stop(self) -> None:
        """Stop the Web chat service and close active client sockets."""
        self._running = False
        active_runs = list(self._active_runs.values())
        for active in active_runs:
            self._start_cancel_for_active_run(active, reason="server_shutdown")
        pending_tasks = {active.task for active in active_runs if not active.task.done()}
        if pending_tasks:
            await asyncio.wait(
                pending_tasks,
                timeout=5.0,
                return_when=asyncio.ALL_COMPLETED,
            )
        self._active_runs.clear()

        connections = [conn for conns in self._sessions.values() for conn in conns]
        for conn in connections:
            with contextlib.suppress(Exception):
                await conn.websocket.close(code=1001, reason="server shutdown")
        self._sessions.clear()
        self._cleanup_pending_uploads()

        if self._server is not None:
            self._server.close()
            await self._server.wait_closed()
            self._server = None

        if self._server_task is not None:
            with contextlib.suppress(asyncio.CancelledError):
                await self._server_task
            self._server_task = None

    async def send(self, message: OutboundMessage) -> None:
        """Send one outbound message to browser clients."""
        payload = self._build_event(message)
        if _should_stream_assistant_payload(payload):
            await self._stream_assistant_message(message.chat_id, payload)
            return
        await self._broadcast(message.chat_id, payload)

    async def _run_server(self) -> None:
        """Run the combined HTTP and WebSocket server."""
        try:
            self._server = await serve(
                self._handle_websocket,
                self.config.host,
                self.config.port,
                process_request=self._process_request,
                max_size=2**20,
            )
            socket = self._server.sockets[0]
            host, port = socket.getsockname()[:2]
            self._host = host
            self._port = port
            self._started.set()
            await self._server.wait_closed()
        except Exception as exc:
            logger.opt(exception=exc).error("CreativeClaw Web chat failed during startup or runtime.")
            raise

    def _build_event(self, message: OutboundMessage) -> dict[str, Any]:
        """Convert one outbound message into a browser payload."""
        metadata = dict(message.metadata or {})
        payload_type = "assistant_message"
        if metadata.get("display_style") == "progress":
            payload_type = "progress"
        elif metadata.get("display_style") == "assistant_delta":
            payload_type = "assistant_delta"
        elif str(message.text or "").startswith("Error:"):
            payload_type = "error"

        payload = {
            "type": payload_type,
            "content": message.text,
            "delta": message.text if payload_type == "assistant_delta" else "",
            "format": "markdown",
            "artifacts": self._build_artifacts(message.artifact_paths),
            "metadata": metadata,
        }
        active = self._active_runs.get(message.chat_id)
        if active is not None:
            runtime_session_id = str(metadata.get("session_id") or "").strip()
            if runtime_session_id:
                active.runtime_session_id = runtime_session_id
            payload["runId"] = active.run_id
        return payload

    async def _stream_assistant_message(self, session_id: str, payload: dict[str, Any]) -> None:
        """Broadcast one assistant message as text deltas followed by the final payload."""
        text = str(payload.get("content") or "")
        stream_base = {
            "format": payload.get("format", "markdown"),
            "metadata": dict(payload.get("metadata") or {}),
        }
        if payload.get("runId"):
            stream_base["runId"] = payload["runId"]
        for chunk in _chunk_text(text, ASSISTANT_STREAM_CHUNK_SIZE):
            await self._broadcast(
                session_id,
                {
                    **stream_base,
                    "type": "assistant_delta",
                    "delta": chunk,
                },
            )
            await asyncio.sleep(0)
        final_payload = dict(payload)
        final_payload["streamComplete"] = True
        await self._broadcast(session_id, final_payload)

    def _build_artifacts(self, artifact_paths: list[str]) -> list[dict[str, Any]]:
        """Build browser-facing artifact metadata for one outbound message."""
        artifacts: list[dict[str, Any]] = []
        seen_urls: set[str] = set()
        for raw_path in artifact_paths:
            cleaned_path = str(raw_path or "").strip()
            if not cleaned_path:
                continue
            try:
                resolved = resolve_workspace_path(cleaned_path)
            except Exception:
                continue
            if not resolved.exists() or not resolved.is_file():
                continue
            relative_path = workspace_relative_path(resolved)
            url = f"/workspace/{quote(relative_path)}"
            if url in seen_urls:
                continue
            seen_urls.add(url)
            artifact_kind = _artifact_kind_for_workspace_file(relative_path, resolved)
            artifact = {
                "name": resolved.name,
                "path": relative_path,
                "url": url,
                "isImage": looks_like_image(resolved),
                "is3D": _looks_like_3d_model(resolved.name),
                "mimeType": _guess_content_type(resolved.name),
                "sizeBytes": resolved.stat().st_size,
            }
            if artifact_kind:
                artifact["artifactKind"] = artifact_kind
            artifacts.append(
                artifact
            )
        return artifacts

    async def _broadcast(self, session_id: str, payload: dict[str, Any]) -> None:
        """Broadcast one payload to all browser tabs in the same session."""
        connections = list(self._sessions.get(session_id, []))
        if not connections:
            return
        encoded = json.dumps(payload, ensure_ascii=False)
        stale: list[_ClientConnection] = []
        for conn in connections:
            try:
                await conn.websocket.send(encoded)
            except ConnectionClosed:
                stale.append(conn)
        for conn in stale:
            self._drop_connection(conn)

    async def _send_to(self, conn: _ClientConnection, payload: dict[str, Any]) -> None:
        """Send one JSON payload to a single browser client."""
        await conn.websocket.send(json.dumps(payload, ensure_ascii=False))

    async def _handle_websocket(self, websocket: WebSocketServerProtocol) -> None:
        """Handle one browser websocket session."""
        parsed = urlparse(websocket.path)
        if parsed.path != "/ws":
            await websocket.close(code=1008, reason="invalid path")
            return

        query = parse_qs(parsed.query)
        session_id = (query.get("session_id") or ["default"])[0].strip() or "default"
        client_id = self._next_client_id()
        conn = _ClientConnection(websocket=websocket, session_id=session_id, client_id=client_id)
        self._sessions.setdefault(session_id, []).append(conn)
        await websocket.send(
            json.dumps(
                {
                    "type": "ready",
                    "sessionId": session_id,
                    "clientId": client_id,
                    "title": self.config.title,
                },
                ensure_ascii=False,
            )
        )
        try:
            async for raw in websocket:
                await self._handle_client_message(conn, raw)
        except ConnectionClosed:
            pass
        finally:
            self._drop_connection(conn)

    async def _handle_client_message(self, conn: _ClientConnection, raw: str) -> None:
        """Normalize one browser-originated message into an inbound runtime event."""
        try:
            payload = json.loads(raw)
        except json.JSONDecodeError:
            await conn.websocket.send(
                json.dumps({"type": "error", "message": "Invalid JSON payload"}, ensure_ascii=False)
            )
            return

        event_type = payload.get("type")
        if event_type == "upload_start":
            await self._handle_upload_start(conn, payload)
            return
        if event_type == "upload_chunk":
            await self._handle_upload_chunk(conn, payload)
            return
        if event_type == "upload_finish":
            await self._handle_upload_finish(conn, payload)
            return
        if event_type == "upload_cancel":
            await self._handle_upload_cancel(conn, payload)
            return
        if event_type == "stop":
            await self._handle_stop(conn, payload)
            return

        if event_type != "chat":
            await self._send_to(conn, {"type": "error", "message": "Unsupported event type"})
            return

        await self._dispatch_chat_task(conn, payload)

    async def _dispatch_chat_task(self, conn: _ClientConnection, payload: dict[str, Any]) -> None:
        """Start one Web Chat run without blocking the websocket read loop."""
        existing = self._active_runs.get(conn.session_id)
        if existing is not None and not existing.task.done():
            await self._send_to(
                conn,
                {
                    "type": "error",
                    "code": "task_running",
                    "runId": existing.run_id,
                    "message": "当前任务仍在运行，请先停止后再提交。",
                },
            )
            return

        content = str(payload.get("content") or "").strip()
        if not content:
            await self._send_to(conn, {"type": "error", "message": "Message content is required"})
            return

        try:
            attachments = self._attachments_from_chat_payload(payload)
        except ValueError as exc:
            await self._send_to(conn, {"type": "error", "message": str(exc)})
            return

        run_id = str(payload.get("runId") or "").strip() or uuid.uuid4().hex
        get_cancellation_manager().register_run(
            run_id=run_id,
            channel=self.name,
            chat_id=conn.session_id,
        )
        task = asyncio.create_task(
            self._run_chat_task(conn, content, run_id, attachments),
            name=f"web-chat-{conn.session_id}-{run_id}",
        )
        active = _ActiveRun(task=task, run_id=run_id, session_id=conn.session_id)
        self._active_runs[conn.session_id] = active

        def _cleanup(_task: asyncio.Task[None], *, ref: _ActiveRun = active) -> None:
            if self._active_runs.get(ref.session_id) is ref:
                self._active_runs.pop(ref.session_id, None)

        task.add_done_callback(_cleanup)
        await self._broadcast(conn.session_id, {"type": "task_started", "runId": run_id})

    async def _run_chat_task(
        self,
        conn: _ClientConnection,
        content: str,
        run_id: str,
        attachments: list[MessageAttachment] | None = None,
    ) -> None:
        """Run one chat task and emit lifecycle events."""
        reason = "completed"
        try:
            await self.inbound_handler(
                InboundMessage(
                    channel=self.name,
                    sender_id=conn.client_id,
                    chat_id=conn.session_id,
                    text=content,
                    attachments=list(attachments or []),
                    metadata={"client_id": conn.client_id, "run_id": run_id},
                )
            )
        except asyncio.CancelledError:
            reason = "cancelled"
        except Exception as exc:
            reason = "error"
            logger.opt(exception=exc).warning(
                "Web chat task failed: session_id={} run_id={}",
                conn.session_id,
                run_id,
            )
        finally:
            active = self._active_runs.get(conn.session_id)
            if reason == "cancelled" and active is not None and active.cleanup_task is not None:
                try:
                    summary = await asyncio.wait_for(active.cleanup_task, timeout=3.0)
                    logger.info(
                        "Cancel cleanup finished: session_id={} run_id={} found={} killed={} failed={}",
                        conn.session_id,
                        run_id,
                        summary.found if summary else 0,
                        summary.killed if summary else 0,
                        summary.failed if summary else 0,
                    )
                except asyncio.TimeoutError:
                    logger.error("Cancel cleanup timed out: session_id={} run_id={}", conn.session_id, run_id)
                except Exception as cleanup_exc:
                    logger.opt(exception=cleanup_exc).warning(
                        "Cancel cleanup failed: session_id={} run_id={}",
                        conn.session_id,
                        run_id,
                    )

            active = self._active_runs.get(conn.session_id)
            if active is not None and active.run_id == run_id:
                self._active_runs.pop(conn.session_id, None)

            get_cancellation_manager().complete_run(run_id)
            with contextlib.suppress(Exception):
                await self._broadcast(
                    conn.session_id,
                    {"type": "task_finished", "runId": run_id, "reason": reason},
                )

    def _attachments_from_chat_payload(self, payload: dict[str, Any]) -> list[MessageAttachment]:
        """Normalize browser-uploaded attachment records from one chat payload."""
        raw_attachments = payload.get("attachments") or []
        if not isinstance(raw_attachments, list):
            raise ValueError("Chat attachments must be a list.")
        if len(raw_attachments) > 20:
            raise ValueError("Too many attachments in one message.")

        attachments: list[MessageAttachment] = []
        upload_root = UPLOAD_ROOT.resolve()
        for index, item in enumerate(raw_attachments, start=1):
            if not isinstance(item, dict):
                raise ValueError(f"Attachment {index} is not valid.")
            raw_path = str(item.get("path") or "").strip()
            if not raw_path:
                raise ValueError(f"Attachment {index} is missing a path.")
            try:
                resolved_path = Path(raw_path).resolve()
                resolved_path.relative_to(upload_root)
            except Exception as exc:
                raise ValueError(f"Attachment {index} is not a valid uploaded file.") from exc
            if not resolved_path.is_file():
                raise ValueError(f"Attachment {index} does not exist.")

            attachments.append(
                MessageAttachment(
                    path=str(resolved_path),
                    name=str(item.get("name") or resolved_path.name).strip() or resolved_path.name,
                    mime_type=str(item.get("mimeType") or "").strip(),
                    description=str(item.get("description") or "").strip(),
                )
            )
        return attachments

    async def _handle_stop(self, conn: _ClientConnection, payload: dict[str, Any]) -> None:
        """Stop the currently active run for one Web Chat session."""
        requested_run_id = str(payload.get("runId") or "").strip()
        active = self._active_runs.get(conn.session_id)
        if active is None or active.task.done():
            await self._send_to(
                conn,
                {"type": "task_finished", "runId": requested_run_id, "reason": "cancelled"},
            )
            return

        if requested_run_id and requested_run_id != active.run_id:
            await self._send_to(
                conn,
                {
                    "type": "task_stop_ignored",
                    "runId": requested_run_id,
                    "currentRunId": active.run_id,
                },
            )
            return

        if active.stopping:
            return

        active.stopping = True
        await self._broadcast(conn.session_id, {"type": "task_stopping", "runId": active.run_id})
        self._start_cancel_for_active_run(active, reason="user_stop")

    def _start_cancel_for_active_run(self, active: _ActiveRun, *, reason: str) -> None:
        """Start best-effort cancellation for one active run."""
        if active.cleanup_task is None or active.cleanup_task.done():
            active.cleanup_task = asyncio.create_task(
                self._cancel_background_work(active, reason=reason),
                name=f"web-cancel-{active.session_id}-{active.run_id}",
            )
        if not active.task.done():
            active.task.cancel()

    async def _cancel_background_work(self, active: _ActiveRun, *, reason: str) -> ProcessKillSummary | None:
        """Best-effort cancel local background work for one active run."""
        cancellation = get_cancellation_manager()
        summary = await asyncio.to_thread(
            cancellation.request_cancel_by_run_id,
            active.run_id,
            reason,
        )
        if summary is not None:
            return summary
        if active.runtime_session_id:
            return await asyncio.to_thread(
                cancellation.request_cancel_by_session,
                active.runtime_session_id,
                reason,
            )
        logger.info(
            "Cancel requested before runtime session was known: session_id={} run_id={}",
            active.session_id,
            active.run_id,
        )
        return None

    async def _handle_upload_start(self, conn: _ClientConnection, payload: dict[str, Any]) -> None:
        """Initialize a streamed browser file upload."""
        upload_id = str(payload.get("uploadId") or "").strip()
        original_name = str(payload.get("name") or "attachment").strip()
        mime_type = str(payload.get("mimeType") or "").strip()
        try:
            expected_size = int(payload.get("size") or 0)
        except (TypeError, ValueError):
            expected_size = 0

        if not upload_id or expected_size < 0 or expected_size > UPLOAD_SIZE_LIMIT:
            await self._send_upload_error(conn, upload_id, "Upload is missing an id or exceeds the size limit.")
            return

        key = self._upload_key(conn, upload_id)
        self._cleanup_pending_upload(key)
        safe_session = _safe_upload_segment(conn.session_id, fallback="session")
        safe_upload_id = _safe_upload_segment(upload_id, fallback=uuid.uuid4().hex)
        safe_name = _safe_upload_segment(Path(original_name).name, fallback="attachment")
        upload_dir = UPLOAD_ROOT / safe_session / safe_upload_id
        upload_dir.mkdir(parents=True, exist_ok=True)
        upload_path = upload_dir / safe_name
        upload_path.write_bytes(b"")
        self._pending_uploads[key] = _PendingUpload(
            path=upload_path,
            original_name=original_name,
            mime_type=mime_type,
            expected_size=expected_size,
        )
        await conn.websocket.send(
            json.dumps(
                {
                    "type": "upload_started",
                    "uploadId": upload_id,
                    "name": original_name,
                    "size": expected_size,
                },
                ensure_ascii=False,
            )
        )

    async def _handle_upload_chunk(self, conn: _ClientConnection, payload: dict[str, Any]) -> None:
        """Append one base64 encoded upload chunk to a pending file."""
        upload_id = str(payload.get("uploadId") or "").strip()
        key = self._upload_key(conn, upload_id)
        pending = self._pending_uploads.get(key)
        if pending is None:
            await self._send_upload_error(conn, upload_id, "Upload was not started.")
            return

        try:
            chunk = base64.b64decode(str(payload.get("data") or ""), validate=True)
        except Exception:
            self._cleanup_pending_upload(key)
            await self._send_upload_error(conn, upload_id, "Upload chunk was not valid base64.")
            return

        next_size = pending.received_size + len(chunk)
        if next_size > pending.expected_size or next_size > UPLOAD_SIZE_LIMIT:
            self._cleanup_pending_upload(key)
            await self._send_upload_error(conn, upload_id, "Upload exceeded the expected size.")
            return

        with pending.path.open("ab") as file_obj:
            file_obj.write(chunk)
        pending.received_size = next_size
        await conn.websocket.send(
            json.dumps(
                {
                    "type": "upload_chunk_received",
                    "uploadId": upload_id,
                    "received": pending.received_size,
                },
                ensure_ascii=False,
            )
        )

    async def _handle_upload_finish(self, conn: _ClientConnection, payload: dict[str, Any]) -> None:
        """Finalize a streamed browser file upload and return the local path."""
        upload_id = str(payload.get("uploadId") or "").strip()
        key = self._upload_key(conn, upload_id)
        pending = self._pending_uploads.pop(key, None)
        if pending is None:
            await self._send_upload_error(conn, upload_id, "Upload was not started.")
            return
        if pending.received_size != pending.expected_size:
            self._cleanup_upload_file(pending.path)
            await self._send_upload_error(conn, upload_id, "Upload finished before all bytes were received.")
            return

        await conn.websocket.send(
            json.dumps(
                {
                    "type": "upload_complete",
                    "uploadId": upload_id,
                    "name": pending.original_name,
                    "path": str(pending.path),
                    "size": pending.received_size,
                    "mimeType": pending.mime_type,
                },
                ensure_ascii=False,
            )
        )

    async def _handle_upload_cancel(self, conn: _ClientConnection, payload: dict[str, Any]) -> None:
        """Cancel a pending browser file upload."""
        upload_id = str(payload.get("uploadId") or "").strip()
        self._cleanup_pending_upload(self._upload_key(conn, upload_id))
        await conn.websocket.send(
            json.dumps({"type": "upload_cancelled", "uploadId": upload_id}, ensure_ascii=False)
        )

    async def _send_upload_error(self, conn: _ClientConnection, upload_id: str, message: str) -> None:
        """Send one upload-scoped error to a browser client."""
        await conn.websocket.send(
            json.dumps({"type": "upload_error", "uploadId": upload_id, "message": message}, ensure_ascii=False)
        )

    def _upload_key(self, conn: _ClientConnection, upload_id: str) -> str:
        """Return the internal key for a client-scoped upload."""
        return f"{conn.client_id}:{upload_id}"

    def _cleanup_pending_upload(self, key: str) -> None:
        """Remove one pending upload and its partially written file."""
        pending = self._pending_uploads.pop(key, None)
        if pending is not None:
            self._cleanup_upload_file(pending.path)

    def _cleanup_pending_uploads(self, *, client_id: str | None = None) -> None:
        """Remove pending upload files, optionally only for one client."""
        keys = list(self._pending_uploads)
        for key in keys:
            if client_id is None or key.startswith(f"{client_id}:"):
                self._cleanup_pending_upload(key)

    def _cleanup_upload_file(self, path: Path) -> None:
        """Delete one partial upload file and its empty parent directory."""
        with contextlib.suppress(FileNotFoundError):
            path.unlink()
        with contextlib.suppress(OSError):
            path.parent.rmdir()

    async def _process_request(
        self,
        path: str,
        _request_headers: Headers,
    ) -> tuple[HTTPStatus, list[tuple[str, str]], bytes] | None:
        """Serve static assets and workspace files from the websocket server."""
        parsed = urlparse(path)
        if parsed.path == "/ws":
            return None

        if parsed.path == "/api/design-systems":
            return _json_response({"designSystems": [item.to_dict() for item in list_design_systems()]})

        if parsed.path.startswith("/api/design-systems/"):
            return self._serve_design_system_api(parsed.path.removeprefix("/api/design-systems/"))

        if parsed.path.startswith("/workspace-3d-package/manifest/"):
            return self._serve_model_package_manifest(
                parsed.path.removeprefix("/workspace-3d-package/manifest/")
            )

        if parsed.path.startswith("/workspace-3d-package/file/"):
            return self._serve_model_package_file(
                parsed.path.removeprefix("/workspace-3d-package/file/"),
                parse_qs(parsed.query).get("entry", [""])[0],
            )

        if parsed.path.startswith("/workspace-preview/"):
            return self._serve_workspace_preview(parsed.path.removeprefix("/workspace-preview/"))

        if parsed.path.startswith("/workspace/"):
            return self._serve_workspace_asset(parsed.path.removeprefix("/workspace/"))

        resolved_path = _normalize_static_path(parsed.path)
        if not self._asset_exists(resolved_path):
            return (
                HTTPStatus.NOT_FOUND,
                [("Content-Type", "text/plain; charset=utf-8")],
                b"Not Found",
            )

        body = self._read_asset(resolved_path)
        headers = [
            ("Content-Type", self._content_type_header(resolved_path)),
            ("Cache-Control", "no-cache"),
        ]
        if resolved_path == INDEX_FILE:
            body = body.replace(b"__CREATIVE_CLAW_TITLE__", self.config.title.encode("utf-8"))
        return HTTPStatus.OK, headers, body

    def _serve_design_system_api(self, raw_path: str) -> tuple[HTTPStatus, list[tuple[str, str]], bytes]:
        """Serve design system metadata and preview HTML."""
        parts = [unquote(part) for part in raw_path.split("/") if part]
        if len(parts) != 2:
            return _json_response({"error": "not found"}, status=HTTPStatus.NOT_FOUND)
        system_id, action = parts
        if action == "detail":
            body = read_design_system(system_id)
            if body is None:
                return _json_response({"error": "design system not found"}, status=HTTPStatus.NOT_FOUND)
            return _json_response({"id": system_id, "body": body})

        if action in {"preview", "showcase", "preview-dark"}:
            preview_path = resolve_design_system_preview(system_id, dark=action == "preview-dark")
            if preview_path is None:
                return (
                    HTTPStatus.NOT_FOUND,
                    [("Content-Type", "text/plain; charset=utf-8")],
                    b"Not Found",
                )
            return (
                HTTPStatus.OK,
                [("Content-Type", "text/html; charset=utf-8"), ("Cache-Control", "no-cache")],
                preview_path.read_bytes(),
            )

        return _json_response({"error": "not found"}, status=HTTPStatus.NOT_FOUND)

    def _serve_model_package_manifest(self, raw_relative_path: str) -> tuple[HTTPStatus, list[tuple[str, str]], bytes]:
        """Serve metadata for one zipped 3D model package."""
        package_path = self._resolve_workspace_file(raw_relative_path)
        if package_path is None or package_path.suffix.lower() != ".zip":
            return _json_response({"error": "model package not found"}, status=HTTPStatus.NOT_FOUND)

        try:
            with zipfile.ZipFile(package_path) as archive:
                entries = self._model_package_entries(archive)
        except zipfile.BadZipFile:
            return _json_response({"error": "invalid zip package"}, status=HTTPStatus.UNSUPPORTED_MEDIA_TYPE)
        except ValueError as exc:
            return _json_response({"error": str(exc)}, status=HTTPStatus.REQUEST_ENTITY_TOO_LARGE)

        model_entries = [
            entry for entry in entries if Path(str(entry["name"])).suffix.lower() in MODEL_PACKAGE_EXTENSIONS
        ]
        if not model_entries:
            return _json_response(
                {"error": "zip package does not contain a supported 3D model"},
                status=HTTPStatus.UNSUPPORTED_MEDIA_TYPE,
            )

        selected = sorted(model_entries, key=_zip_entry_sort_key)[0]
        relative_path = workspace_relative_path(package_path)
        file_url = f"/workspace-3d-package/file/{quote(relative_path)}"
        model_entry = str(selected["name"])
        model_url = f"{file_url}?entry={quote(model_entry, safe='')}"
        model_directory = str(PurePosixPath(model_entry).parent)
        if model_directory == ".":
            model_directory = ""

        return _json_response(
            {
                "name": package_path.name,
                "path": relative_path,
                "packageSizeBytes": package_path.stat().st_size,
                "fileUrl": file_url,
                "modelEntry": model_entry,
                "modelUrl": model_url,
                "modelDirectory": model_directory,
                "modelSizeBytes": selected["sizeBytes"],
                "entries": entries[:200],
                "entryCount": len(entries),
            }
        )

    def _serve_model_package_file(
        self,
        raw_relative_path: str,
        raw_entry: str,
    ) -> tuple[HTTPStatus, list[tuple[str, str]], bytes]:
        """Serve one safe entry from a zipped 3D model package."""
        package_path = self._resolve_workspace_file(raw_relative_path)
        entry_name = _normalize_zip_entry(raw_entry)
        if package_path is None or package_path.suffix.lower() != ".zip" or entry_name is None:
            return (
                HTTPStatus.NOT_FOUND,
                [("Content-Type", "text/plain; charset=utf-8")],
                b"Not Found",
            )

        try:
            with zipfile.ZipFile(package_path) as archive:
                info = self._find_zip_entry(archive, entry_name)
                if info is None or info.is_dir():
                    return (
                        HTTPStatus.NOT_FOUND,
                        [("Content-Type", "text/plain; charset=utf-8")],
                        b"Not Found",
                    )
                if info.file_size > MODEL_PACKAGE_MAX_ENTRY_BYTES:
                    return _json_response(
                        {"error": "zip entry is too large to preview inline"},
                        status=HTTPStatus.REQUEST_ENTITY_TOO_LARGE,
                    )
                body = archive.read(info)
        except zipfile.BadZipFile:
            return _json_response({"error": "invalid zip package"}, status=HTTPStatus.UNSUPPORTED_MEDIA_TYPE)

        return (
            HTTPStatus.OK,
            [("Content-Type", self._content_type_header(entry_name)), ("Cache-Control", "no-cache")],
            body,
        )

    def _model_package_entries(self, archive: zipfile.ZipFile) -> list[dict[str, Any]]:
        """Return safe file entries from one zip archive."""
        infos = [info for info in archive.infolist() if not info.is_dir()]
        if len(infos) > MODEL_PACKAGE_MAX_ENTRIES:
            raise ValueError("zip package contains too many files to preview inline")

        entries: list[dict[str, Any]] = []
        seen: set[str] = set()
        for info in infos:
            name = _normalize_zip_entry(info.filename)
            if name is None or name in seen:
                continue
            seen.add(name)
            entries.append(
                {
                    "name": name,
                    "sizeBytes": int(info.file_size),
                    "mimeType": _guess_content_type(name),
                    "isModel": Path(name).suffix.lower() in MODEL_PACKAGE_EXTENSIONS,
                }
            )
        return entries

    def _find_zip_entry(self, archive: zipfile.ZipFile, entry_name: str) -> zipfile.ZipInfo | None:
        """Find one zip entry by normalized safe name."""
        for info in archive.infolist():
            if _normalize_zip_entry(info.filename) == entry_name:
                return info
        return None

    def _serve_workspace_asset(self, raw_relative_path: str) -> tuple[HTTPStatus, list[tuple[str, str]], bytes]:
        """Serve one file from the CreativeClaw workspace."""
        normalized = _normalize_workspace_relative_path(raw_relative_path)
        if normalized is None:
            return (
                HTTPStatus.NOT_FOUND,
                [("Content-Type", "text/plain; charset=utf-8")],
                b"Not Found",
            )
        try:
            resolved = resolve_workspace_path(normalized)
        except Exception:
            return (
                HTTPStatus.NOT_FOUND,
                [("Content-Type", "text/plain; charset=utf-8")],
                b"Not Found",
            )
        if not resolved.exists() or not resolved.is_file():
            return (
                HTTPStatus.NOT_FOUND,
                [("Content-Type", "text/plain; charset=utf-8")],
                b"Not Found",
            )
        return (
            HTTPStatus.OK,
            [("Content-Type", self._content_type_header(resolved.name)), ("Cache-Control", "no-cache")],
            resolved.read_bytes(),
        )

    def _serve_workspace_preview(self, raw_relative_path: str) -> tuple[HTTPStatus, list[tuple[str, str]], bytes]:
        """Serve one browser-renderable preview for a workspace file."""
        resolved = self._resolve_workspace_file(raw_relative_path)
        if resolved is None:
            return (
                HTTPStatus.NOT_FOUND,
                [("Content-Type", "text/plain; charset=utf-8")],
                b"Not Found",
            )
        suffix = resolved.suffix.lower()
        if suffix == ".pdf":
            try:
                return _html_response(_render_pdf_preview_html(resolved))
            except Exception as exc:
                logger.opt(exception=exc).warning("Failed to render PDF preview for {}", resolved)
                return _html_response(
                    _simple_preview_error(PDF_PREVIEW_ERROR_TITLE, f"Could not render {resolved.name}: {exc}"),
                    status=HTTPStatus.INTERNAL_SERVER_ERROR,
                )
        if suffix != ".pptx":
            return _html_response(
                _simple_preview_error(
                    "Preview unsupported",
                    f"No inline preview is available for {resolved.name}. Open the original file instead.",
                ),
                status=HTTPStatus.UNSUPPORTED_MEDIA_TYPE,
            )

        try:
            return _html_response(_render_pptx_preview_html(resolved))
        except Exception as exc:
            logger.opt(exception=exc).warning("Failed to render PPTX preview for {}", resolved)
            return _html_response(
                _simple_preview_error(PPTX_PREVIEW_ERROR_TITLE, f"Could not render {resolved.name}: {exc}"),
                status=HTTPStatus.INTERNAL_SERVER_ERROR,
            )

    def _resolve_workspace_file(self, raw_relative_path: str) -> Path | None:
        """Resolve one normalized workspace file path for browser routes."""
        normalized = _normalize_workspace_relative_path(raw_relative_path)
        if normalized is None:
            return None
        try:
            resolved = resolve_workspace_path(normalized)
        except Exception:
            return None
        if not resolved.exists() or not resolved.is_file():
            return None
        return resolved

    def _asset_exists(self, relative_path: str) -> bool:
        """Return whether one packaged static asset exists."""
        return resources.files(STATIC_PACKAGE).joinpath(relative_path).is_file()

    def _read_asset(self, relative_path: str) -> bytes:
        """Read one packaged static asset."""
        return resources.files(STATIC_PACKAGE).joinpath(relative_path).read_bytes()

    def _content_type_header(self, filename: str) -> str:
        """Build one response content type header."""
        content_type = _guess_content_type(filename)
        if filename.endswith((".html", ".js", ".css", ".json", ".svg")):
            return f"{content_type}; charset=utf-8"
        return content_type

    def _next_client_id(self) -> str:
        """Generate the next browser client identifier."""
        self._client_seq += 1
        return f"web-client-{self._client_seq}"

    def _drop_connection(self, conn: _ClientConnection) -> None:
        """Remove one closed browser connection from the session registry."""
        self._cleanup_pending_uploads(client_id=conn.client_id)
        session_connections = self._sessions.get(conn.session_id)
        if not session_connections:
            return
        self._sessions[conn.session_id] = [item for item in session_connections if item is not conn]
        if not self._sessions[conn.session_id]:
            self._sessions.pop(conn.session_id, None)
