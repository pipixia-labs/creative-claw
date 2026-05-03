"""Conversion helpers for the AnythingToMD expert."""

from __future__ import annotations

import base64
import mimetypes
import os
import re
import shutil
import zipfile
from dataclasses import dataclass
from datetime import date, datetime, time
from pathlib import Path
from typing import Any
from urllib.parse import unquote, urljoin, urlparse

from src.agents.experts.basic_operations_helpers import (
    build_error_output,
    normalize_optional_int,
)
from src.runtime.workspace import (
    build_generated_output_path,
    build_workspace_file_record,
    resolve_workspace_path,
    workspace_relative_path,
)

_EXPERT_NAME = "AnythingToMD"

_PPT_FORMATS = {".pptx", ".pptm", ".ppsx", ".ppsm", ".potx", ".potm"}
_EXCEL_FORMATS = {".xlsx", ".xlsm"}
_DOC_FORMATS = {".docx", ".html", ".htm"}
_PLAIN_TEXT_FORMATS = {".txt", ".md", ".markdown", ".csv", ".json", ".yaml", ".yml", ".toml", ".xml"}
_SUPPORTED_PRIMARY_FORMATS = (
    {".pdf"} | _PPT_FORMATS | _EXCEL_FORMATS | _DOC_FORMATS | _PLAIN_TEXT_FORMATS
)


@dataclass(frozen=True, slots=True)
class ConversionResult:
    """A successful Markdown conversion result."""

    markdown: str
    method: str
    source_label: str


def convert_anything_to_markdown(parameters: dict[str, Any]) -> dict[str, Any]:
    """Convert one workspace file or URL to Markdown and return runtime output."""
    try:
        output_path = str(parameters.get("output_path", "") or "").strip()
        url = str(parameters.get("url", "") or "").strip()
        if url:
            output_file = _resolve_output_path(parameters, output_path, source_name=_url_stem(url))
            primary_error = ""
            try:
                conversion = _convert_url_primary(url, output_file)
            except Exception as exc:
                primary_error = f"{type(exc).__name__}: {exc}"
                conversion = None
            if conversion is None:
                conversion = _convert_with_markitdown(url)
            if conversion is None:
                detail = f" Primary error: {primary_error}." if primary_error else ""
                raise ValueError(
                    "URL conversion failed. Install markitdown or provide a supported HTML page."
                    + detail
                )
            return _write_success_output(parameters, output_file, conversion, input_paths=[], source_url=url)

        input_path = str(parameters.get("input_path", "") or "").strip()
        if not input_path:
            raise ValueError("input_path or url is required.")
        resolved_input = resolve_workspace_path(input_path)
        if not resolved_input.is_file():
            raise FileNotFoundError(f"Input file not found: {input_path}")

        output_file = _resolve_output_path(parameters, output_path, source_name=resolved_input.stem)
        primary_error = ""
        try:
            conversion = _convert_file_primary(
                resolved_input,
                output_file,
                max_rows=normalize_optional_int(parameters.get("max_rows"), "max_rows") or 0,
                max_cols=normalize_optional_int(parameters.get("max_cols"), "max_cols") or 0,
                doc2x_api_key=str(parameters.get("doc2x_api_key", "") or "").strip(),
                doc2x_formula_level=parameters.get("doc2x_formula_level", parameters.get("formula_level")),
            )
        except Exception as exc:
            primary_error = f"{type(exc).__name__}: {exc}"
            conversion = None
        if conversion is None:
            conversion = _convert_with_markitdown(resolved_input)
        if conversion is None:
            supported = ", ".join(sorted(_SUPPORTED_PRIMARY_FORMATS))
            detail = f" Primary error: {primary_error}." if primary_error else ""
            raise ValueError(
                f"Unsupported or failed source format: {resolved_input.suffix.lower()}. "
                f"Primary supports: {supported}. MarkItDown fallback is unavailable or failed."
                + detail
            )
        return _write_success_output(parameters, output_file, conversion, input_paths=[input_path], source_url="")
    except Exception as exc:
        return build_error_output(_EXPERT_NAME, f"{_EXPERT_NAME} failed: {exc}")


def _resolve_output_path(parameters: dict[str, Any], output_path: str, *, source_name: str) -> Path:
    """Resolve the Markdown output path inside the workspace."""
    if output_path:
        resolved = resolve_workspace_path(output_path)
        if resolved.suffix.lower() != ".md":
            resolved = resolved.with_suffix(".md")
        resolved.parent.mkdir(parents=True, exist_ok=True)
        return resolved

    destination = build_generated_output_path(
        session_id=str(parameters.get("__session_id", "") or "default"),
        turn_index=int(parameters.get("__turn_index", 0) or 0),
        step=int(parameters.get("__step", 0) or 0),
        output_type=f"anything_to_md_{_safe_stem(source_name)}",
        index=0,
        extension=".md",
    )
    destination.parent.mkdir(parents=True, exist_ok=True)
    return destination


def _write_success_output(
    parameters: dict[str, Any],
    output_file: Path,
    conversion: ConversionResult,
    *,
    input_paths: list[str],
    source_url: str,
) -> dict[str, Any]:
    """Write the Markdown file and build a normalized success payload."""
    markdown = conversion.markdown.rstrip() + "\n"
    output_file.write_text(markdown, encoding="utf-8")
    relative_output_path = workspace_relative_path(output_file)
    artifact_name = output_file.name
    turn_index = int(parameters.get("__turn_index", 0) or 0)
    step = int(parameters.get("__step", 0) or 0)
    expert_step = int(parameters.get("__expert_step", 0) or 0)
    return {
        "status": "success",
        "message": f"{_EXPERT_NAME} converted source to Markdown: {artifact_name}.",
        "message_for_user": f"{_EXPERT_NAME} converted source to Markdown.",
        "output_text": markdown,
        "results": {
            "input_paths": input_paths,
            "source_url": source_url,
            "output_path": relative_output_path,
            "method": conversion.method,
            "source_label": conversion.source_label,
        },
        "output_files": [
            build_workspace_file_record(
                output_file,
                description=f"{_EXPERT_NAME} Markdown output generated from {conversion.source_label}.",
                source="expert",
                name=artifact_name,
                turn=turn_index,
                step=step,
                expert_step=expert_step,
            )
        ],
    }


def _convert_file_primary(
    input_file: Path,
    output_file: Path,
    *,
    max_rows: int,
    max_cols: int,
    doc2x_api_key: str,
    doc2x_formula_level: Any,
) -> ConversionResult | None:
    """Run the source_to_md-style primary converter for one local file."""
    suffix = input_file.suffix.lower()
    if suffix in _PLAIN_TEXT_FORMATS:
        return ConversionResult(
            markdown=_convert_plain_text(input_file),
            method="primary:plain_text",
            source_label=input_file.name,
        )
    if suffix in _PPT_FORMATS:
        return ConversionResult(
            markdown=_convert_pptx(input_file, output_file),
            method="primary:pptx",
            source_label=input_file.name,
        )
    if suffix in _EXCEL_FORMATS:
        return ConversionResult(
            markdown=_convert_excel(input_file, max_rows=max_rows, max_cols=max_cols),
            method="primary:excel",
            source_label=input_file.name,
        )
    if suffix == ".pdf":
        doc2x_result = _convert_pdf_with_doc2x_v3(
            input_file,
            output_file,
            api_key=doc2x_api_key,
            formula_level=doc2x_formula_level,
        )
        if doc2x_result is not None:
            return doc2x_result
        markdown = _convert_pdf(input_file, output_file)
        return ConversionResult(markdown=markdown, method="primary:pdf", source_label=input_file.name)
    if suffix == ".docx":
        markdown = _convert_docx(input_file, output_file)
        return ConversionResult(markdown=markdown, method="primary:docx", source_label=input_file.name)
    if suffix in {".html", ".htm"}:
        markdown = _convert_html(input_file, output_file)
        return ConversionResult(markdown=markdown, method="primary:html", source_label=input_file.name)
    return None


def _convert_with_markitdown(source: str | Path) -> ConversionResult | None:
    """Use the optional MarkItDown package as a fallback converter."""
    try:
        from markitdown import MarkItDown
    except Exception:
        return None
    try:
        result = MarkItDown().convert(source)
    except Exception:
        return None
    markdown = str(getattr(result, "text_content", "") or "").strip()
    if not markdown:
        return None
    return ConversionResult(markdown=markdown, method="fallback:markitdown", source_label=str(source))


def _convert_plain_text(input_file: Path) -> str:
    """Convert a text-like file to Markdown without semantic changes."""
    text = input_file.read_text(encoding="utf-8", errors="replace")
    if input_file.suffix.lower() in {".md", ".markdown"}:
        return text
    return "\n".join([f"# {input_file.stem}", "", "```", text.rstrip(), "```"])


def _convert_pptx(input_file: Path, output_file: Path) -> str:
    """Convert a PowerPoint file into Markdown text, tables, notes, and images."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for PowerPoint conversion.") from exc

    presentation = Presentation(str(input_file))
    asset_dir = output_file.parent / f"{output_file.stem}_files"
    lines = [
        f"# {input_file.stem}",
        "",
        f"- Source: `{input_file.name}`",
        f"- Total slides: {len(presentation.slides)}",
        "",
    ]
    image_count = 0
    asset_dir_used = False
    for slide_index, slide in enumerate(presentation.slides, 1):
        lines.extend([f"## Slide {slide_index}", ""])
        blocks: list[str] = []
        for shape in _iter_ppt_leaf_shapes(slide.shapes, MSO_SHAPE_TYPE):
            if getattr(shape, "has_table", False):
                table_md = _ppt_table_to_markdown(shape.table)
                if table_md:
                    blocks.append(table_md)
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                asset_dir.mkdir(parents=True, exist_ok=True)
                image_count += 1
                filename = _save_ppt_picture(shape, asset_dir, slide_index, image_count)
                if filename:
                    asset_dir_used = True
                    blocks.append(f"![Slide {slide_index} Image {image_count}]({asset_dir.name}/{filename})")
                continue
            if getattr(shape, "has_text_frame", False):
                text_md = _ppt_text_frame_to_markdown(shape.text_frame)
                if text_md:
                    blocks.append(text_md)
        lines.extend([("\n\n".join(blocks) if blocks else "_No extractable text content._"), ""])
        notes = _extract_ppt_notes(slide, MSO_SHAPE_TYPE)
        if notes:
            lines.extend(["### Speaker Notes", "", notes, ""])
    if not asset_dir_used and asset_dir.exists():
        shutil.rmtree(asset_dir)
    return "\n".join(lines).strip()


def _iter_ppt_leaf_shapes(shapes: object, shape_type_enum: Any) -> list[Any]:
    """Return PowerPoint shapes in stable top-left reading order."""
    items: list[tuple[int, int, Any]] = []
    for shape in shapes:
        if shape.shape_type == shape_type_enum.GROUP:
            for child in _iter_ppt_leaf_shapes(shape.shapes, shape_type_enum):
                items.append((int(getattr(child, "top", 0) or 0), int(getattr(child, "left", 0) or 0), child))
            continue
        items.append((int(getattr(shape, "top", 0) or 0), int(getattr(shape, "left", 0) or 0), shape))
    return [shape for _, _, shape in sorted(items, key=lambda item: (item[0], item[1]))]


def _ppt_text_frame_to_markdown(text_frame: object) -> str:
    """Convert a PowerPoint text frame to simple Markdown."""
    paragraphs = []
    visible = [paragraph for paragraph in text_frame.paragraphs if _normalize_text(paragraph.text)]
    list_like = len(visible) > 1 or any(paragraph.level > 0 for paragraph in visible)
    for paragraph in visible:
        text = _normalize_text(paragraph.text)
        if list_like:
            paragraphs.append(f"{'  ' * max(paragraph.level, 0)}- {text}")
        else:
            paragraphs.append(text)
    return "\n".join(paragraphs) if list_like else "\n\n".join(paragraphs)


def _ppt_table_to_markdown(table: object) -> str:
    """Convert a PowerPoint table to Markdown."""
    rows = [[_escape_table_cell(cell.text) for cell in row.cells] for row in table.rows]
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    rows = [row + [" "] * (width - len(row)) for row in rows]
    lines = ["| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
    lines.extend("| " + " | ".join(row) + " |" for row in rows[1:])
    return "\n".join(lines)


def _save_ppt_picture(shape: object, asset_dir: Path, slide_index: int, image_index: int) -> str | None:
    """Save one PowerPoint picture and return its asset filename."""
    try:
        image = shape.image
    except Exception:
        return None
    ext = (image.ext or "png").lower()
    filename = f"slide_{slide_index:02d}_image_{image_index:02d}.{ext}"
    (asset_dir / filename).write_bytes(image.blob)
    return filename


def _extract_ppt_notes(slide: object, shape_type_enum: Any) -> str:
    """Extract speaker notes from one PowerPoint slide."""
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return ""
    blocks = []
    for shape in _iter_ppt_leaf_shapes(notes_slide.shapes, shape_type_enum):
        if getattr(shape, "has_text_frame", False):
            text = _ppt_text_frame_to_markdown(shape.text_frame)
            if text:
                blocks.append(text)
    return "\n\n".join(blocks).strip()


def _convert_excel(input_file: Path, *, max_rows: int, max_cols: int) -> str:
    """Convert visible workbook sheets to Markdown tables."""
    try:
        from openpyxl import load_workbook
        from openpyxl.utils import get_column_letter
    except ImportError as exc:
        raise RuntimeError("openpyxl is required for Excel conversion.") from exc
    if max_rows < 0 or max_cols < 0:
        raise ValueError("max_rows and max_cols must be zero or positive integers.")

    workbook = load_workbook(input_file, data_only=True, read_only=False)
    visible_sheets = [sheet for sheet in workbook.worksheets if sheet.sheet_state == "visible"]
    lines = [
        f"# Spreadsheet Source: {input_file.name}",
        "",
        "## Workbook Summary",
        "",
        f"- Sheets: {len(workbook.worksheets)}",
        f"- Visible sheets: {', '.join(sheet.title for sheet in visible_sheets) or 'None'}",
        "",
        "> Note: Formula cells are exported as cached values. This converter does not recalculate formulas.",
        "",
    ]
    for worksheet in visible_sheets:
        rows = list(worksheet.iter_rows(values_only=True))
        bounds = _excel_content_bounds(rows)
        lines.extend([f"## Sheet: {worksheet.title}", ""])
        if bounds is None:
            lines.extend(["_No content found._", ""])
            continue
        min_row, min_col, max_row, max_col = bounds
        row_limit = min(max_row, min_row + max_rows - 1) if max_rows else max_row
        col_limit = min(max_col, min_col + max_cols - 1) if max_cols else max_col
        used_range = f"{get_column_letter(min_col + 1)}{min_row + 1}:{get_column_letter(max_col + 1)}{max_row + 1}"
        lines.extend([
            f"- Used range: {used_range}",
            f"- Rows: {max_row - min_row + 1}",
            f"- Columns: {max_col - min_col + 1}",
            "",
        ])
        table_rows = [
            [rows[row_index][col_index] for col_index in range(min_col, col_limit + 1)]
            for row_index in range(min_row, row_limit + 1)
        ]
        lines.extend([_rows_to_markdown_table(table_rows), ""])
    return "\n".join(lines).rstrip()


def _excel_content_bounds(rows: list[tuple[Any, ...]]) -> tuple[int, int, int, int] | None:
    """Return row/column bounds for non-empty spreadsheet content."""
    min_row = min_col = max_row = max_col = None
    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            if _is_empty_cell(value):
                continue
            min_row = row_index if min_row is None else min(min_row, row_index)
            max_row = row_index if max_row is None else max(max_row, row_index)
            min_col = col_index if min_col is None else min(min_col, col_index)
            max_col = col_index if max_col is None else max(max_col, col_index)
    if min_row is None or min_col is None or max_row is None or max_col is None:
        return None
    return min_row, min_col, max_row, max_col


def _rows_to_markdown_table(rows: list[list[Any]]) -> str:
    """Convert cell values to a Markdown table."""
    if not rows:
        return "_No tabular content found._"
    width = max(len(row) for row in rows)
    formatted = [[_format_cell(value) for value in row + [""] * (width - len(row))] for row in rows]
    lines = ["| " + " | ".join(formatted[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
    lines.extend("| " + " | ".join(row) + " |" for row in formatted[1:])
    return "\n".join(lines)


def _convert_pdf(input_file: Path, output_file: Path) -> str:
    """Convert a PDF to Markdown using PyMuPDF when available."""
    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError("PyMuPDF is required for primary PDF conversion.") from exc
    doc = fitz.open(str(input_file))
    lines = [f"# {input_file.stem}", ""]
    asset_dir = output_file.parent / f"{output_file.stem}_files"
    image_count = 0
    asset_dir_used = False
    for page_number, page in enumerate(doc, 1):
        if page_number > 1:
            lines.extend(["", f"<!-- Page {page_number} -->", ""])
        text = page.get_text("text").strip()
        if text:
            lines.extend([text, ""])
        for image_index, image in enumerate(page.get_images(full=True), 1):
            xref = image[0]
            try:
                pix = fitz.Pixmap(doc, xref)
                if pix.width < 100 or pix.height < 100:
                    continue
                asset_dir.mkdir(parents=True, exist_ok=True)
                image_count += 1
                filename = f"page_{page_number:03d}_image_{image_index:02d}.png"
                output_image = asset_dir / filename
                if pix.alpha:
                    pix.save(str(output_image))
                else:
                    pix.save(str(output_image))
                asset_dir_used = True
                lines.append(f"![Page {page_number} Image {image_count}]({asset_dir.name}/{filename})")
            except Exception:
                continue
    if not asset_dir_used and asset_dir.exists():
        shutil.rmtree(asset_dir)
    return "\n".join(lines).strip()


def _convert_pdf_with_doc2x_v3(
    input_file: Path,
    output_file: Path,
    *,
    api_key: str,
    formula_level: Any,
) -> ConversionResult | None:
    """Convert PDF to Markdown with Doc2X v3 through the optional pdfdeal SDK."""
    resolved_api_key = api_key or os.environ.get("DOC2X_API_KEY", "")
    if not resolved_api_key:
        return None
    formula_level_value = normalize_optional_int(formula_level, "doc2x_formula_level")
    try:
        from pdfdeal import Doc2X
    except Exception:
        return None

    output_dir = output_file.parent / f"{output_file.stem}_doc2x"
    output_dir.mkdir(parents=True, exist_ok=True)
    try:
        client = Doc2X(apikey=resolved_api_key, debug=False, thread=1)
        success, failed, has_error = client.pdf2file(
            pdf_file=str(input_file),
            output_path=str(output_dir),
            output_names=[output_file.name],
            output_format="md",
            model="v3-2026",
            formula_level=0 if formula_level_value is None else formula_level_value,
        )
    except Exception:
        return None

    if has_error or not success:
        return None
    markdown_path = _resolve_doc2x_markdown_path(success[0])
    if markdown_path is None:
        return None
    markdown = markdown_path.read_text(encoding="utf-8", errors="replace").strip()
    if not markdown:
        return None
    return ConversionResult(
        markdown=markdown,
        method="primary:doc2x_v3",
        source_label=input_file.name,
    )


def _resolve_doc2x_markdown_path(result_path: Any) -> Path | None:
    """Return a Markdown file path from pdfdeal output, extracting zip archives when needed."""
    path = Path(str(result_path or "")).expanduser()
    if not path.exists():
        return None
    if path.suffix.lower() in {".md", ".markdown"}:
        return path
    if path.suffix.lower() != ".zip":
        return None

    extract_dir = path.with_suffix("")
    extract_dir.mkdir(parents=True, exist_ok=True)
    try:
        _extract_zip_safely(path, extract_dir)
    except Exception:
        return None
    markdown_files = sorted(
        [item for item in extract_dir.rglob("*") if item.is_file() and item.suffix.lower() in {".md", ".markdown"}],
        key=lambda item: len(item.parts),
    )
    return markdown_files[0] if markdown_files else None


def _extract_zip_safely(zip_path: Path, destination: Path) -> None:
    """Extract a zip archive while rejecting absolute or parent-traversal paths."""
    destination_root = destination.resolve()
    with zipfile.ZipFile(zip_path) as archive:
        for member in archive.infolist():
            target = (destination / member.filename).resolve()
            try:
                target.relative_to(destination_root)
            except ValueError:
                continue
            if member.is_dir():
                target.mkdir(parents=True, exist_ok=True)
                continue
            target.parent.mkdir(parents=True, exist_ok=True)
            with archive.open(member) as source, target.open("wb") as sink:
                shutil.copyfileobj(source, sink)


def _convert_docx(input_file: Path, output_file: Path) -> str:
    """Convert DOCX to Markdown and extract embedded images."""
    try:
        import mammoth
    except ImportError as exc:
        raise RuntimeError("mammoth is required for primary DOCX conversion.") from exc
    asset_dir = output_file.parent / f"{output_file.stem}_files"
    rel_asset_dir = asset_dir.name
    counter = {"value": 0}

    def save_image(image: Any) -> dict[str, str]:
        counter["value"] += 1
        ext = mimetypes.guess_extension(image.content_type) or ".bin"
        if ext == ".jpe":
            ext = ".jpg"
        asset_dir.mkdir(parents=True, exist_ok=True)
        filename = f"image_{counter['value']:03d}{ext}"
        with image.open() as stream:
            (asset_dir / filename).write_bytes(stream.read())
        return {"src": f"{rel_asset_dir}/{filename}"}

    with input_file.open("rb") as file:
        result = mammoth.convert_to_markdown(file, convert_image=mammoth.images.img_element(save_image))
    markdown = _html_img_to_md(result.value)
    if asset_dir.exists() and not any(asset_dir.iterdir()):
        asset_dir.rmdir()
    return markdown


def _convert_html(input_file: Path, output_file: Path) -> str:
    """Convert an HTML file to Markdown and copy local/data images."""
    raw_html = input_file.read_text(encoding="utf-8", errors="replace")
    return _html_to_markdown(raw_html, base_dir=input_file.parent, output_file=output_file, source_url="")


def _convert_url_primary(url: str, output_file: Path) -> ConversionResult | None:
    """Fetch an HTML URL and convert its main content to Markdown."""
    try:
        import requests
    except ImportError:
        return None
    try:
        response = requests.get(
            url,
            headers={"User-Agent": "Mozilla/5.0 CreativeClaw AnythingToMD"},
            timeout=30,
        )
        response.raise_for_status()
    except Exception:
        return None
    content_type = response.headers.get("content-type", "").lower()
    if "html" not in content_type and not response.text.lstrip().startswith("<"):
        return None
    markdown = _html_to_markdown(response.text, base_dir=output_file.parent, output_file=output_file, source_url=url)
    return ConversionResult(markdown=markdown, method="primary:web", source_label=url)


def _html_to_markdown(raw_html: str, *, base_dir: Path, output_file: Path, source_url: str) -> str:
    """Convert HTML to Markdown using markdownify and BeautifulSoup."""
    try:
        from bs4 import BeautifulSoup
        from markdownify import markdownify
    except ImportError as exc:
        raise RuntimeError("beautifulsoup4 and markdownify are required for HTML conversion.") from exc
    soup = BeautifulSoup(raw_html, "html.parser")
    for tag in soup(["head", "style", "script", "noscript", "nav", "header", "footer", "aside"]):
        tag.decompose()
    content = _find_main_content(soup)
    asset_dir = output_file.parent / f"{output_file.stem}_files"
    rel_asset_dir = asset_dir.name
    _rewrite_html_images(content, base_dir=base_dir, asset_dir=asset_dir, rel_asset_dir=rel_asset_dir, source_url=source_url)
    markdown = markdownify(str(content or soup), heading_style="ATX", bullets="-")
    return re.sub(r"\n{3,}", "\n\n", markdown).strip()


def _find_main_content(soup: Any) -> Any:
    """Find a likely main content container in a parsed HTML document."""
    candidates = []
    for selector in ("article", "main", "#content", "#article", ".article-content", ".content", ".main-content"):
        candidates.extend(soup.select(selector))
    candidates.append(soup.body or soup)
    return max(candidates, key=lambda element: len(element.get_text(strip=True)), default=soup)


def _rewrite_html_images(
    content: Any,
    *,
    base_dir: Path,
    asset_dir: Path,
    rel_asset_dir: str,
    source_url: str,
) -> None:
    """Extract local, data URI, and URL images from HTML into the Markdown asset directory."""
    if content is None:
        return
    index = 0
    for img in content.find_all("img"):
        src = img.get("src") or img.get("data-src") or img.get("data-original") or ""
        if not src:
            continue
        index += 1
        filename = _copy_or_download_html_image(src, base_dir, asset_dir, index, source_url=source_url)
        if filename:
            img["src"] = f"{rel_asset_dir}/{filename}"


def _copy_or_download_html_image(
    src: str,
    base_dir: Path,
    asset_dir: Path,
    index: int,
    *,
    source_url: str,
) -> str | None:
    """Persist one HTML image source if possible and return its local filename."""
    asset_dir.mkdir(parents=True, exist_ok=True)
    if src.startswith("data:"):
        match = re.match(r"data:(?P<mime>[^;]+);base64,(?P<data>.+)", src)
        if not match:
            return None
        ext = mimetypes.guess_extension(match.group("mime")) or ".bin"
        filename = f"image_{index:03d}{ext}"
        (asset_dir / filename).write_bytes(base64.b64decode(match.group("data")))
        return filename

    parsed = urlparse(src)
    if parsed.scheme in {"http", "https"} or (source_url and src.startswith(("/", "//"))):
        try:
            import requests
            absolute_url = urljoin(source_url, src)
            response = requests.get(absolute_url, timeout=15)
            response.raise_for_status()
        except Exception:
            return None
        ext = Path(urlparse(absolute_url).path).suffix
        if not ext:
            ext = mimetypes.guess_extension(response.headers.get("content-type", "").split(";")[0]) or ".bin"
        filename = f"image_{index:03d}{ext}"
        (asset_dir / filename).write_bytes(response.content)
        return filename

    local_path = Path(unquote(parsed.path if parsed.scheme == "file" else src))
    if not local_path.is_absolute():
        local_path = (base_dir / local_path).resolve()
    if not local_path.is_file():
        return None
    filename = f"image_{index:03d}{local_path.suffix or '.bin'}"
    shutil.copy2(local_path, asset_dir / filename)
    return filename


def _html_img_to_md(markdown_content: str) -> str:
    """Convert leftover image HTML tags in Markdown to Markdown image syntax."""
    return re.sub(
        r'<img\s[^>]*?src="(?P<src>[^"]+)"[^>]*?(?:alt="(?P<alt>[^"]*)")?[^>]*?/?>',
        lambda match: f"![{match.group('alt') or Path(match.group('src')).stem}]({match.group('src')})",
        markdown_content,
    )


def _format_cell(value: Any) -> str:
    """Format one spreadsheet cell for Markdown."""
    if _is_empty_cell(value):
        return ""
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    if isinstance(value, datetime):
        return value.isoformat(sep=" ", timespec="seconds")
    if isinstance(value, date):
        return value.isoformat()
    if isinstance(value, time):
        return value.isoformat(timespec="seconds")
    if isinstance(value, float):
        return _escape_table_cell(f"{value:g}")
    return _escape_table_cell(str(value))


def _is_empty_cell(value: Any) -> bool:
    """Return whether a spreadsheet cell should count as empty."""
    return value is None or (isinstance(value, str) and value.strip() == "")


def _normalize_text(value: str) -> str:
    """Collapse whitespace while preserving paragraph boundaries elsewhere."""
    value = value.replace("\r\n", "\n").replace("\r", "\n")
    lines = [re.sub(r"\s+", " ", line).strip() for line in value.split("\n")]
    return "\n".join(line for line in lines if line)


def _escape_table_cell(value: str) -> str:
    """Escape Markdown table syntax inside a cell."""
    return _normalize_text(value).replace("|", r"\|") or " "


def _url_stem(url: str) -> str:
    """Derive a safe output stem from a URL."""
    parsed = urlparse(url)
    candidate = Path(parsed.path).stem or parsed.netloc or "web_page"
    return _safe_stem(candidate)


def _safe_stem(value: str) -> str:
    """Sanitize a string for a generated filename segment."""
    cleaned = re.sub(r"[^A-Za-z0-9_-]+", "_", str(value or "").strip()).strip("_")
    return cleaned[:80] or "source"
