"""Native SVG to PPTX conversion helpers for the PPT SVG route.

The converter intentionally supports a controlled SVG subset that maps to
editable PowerPoint DrawingML. Unsupported visual features fail before the
public PPTX is replaced.
"""

from __future__ import annotations

import base64
from dataclasses import dataclass, field
import hashlib
import io
import math
import mimetypes
from pathlib import Path
import re
import shutil
import tempfile
import zipfile
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.util import Inches

from src.productions.ppt.schemas import PptSvgExecutionPlan
from src.runtime.workspace import resolve_workspace_path

SVG_NS = "http://www.w3.org/2000/svg"
XLINK_NS = "http://www.w3.org/1999/xlink"
EMU_PER_INCH = 914400
ANGLE_UNIT = 60000

SUPPORTED_SVG_TAGS = frozenset(
    {
        "svg",
        "g",
        "defs",
        "title",
        "desc",
        "metadata",
        "linearGradient",
        "radialGradient",
        "stop",
        "filter",
        "feGaussianBlur",
        "feOffset",
        "feFlood",
        "feComposite",
        "feMerge",
        "feMergeNode",
        "feDropShadow",
        "feFuncA",
        "marker",
        "clipPath",
        "rect",
        "circle",
        "ellipse",
        "line",
        "path",
        "polygon",
        "polyline",
        "text",
        "tspan",
        "image",
    }
)
CONVERTIBLE_VISUAL_TAGS = frozenset(
    {
        "g",
        "rect",
        "circle",
        "ellipse",
        "line",
        "path",
        "polygon",
        "polyline",
        "text",
        "image",
    }
)
NON_VISUAL_TAGS = frozenset(
    {
        "defs",
        "title",
        "desc",
        "metadata",
        "linearGradient",
        "radialGradient",
        "stop",
        "filter",
        "feGaussianBlur",
        "feOffset",
        "feFlood",
        "feComposite",
        "feMerge",
        "feMergeNode",
        "feDropShadow",
        "feFuncA",
        "marker",
        "clipPath",
    }
)
FORBIDDEN_SVG_TAGS = frozenset(
    {
        "style",
        "foreignObject",
        "mask",
        "script",
        "iframe",
        "video",
        "audio",
        "symbol",
        "use",
        "textPath",
        "set",
    }
)
FORBIDDEN_SVG_ATTRS = frozenset({"class", "style", "onclick", "onload", "onmouseover", "onerror"})
SUPPORTED_IMAGE_EXTENSIONS = frozenset({"png", "jpg", "jpeg", "gif", "webp", "bmp", "tif", "tiff", "svg"})
SUPPORTED_PATH_COMMANDS = frozenset({"M", "L", "H", "V", "C", "S", "Q", "T", "A", "Z"})
INHERITABLE_ATTRS = frozenset(
    {
        "fill",
        "stroke",
        "stroke-width",
        "stroke-dasharray",
        "stroke-linecap",
        "stroke-linejoin",
        "opacity",
        "fill-opacity",
        "stroke-opacity",
        "font-family",
        "font-size",
        "font-weight",
        "font-style",
        "text-anchor",
        "text-decoration",
        "letter-spacing",
        "marker-start",
        "marker-end",
    }
)
CONTENT_TYPES = {
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "gif": "image/gif",
    "webp": "image/webp",
    "bmp": "image/bmp",
    "tif": "image/tiff",
    "tiff": "image/tiff",
    "svg": "image/svg+xml",
}
DASH_PRESETS = {
    "1 2": "sysDot",
    "2 2": "sysDot",
    "4 4": "dash",
    "6 3": "dash",
    "8 4": "lgDash",
    "8 4 2 4": "lgDashDot",
}


class PptSvgNativeConversionError(RuntimeError):
    """Raised when SVG cannot be safely converted to editable DrawingML."""


@dataclass(frozen=True)
class PptSvgNativeExportResult:
    """Native SVG export result returned to the route layer."""

    pptx_path: Path
    conversion_report: dict
    warnings: list[str] = field(default_factory=list)


@dataclass
class PathCommand:
    """One normalized path command."""

    cmd: str
    args: list[float] = field(default_factory=list)


@dataclass
class SvgConvertContext:
    """Mutable state for one SVG page conversion."""

    slide_num: int
    svg_dir: Path
    unit_x: float
    unit_y: float
    slide_width_emu: int
    slide_height_emu: int
    defs: dict[str, ET.Element] = field(default_factory=dict)
    id_counter: int = 2
    rel_counter: int = 2
    translate_x: float = 0.0
    translate_y: float = 0.0
    scale_x: float = 1.0
    scale_y: float = 1.0
    inherited_styles: dict[str, str] = field(default_factory=dict)
    media_files: dict[str, bytes] = field(default_factory=dict)
    rel_entries: list[dict[str, str]] = field(default_factory=list)
    shape_count: int = 0
    text_count: int = 0
    image_count: int = 0

    def next_id(self) -> int:
        """Allocate a PowerPoint shape id."""
        shape_id = self.id_counter
        self.id_counter += 1
        return shape_id

    def next_rel_id(self) -> str:
        """Allocate a slide relationship id."""
        rel_id = f"rId{self.rel_counter}"
        self.rel_counter += 1
        return rel_id

    def child(
        self,
        *,
        dx: float = 0.0,
        dy: float = 0.0,
        sx: float = 1.0,
        sy: float = 1.0,
        style_overrides: dict[str, str] | None = None,
    ) -> "SvgConvertContext":
        """Create a child context for a nested group."""
        inherited = dict(self.inherited_styles)
        inherited.update(style_overrides or {})
        return SvgConvertContext(
            slide_num=self.slide_num,
            svg_dir=self.svg_dir,
            unit_x=self.unit_x,
            unit_y=self.unit_y,
            slide_width_emu=self.slide_width_emu,
            slide_height_emu=self.slide_height_emu,
            defs=self.defs,
            id_counter=self.id_counter,
            rel_counter=self.rel_counter,
            translate_x=self.translate_x + dx * self.scale_x,
            translate_y=self.translate_y + dy * self.scale_y,
            scale_x=self.scale_x * sx,
            scale_y=self.scale_y * sy,
            inherited_styles=inherited,
            media_files=self.media_files,
            rel_entries=self.rel_entries,
            shape_count=self.shape_count,
            text_count=self.text_count,
            image_count=self.image_count,
        )

    def sync_from_child(self, child: "SvgConvertContext") -> None:
        """Merge mutable counters from a child conversion context."""
        self.id_counter = child.id_counter
        self.rel_counter = child.rel_counter
        self.shape_count = child.shape_count
        self.text_count = child.text_count
        self.image_count = child.image_count


def export_svg_pages_to_native_pptx(
    *,
    svg_page_paths: list[str],
    pptx_path: Path,
    execution_plan: PptSvgExecutionPlan,
) -> PptSvgNativeExportResult:
    """Convert SVG pages to an editable PPTX by writing native DrawingML XML."""
    if not svg_page_paths:
        raise PptSvgNativeConversionError("No SVG pages were provided for PPTX export.")

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    width_emu, height_emu = _slide_dimensions(execution_plan)
    media_extensions: set[str] = set()
    pages_report: list[dict] = []

    with tempfile.TemporaryDirectory(prefix="ppt_svg_native_", dir=str(pptx_path.parent)) as temp_dir_name:
        temp_dir = Path(temp_dir_name)
        base_pptx = temp_dir / "base.pptx"
        extract_dir = temp_dir / "pptx_content"

        presentation = Presentation()
        presentation.slide_width = width_emu
        presentation.slide_height = height_emu
        blank_layout = presentation.slide_layouts[6]
        for _ in svg_page_paths:
            presentation.slides.add_slide(blank_layout)
        presentation.save(str(base_pptx))

        with zipfile.ZipFile(base_pptx, "r") as package:
            package.extractall(extract_dir)

        media_dir = extract_dir / "ppt" / "media"
        media_dir.mkdir(parents=True, exist_ok=True)

        for slide_num, raw_svg_path in enumerate(svg_page_paths, start=1):
            svg_path = _resolve_svg_path(raw_svg_path)
            slide_xml, ctx = convert_svg_file_to_slide_xml(
                svg_path=svg_path,
                slide_num=slide_num,
                execution_plan=execution_plan,
                width_emu=width_emu,
                height_emu=height_emu,
            )

            slide_xml_path = extract_dir / "ppt" / "slides" / f"slide{slide_num}.xml"
            slide_xml_path.write_text(slide_xml, encoding="utf-8")

            media_name_map: dict[str, str] = {}
            for media_name, media_bytes in ctx.media_files.items():
                ext = media_name.rsplit(".", 1)[-1].lower()
                digest = hashlib.sha256(media_bytes).hexdigest()[:16]
                final_name = f"image_{digest}.{ext}"
                final_path = media_dir / final_name
                if not final_path.exists():
                    final_path.write_bytes(media_bytes)
                media_name_map[media_name] = final_name
                media_extensions.add(ext)

            rels_dir = extract_dir / "ppt" / "slides" / "_rels"
            rels_dir.mkdir(parents=True, exist_ok=True)
            rels_path = rels_dir / f"slide{slide_num}.xml.rels"
            rels_path.write_text(_build_slide_rels_xml(ctx.rel_entries, media_name_map), encoding="utf-8")

            pages_report.append(
                {
                    "slide_number": slide_num,
                    "status": "svg_to_native_drawingml",
                    "editable_level": "high",
                    "shape_count": ctx.shape_count,
                    "text_count": ctx.text_count,
                    "image_count": ctx.image_count,
                    "warnings": [],
                    "errors": [],
                }
            )

        _update_content_types(extract_dir / "[Content_Types].xml", media_extensions)

        result_pptx = temp_dir / "result.pptx"
        with zipfile.ZipFile(result_pptx, "w", zipfile.ZIP_DEFLATED) as package:
            for file_path in extract_dir.rglob("*"):
                if file_path.is_file():
                    package.write(file_path, file_path.relative_to(extract_dir))
        shutil.move(str(result_pptx), str(pptx_path))

    report = {
        "engine": "native_drawingml_svg_converter",
        "requested_strategy": "svg_to_native_drawingml_pptx",
        "final_strategy": "svg_to_native_drawingml_pptx",
        "ok": pptx_path.exists(),
        "fallback_used": False,
        "editable_level": "high",
        "warnings": [],
        "errors": [],
        "pages": pages_report,
    }
    return PptSvgNativeExportResult(pptx_path=pptx_path, conversion_report=report)


def convert_svg_file_to_slide_xml(
    *,
    svg_path: Path,
    slide_num: int,
    execution_plan: PptSvgExecutionPlan,
    width_emu: int,
    height_emu: int,
) -> tuple[str, SvgConvertContext]:
    """Convert one SVG file into complete slide XML plus media relationships."""
    svg_text = svg_path.read_text(encoding="utf-8")
    issues = validate_svg_content(
        svg_text,
        execution_plan=execution_plan,
        svg_dir=svg_path.parent,
        path_label=str(svg_path),
    )
    errors = [issue for issue in issues if issue.get("severity") == "error"]
    if errors:
        detail = "; ".join(str(issue.get("message", "")) for issue in errors[:5])
        raise PptSvgNativeConversionError(f"{svg_path.name}: {detail}")

    root = ET.fromstring(svg_text)
    defs = _collect_defs(root)
    ctx = SvgConvertContext(
        slide_num=slide_num,
        svg_dir=svg_path.parent,
        unit_x=width_emu / max(1, execution_plan.canvas_width),
        unit_y=height_emu / max(1, execution_plan.canvas_height),
        slide_width_emu=width_emu,
        slide_height_emu=height_emu,
        defs=defs,
    )

    shapes_xml: list[str] = []
    for child in list(root):
        if _local_name(child.tag) in NON_VISUAL_TAGS:
            continue
        converted = _convert_element(child, ctx)
        if converted:
            shapes_xml.extend(converted)

    slide_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
       xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm>
          <a:off x="0" y="0"/>
          <a:ext cx="0" cy="0"/>
          <a:chOff x="0" y="0"/>
          <a:chExt cx="0" cy="0"/>
        </a:xfrm>
      </p:grpSpPr>
      {"".join(shapes_xml)}
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sld>'''
    return slide_xml, ctx


def validate_svg_file(
    svg_path: Path,
    *,
    execution_plan: PptSvgExecutionPlan,
) -> list[dict]:
    """Validate one SVG file against the native converter subset."""
    return validate_svg_content(
        svg_path.read_text(encoding="utf-8"),
        execution_plan=execution_plan,
        svg_dir=svg_path.parent,
        path_label=str(svg_path),
    )


def validate_svg_content(
    svg_content: str,
    *,
    execution_plan: PptSvgExecutionPlan,
    svg_dir: Path | None = None,
    path_label: str = "",
) -> list[dict]:
    """Validate raw SVG text against the native converter contract."""
    issues: list[dict] = []
    try:
        root = ET.fromstring(svg_content)
    except Exception as exc:
        return [
            {
                "severity": "error",
                "code": "invalid_svg_xml",
                "message": str(exc),
                "path": path_label,
            }
        ]

    if _local_name(root.tag) != "svg":
        issues.append(_issue("missing_svg_root", "Root tag is not svg.", path_label))
        return issues

    expected_viewbox = f"0 0 {execution_plan.canvas_width} {execution_plan.canvas_height}"
    defs = _collect_defs(root)
    if str(root.attrib.get("viewBox") or "").strip() != expected_viewbox:
        issues.append(
            _issue(
                "viewbox_mismatch",
                f"Expected viewBox `{expected_viewbox}`.",
                path_label,
            )
        )

    def walk(elem: ET.Element, ancestry: list[str], in_defs: bool = False) -> None:
        tag = _local_name(elem.tag)
        current_path = "/".join([*ancestry, tag])
        if tag in FORBIDDEN_SVG_TAGS or tag.startswith("animate"):
            issues.append(
                _issue("forbidden_svg_feature", f"Forbidden SVG tag <{tag}>.", path_label, current_path)
            )
        for attr_name, attr_value in elem.attrib.items():
            clean_attr = _local_name(attr_name)
            if clean_attr in FORBIDDEN_SVG_ATTRS:
                issues.append(
                    _issue(
                        "forbidden_svg_attribute",
                        f"Forbidden SVG attribute `{clean_attr}` on <{tag}>.",
                        path_label,
                        current_path,
                    )
                )
            if clean_attr in {"fill", "stroke"}:
                _validate_paint_value(issues, attr_value, path_label, current_path, clean_attr, defs)
            if clean_attr in {"marker-start", "marker-end"}:
                _validate_marker_ref(issues, elem, clean_attr, defs, path_label, current_path)
            if clean_attr == "clip-path":
                _validate_clip_path_ref(issues, elem, defs, path_label, current_path)
            if clean_attr == "filter":
                _validate_filter_ref(issues, elem, defs, path_label, current_path)
            if clean_attr == "transform" and not _is_supported_transform(str(attr_value)):
                issues.append(
                    _issue(
                        "unsupported_transform",
                        f"Unsupported transform `{attr_value}`.",
                        path_label,
                        current_path,
                    )
                )
            if clean_attr == "transform" and tag == "g" and "rotate(" in str(attr_value):
                issues.append(
                    _issue(
                        "unsupported_transform",
                        "Group rotate transform is not supported by the native converter.",
                        path_label,
                        current_path,
                    )
                )

        if not in_defs and tag not in SUPPORTED_SVG_TAGS:
            issues.append(_issue("unsupported_svg_tag", f"Unsupported SVG tag <{tag}>.", path_label, current_path))
        if not in_defs and tag == "path":
            unsupported = sorted(_unsupported_path_commands(str(elem.attrib.get("d") or "")))
            if unsupported:
                issues.append(
                    _issue(
                        "unsupported_svg_path_command",
                        f"Unsupported SVG path command(s): {', '.join(unsupported)}.",
                        path_label,
                        current_path,
                    )
                )
            else:
                try:
                    normalized_commands = _normalize_path_commands(str(elem.attrib.get("d") or ""))
                except PptSvgNativeConversionError as exc:
                    issues.append(_issue("malformed_svg_path", str(exc), path_label, current_path))
                else:
                    if not normalized_commands:
                        issues.append(
                            _issue(
                                "malformed_svg_path",
                                "SVG path has no drawable commands.",
                                path_label,
                                current_path,
                            )
                        )
        if not in_defs and tag == "image":
            _validate_image_href(issues, elem, svg_dir=svg_dir, path_label=path_label, xml_path=current_path)

        child_in_defs = in_defs or tag == "defs"
        for child in list(elem):
            walk(child, [*ancestry, tag], child_in_defs)

    walk(root, [], False)
    return issues


def _convert_element(elem: ET.Element, ctx: SvgConvertContext) -> list[str]:
    tag = _local_name(elem.tag)
    if tag in NON_VISUAL_TAGS or tag == "tspan":
        return []
    if tag == "g":
        return _convert_group(elem, ctx)
    if tag == "rect":
        return [_convert_rect(elem, ctx)]
    if tag == "circle":
        return [_convert_circle(elem, ctx)]
    if tag == "ellipse":
        return [_convert_ellipse(elem, ctx)]
    if tag == "line":
        return [_convert_line(elem, ctx)]
    if tag == "path":
        return [_convert_path(elem, ctx)]
    if tag == "polygon":
        return [_convert_polygon(elem, ctx, closed=True)]
    if tag == "polyline":
        return [_convert_polygon(elem, ctx, closed=False)]
    if tag == "text":
        converted = _convert_text(elem, ctx)
        return [converted] if converted else []
    if tag == "image":
        converted = _convert_image(elem, ctx)
        return [converted] if converted else []
    raise PptSvgNativeConversionError(f"Unsupported visual SVG element <{tag}>.")


def _convert_group(elem: ET.Element, ctx: SvgConvertContext) -> list[str]:
    dx, dy, sx, sy = _parse_transform(str(elem.attrib.get("transform") or ""))[:4]
    child_ctx = ctx.child(dx=dx, dy=dy, sx=sx, sy=sy, style_overrides=_extract_styles(elem))
    shapes: list[str] = []
    for child in list(elem):
        shapes.extend(_convert_element(child, child_ctx))
    ctx.sync_from_child(child_ctx)
    group_id = str(elem.get("id") or "").strip()
    if not group_id or not shapes:
        return shapes

    shape_id = ctx.next_id()
    ctx.shape_count += 1
    group_effect = _effect_xml(elem, ctx)
    return [
        f'''<p:grpSp>
<p:nvGrpSpPr>
<p:cNvPr id="{shape_id}" name="{_xml_escape(group_id)}"/>
<p:cNvGrpSpPr/><p:nvPr/>
</p:nvGrpSpPr>
<p:grpSpPr>
<a:xfrm>
<a:off x="0" y="0"/>
<a:ext cx="{ctx.slide_width_emu}" cy="{ctx.slide_height_emu}"/>
<a:chOff x="0" y="0"/>
<a:chExt cx="{ctx.slide_width_emu}" cy="{ctx.slide_height_emu}"/>
</a:xfrm>
{group_effect}
</p:grpSpPr>
{"".join(shapes)}
</p:grpSp>'''
    ]


def _convert_rect(elem: ET.Element, ctx: SvgConvertContext) -> str:
    x = _num(elem.get("x"))
    y = _num(elem.get("y"))
    width = _num(elem.get("width"))
    height = _num(elem.get("height"))
    if width <= 0 or height <= 0:
        return ""
    shape_id = ctx.next_id()
    ctx.shape_count += 1
    rx = _num(elem.get("rx"))
    ry = _num(elem.get("ry"), rx)
    geom = '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
    if rx > 0 or ry > 0:
        radius = min(max(rx, ry), min(width, height) / 2)
        adj = max(0, min(50000, int(round(radius / max(1, min(width, height)) * 100000))))
        geom = f'<a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val {adj}"/></a:avLst></a:prstGeom>'
    return _wrap_shape(
        shape_id=shape_id,
        name=f"Rectangle {shape_id}",
        off_x=ctx_x(ctx, x),
        off_y=ctx_y(ctx, y),
        ext_cx=ctx_w(ctx, width),
        ext_cy=ctx_h(ctx, height),
        geom_xml=geom,
        fill_xml=_fill_xml(elem, ctx),
        line_xml=_line_xml(elem, ctx),
        effect_xml=_effect_xml(elem, ctx),
        rot=_rotation_emu(elem),
    )


def _convert_circle(elem: ET.Element, ctx: SvgConvertContext) -> str:
    cx = _num(elem.get("cx"))
    cy = _num(elem.get("cy"))
    radius = _num(elem.get("r"))
    if radius <= 0:
        return ""
    shape_id = ctx.next_id()
    ctx.shape_count += 1
    return _wrap_shape(
        shape_id=shape_id,
        name=f"Ellipse {shape_id}",
        off_x=ctx_x(ctx, cx - radius),
        off_y=ctx_y(ctx, cy - radius),
        ext_cx=ctx_w(ctx, radius * 2),
        ext_cy=ctx_h(ctx, radius * 2),
        geom_xml='<a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>',
        fill_xml=_fill_xml(elem, ctx),
        line_xml=_line_xml(elem, ctx),
        effect_xml=_effect_xml(elem, ctx),
        rot=_rotation_emu(elem),
    )


def _convert_ellipse(elem: ET.Element, ctx: SvgConvertContext) -> str:
    cx = _num(elem.get("cx"))
    cy = _num(elem.get("cy"))
    rx = _num(elem.get("rx"))
    ry = _num(elem.get("ry"))
    if rx <= 0 or ry <= 0:
        return ""
    shape_id = ctx.next_id()
    ctx.shape_count += 1
    return _wrap_shape(
        shape_id=shape_id,
        name=f"Ellipse {shape_id}",
        off_x=ctx_x(ctx, cx - rx),
        off_y=ctx_y(ctx, cy - ry),
        ext_cx=ctx_w(ctx, rx * 2),
        ext_cy=ctx_h(ctx, ry * 2),
        geom_xml='<a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>',
        fill_xml=_fill_xml(elem, ctx),
        line_xml=_line_xml(elem, ctx),
        effect_xml=_effect_xml(elem, ctx),
        rot=_rotation_emu(elem),
    )


def _convert_line(elem: ET.Element, ctx: SvgConvertContext) -> str:
    if _attr(elem, ctx, "marker-start") or _attr(elem, ctx, "marker-end"):
        x1 = _num(elem.get("x1"))
        y1 = _num(elem.get("y1"))
        x2 = _num(elem.get("x2"))
        y2 = _num(elem.get("y2"))
        shape_id = ctx.next_id()
        ctx.shape_count += 1
        min_x = min(ctx_x(ctx, x1), ctx_x(ctx, x2))
        min_y = min(ctx_y(ctx, y1), ctx_y(ctx, y2))
        ext_cx = max(1, abs(ctx_x(ctx, x2) - ctx_x(ctx, x1)))
        ext_cy = max(1, abs(ctx_y(ctx, y2) - ctx_y(ctx, y1)))
        flip_attr = ""
        if x1 > x2 and y1 > y2:
            flip_attr = ' flipH="1" flipV="1"'
        elif x1 > x2:
            flip_attr = ' flipH="1"'
        elif y1 > y2:
            flip_attr = ' flipV="1"'
        rot = _rotation_emu(elem)
        rot_attr = f' rot="{rot}"' if rot else ""
        return f'''<p:sp>
<p:nvSpPr>
<p:cNvPr id="{shape_id}" name="Line {shape_id}"/>
<p:cNvSpPr/><p:nvPr/>
</p:nvSpPr>
<p:spPr>
<a:xfrm{flip_attr}{rot_attr}><a:off x="{min_x}" y="{min_y}"/><a:ext cx="{ext_cx}" cy="{ext_cy}"/></a:xfrm>
<a:prstGeom prst="line"><a:avLst/></a:prstGeom>
<a:noFill/>
{_line_xml(elem, ctx)}
{_effect_xml(elem, ctx)}
</p:spPr>
</p:sp>'''
    commands = [
        PathCommand("M", [_num(elem.get("x1")), _num(elem.get("y1"))]),
        PathCommand("L", [_num(elem.get("x2")), _num(elem.get("y2"))]),
    ]
    return _custom_path_shape(elem, ctx, commands, name_prefix="Line", force_no_fill=True)


def _convert_path(elem: ET.Element, ctx: SvgConvertContext) -> str:
    commands = _normalize_path_commands(str(elem.attrib.get("d") or ""))
    if not commands:
        return ""
    return _custom_path_shape(elem, ctx, commands, name_prefix="Freeform")


def _convert_polygon(elem: ET.Element, ctx: SvgConvertContext, *, closed: bool) -> str:
    points = _parse_points(str(elem.attrib.get("points") or ""))
    if len(points) < 2:
        return ""
    commands = [PathCommand("M", [points[0][0], points[0][1]])]
    commands.extend(PathCommand("L", [x, y]) for x, y in points[1:])
    if closed:
        commands.append(PathCommand("Z", []))
    return _custom_path_shape(elem, ctx, commands, name_prefix="Polygon" if closed else "Polyline", force_no_fill=not closed)


def _custom_path_shape(
    elem: ET.Element,
    ctx: SvgConvertContext,
    commands: list[PathCommand],
    *,
    name_prefix: str,
    force_no_fill: bool = False,
) -> str:
    path_xml, min_x, min_y, width, height = _path_commands_to_drawingml(commands, ctx)
    if not path_xml:
        return ""
    shape_id = ctx.next_id()
    ctx.shape_count += 1
    w_emu = max(1, int(width * ctx.unit_x))
    h_emu = max(1, int(height * ctx.unit_y))
    geom = f'''<a:custGeom>
<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>
<a:rect l="l" t="t" r="r" b="b"/>
<a:pathLst><a:path w="{w_emu}" h="{h_emu}">
{path_xml}
</a:path></a:pathLst>
</a:custGeom>'''
    return _wrap_shape(
        shape_id=shape_id,
        name=f"{name_prefix} {shape_id}",
        off_x=int(min_x * ctx.unit_x),
        off_y=int(min_y * ctx.unit_y),
        ext_cx=w_emu,
        ext_cy=h_emu,
        geom_xml=geom,
        fill_xml="<a:noFill/>" if force_no_fill else _fill_xml(elem, ctx),
        line_xml=_line_xml(elem, ctx),
        effect_xml=_effect_xml(elem, ctx),
        rot=_rotation_emu(elem),
    )


def _convert_text(elem: ET.Element, ctx: SvgConvertContext) -> str:
    blocks = _collect_text_blocks(elem, ctx)
    if not blocks:
        return ""
    return "\n".join(_text_shape_xml(elem, ctx, block) for block in blocks)


def _convert_image(elem: ET.Element, ctx: SvgConvertContext) -> str:
    href = elem.get("href") or elem.get(f"{{{XLINK_NS}}}href")
    if not href:
        return ""
    img_format, img_bytes = _read_image_bytes(str(href), ctx.svg_dir)
    image_index = len(ctx.media_files) + 1
    media_name = f"s{ctx.slide_num}_img{image_index}.{img_format}"
    ctx.media_files[media_name] = img_bytes
    rel_id = ctx.next_rel_id()
    ctx.rel_entries.append(
        {
            "id": rel_id,
            "type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
            "target": f"../media/{media_name}",
        }
    )
    x = _num(elem.get("x"))
    y = _num(elem.get("y"))
    width = _num(elem.get("width"))
    height = _num(elem.get("height"))
    if width <= 0 or height <= 0:
        return ""
    source_size = _read_image_size(img_bytes)
    src_rect_xml = _image_src_rect_xml(elem, source_size, width, height)
    if not src_rect_xml:
        x, y, width, height = _image_meet_frame(elem, source_size, x, y, width, height)
    clip_geom_xml = _clip_geometry_xml(elem, ctx, base_x=x, base_y=y, width=width, height=height)
    shape_id = ctx.next_id()
    ctx.image_count += 1
    rot = _rotation_emu(elem)
    rot_attr = f' rot="{rot}"' if rot else ""
    return f'''<p:pic>
<p:nvPicPr>
<p:cNvPr id="{shape_id}" name="Image {shape_id}"/>
<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
<p:nvPr/>
</p:nvPicPr>
<p:blipFill>
<a:blip r:embed="{rel_id}"/>
{src_rect_xml}
<a:stretch><a:fillRect/></a:stretch>
</p:blipFill>
<p:spPr>
<a:xfrm{rot_attr}><a:off x="{ctx_x(ctx, x)}" y="{ctx_y(ctx, y)}"/><a:ext cx="{ctx_w(ctx, width)}" cy="{ctx_h(ctx, height)}"/></a:xfrm>
{clip_geom_xml}
{_effect_xml(elem, ctx)}
</p:spPr>
</p:pic>'''


def _wrap_shape(
    *,
    shape_id: int,
    name: str,
    off_x: int,
    off_y: int,
    ext_cx: int,
    ext_cy: int,
    geom_xml: str,
    fill_xml: str,
    line_xml: str,
    effect_xml: str = "",
    rot: int = 0,
) -> str:
    rot_attr = f' rot="{rot}"' if rot else ""
    return f'''<p:sp>
<p:nvSpPr>
<p:cNvPr id="{shape_id}" name="{_xml_escape(name)}"/>
<p:cNvSpPr/><p:nvPr/>
</p:nvSpPr>
<p:spPr>
<a:xfrm{rot_attr}><a:off x="{off_x}" y="{off_y}"/><a:ext cx="{max(1, ext_cx)}" cy="{max(1, ext_cy)}"/></a:xfrm>
{geom_xml}
{fill_xml}
{line_xml}
{effect_xml}
</p:spPr>
</p:sp>'''


def _fill_xml(elem: ET.Element, ctx: SvgConvertContext) -> str:
    fill = str(_attr(elem, ctx, "fill") or "#000000").strip()
    if fill.lower() == "none":
        return "<a:noFill/>"
    opacity = _combined_opacity(elem, ctx, "fill-opacity")
    return _paint_xml(fill, opacity, ctx.defs, fallback_no_fill=True)


def _line_xml(elem: ET.Element, ctx: SvgConvertContext) -> str:
    stroke = str(_attr(elem, ctx, "stroke") or "").strip()
    if not stroke or stroke.lower() == "none":
        return "<a:ln><a:noFill/></a:ln>"
    paint_xml = _paint_xml(stroke, _combined_opacity(elem, ctx, "stroke-opacity"), ctx.defs, fallback_no_fill=True)
    if paint_xml == "<a:noFill/>":
        return "<a:ln><a:noFill/></a:ln>"
    width_px = max(0.25, _num(_attr(elem, ctx, "stroke-width"), 1.0))
    dash = _dash_xml(str(_attr(elem, ctx, "stroke-dasharray") or ""))
    cap = {"butt": "flat", "round": "rnd", "square": "sq"}.get(str(_attr(elem, ctx, "stroke-linecap") or "").lower())
    cap_attr = f' cap="{cap}"' if cap else ""
    join = _line_join_xml(str(_attr(elem, ctx, "stroke-linejoin") or ""))
    marker_start = _line_end_xml(elem, ctx, "marker-start")
    marker_end = _line_end_xml(elem, ctx, "marker-end")
    return (
        f'<a:ln w="{max(1, int(width_px * ctx.unit_x))}"{cap_attr}>'
        f"{paint_xml}"
        f"{dash}"
        f"{join}"
        f"{marker_start}"
        f"{marker_end}"
        f"</a:ln>"
    )


def _text_paragraph_xml(
    *,
    line: str,
    font_size: float,
    font_family: str,
    bold: bool,
    color: str,
    align: str,
) -> str:
    color_hex = _normalize_hex_color(color) or "#000000"
    bold_attr = ' b="1"' if bold else ""
    size = max(100, int(round(font_size * 75)))
    escaped_font = _xml_escape(font_family)
    return f'''<a:p>
<a:pPr algn="{align}"/>
<a:r>
<a:rPr lang="en-US" sz="{size}"{bold_attr}>
<a:solidFill><a:srgbClr val="{color_hex[1:]}"/></a:solidFill>
<a:latin typeface="{escaped_font}"/>
<a:ea typeface="{escaped_font}"/>
<a:cs typeface="{escaped_font}"/>
</a:rPr>
<a:t>{_xml_escape(line)}</a:t>
</a:r>
</a:p>'''


def _collect_text_lines(elem: ET.Element) -> list[str]:
    lines: list[str] = []
    direct = (elem.text or "").strip()
    if direct:
        lines.append(direct)
    for child in list(elem):
        if _local_name(child.tag) != "tspan":
            continue
        text = "".join(child.itertext()).strip()
        if not text:
            continue
        dy = _num(child.get("dy"), 0.0)
        has_position_break = bool(child.get("x") or child.get("y") or abs(dy) > 0.01)
        if not lines or has_position_break:
            lines.append(text)
        else:
            lines[-1] += text
    if not lines:
        combined = "".join(elem.itertext()).strip()
        if combined:
            lines.append(combined)
    return lines


def _collect_text_blocks(elem: ET.Element, ctx: SvgConvertContext) -> list[dict[str, object]]:
    base_x = _num(elem.get("x"))
    base_y = _num(elem.get("y"))
    current_y = base_y
    blocks: list[dict[str, object]] = []
    current_block: dict[str, object] = {"x": base_x, "y": base_y, "runs": []}

    def append_run(text: object, source: ET.Element) -> None:
        clean = _normalize_text_space(str(text or ""))
        if not clean:
            return
        runs = current_block["runs"]
        assert isinstance(runs, list)
        runs.append((clean, source))

    append_run(elem.text, elem)
    for child in list(elem):
        if _local_name(child.tag) != "tspan":
            continue
        dy = _num(child.get("dy"), 0.0)
        has_position = bool(child.get("x") or child.get("y") or abs(dy) > 0.01)
        if has_position:
            runs = current_block["runs"]
            if isinstance(runs, list) and runs:
                blocks.append(current_block)
            current_y = _num(child.get("y"), current_y + dy)
            current_block = {
                "x": _num(child.get("x"), base_x),
                "y": current_y,
                "runs": [],
            }
        append_run("".join(child.itertext()), child)
        append_run(child.tail, elem)

    runs = current_block["runs"]
    if isinstance(runs, list) and runs:
        blocks.append(current_block)
    return blocks


def _text_shape_xml(base_elem: ET.Element, ctx: SvgConvertContext, block: dict[str, object]) -> str:
    runs = block.get("runs") or []
    if not isinstance(runs, list) or not runs:
        return ""

    max_font_size = max(
        max(8.0, _num(_text_attr(source, base_elem, ctx, "font-size"), 20.0))
        for _, source in runs
    )
    text_anchor = str(_text_attr(base_elem, base_elem, ctx, "text-anchor") or "start").strip()
    data_width = max(80.0, _num(base_elem.get("data-width"), 720.0))
    line_height = max_font_size * 1.35
    box_w = data_width
    box_h = max(line_height, max_font_size + 10)
    x = float(block.get("x") or 0.0)
    y = max(0.0, float(block.get("y") or 0.0) - max_font_size)
    if text_anchor == "middle":
        x -= box_w / 2
    elif text_anchor == "end":
        x -= box_w
    align = {"middle": "ctr", "end": "r"}.get(text_anchor, "l")

    shape_id = ctx.next_id()
    ctx.text_count += 1
    run_xml = "".join(_text_run_xml(text, source, base_elem, ctx) for text, source in runs)
    rot = _rotation_emu(base_elem)
    rot_attr = f' rot="{rot}"' if rot else ""
    return f'''<p:sp>
<p:nvSpPr>
<p:cNvPr id="{shape_id}" name="TextBox {shape_id}"/>
<p:cNvSpPr txBox="1"/><p:nvPr/>
</p:nvSpPr>
<p:spPr>
<a:xfrm{rot_attr}><a:off x="{ctx_x(ctx, x)}" y="{ctx_y(ctx, y)}"/><a:ext cx="{ctx_w(ctx, box_w)}" cy="{ctx_h(ctx, box_h)}"/></a:xfrm>
<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
<a:noFill/><a:ln><a:noFill/></a:ln>
{_effect_xml(base_elem, ctx)}
</p:spPr>
<p:txBody>
<a:bodyPr wrap="none" lIns="0" tIns="0" rIns="0" bIns="0" anchor="t"><a:spAutoFit/></a:bodyPr>
<a:lstStyle/>
<a:p><a:pPr algn="{align}"/>{run_xml}</a:p>
</p:txBody>
</p:sp>'''


def _text_run_xml(text: str, source: ET.Element, base_elem: ET.Element, ctx: SvgConvertContext) -> str:
    font_size = max(8.0, _num(_text_attr(source, base_elem, ctx, "font-size"), 20.0))
    font_family = _text_attr(source, base_elem, ctx, "font-family") or "Aptos"
    font_weight = str(_text_attr(source, base_elem, ctx, "font-weight") or "").strip().lower()
    font_style = str(_text_attr(source, base_elem, ctx, "font-style") or "").strip().lower()
    decoration = str(_text_attr(source, base_elem, ctx, "text-decoration") or "").strip().lower()
    letter_spacing = _num(_text_attr(source, base_elem, ctx, "letter-spacing"), 0.0)
    bold_attr = ' b="1"' if font_weight in {"700", "800", "900", "bold", "bolder"} else ""
    italic_attr = ' i="1"' if font_style == "italic" else ""
    underline_attr = ' u="sng"' if "underline" in decoration else ""
    strike_attr = ' strike="sngStrike"' if "line-through" in decoration else ""
    spacing_attr = f' spc="{int(round(letter_spacing * 1000))}"' if abs(letter_spacing) > 0.001 else ""
    size = max(100, int(round(font_size * 75)))
    escaped_font = _xml_escape(font_family)
    return f'''<a:r>
<a:rPr lang="en-US" sz="{size}"{bold_attr}{italic_attr}{underline_attr}{strike_attr}{spacing_attr}>
{_text_outline_xml(source, base_elem, ctx)}
{_text_fill_xml(source, base_elem, ctx)}
<a:latin typeface="{escaped_font}"/>
<a:ea typeface="{escaped_font}"/>
<a:cs typeface="{escaped_font}"/>
</a:rPr>
<a:t>{_xml_escape(text)}</a:t>
</a:r>'''


def _text_fill_xml(source: ET.Element, base_elem: ET.Element, ctx: SvgConvertContext) -> str:
    fill = _text_attr(source, base_elem, ctx, "fill") or "#000000"
    opacity = _text_opacity(source, base_elem, ctx, "fill-opacity")
    return _paint_xml(fill, opacity, ctx.defs, fallback_no_fill=False)


def _text_outline_xml(source: ET.Element, base_elem: ET.Element, ctx: SvgConvertContext) -> str:
    stroke = str(_text_attr(source, base_elem, ctx, "stroke") or "").strip()
    if not stroke or stroke.lower() == "none":
        return ""
    paint_xml = _paint_xml(stroke, _text_opacity(source, base_elem, ctx, "stroke-opacity"), ctx.defs, fallback_no_fill=True)
    if paint_xml == "<a:noFill/>":
        return ""
    width_px = max(0.25, _num(_text_attr(source, base_elem, ctx, "stroke-width"), 1.0))
    return f'<a:ln w="{max(1, int(width_px * ctx.unit_x))}">{paint_xml}</a:ln>'


def _text_attr(source: ET.Element, base_elem: ET.Element, ctx: SvgConvertContext, name: str) -> str:
    return str(source.attrib.get(name) or base_elem.attrib.get(name) or ctx.inherited_styles.get(name) or "").strip()


def _text_opacity(source: ET.Element, base_elem: ET.Element, ctx: SvgConvertContext, specific_attr: str) -> float:
    opacity = _num(_text_attr(source, base_elem, ctx, "opacity"), 1.0)
    specific = _num(_text_attr(source, base_elem, ctx, specific_attr), 1.0)
    return max(0.0, min(1.0, opacity * specific))


def _normalize_text_space(text: str) -> str:
    return re.sub(r"\s+", " ", str(text or "")).strip()


def _path_commands_to_drawingml(commands: list[PathCommand], ctx: SvgConvertContext) -> tuple[str, float, float, float, float]:
    transformed_points: list[tuple[float, float]] = []
    for command in commands:
        if command.cmd in {"M", "L"}:
            transformed_points.append(_map_point(ctx, command.args[0], command.args[1]))
        elif command.cmd == "C":
            for index in range(0, 6, 2):
                transformed_points.append(_map_point(ctx, command.args[index], command.args[index + 1]))
    if not transformed_points:
        return "", 0, 0, 0, 0

    min_x = min(point[0] for point in transformed_points)
    min_y = min(point[1] for point in transformed_points)
    max_x = max(point[0] for point in transformed_points)
    max_y = max(point[1] for point in transformed_points)
    width = max(1.0, max_x - min_x)
    height = max(1.0, max_y - min_y)

    parts: list[str] = []
    for command in commands:
        if command.cmd == "M":
            x, y = _map_point(ctx, command.args[0], command.args[1])
            parts.append(f'<a:moveTo><a:pt x="{int((x - min_x) * ctx.unit_x)}" y="{int((y - min_y) * ctx.unit_y)}"/></a:moveTo>')
        elif command.cmd == "L":
            x, y = _map_point(ctx, command.args[0], command.args[1])
            parts.append(f'<a:lnTo><a:pt x="{int((x - min_x) * ctx.unit_x)}" y="{int((y - min_y) * ctx.unit_y)}"/></a:lnTo>')
        elif command.cmd == "C":
            points = []
            for index in range(0, 6, 2):
                x, y = _map_point(ctx, command.args[index], command.args[index + 1])
                points.append(f'<a:pt x="{int((x - min_x) * ctx.unit_x)}" y="{int((y - min_y) * ctx.unit_y)}"/>')
            parts.append(f'<a:cubicBezTo>{"".join(points)}</a:cubicBezTo>')
        elif command.cmd == "Z":
            parts.append("<a:close/>")
    return "\n".join(parts), min_x, min_y, width, height


def _normalize_path_commands(path_data: str) -> list[PathCommand]:
    raw = _path_to_absolute(_parse_path(path_data))
    commands: list[PathCommand] = []
    current_x = current_y = 0.0
    start_x = start_y = 0.0
    last_cubic_control: tuple[float, float] | None = None
    last_quad_control: tuple[float, float] | None = None
    last_cmd = ""

    for item in raw:
        cmd = item.cmd.upper()
        args = item.args
        if cmd == "M":
            x, y = args[0], args[1]
            commands.append(PathCommand("M", [x, y]))
            current_x, current_y = x, y
            start_x, start_y = x, y
            last_cubic_control = None
            last_quad_control = None
        elif cmd == "L":
            x, y = args[0], args[1]
            commands.append(PathCommand("L", [x, y]))
            current_x, current_y = x, y
            last_cubic_control = None
            last_quad_control = None
        elif cmd == "H":
            x = args[0]
            commands.append(PathCommand("L", [x, current_y]))
            current_x = x
            last_cubic_control = None
            last_quad_control = None
        elif cmd == "V":
            y = args[0]
            commands.append(PathCommand("L", [current_x, y]))
            current_y = y
            last_cubic_control = None
            last_quad_control = None
        elif cmd == "C":
            commands.append(PathCommand("C", list(args)))
            last_cubic_control = (args[2], args[3])
            last_quad_control = None
            current_x, current_y = args[4], args[5]
        elif cmd == "S":
            c1x, c1y = (
                _reflect_control_point(last_cubic_control, current_x, current_y)
                if last_cmd in {"C", "S"} and last_cubic_control is not None
                else (current_x, current_y)
            )
            commands.append(PathCommand("C", [c1x, c1y, args[0], args[1], args[2], args[3]]))
            last_cubic_control = (args[0], args[1])
            last_quad_control = None
            current_x, current_y = args[2], args[3]
        elif cmd == "Q":
            cubic = _quad_to_cubic(current_x, current_y, args[0], args[1], args[2], args[3])
            commands.append(PathCommand("C", cubic))
            last_cubic_control = None
            last_quad_control = (args[0], args[1])
            current_x, current_y = args[2], args[3]
        elif cmd == "T":
            qx, qy = (
                _reflect_control_point(last_quad_control, current_x, current_y)
                if last_cmd in {"Q", "T"} and last_quad_control is not None
                else (current_x, current_y)
            )
            cubic = _quad_to_cubic(current_x, current_y, qx, qy, args[0], args[1])
            commands.append(PathCommand("C", cubic))
            last_cubic_control = None
            last_quad_control = (qx, qy)
            current_x, current_y = args[0], args[1]
        elif cmd == "A":
            segments = _arc_to_cubic_beziers(
                current_x,
                current_y,
                args[0],
                args[1],
                args[2],
                bool(int(args[3])),
                bool(int(args[4])),
                args[5],
                args[6],
            )
            if not segments:
                commands.append(PathCommand("L", [args[5], args[6]]))
            else:
                for segment in segments:
                    commands.append(PathCommand("C", segment))
            current_x, current_y = args[5], args[6]
            last_cubic_control = None
            last_quad_control = None
        elif cmd == "Z":
            commands.append(PathCommand("Z", []))
            current_x, current_y = start_x, start_y
            last_cubic_control = None
            last_quad_control = None
        else:
            raise PptSvgNativeConversionError(f"Unsupported path command `{cmd}`.")
        last_cmd = cmd
    return commands


def _path_to_absolute(commands: list[PathCommand]) -> list[PathCommand]:
    absolute: list[PathCommand] = []
    current_x = current_y = 0.0
    start_x = start_y = 0.0
    for item in commands:
        cmd = item.cmd
        args = list(item.args)
        upper = cmd.upper()
        relative = cmd.islower()
        if upper == "M":
            x = args[0] + (current_x if relative else 0.0)
            y = args[1] + (current_y if relative else 0.0)
            absolute.append(PathCommand("M", [x, y]))
            current_x, current_y = x, y
            start_x, start_y = x, y
        elif upper == "L":
            x = args[0] + (current_x if relative else 0.0)
            y = args[1] + (current_y if relative else 0.0)
            absolute.append(PathCommand("L", [x, y]))
            current_x, current_y = x, y
        elif upper == "H":
            x = args[0] + (current_x if relative else 0.0)
            absolute.append(PathCommand("H", [x]))
            current_x = x
        elif upper == "V":
            y = args[0] + (current_y if relative else 0.0)
            absolute.append(PathCommand("V", [y]))
            current_y = y
        elif upper == "C":
            values = [
                args[0] + (current_x if relative else 0.0),
                args[1] + (current_y if relative else 0.0),
                args[2] + (current_x if relative else 0.0),
                args[3] + (current_y if relative else 0.0),
                args[4] + (current_x if relative else 0.0),
                args[5] + (current_y if relative else 0.0),
            ]
            absolute.append(PathCommand("C", values))
            current_x, current_y = values[4], values[5]
        elif upper == "S":
            values = [
                args[0] + (current_x if relative else 0.0),
                args[1] + (current_y if relative else 0.0),
                args[2] + (current_x if relative else 0.0),
                args[3] + (current_y if relative else 0.0),
            ]
            absolute.append(PathCommand("S", values))
            current_x, current_y = values[2], values[3]
        elif upper == "Q":
            values = [
                args[0] + (current_x if relative else 0.0),
                args[1] + (current_y if relative else 0.0),
                args[2] + (current_x if relative else 0.0),
                args[3] + (current_y if relative else 0.0),
            ]
            absolute.append(PathCommand("Q", values))
            current_x, current_y = values[2], values[3]
        elif upper == "T":
            x = args[0] + (current_x if relative else 0.0)
            y = args[1] + (current_y if relative else 0.0)
            absolute.append(PathCommand("T", [x, y]))
            current_x, current_y = x, y
        elif upper == "A":
            x = args[5] + (current_x if relative else 0.0)
            y = args[6] + (current_y if relative else 0.0)
            absolute.append(PathCommand("A", [args[0], args[1], args[2], args[3], args[4], x, y]))
            current_x, current_y = x, y
        elif upper == "Z":
            absolute.append(PathCommand("Z", []))
            current_x, current_y = start_x, start_y
        else:
            raise PptSvgNativeConversionError(f"Unsupported path command `{cmd}`.")
    return absolute


def _parse_path(path_data: str) -> list[PathCommand]:
    tokens = re.findall(r"[A-Za-z]|[-+]?(?:\d+\.?\d*|\.\d+)(?:[eE][-+]?\d+)?", path_data)
    commands: list[PathCommand] = []
    index = 0
    current_cmd = ""
    arg_counts = {"M": 2, "L": 2, "H": 1, "V": 1, "C": 6, "S": 4, "Q": 4, "T": 2, "A": 7, "Z": 0}

    while index < len(tokens):
        token = tokens[index]
        if re.fullmatch(r"[A-Za-z]", token):
            current_cmd = token
            index += 1
        if not current_cmd:
            raise PptSvgNativeConversionError("Path data starts without a command.")
        upper = current_cmd.upper()
        if upper not in arg_counts:
            raise PptSvgNativeConversionError(f"Unsupported path command `{current_cmd}`.")
        count = arg_counts[upper]
        if count == 0:
            commands.append(PathCommand(current_cmd, []))
            current_cmd = ""
            continue
        first_for_move = upper == "M"
        produced = False
        while index + count <= len(tokens) and not re.fullmatch(r"[A-Za-z]", tokens[index]):
            args = [float(tokens[index + offset]) for offset in range(count)]
            cmd = current_cmd
            if first_for_move:
                first_for_move = False
            elif upper == "M":
                cmd = "l" if current_cmd.islower() else "L"
            commands.append(PathCommand(cmd, args))
            index += count
            produced = True
            if index >= len(tokens) or re.fullmatch(r"[A-Za-z]", tokens[index]):
                break
        if not produced:
            raise PptSvgNativeConversionError(f"Path command `{current_cmd}` is missing required arguments.")
    return commands


def _reflect_control_point(
    control: tuple[float, float] | None,
    current_x: float,
    current_y: float,
) -> tuple[float, float]:
    if control is None:
        return current_x, current_y
    return 2 * current_x - control[0], 2 * current_y - control[1]


def _quad_to_cubic(
    start_x: float,
    start_y: float,
    quad_x: float,
    quad_y: float,
    end_x: float,
    end_y: float,
) -> list[float]:
    c1x = start_x + (2.0 / 3.0) * (quad_x - start_x)
    c1y = start_y + (2.0 / 3.0) * (quad_y - start_y)
    c2x = end_x + (2.0 / 3.0) * (quad_x - end_x)
    c2y = end_y + (2.0 / 3.0) * (quad_y - end_y)
    return [c1x, c1y, c2x, c2y, end_x, end_y]


def _arc_to_cubic_beziers(
    start_x: float,
    start_y: float,
    rx: float,
    ry: float,
    x_axis_rotation: float,
    large_arc: bool,
    sweep: bool,
    end_x: float,
    end_y: float,
) -> list[list[float]]:
    if abs(start_x - end_x) < 1e-9 and abs(start_y - end_y) < 1e-9:
        return []
    rx = abs(rx)
    ry = abs(ry)
    if rx < 1e-9 or ry < 1e-9:
        return []

    phi = math.radians(x_axis_rotation % 360)
    cos_phi = math.cos(phi)
    sin_phi = math.sin(phi)
    dx = (start_x - end_x) / 2.0
    dy = (start_y - end_y) / 2.0
    x1p = cos_phi * dx + sin_phi * dy
    y1p = -sin_phi * dx + cos_phi * dy

    radius_scale = (x1p * x1p) / (rx * rx) + (y1p * y1p) / (ry * ry)
    if radius_scale > 1:
        scale = math.sqrt(radius_scale)
        rx *= scale
        ry *= scale

    numerator = rx * rx * ry * ry - rx * rx * y1p * y1p - ry * ry * x1p * x1p
    denominator = rx * rx * y1p * y1p + ry * ry * x1p * x1p
    coefficient = 0.0 if denominator == 0 else math.sqrt(max(0.0, numerator / denominator))
    if large_arc == sweep:
        coefficient = -coefficient
    cxp = coefficient * (rx * y1p / ry)
    cyp = coefficient * (-ry * x1p / rx)
    cx = cos_phi * cxp - sin_phi * cyp + (start_x + end_x) / 2.0
    cy = sin_phi * cxp + cos_phi * cyp + (start_y + end_y) / 2.0

    def vector_angle(u: tuple[float, float], v: tuple[float, float]) -> float:
        dot = u[0] * v[0] + u[1] * v[1]
        det = u[0] * v[1] - u[1] * v[0]
        return math.atan2(det, dot)

    v1 = ((x1p - cxp) / rx, (y1p - cyp) / ry)
    v2 = ((-x1p - cxp) / rx, (-y1p - cyp) / ry)
    theta1 = vector_angle((1.0, 0.0), v1)
    delta_theta = vector_angle(v1, v2)
    if not sweep and delta_theta > 0:
        delta_theta -= 2 * math.pi
    elif sweep and delta_theta < 0:
        delta_theta += 2 * math.pi

    segment_count = max(1, int(math.ceil(abs(delta_theta) / (math.pi / 2))))
    segment_delta = delta_theta / segment_count
    segments: list[list[float]] = []

    def transform(point: tuple[float, float]) -> tuple[float, float]:
        x = point[0]
        y = point[1]
        return (
            cx + rx * (cos_phi * x - sin_phi * y),
            cy + ry * (sin_phi * x + cos_phi * y),
        )

    for index in range(segment_count):
        start_angle = theta1 + index * segment_delta
        end_angle = start_angle + segment_delta
        alpha = 4.0 / 3.0 * math.tan((end_angle - start_angle) / 4.0)
        p1 = (math.cos(start_angle), math.sin(start_angle))
        p2 = (math.cos(end_angle), math.sin(end_angle))
        c1 = (p1[0] - alpha * p1[1], p1[1] + alpha * p1[0])
        c2 = (p2[0] + alpha * p2[1], p2[1] - alpha * p2[0])
        c1x, c1y = transform(c1)
        c2x, c2y = transform(c2)
        px, py = transform(p2)
        segments.append([c1x, c1y, c2x, c2y, px, py])
    if segments:
        segments[-1][4] = end_x
        segments[-1][5] = end_y
    return segments


def _unsupported_path_commands(path_data: str) -> set[str]:
    return {
        token.upper()
        for token in re.findall(r"[A-Za-z]", path_data)
        if token.upper() not in SUPPORTED_PATH_COMMANDS
    }


def _parse_points(points: str) -> list[tuple[float, float]]:
    numbers = [float(item) for item in re.findall(r"[-+]?(?:\d+\.?\d*|\.\d+)", points)]
    return [(numbers[i], numbers[i + 1]) for i in range(0, len(numbers) - 1, 2)]


def _read_image_bytes(href: str, svg_dir: Path) -> tuple[str, bytes]:
    if href.startswith("data:"):
        match = re.match(r"data:image/([A-Za-z0-9.+-]+);base64,(.+)", href, re.DOTALL)
        if not match:
            raise PptSvgNativeConversionError("Unsupported image data URI.")
        ext = _normalize_image_extension(match.group(1))
        return ext, base64.b64decode(match.group(2))
    if re.match(r"^https?://", href, flags=re.IGNORECASE):
        raise PptSvgNativeConversionError("Remote SVG image href is not supported.")

    candidates = [(svg_dir / href).resolve(), (svg_dir.parent / href).resolve()]
    image_path = next((candidate for candidate in candidates if candidate.exists() and candidate.is_file()), None)
    if image_path is None:
        raise PptSvgNativeConversionError(f"Image reference does not exist: {href}")
    ext = _normalize_image_extension(image_path.suffix.lstrip("."))
    return ext, image_path.read_bytes()


def _read_image_size(image_bytes: bytes) -> tuple[int, int] | None:
    try:
        from PIL import Image

        with Image.open(io.BytesIO(image_bytes)) as image:
            return image.size
    except Exception:
        return None


def _image_src_rect_xml(
    elem: ET.Element,
    source_size: tuple[int, int] | None,
    width: float,
    height: float,
) -> str:
    if not source_size:
        return ""
    _, _, mode = _parse_preserve_aspect_ratio(str(elem.get("preserveAspectRatio") or "xMidYMid meet"))
    if mode != "slice":
        return ""
    source_w, source_h = source_size
    if source_w <= 0 or source_h <= 0 or width <= 0 or height <= 0:
        return ""
    align_x, align_y, _ = _parse_preserve_aspect_ratio(str(elem.get("preserveAspectRatio") or "xMidYMid meet"))
    scale = max(width / source_w, height / source_h)
    visible_w = width / scale
    visible_h = height / scale
    crop_x = max(0.0, source_w - visible_w)
    crop_y = max(0.0, source_h - visible_h)
    left = int(round((crop_x * align_x / source_w) * 100000))
    top = int(round((crop_y * align_y / source_h) * 100000))
    right = int(round((crop_x * (1.0 - align_x) / source_w) * 100000))
    bottom = int(round((crop_y * (1.0 - align_y) / source_h) * 100000))
    if not any((left, top, right, bottom)):
        return ""
    return f'<a:srcRect l="{left}" t="{top}" r="{right}" b="{bottom}"/>'


def _image_meet_frame(
    elem: ET.Element,
    source_size: tuple[int, int] | None,
    x: float,
    y: float,
    width: float,
    height: float,
) -> tuple[float, float, float, float]:
    if not source_size:
        return x, y, width, height
    align_x, align_y, mode = _parse_preserve_aspect_ratio(str(elem.get("preserveAspectRatio") or "xMidYMid meet"))
    if mode != "meet":
        return x, y, width, height
    source_w, source_h = source_size
    if source_w <= 0 or source_h <= 0 or width <= 0 or height <= 0:
        return x, y, width, height
    scale = min(width / source_w, height / source_h)
    rendered_w = source_w * scale
    rendered_h = source_h * scale
    return x + (width - rendered_w) * align_x, y + (height - rendered_h) * align_y, rendered_w, rendered_h


def _parse_preserve_aspect_ratio(value: str) -> tuple[float, float, str]:
    clean = str(value or "").strip()
    if not clean:
        clean = "xMidYMid meet"
    if clean == "none":
        return 0.0, 0.0, "none"
    parts = clean.split()
    align = parts[0] if parts else "xMidYMid"
    mode = parts[1] if len(parts) > 1 else "meet"
    align_x = 0.0 if "xMin" in align else 1.0 if "xMax" in align else 0.5
    align_y = 0.0 if "YMin" in align else 1.0 if "YMax" in align else 0.5
    return align_x, align_y, mode if mode in {"meet", "slice"} else "meet"


def _clip_geometry_xml(
    elem: ET.Element,
    ctx: SvgConvertContext,
    *,
    base_x: float,
    base_y: float,
    width: float,
    height: float,
) -> str:
    clip_id = _resolve_url_id(str(_attr(elem, ctx, "clip-path") or ""))
    if not clip_id:
        return '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
    clip = ctx.defs.get(clip_id)
    if clip is None:
        return '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
    shape = next((child for child in list(clip) if _local_name(child.tag) in {"circle", "ellipse", "rect", "path", "polygon"}), None)
    if shape is None:
        return '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
    tag = _local_name(shape.tag)
    if tag in {"circle", "ellipse"}:
        return '<a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>'
    if tag == "rect":
        rx = _num(shape.get("rx"))
        ry = _num(shape.get("ry"), rx)
        commands = _rect_to_path_commands(shape)
        if (rx > 0 or ry > 0) and _clip_rect_covers_box(shape, base_x, base_y, width, height):
            radius = min(max(rx, ry), min(width, height) / 2)
            adj = max(0, min(50000, int(round(radius / max(1, min(width, height)) * 100000))))
            return f'<a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val {adj}"/></a:avLst></a:prstGeom>'
        return _clip_custom_geometry_xml(commands, ctx, base_x=base_x, base_y=base_y, width=width, height=height)
    if tag == "polygon":
        points = _parse_points(str(shape.attrib.get("points") or ""))
        if len(points) < 3:
            return '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        commands = [PathCommand("M", [points[0][0], points[0][1]])]
        commands.extend(PathCommand("L", [x, y]) for x, y in points[1:])
        commands.append(PathCommand("Z", []))
        return _clip_custom_geometry_xml(commands, ctx, base_x=base_x, base_y=base_y, width=width, height=height)
    commands = _normalize_path_commands(str(shape.attrib.get("d") or ""))
    return _clip_custom_geometry_xml(commands, ctx, base_x=base_x, base_y=base_y, width=width, height=height)


def _rect_to_path_commands(elem: ET.Element) -> list[PathCommand]:
    x = _num(elem.get("x"))
    y = _num(elem.get("y"))
    width = _num(elem.get("width"))
    height = _num(elem.get("height"))
    return [
        PathCommand("M", [x, y]),
        PathCommand("L", [x + width, y]),
        PathCommand("L", [x + width, y + height]),
        PathCommand("L", [x, y + height]),
        PathCommand("Z", []),
    ]


def _clip_rect_covers_box(elem: ET.Element, base_x: float, base_y: float, width: float, height: float) -> bool:
    return (
        abs(_num(elem.get("x")) - base_x) < 0.01
        and abs(_num(elem.get("y")) - base_y) < 0.01
        and abs(_num(elem.get("width")) - width) < 0.01
        and abs(_num(elem.get("height")) - height) < 0.01
    )


def _clip_custom_geometry_xml(
    commands: list[PathCommand],
    ctx: SvgConvertContext,
    *,
    base_x: float,
    base_y: float,
    width: float,
    height: float,
) -> str:
    w_emu = ctx_w(ctx, width)
    h_emu = ctx_h(ctx, height)

    def map_clip_point(x: float, y: float) -> tuple[int, int]:
        rel_x = (x - base_x) / max(1.0, width)
        rel_y = (y - base_y) / max(1.0, height)
        return (
            max(0, min(w_emu, int(round(rel_x * w_emu)))),
            max(0, min(h_emu, int(round(rel_y * h_emu)))),
        )

    parts: list[str] = []
    for command in commands:
        if command.cmd == "M":
            x, y = map_clip_point(command.args[0], command.args[1])
            parts.append(f'<a:moveTo><a:pt x="{x}" y="{y}"/></a:moveTo>')
        elif command.cmd == "L":
            x, y = map_clip_point(command.args[0], command.args[1])
            parts.append(f'<a:lnTo><a:pt x="{x}" y="{y}"/></a:lnTo>')
        elif command.cmd == "C":
            points = []
            for index in range(0, 6, 2):
                x, y = map_clip_point(command.args[index], command.args[index + 1])
                points.append(f'<a:pt x="{x}" y="{y}"/>')
            parts.append(f'<a:cubicBezTo>{"".join(points)}</a:cubicBezTo>')
        elif command.cmd == "Z":
            parts.append("<a:close/>")
    if not parts:
        return '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
    return f'''<a:custGeom>
<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>
<a:rect l="l" t="t" r="r" b="b"/>
<a:pathLst><a:path w="{w_emu}" h="{h_emu}">
{"".join(parts)}
</a:path></a:pathLst>
</a:custGeom>'''


def _validate_image_href(
    issues: list[dict],
    elem: ET.Element,
    *,
    svg_dir: Path | None,
    path_label: str,
    xml_path: str,
) -> None:
    href = elem.get("href") or elem.get(f"{{{XLINK_NS}}}href")
    if not href:
        issues.append(_issue("missing_image_href", "SVG image is missing href.", path_label, xml_path))
        return
    if re.match(r"^https?://", str(href), flags=re.IGNORECASE):
        issues.append(_issue("remote_image_href", "Remote image href is not supported.", path_label, xml_path))
        return
    if str(href).startswith("data:"):
        match = re.match(r"data:image/([A-Za-z0-9.+-]+);base64,", str(href))
        ext = _normalize_image_extension(match.group(1)) if match else ""
        if not ext:
            issues.append(_issue("unsupported_image_data_uri", "Unsupported image data URI.", path_label, xml_path))
        return
    ext = _normalize_image_extension(Path(str(href)).suffix.lstrip("."))
    if not ext:
        issues.append(_issue("unsupported_image_format", f"Unsupported image format: {href}", path_label, xml_path))
        return
    if svg_dir is not None:
        candidates = [(svg_dir / str(href)).resolve(), (svg_dir.parent / str(href)).resolve()]
        if not any(candidate.exists() and candidate.is_file() for candidate in candidates):
            issues.append(_issue("missing_image_file", f"Image reference does not exist: {href}", path_label, xml_path))


def _normalize_image_extension(ext: str) -> str:
    clean = str(ext or "").lower().strip()
    if clean == "jpeg":
        return "jpg"
    return clean if clean in SUPPORTED_IMAGE_EXTENSIONS else ""


def _build_slide_rels_xml(rel_entries: list[dict[str, str]], media_name_map: dict[str, str]) -> str:
    extra = ""
    for rel in rel_entries:
        target = rel["target"]
        media_name = target.split("../media/", 1)[-1]
        target = f"../media/{media_name_map.get(media_name, media_name)}"
        extra += (
            f'\n  <Relationship Id="{rel["id"]}" Type="{rel["type"]}" '
            f'Target="{target}"/>'
        )
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>{extra}
</Relationships>'''


def _update_content_types(content_types_path: Path, media_extensions: set[str]) -> None:
    content = content_types_path.read_text(encoding="utf-8")
    additions: list[str] = []
    for ext in sorted(media_extensions):
        if f'Extension="{ext}"' in content:
            continue
        content_type = CONTENT_TYPES.get(ext) or mimetypes.types_map.get(f".{ext}")
        if not content_type:
            raise PptSvgNativeConversionError(f"Unknown media content type for extension: {ext}")
        additions.append(f'  <Default Extension="{ext}" ContentType="{content_type}"/>')
    if additions:
        content = content.replace("</Types>", "\n".join(additions) + "\n</Types>")
        content_types_path.write_text(content, encoding="utf-8")


def _slide_dimensions(execution_plan: PptSvgExecutionPlan) -> tuple[int, int]:
    if execution_plan.aspect_ratio == "4:3":
        return int(Inches(10)), int(Inches(7.5))
    return int(Inches(13.333333)), int(Inches(7.5))


def _resolve_svg_path(path: str) -> Path:
    try:
        resolved = resolve_workspace_path(path)
    except Exception:
        resolved = Path(path)
    if not resolved.exists() or not resolved.is_file():
        raise PptSvgNativeConversionError(f"SVG file does not exist: {path}")
    return resolved


def _extract_styles(elem: ET.Element) -> dict[str, str]:
    return {
        key: str(value)
        for key, value in elem.attrib.items()
        if _local_name(key) in INHERITABLE_ATTRS
    }


def _collect_defs(root: ET.Element) -> dict[str, ET.Element]:
    defs: dict[str, ET.Element] = {}
    for elem in root.iter():
        elem_id = str(elem.attrib.get("id") or "").strip()
        if elem_id:
            defs[elem_id] = elem
    return defs


def _attr(elem: ET.Element, ctx: SvgConvertContext, name: str) -> str:
    return str(elem.attrib.get(name) or ctx.inherited_styles.get(name) or "").strip()


def _combined_opacity(elem: ET.Element, ctx: SvgConvertContext, specific_attr: str) -> float:
    opacity = _num(_attr(elem, ctx, "opacity"), 1.0)
    specific = _num(_attr(elem, ctx, specific_attr), 1.0)
    return max(0.0, min(1.0, opacity * specific))


def _paint_xml(value: str, opacity: float, defs: dict[str, ET.Element], *, fallback_no_fill: bool) -> str:
    clean = str(value or "").strip()
    if not clean or clean.lower() == "none":
        return "<a:noFill/>" if fallback_no_fill else _solid_fill_xml("#000000", opacity)
    paint_id = _resolve_url_id(clean)
    if paint_id:
        definition = defs.get(paint_id)
        if definition is not None and _local_name(definition.tag) in {"linearGradient", "radialGradient"}:
            return _gradient_fill_xml(definition, opacity)
        return "<a:noFill/>" if fallback_no_fill else _solid_fill_xml("#000000", opacity)
    color = _normalize_hex_color(clean)
    if not color:
        return "<a:noFill/>" if fallback_no_fill else _solid_fill_xml("#000000", opacity)
    return _solid_fill_xml(color, opacity)


def _solid_fill_xml(color: str, opacity: float) -> str:
    color_hex = _normalize_hex_color(color) or "#000000"
    alpha = _alpha_xml(opacity)
    return f'<a:solidFill><a:srgbClr val="{color_hex[1:]}">{alpha}</a:srgbClr></a:solidFill>'


def _gradient_fill_xml(gradient: ET.Element, opacity: float) -> str:
    stops = _gradient_stops_xml(gradient, opacity)
    if not stops:
        return _solid_fill_xml("#000000", opacity)
    tag = _local_name(gradient.tag)
    if tag == "radialGradient":
        gradient_mode = '<a:path path="circle"><a:fillToRect l="0" t="0" r="0" b="0"/></a:path>'
    else:
        gradient_mode = f'<a:lin ang="{_linear_gradient_angle(gradient)}" scaled="1"/>'
    return f'<a:gradFill flip="none" rotWithShape="1"><a:gsLst>{"".join(stops)}</a:gsLst>{gradient_mode}</a:gradFill>'


def _gradient_stops_xml(gradient: ET.Element, opacity: float) -> list[str]:
    stops: list[str] = []
    for stop in [child for child in list(gradient) if _local_name(child.tag) == "stop"]:
        color = _normalize_hex_color(str(stop.get("stop-color") or "#000000")) or "#000000"
        stop_opacity = max(0.0, min(1.0, _num(stop.get("stop-opacity"), 1.0) * opacity))
        offset = _gradient_offset(stop.get("offset"))
        stops.append(
            f'<a:gs pos="{offset}"><a:srgbClr val="{color[1:]}">{_alpha_xml(stop_opacity)}</a:srgbClr></a:gs>'
        )
    if len(stops) == 1:
        first_stop = next(child for child in list(gradient) if _local_name(child.tag) == "stop")
        color = _normalize_hex_color(str(first_stop.get("stop-color") or "#000000")) or "#000000"
        stop_opacity = max(0.0, min(1.0, _num(first_stop.get("stop-opacity"), 1.0) * opacity))
        stops.append(
            f'<a:gs pos="100000"><a:srgbClr val="{color[1:]}">{_alpha_xml(stop_opacity)}</a:srgbClr></a:gs>'
        )
    return stops


def _gradient_offset(value: object) -> int:
    clean = str(value or "0").strip()
    if clean.endswith("%"):
        return max(0, min(100000, int(round(_num(clean[:-1]) * 1000))))
    return max(0, min(100000, int(round(_num(clean) * 100000))))


def _linear_gradient_angle(gradient: ET.Element) -> int:
    x1 = _gradient_coord(gradient.get("x1"), 0.0)
    y1 = _gradient_coord(gradient.get("y1"), 0.0)
    x2 = _gradient_coord(gradient.get("x2"), 1.0)
    y2 = _gradient_coord(gradient.get("y2"), 0.0)
    angle = math.degrees(math.atan2(y2 - y1, x2 - x1)) % 360
    return int(round(angle * ANGLE_UNIT))


def _gradient_coord(value: object, default: float) -> float:
    clean = str(value or "").strip()
    if not clean:
        return default
    if clean.endswith("%"):
        return _num(clean[:-1]) / 100.0
    return _num(clean, default)


def _alpha_xml(opacity: float) -> str:
    clean_opacity = max(0.0, min(1.0, opacity))
    return f'<a:alpha val="{int(round(clean_opacity * 100000))}"/>' if clean_opacity < 1 else ""


def _resolve_url_id(value: str) -> str:
    match = re.fullmatch(r"url\(\s*#([^\s)]+)\s*\)", str(value or "").strip())
    return match.group(1) if match else ""


def _line_join_xml(value: str) -> str:
    clean = str(value or "").strip().lower()
    if clean == "round":
        return "<a:round/>"
    if clean == "bevel":
        return "<a:bevel/>"
    if clean == "miter":
        return "<a:miter/>"
    return ""


def _line_end_xml(elem: ET.Element, ctx: SvgConvertContext, attr_name: str) -> str:
    marker_id = _resolve_url_id(str(_attr(elem, ctx, attr_name) or ""))
    if not marker_id:
        return ""
    marker = ctx.defs.get(marker_id)
    marker_type = _classify_marker(marker) if marker is not None else ""
    if not marker_type:
        return ""
    tag = "headEnd" if attr_name == "marker-start" else "tailEnd"
    width = _marker_size_bucket(_num(marker.get("markerWidth"), 3.0) if marker is not None else 3.0)
    length = _marker_size_bucket(_num(marker.get("markerHeight"), 3.0) if marker is not None else 3.0)
    return f'<a:{tag} type="{marker_type}" w="{width}" len="{length}"/>'


def _classify_marker(marker: ET.Element | None) -> str:
    if marker is None or _local_name(marker.tag) != "marker":
        return ""
    children = [child for child in list(marker) if _local_name(child.tag) in {"path", "polygon", "circle", "ellipse"}]
    if any(_local_name(child.tag) in {"circle", "ellipse"} for child in children):
        return "oval"
    for child in children:
        tag = _local_name(child.tag)
        if tag == "polygon":
            points = _parse_points(str(child.attrib.get("points") or ""))
            return "diamond" if len(points) == 4 else "triangle"
        if tag == "path":
            path_data = str(child.attrib.get("d") or "")
            numbers = re.findall(r"[-+]?(?:\d+\.?\d*|\.\d+)", path_data)
            return "diamond" if len(numbers) >= 8 and "Z" in path_data.upper() else "triangle"
    return ""


def _marker_size_bucket(size: float) -> str:
    if size <= 3:
        return "sm"
    if size >= 8:
        return "lg"
    return "med"


def _effect_xml(elem: ET.Element, ctx: SvgConvertContext) -> str:
    filter_id = _resolve_url_id(str(_attr(elem, ctx, "filter") or ""))
    if not filter_id:
        return ""
    filter_elem = ctx.defs.get(filter_id)
    if filter_elem is None or _local_name(filter_elem.tag) != "filter":
        return ""
    shadow = _filter_shadow(filter_elem, ctx)
    if shadow:
        return f"<a:effectLst>{shadow}</a:effectLst>"
    glow = _filter_glow(filter_elem, ctx)
    return f"<a:effectLst>{glow}</a:effectLst>" if glow else ""


def _filter_shadow(filter_elem: ET.Element, ctx: SvgConvertContext) -> str:
    drop_shadow = next((child for child in list(filter_elem) if _local_name(child.tag) == "feDropShadow"), None)
    if drop_shadow is not None:
        dx = _num(drop_shadow.get("dx"), 2.0)
        dy = _num(drop_shadow.get("dy"), 2.0)
        blur = _num(drop_shadow.get("stdDeviation"), 4.0)
        color = _filter_color(drop_shadow)
        opacity = _num(drop_shadow.get("flood-opacity"), 0.35)
        return _shadow_xml(dx, dy, blur, color, opacity, ctx)

    blur_node = next((child for child in list(filter_elem) if _local_name(child.tag) == "feGaussianBlur"), None)
    offset_node = next((child for child in list(filter_elem) if _local_name(child.tag) == "feOffset"), None)
    if blur_node is None or offset_node is None:
        return ""
    flood_node = next((child for child in list(filter_elem) if _local_name(child.tag) == "feFlood"), None)
    dx = _num(offset_node.get("dx"), 2.0)
    dy = _num(offset_node.get("dy"), 2.0)
    blur = _num(blur_node.get("stdDeviation"), 4.0)
    color = _filter_color(flood_node) if flood_node is not None else "#000000"
    opacity = _num(flood_node.get("flood-opacity"), 0.35) if flood_node is not None else 0.35
    return _shadow_xml(dx, dy, blur, color, opacity, ctx)


def _filter_glow(filter_elem: ET.Element, ctx: SvgConvertContext) -> str:
    blur_node = next((child for child in list(filter_elem) if _local_name(child.tag) == "feGaussianBlur"), None)
    if blur_node is None:
        return ""
    radius = max(1, int(_num(blur_node.get("stdDeviation"), 4.0) * (ctx.unit_x + ctx.unit_y) / 2))
    flood_node = next((child for child in list(filter_elem) if _local_name(child.tag) == "feFlood"), None)
    color = _filter_color(flood_node) if flood_node is not None else "#000000"
    opacity = _num(flood_node.get("flood-opacity"), 0.35) if flood_node is not None else 0.35
    color_hex = _normalize_hex_color(color) or "#000000"
    return f'<a:glow rad="{radius}"><a:srgbClr val="{color_hex[1:]}">{_alpha_xml(opacity)}</a:srgbClr></a:glow>'


def _shadow_xml(dx: float, dy: float, blur: float, color: str, opacity: float, ctx: SvgConvertContext) -> str:
    avg_unit = (ctx.unit_x + ctx.unit_y) / 2.0
    blur_rad = max(1, int(abs(blur) * avg_unit))
    distance = max(1, int(math.hypot(dx * ctx.unit_x, dy * ctx.unit_y)))
    direction = int((math.degrees(math.atan2(dy, dx)) % 360) * ANGLE_UNIT)
    color_hex = _normalize_hex_color(color) or "#000000"
    return (
        f'<a:outerShdw blurRad="{blur_rad}" dist="{distance}" dir="{direction}" algn="ctr" rotWithShape="0">'
        f'<a:srgbClr val="{color_hex[1:]}">{_alpha_xml(opacity)}</a:srgbClr>'
        f"</a:outerShdw>"
    )


def _filter_color(elem: ET.Element | None) -> str:
    if elem is None:
        return "#000000"
    return _normalize_hex_color(str(elem.get("flood-color") or "#000000")) or "#000000"


def _rotation_emu(elem: ET.Element) -> int:
    transform = str(elem.attrib.get("transform") or "")
    match = re.search(r"rotate\(\s*([-+]?(?:\d+\.?\d*|\.\d+))", transform)
    if not match:
        return 0
    return int(round(float(match.group(1)) * ANGLE_UNIT)) % (360 * ANGLE_UNIT)


def _validate_paint_value(
    issues: list[dict],
    value: str,
    path_label: str,
    xml_path: str,
    attr_name: str,
    defs: dict[str, ET.Element],
) -> None:
    clean = str(value or "").strip()
    if not clean or clean.lower() == "none":
        return
    paint_id = _resolve_url_id(clean)
    if clean.startswith("url("):
        definition = defs.get(paint_id)
        if definition is None or _local_name(definition.tag) not in {"linearGradient", "radialGradient"}:
            issues.append(_issue("unsupported_paint_server", f"`{attr_name}` references an unsupported paint server.", path_label, xml_path))
            return
        stops = [child for child in list(definition) if _local_name(child.tag) == "stop"]
        if not stops:
            issues.append(_issue("unsupported_paint_server", f"`{attr_name}` gradient has no stops.", path_label, xml_path))
            return
        for stop in stops:
            if not _normalize_hex_color(str(stop.get("stop-color") or "#000000")):
                issues.append(_issue("unsupported_color", f"Gradient stop-color must be #RGB or #RRGGBB.", path_label, xml_path))
                return
        return
    if not _normalize_hex_color(clean):
        issues.append(_issue("unsupported_color", f"`{attr_name}` must be #RGB or #RRGGBB, got `{clean}`.", path_label, xml_path))


def _validate_marker_ref(
    issues: list[dict],
    elem: ET.Element,
    attr_name: str,
    defs: dict[str, ET.Element],
    path_label: str,
    xml_path: str,
) -> None:
    marker_id = _resolve_url_id(str(elem.attrib.get(attr_name) or ""))
    if not marker_id:
        return
    tag = _local_name(elem.tag)
    if tag not in {"line", "path"}:
        issues.append(_issue("unsupported_marker", f"`{attr_name}` is only supported on line/path.", path_label, xml_path))
        return
    marker = defs.get(marker_id)
    if marker is None or _local_name(marker.tag) != "marker" or not _classify_marker(marker):
        issues.append(_issue("unsupported_marker", f"`{attr_name}` references an unsupported marker.", path_label, xml_path))


def _validate_clip_path_ref(
    issues: list[dict],
    elem: ET.Element,
    defs: dict[str, ET.Element],
    path_label: str,
    xml_path: str,
) -> None:
    clip_id = _resolve_url_id(str(elem.attrib.get("clip-path") or ""))
    if not clip_id:
        return
    if _local_name(elem.tag) != "image":
        issues.append(_issue("unsupported_clip_path", "`clip-path` is only supported on image elements.", path_label, xml_path))
        return
    clip = defs.get(clip_id)
    if clip is None or _local_name(clip.tag) != "clipPath":
        issues.append(_issue("unsupported_clip_path", "`clip-path` references an unsupported clipPath.", path_label, xml_path))
        return
    shape = next((child for child in list(clip) if _local_name(child.tag) in {"circle", "ellipse", "rect", "path", "polygon"}), None)
    if shape is None:
        issues.append(_issue("unsupported_clip_path", "clipPath must contain circle, ellipse, rect, path, or polygon.", path_label, xml_path))
        return
    tag = _local_name(shape.tag)
    if tag == "path":
        try:
            _normalize_path_commands(str(shape.attrib.get("d") or ""))
        except PptSvgNativeConversionError as exc:
            issues.append(_issue("malformed_svg_path", str(exc), path_label, xml_path))
    elif tag == "polygon" and len(_parse_points(str(shape.attrib.get("points") or ""))) < 3:
        issues.append(_issue("unsupported_clip_path", "clipPath polygon must have at least 3 points.", path_label, xml_path))


def _validate_filter_ref(
    issues: list[dict],
    elem: ET.Element,
    defs: dict[str, ET.Element],
    path_label: str,
    xml_path: str,
) -> None:
    filter_id = _resolve_url_id(str(elem.attrib.get("filter") or ""))
    if not filter_id:
        return
    filter_elem = defs.get(filter_id)
    if filter_elem is None or _local_name(filter_elem.tag) != "filter":
        issues.append(_issue("unsupported_filter", "`filter` references an unsupported filter.", path_label, xml_path))
        return
    child_tags = {_local_name(child.tag) for child in list(filter_elem)}
    if "feDropShadow" in child_tags:
        return
    if "feGaussianBlur" in child_tags:
        unsupported = child_tags - {"feGaussianBlur", "feOffset", "feFlood", "feComposite", "feMerge", "feMergeNode", "feFuncA"}
        if not unsupported:
            return
    issues.append(_issue("unsupported_filter", "Only drop shadow or blur/glow filters are supported.", path_label, xml_path))


def _issue(code: str, message: str, path_label: str = "", xml_path: str = "") -> dict:
    issue = {"severity": "error", "code": code, "message": message}
    if path_label:
        issue["path"] = path_label
    if xml_path:
        issue["xml_path"] = xml_path
    return issue


def _parse_transform(transform: str) -> tuple[float, float, float, float]:
    if not transform:
        return 0.0, 0.0, 1.0, 1.0
    dx = dy = 0.0
    sx = sy = 1.0
    translate = re.search(r"translate\(\s*([-\d.]+)(?:[\s,]+([-\d.]+))?\s*\)", transform)
    if translate:
        dx = float(translate.group(1))
        dy = float(translate.group(2) or 0)
    scale = re.search(r"scale\(\s*([-\d.]+)(?:[\s,]+([-\d.]+))?\s*\)", transform)
    if scale:
        sx = float(scale.group(1))
        sy = float(scale.group(2) or sx)
    return dx, dy, sx, sy


def _is_supported_transform(transform: str) -> bool:
    clean = str(transform or "").strip()
    if not clean:
        return True
    remaining = re.sub(r"translate\(\s*[-\d.]+(?:[\s,]+[-\d.]+)?\s*\)", "", clean)
    remaining = re.sub(r"scale\(\s*[-\d.]+(?:[\s,]+[-\d.]+)?\s*\)", "", remaining)
    remaining = re.sub(r"rotate\(\s*[-\d.]+(?:[\s,]+[-\d.]+){0,2}\s*\)", "", remaining)
    return not remaining.strip()


def _dash_xml(value: str) -> str:
    clean = re.sub(r"\s+", " ", str(value or "").strip().replace(",", " "))
    if not clean or clean.lower() == "none":
        return ""
    preset = DASH_PRESETS.get(clean)
    if preset:
        return f'<a:prstDash val="{preset}"/>'
    return ""


def ctx_x(ctx: SvgConvertContext, value: float) -> int:
    return int((value * ctx.scale_x + ctx.translate_x) * ctx.unit_x)


def ctx_y(ctx: SvgConvertContext, value: float) -> int:
    return int((value * ctx.scale_y + ctx.translate_y) * ctx.unit_y)


def ctx_w(ctx: SvgConvertContext, value: float) -> int:
    return max(1, int(value * ctx.scale_x * ctx.unit_x))


def ctx_h(ctx: SvgConvertContext, value: float) -> int:
    return max(1, int(value * ctx.scale_y * ctx.unit_y))


def _map_point(ctx: SvgConvertContext, x: float, y: float) -> tuple[float, float]:
    return x * ctx.scale_x + ctx.translate_x, y * ctx.scale_y + ctx.translate_y


def _num(value: object, default: float = 0.0) -> float:
    raw = str(value or "").strip()
    if not raw:
        return default
    match = re.match(r"[-+]?\d+(?:\.\d+)?", raw)
    return float(match.group(0)) if match else default


def _local_name(tag: object) -> str:
    return str(tag or "").rsplit("}", 1)[-1]


def _normalize_hex_color(color: str) -> str:
    clean = str(color or "").strip()
    if re.fullmatch(r"#[0-9a-fA-F]{6}", clean):
        return clean.upper()
    if re.fullmatch(r"#[0-9a-fA-F]{3}", clean):
        return "#" + "".join(char * 2 for char in clean[1:]).upper()
    return ""


def _xml_escape(value: object) -> str:
    return (
        str(value or "")
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
        .replace("'", "&apos;")
    )
