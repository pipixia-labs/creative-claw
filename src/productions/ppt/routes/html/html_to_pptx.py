"""Structured HTML-to-PPTX conversion for the PPT HTML route."""

from __future__ import annotations

from dataclasses import dataclass, field
import math
import re
from pathlib import Path
from urllib.parse import unquote, urlparse

from bs4 import BeautifulSoup, Tag
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from src.productions.ppt.routes.html.browser_html_to_pptx import convert_html_pages_with_browser
from src.productions.ppt.schemas import HtmlTemplatePackage
from src.runtime.workspace import resolve_workspace_path

_TEXT_TAGS = {"h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "span", "strong", "em", "label"}
_CONTAINER_TAGS = {"div", "section", "article", "aside", "main", "header", "footer"}
_UNSUPPORTED_TAGS = {"script", "style", "canvas", "video", "audio", "iframe", "object", "embed"}
_CRITICAL_UNSUPPORTED_STYLE_TOKENS = ("linear-gradient", "radial-gradient", "conic-gradient", "filter:", "backdrop-filter")
_WARN_UNSUPPORTED_STYLE_TOKENS = ("box-shadow", "text-shadow", "clip-path", "mask", "transform")
_DEFAULT_FONT = "Aptos"
_DEFAULT_TEXT_COLOR = RGBColor(23, 32, 51)
_LINE_HEIGHT = 1.12
_MIN_SAFE_MARGIN_IN = 0.03
_TEXT_WIDTH_FACTOR_CJK = 0.92
_TEXT_WIDTH_FACTOR_LATIN = 0.52
_NAMED_COLORS = {
    "black": RGBColor(0, 0, 0),
    "white": RGBColor(255, 255, 255),
    "red": RGBColor(220, 38, 38),
    "green": RGBColor(22, 163, 74),
    "blue": RGBColor(37, 99, 235),
    "yellow": RGBColor(250, 204, 21),
    "gray": RGBColor(107, 114, 128),
    "grey": RGBColor(107, 114, 128),
}


@dataclass(frozen=True)
class HtmlToPptxConversionResult:
    """Result returned by the structured HTML-to-PPTX converter."""

    ok: bool
    pptx_path: Path
    strategy: str = "html_to_pptx"
    engine: str = "python_structured_html"
    warnings: list[str] = field(default_factory=list)
    errors: list[str] = field(default_factory=list)
    preflight_report: dict[str, object] = field(default_factory=dict)


@dataclass(frozen=True)
class PptxBox:
    """A slide-space box stored in pixels."""

    x: float
    y: float
    width: float
    height: float

    def to_inches(self, template: HtmlTemplatePackage, slide_width: float, slide_height: float) -> tuple[float, float, float, float]:
        """Convert the box to PowerPoint inches."""
        return (
            max(0.0, self.x / template.viewport_width * slide_width),
            max(0.0, self.y / template.viewport_height * slide_height),
            max(0.01, self.width / template.viewport_width * slide_width),
            max(0.01, self.height / template.viewport_height * slide_height),
        )

    @property
    def right(self) -> float:
        """Return the right edge in pixels."""
        return self.x + self.width

    @property
    def bottom(self) -> float:
        """Return the bottom edge in pixels."""
        return self.y + self.height

    def overlaps(self, other: PptxBox) -> bool:
        """Return whether two boxes overlap."""
        return not (
            self.right <= other.x
            or other.right <= self.x
            or self.bottom <= other.y
            or other.bottom <= self.y
        )


@dataclass(frozen=True)
class PptxElement:
    """One extracted PPTX-renderable element."""

    kind: str
    box: PptxBox
    tag_name: str = ""
    text: str = ""
    src: str = ""
    style: dict[str, str] = field(default_factory=dict)
    z_index: int = 0


@dataclass(frozen=True)
class PptxSlideModel:
    """Intermediate model for one HTML slide."""

    slide_number: int
    width_px: int
    height_px: int
    elements: list[PptxElement]
    background: RGBColor | None = None
    warnings: list[str] = field(default_factory=list)
    errors: list[str] = field(default_factory=list)


@dataclass(frozen=True)
class PreflightFinding:
    """One preflight QA finding."""

    severity: str
    slide_number: int
    code: str
    message: str

    def as_text(self) -> str:
        """Return a compact user-facing message."""
        return f"slide {self.slide_number}: {self.code}: {self.message}"


@dataclass(frozen=True)
class PreflightReport:
    """Preflight result for one conversion run."""

    ok: bool
    findings: list[PreflightFinding]

    def warnings(self) -> list[str]:
        """Return warning messages."""
        return [finding.as_text() for finding in self.findings if finding.severity == "warning"]

    def errors(self) -> list[str]:
        """Return fatal error messages."""
        return [finding.as_text() for finding in self.findings if finding.severity == "error"]

    def model_dump(self) -> dict[str, object]:
        """Return a JSON-safe report."""
        return {
            "ok": self.ok,
            "findings": [
                {
                    "severity": finding.severity,
                    "slide_number": finding.slide_number,
                    "code": finding.code,
                    "message": finding.message,
                }
                for finding in self.findings
            ],
        }


def convert_html_pages_to_pptx(
    *,
    html_pages: list[str],
    pptx_path: Path,
    template: HtmlTemplatePackage,
) -> HtmlToPptxConversionResult:
    """Convert per-slide HTML fragments into an editable PPTX.

    The converter supports a deliberate subset of HTML: semantic text tags,
    local images, simple shapes, lines, and explicit geometry. It fails closed
    when high-risk HTML would produce a misleading or broken editable deck.
    """
    clean_pages = [str(page or "").strip() for page in html_pages if str(page or "").strip()]
    if not clean_pages:
        return HtmlToPptxConversionResult(
            ok=False,
            pptx_path=pptx_path,
            errors=["No generated HTML pages were available for structured PPTX conversion."],
        )

    browser_result = convert_html_pages_with_browser(
        html_pages=clean_pages,
        pptx_path=pptx_path,
        template=template,
    )
    browser_preflight_report = _browser_report_to_preflight_report(browser_result)
    if browser_result.ok:
        return HtmlToPptxConversionResult(
            ok=True,
            pptx_path=pptx_path,
            engine=browser_result.engine,
            warnings=browser_result.warnings,
            preflight_report=browser_preflight_report,
        )
    if not browser_result.unavailable:
        return HtmlToPptxConversionResult(
            ok=False,
            pptx_path=pptx_path,
            engine=browser_result.engine,
            warnings=browser_result.warnings,
            errors=browser_result.errors,
            preflight_report=browser_preflight_report,
        )

    models = extract_html_slide_models(html_pages=clean_pages, template=template)
    report = preflight_html_slide_models(models=models, template=template)
    if not report.ok:
        return HtmlToPptxConversionResult(
            ok=False,
            pptx_path=pptx_path,
            engine="python_structured_html",
            warnings=report.warnings(),
            errors=[*browser_result.errors, *report.errors()],
            preflight_report=_merge_unavailable_browser_report(
                browser_preflight_report,
                report.model_dump(),
            ),
        )

    try:
        render_html_slide_models_to_pptx(
            models=models,
            pptx_path=pptx_path,
            template=template,
        )
    except Exception as exc:
        return HtmlToPptxConversionResult(
            ok=False,
            pptx_path=pptx_path,
            engine="python_structured_html",
            warnings=report.warnings(),
            errors=[
                *browser_result.errors,
                *report.errors(),
                f"Structured HTML-to-PPTX render error: {type(exc).__name__}: {exc}",
            ],
            preflight_report=_merge_unavailable_browser_report(
                browser_preflight_report,
                report.model_dump(),
            ),
        )

    return HtmlToPptxConversionResult(
        ok=True,
        pptx_path=pptx_path,
        engine="python_structured_html",
        warnings=[
            "Browser HTML-to-PPTX converter was unavailable; used Python structured converter fallback.",
            *report.warnings(),
        ],
        preflight_report=_merge_unavailable_browser_report(
            browser_preflight_report,
            report.model_dump(),
        ),
    )


def _browser_report_to_preflight_report(browser_result: object) -> dict[str, object]:
    """Convert the browser converter report into the route preflight shape."""
    raw_report = dict(getattr(browser_result, "report", {}) or {})
    findings: list[dict[str, object]] = []
    for page in raw_report.get("pages") or []:
        try:
            slide_number = int(page.get("slideNumber") or page.get("slide_number") or 0)
        except (AttributeError, TypeError, ValueError):
            slide_number = 0
        if slide_number <= 0:
            continue
        for warning in page.get("warnings") or []:
            findings.append(
                {
                    "severity": "warning",
                    "slide_number": slide_number,
                    "code": "browser_converter_warning",
                    "message": str(warning),
                }
            )
        for error in page.get("errors") or []:
            findings.append(
                {
                    "severity": "error",
                    "slide_number": slide_number,
                    "code": "browser_converter_error",
                    "message": str(error),
                }
            )
    for error in getattr(browser_result, "errors", []) or []:
        findings.append(
            {
                "severity": "error",
                "slide_number": 0,
                "code": "browser_converter_error",
                "message": str(error),
            }
        )
    return {
        "ok": bool(getattr(browser_result, "ok", False)),
        "engine": getattr(browser_result, "engine", "node_playwright_pptxgenjs"),
        "browser_report": raw_report,
        "findings": findings,
    }


def _merge_unavailable_browser_report(
    browser_report: dict[str, object],
    python_report: dict[str, object],
) -> dict[str, object]:
    """Merge unavailable-browser context with the Python fallback report."""
    merged_findings = [
        *(browser_report.get("findings") or []),
        *(python_report.get("findings") or []),
    ]
    return {
        "ok": python_report.get("ok", False),
        "engine": "python_structured_html",
        "browser_report": browser_report.get("browser_report") or {},
        "python_report": python_report,
        "findings": merged_findings,
    }


def extract_html_slide_models(
    *,
    html_pages: list[str],
    template: HtmlTemplatePackage,
) -> list[PptxSlideModel]:
    """Extract intermediate slide models from generated HTML fragments."""
    return [
        _extract_single_slide_model(
            html_fragment=html_fragment,
            slide_number=index,
            template=template,
        )
        for index, html_fragment in enumerate(html_pages, start=1)
    ]


def preflight_html_slide_models(
    *,
    models: list[PptxSlideModel],
    template: HtmlTemplatePackage,
) -> PreflightReport:
    """Validate extracted models before writing a PPTX file."""
    findings: list[PreflightFinding] = []
    for model in models:
        for warning in model.warnings:
            findings.append(PreflightFinding("warning", model.slide_number, "html_warning", warning))
        for error in model.errors:
            findings.append(PreflightFinding("error", model.slide_number, "html_error", error))
        if not model.elements:
            findings.append(
                PreflightFinding(
                    "error",
                    model.slide_number,
                    "empty_slide",
                    "No convertible text, image, shape, or line elements were found.",
                )
            )
            continue

        text_count = len([element for element in model.elements if element.kind == "text" and element.text.strip()])
        if text_count == 0:
            findings.append(
                PreflightFinding(
                    "error",
                    model.slide_number,
                    "no_editable_text",
                    "The slide has no editable text element, so PPTX output would lose core content.",
                )
            )

        for element in model.elements:
            findings.extend(_preflight_element(element, model=model, template=template))
        findings.extend(_preflight_overlaps(model))

    return PreflightReport(
        ok=not any(finding.severity == "error" for finding in findings),
        findings=findings,
    )


def render_html_slide_models_to_pptx(
    *,
    models: list[PptxSlideModel],
    pptx_path: Path,
    template: HtmlTemplatePackage,
) -> None:
    """Render extracted slide models into an editable PPTX."""
    prs = Presentation()
    slide_width, slide_height = _configure_slide_size(prs, template.aspect_ratio)
    blank_layout = prs.slide_layouts[6]
    for model in models:
        slide = prs.slides.add_slide(blank_layout)
        if model.background is not None:
            _set_slide_background(slide, model.background)
        for element in sorted(model.elements, key=lambda item: item.z_index):
            box = element.box.to_inches(template, slide_width, slide_height)
            if element.kind == "text":
                _add_text(slide, element.text, box, style=element.style, tag_name=element.tag_name)
            elif element.kind == "image":
                _add_image(slide, element.src, box)
            elif element.kind == "line":
                _add_line(slide, box, style=element.style)
            elif element.kind == "shape":
                _add_shape(slide, box, style=element.style)
    prs.save(pptx_path)


def _extract_single_slide_model(
    *,
    html_fragment: str,
    slide_number: int,
    template: HtmlTemplatePackage,
) -> PptxSlideModel:
    """Extract one slide model from one HTML fragment."""
    soup = BeautifulSoup(html_fragment, "html.parser")
    root = _select_slide_root(soup)
    warnings = _collect_unsupported_html_warnings(root)
    errors = _collect_unsupported_html_errors(root)
    root_style = _parse_style(str(root.get("style") or ""))
    background = _parse_css_color(root_style.get("background-color") or root_style.get("background"))
    elements: list[PptxElement] = []

    for element in root.find_all(True):
        if not isinstance(element, Tag) or element is root:
            continue
        tag_name = str(element.name or "").lower()
        if tag_name in _UNSUPPORTED_TAGS:
            continue
        style = _parse_style(str(element.get("style") or ""))
        box = _style_box(style, template=template)
        if box is None:
            continue
        z_index = _safe_int(style.get("z-index"), default=0)

        if tag_name == "img":
            elements.append(PptxElement(kind="image", box=box, tag_name=tag_name, src=str(element.get("src") or ""), style=style, z_index=z_index))
            continue
        if _looks_like_line(box, style=style):
            elements.append(PptxElement(kind="line", box=box, tag_name=tag_name, style=style, z_index=z_index))
            continue

        flow_elements = _extract_positioned_flow_text_elements(
            element,
            container_box=box,
            template=template,
            base_z_index=z_index,
        )
        if flow_elements:
            if tag_name in _CONTAINER_TAGS and _has_shape_style(style):
                elements.append(PptxElement(kind="shape", box=box, tag_name=tag_name, style=style, z_index=z_index))
            elements.extend(flow_elements)
            continue

        text = _element_text(element)
        if text and _is_text_candidate(element):
            elements.append(PptxElement(kind="text", box=box, tag_name=tag_name, text=text, style=style, z_index=z_index))
            continue
        if tag_name in _CONTAINER_TAGS and not text and _has_shape_style(style):
            elements.append(PptxElement(kind="shape", box=box, tag_name=tag_name, style=style, z_index=z_index))

    if not elements:
        elements = _extract_sequential_text_model(root=root, template=template)
        if elements:
            warnings.append("Slide used sequential HTML-to-PPTX fallback because elements lacked explicit geometry.")
    return PptxSlideModel(
        slide_number=slide_number,
        width_px=template.viewport_width,
        height_px=template.viewport_height,
        elements=elements,
        background=background,
        warnings=warnings,
        errors=errors,
    )


def _extract_sequential_text_model(*, root: Tag, template: HtmlTemplatePackage) -> list[PptxElement]:
    """Extract a conservative sequential text layout for simple unpositioned HTML."""
    elements: list[PptxElement] = []
    y = template.viewport_height * 0.08
    x = template.viewport_width * 0.075
    width = template.viewport_width * 0.85
    for element in root.find_all(True):
        if not _is_text_candidate(element):
            continue
        text = _element_text(element)
        if not text:
            continue
        style = _parse_style(str(element.get("style") or ""))
        tag_name = str(element.name or "p").lower()
        height = _default_text_height_px(tag_name)
        elements.append(
            PptxElement(
                kind="text",
                box=PptxBox(x=x, y=y, width=width, height=height),
                tag_name=tag_name,
                text=text,
                style=style,
            )
        )
        y += height + 18
        if y > template.viewport_height * 0.9:
            break
    return elements


def _extract_positioned_flow_text_elements(
    container: Tag,
    *,
    container_box: PptxBox,
    template: HtmlTemplatePackage,
    base_z_index: int,
) -> list[PptxElement]:
    """Expand unpositioned text descendants inside a positioned HTML container."""
    tag_name = str(container.name or "").lower()
    if tag_name not in _CONTAINER_TAGS or _is_text_candidate(container):
        return []

    elements: list[PptxElement] = []
    container_style = _parse_style(str(container.get("style") or ""))
    padding_top = _css_edge_px(container_style, "padding", "top", template.viewport_height)
    padding_right = _css_edge_px(container_style, "padding", "right", template.viewport_width)
    padding_left = _css_edge_px(container_style, "padding", "left", template.viewport_width)
    content_x = container_box.x + padding_left
    content_y = container_box.y + padding_top
    content_width = max(1.0, container_box.width - padding_left - padding_right)
    cursor_y = content_y

    for child in container.find_all(True, recursive=False):
        if not isinstance(child, Tag) or str(child.name or "").lower() in _UNSUPPORTED_TAGS:
            continue
        child_style = _parse_style(str(child.get("style") or ""))
        if _style_box(child_style, template=template) is not None:
            continue
        text = _element_text(child)
        if not text:
            continue
        primary = _primary_flow_text_element(child) or child
        primary_style = _parse_style(str(primary.get("style") or ""))
        primary_tag = str(primary.name or child.name or "p").lower()
        margin_top = _css_edge_px(child_style, "margin", "top", template.viewport_height)
        margin_right = _css_edge_px(child_style, "margin", "right", template.viewport_width)
        margin_bottom = _css_edge_px(child_style, "margin", "bottom", template.viewport_height)
        margin_left = _css_edge_px(child_style, "margin", "left", template.viewport_width)
        item_x = content_x + margin_left
        item_y = cursor_y + margin_top + _flow_child_text_offset_y(child, primary)
        item_width = max(1.0, content_width - margin_left - margin_right)
        item_height = _estimate_flow_text_height_px(
            text,
            style=primary_style,
            tag_name=primary_tag,
            width_px=item_width,
            template=template,
        )
        if item_y + item_height > template.viewport_height:
            item_height = max(1.0, template.viewport_height - item_y)
        elements.append(
            PptxElement(
                kind="text",
                box=PptxBox(x=item_x, y=item_y, width=item_width, height=item_height),
                tag_name=primary_tag,
                text=text,
                style=primary_style,
                z_index=base_z_index + 1,
            )
        )
        cursor_y = item_y + item_height + margin_bottom

    return elements


def _primary_flow_text_element(element: Tag) -> Tag | None:
    """Return the descendant whose style best represents a flow text group."""
    candidates = [
        descendant
        for descendant in element.find_all(True)
        if isinstance(descendant, Tag)
        and _element_text(descendant)
        and _is_text_candidate(descendant)
    ]
    if not candidates:
        return element if _is_text_candidate(element) and _element_text(element) else None
    for candidate in candidates:
        text = _element_text(candidate)
        if str(candidate.name or "").lower() in _TEXT_TAGS and not _looks_like_marker_text(text):
            return candidate
    return candidates[0]


def _looks_like_marker_text(text: str) -> bool:
    """Return whether text is probably a numeric marker rather than body copy."""
    clean = str(text or "").strip()
    return bool(re.fullmatch(r"[\dA-Za-z一二三四五六七八九十]+[.)、]?", clean))


def _flow_child_text_offset_y(child: Tag, primary: Tag) -> float:
    """Return a small vertical offset for flex-style rows."""
    child_style = _parse_style(str(child.get("style") or ""))
    display = str(child_style.get("display") or "").lower()
    if display != "flex":
        return 0.0
    primary_style = _parse_style(str(primary.get("style") or ""))
    return _css_edge_px(primary_style, "margin", "top", 720)


def _estimate_flow_text_height_px(
    text: str,
    *,
    style: dict[str, str],
    tag_name: str,
    width_px: float,
    template: HtmlTemplatePackage,
) -> float:
    """Estimate a flow text box height in slide pixels."""
    font_size_pt = _font_size_pt(style, tag_name=tag_name)
    font_size_px = font_size_pt / 0.75
    slide_width_in = 13.333 if template.aspect_ratio == "16:9" else 10.0
    usable_width_pt = max(1.0, width_px / template.viewport_width * slide_width_in * 72 - 8)
    estimated_lines = max(1, math.ceil(_estimate_text_width_pt(text, font_size_pt) / usable_width_pt))
    line_height_px = _css_line_height_px(style, font_size_px)
    return max(_default_text_height_px(tag_name), estimated_lines * line_height_px + 8)


def _preflight_element(
    element: PptxElement,
    *,
    model: PptxSlideModel,
    template: HtmlTemplatePackage,
) -> list[PreflightFinding]:
    """Run element-level QA checks."""
    findings: list[PreflightFinding] = []
    if element.box.x < 0 or element.box.y < 0 or element.box.right > model.width_px or element.box.bottom > model.height_px:
        findings.append(
            PreflightFinding(
                "error",
                model.slide_number,
                "element_out_of_bounds",
                f"{element.kind} element extends outside the slide viewport.",
            )
        )
    if element.kind == "text":
        risk = _estimate_text_fit_risk(element, template=template)
        if risk == "error":
            findings.append(
                PreflightFinding(
                    "error",
                    model.slide_number,
                    "text_overflow_risk",
                    f"Text may not fit in its PPTX box: {element.text[:80]}",
                )
            )
        elif risk == "warning":
            findings.append(
                PreflightFinding(
                    "warning",
                    model.slide_number,
                    "text_tight_fit",
                    f"Text is close to the PPTX box height limit: {element.text[:80]}",
                )
            )
    if element.kind == "image" and _resolve_image_path(element.src) is None:
        findings.append(
            PreflightFinding(
                "error",
                model.slide_number,
                "image_missing",
                f"Image src is not a local readable file: {element.src}",
            )
        )
    if element.kind in {"shape", "line"} and _style_has_critical_unsupported_css(element.style):
        findings.append(
            PreflightFinding(
                "error",
                model.slide_number,
                "unsupported_css",
                "The element uses CSS that cannot be mapped into editable PPTX safely.",
            )
        )
    _ = template
    return findings


def _preflight_overlaps(model: PptxSlideModel) -> list[PreflightFinding]:
    """Warn when independent editable text boxes overlap other text boxes."""
    findings: list[PreflightFinding] = []
    text_elements = [element for element in model.elements if element.kind == "text"]
    for index, left in enumerate(text_elements):
        for right in text_elements[index + 1 :]:
            if left.box.overlaps(right.box):
                findings.append(
                    PreflightFinding(
                        "warning",
                        model.slide_number,
                        "text_overlap_risk",
                        "Two editable text boxes overlap in the extracted model.",
                    )
                )
                return findings
    return findings


def _collect_unsupported_html_warnings(root: Tag) -> list[str]:
    """Collect non-fatal unsupported HTML/CSS warnings."""
    warnings: list[str] = []
    for element in root.find_all(True):
        style_text = str(element.get("style") or "").lower()
        if any(token in style_text for token in _WARN_UNSUPPORTED_STYLE_TOKENS):
            warnings.append("Unsupported CSS detail was ignored during HTML-to-PPTX conversion.")
            break
    return warnings


def _collect_unsupported_html_errors(root: Tag) -> list[str]:
    """Collect fatal unsupported HTML/CSS findings."""
    errors: list[str] = []
    unsupported_tags = sorted({str(tag.name).lower() for tag in root.find_all(_UNSUPPORTED_TAGS)})
    if unsupported_tags:
        errors.append(f"Unsupported tags are present: {', '.join(unsupported_tags)}.")
    for element in root.find_all(True):
        style_text = str(element.get("style") or "").lower()
        if any(token in style_text for token in _CRITICAL_UNSUPPORTED_STYLE_TOKENS):
            errors.append("Gradient/filter CSS cannot be converted into editable PPTX safely.")
            break
    return errors


def _configure_slide_size(prs: Presentation, aspect_ratio: str) -> tuple[float, float]:
    """Configure slide size and return dimensions in inches."""
    if aspect_ratio == "4:3":
        width, height = 10.0, 7.5
    else:
        width, height = 13.333, 7.5
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    return width, height


def _select_slide_root(soup: BeautifulSoup) -> Tag:
    """Return the most likely slide root from one HTML fragment."""
    section = soup.find("section")
    if isinstance(section, Tag):
        return section
    body = soup.find("body")
    if isinstance(body, Tag):
        return body
    return soup


def _parse_style(style_text: str) -> dict[str, str]:
    """Parse an inline CSS style attribute into a lowercase dictionary."""
    result: dict[str, str] = {}
    for item in style_text.split(";"):
        if ":" not in item:
            continue
        key, value = item.split(":", 1)
        result[key.strip().lower()] = value.strip()
    return result


def _style_box(style: dict[str, str], *, template: HtmlTemplatePackage) -> PptxBox | None:
    """Return an element box in slide pixels when inline geometry is explicit."""
    left = _css_length_to_px(style.get("left") or style.get("x"), template.viewport_width)
    top = _css_length_to_px(style.get("top") or style.get("y"), template.viewport_height)
    width = _css_length_to_px(style.get("width"), template.viewport_width)
    height = _css_length_to_px(style.get("height"), template.viewport_height)
    right = _css_length_to_px(style.get("right"), template.viewport_width)
    bottom = _css_length_to_px(style.get("bottom"), template.viewport_height)
    margin_left = _css_length_to_px(style.get("margin-left"), template.viewport_width) or 0
    margin_top = _css_length_to_px(style.get("margin-top"), template.viewport_height) or 0

    if left is None and right is not None and width is not None:
        left = template.viewport_width - right - width
    if top is None and bottom is not None and height is not None:
        top = template.viewport_height - bottom - height
    if width is None and left is not None and right is not None:
        width = template.viewport_width - left - right
    if height is None and top is not None and bottom is not None:
        height = template.viewport_height - top - bottom
    if left is None or top is None or width is None:
        return None
    if height is None:
        height = _css_length_to_px(style.get("min-height"), template.viewport_height) or 72.0
    return PptxBox(
        x=left + margin_left,
        y=top + margin_top,
        width=max(1.0, width),
        height=max(1.0, height),
    )


def _css_length_to_px(value: str | None, axis_px: int) -> float | None:
    """Parse a simple CSS length into pixels."""
    clean = str(value or "").strip().lower()
    if not clean or clean in {"auto", "initial", "inherit"}:
        return None
    if clean.endswith("%"):
        try:
            return float(clean[:-1]) / 100.0 * axis_px
        except ValueError:
            return None
    match = re.match(r"^-?\d+(?:\.\d+)?", clean)
    if match is None:
        return None
    number = float(match.group(0))
    if clean.endswith("rem") or clean.endswith("em"):
        return number * 16.0
    if clean.endswith("pt"):
        return number / 0.75
    if clean.endswith("in"):
        return number * 96.0
    return number


def _element_text(element: Tag) -> str:
    """Return normalized display text for a convertible text element."""
    return " ".join(element.get_text(" ", strip=True).split())


def _is_text_candidate(element: Tag) -> bool:
    """Return whether an element should be converted as editable text."""
    tag_name = str(element.name or "").lower()
    if tag_name in _TEXT_TAGS:
        return True
    if tag_name not in {"div", "label"}:
        return False
    if element.find(True, recursive=False):
        return False
    return bool(" ".join(element.find_all(string=True, recursive=False)).strip())


def _css_edge_px(style: dict[str, str], prefix: str, edge: str, axis_px: int) -> float:
    """Return one CSS margin/padding edge in pixels."""
    explicit = _css_length_to_px(style.get(f"{prefix}-{edge}"), axis_px)
    if explicit is not None:
        return explicit
    shorthand = str(style.get(prefix) or "").strip()
    if not shorthand:
        return 0.0
    parts = shorthand.split()
    if len(parts) == 1:
        token = parts[0]
    elif len(parts) == 2:
        token = parts[0] if edge in {"top", "bottom"} else parts[1]
    elif len(parts) == 3:
        token = parts[0] if edge == "top" else parts[2] if edge == "bottom" else parts[1]
    else:
        token = {"top": parts[0], "right": parts[1], "bottom": parts[2], "left": parts[3]}[edge]
    return _css_length_to_px(token, axis_px) or 0.0


def _css_line_height_px(style: dict[str, str], font_size_px: float) -> float:
    """Return CSS line-height in pixels."""
    raw = str(style.get("line-height") or "").strip().lower()
    if not raw:
        return font_size_px * _LINE_HEIGHT
    try:
        return max(font_size_px, float(raw) * font_size_px)
    except ValueError:
        parsed = _css_length_to_px(raw, 720)
        return max(font_size_px, parsed) if parsed is not None else font_size_px * _LINE_HEIGHT


def _has_shape_style(style: dict[str, str]) -> bool:
    """Return whether an element has a simple visual style worth converting."""
    return bool(
        _parse_css_color(style.get("background-color") or style.get("background"))
        or _parse_css_color(style.get("border-color"))
        or style.get("border")
    )


def _looks_like_line(box: PptxBox, *, style: dict[str, str]) -> bool:
    """Return whether a thin HTML element should become a PPTX line."""
    has_line_style = bool(style.get("border") or style.get("border-top") or style.get("border-left") or style.get("background-color"))
    return has_line_style and (box.width <= 3 or box.height <= 3)


def _estimate_text_fit_risk(element: PptxElement, *, template: HtmlTemplatePackage) -> str:
    """Estimate whether text can fit inside its PPTX box."""
    font_size = _font_size_pt(element.style, tag_name=element.tag_name)
    box_width_in = max(0.01, element.box.width / template.viewport_width * (13.333 if template.aspect_ratio == "16:9" else 10.0))
    box_height_in = max(0.01, element.box.height / template.viewport_height * 7.5)
    usable_width_pt = max(1.0, box_width_in * 72 - 8)
    usable_height_pt = max(1.0, box_height_in * 72 - 4)
    estimated_line_width_pt = _estimate_text_width_pt(element.text, font_size)
    estimated_lines = max(1, math.ceil(estimated_line_width_pt / usable_width_pt))
    needed_height_pt = estimated_lines * font_size * _LINE_HEIGHT
    ratio = usable_height_pt / needed_height_pt
    if ratio < 0.9:
        return "error"
    if ratio < 1.1:
        return "warning"
    return "ok"


def _estimate_text_width_pt(text: str, font_size: float) -> float:
    """Estimate text width in points for preflight only."""
    width = 0.0
    for char in text:
        if "\u4e00" <= char <= "\u9fff":
            width += font_size * _TEXT_WIDTH_FACTOR_CJK
        elif char.isspace():
            width += font_size * 0.28
        else:
            width += font_size * _TEXT_WIDTH_FACTOR_LATIN
    return width


def _style_has_critical_unsupported_css(style: dict[str, str]) -> bool:
    """Return whether style contains critical unsupported CSS."""
    style_text = ";".join(f"{key}:{value}" for key, value in style.items()).lower()
    return any(token in style_text for token in _CRITICAL_UNSUPPORTED_STYLE_TOKENS)


def _add_shape(
    slide: object,
    box: tuple[float, float, float, float],
    *,
    style: dict[str, str],
) -> None:
    """Add a simple rectangle or rounded rectangle shape."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if _css_length_to_px(style.get("border-radius"), 1280) else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(box[0]), Inches(box[1]), Inches(box[2]), Inches(box[3]))
    fill_color = _parse_css_color(style.get("background-color") or style.get("background"))
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    border_color = _parse_css_color(style.get("border-color") or _border_color_from_shorthand(style.get("border", "")))
    if border_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(_border_width_pt(style))


def _add_line(
    slide: object,
    box: tuple[float, float, float, float],
    *,
    style: dict[str, str],
) -> None:
    """Add a horizontal or vertical line."""
    x, y, width, height = box
    if width >= height:
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + width), Inches(y))
    else:
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x), Inches(y + height))
    connector.line.color.rgb = _parse_css_color(style.get("background-color") or style.get("border-color") or _border_color_from_shorthand(style.get("border", ""))) or _DEFAULT_TEXT_COLOR
    connector.line.width = Pt(max(0.5, _border_width_pt(style)))


def _add_text(
    slide: object,
    text: str,
    box: tuple[float, float, float, float],
    *,
    style: dict[str, str],
    tag_name: str,
) -> None:
    """Add one editable text box."""
    shape = slide.shapes.add_textbox(Inches(box[0]), Inches(box[1]), Inches(box[2]), Inches(box[3]))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    margin_in = max(_MIN_SAFE_MARGIN_IN, (_css_length_to_px(style.get("padding"), 1280) or 4) / 96.0)
    frame.margin_left = Inches(margin_in)
    frame.margin_right = Inches(margin_in)
    frame.margin_top = Inches(min(margin_in, 0.08))
    frame.margin_bottom = Inches(min(margin_in, 0.08))
    paragraph = frame.paragraphs[0]
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = _LINE_HEIGHT
    align = _paragraph_align(style.get("text-align", ""))
    if align is not None:
        paragraph.alignment = align
    if tag_name == "li":
        paragraph.level = 0
        paragraph.text = ""
        run = paragraph.add_run()
        run.text = f"• {text}"
    else:
        run = paragraph.add_run()
        run.text = text
    run.font.name = _font_family(style)
    run.font.size = Pt(_font_size_pt(style, tag_name=tag_name))
    run.font.bold = _is_bold(style, tag_name=tag_name)
    run.font.italic = str(style.get("font-style") or "").lower() == "italic" or tag_name == "em"
    run.font.color.rgb = _parse_css_color(style.get("color")) or _DEFAULT_TEXT_COLOR


def _add_image(
    slide: object,
    src: str,
    box: tuple[float, float, float, float],
) -> None:
    """Add one local image fitted into a box."""
    image_path = _resolve_image_path(src)
    if image_path is None:
        raise ValueError(f"Image src is not a local readable file: {src}")
    x, y, width, height = _fit_image_box(image_path, box)
    slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(width), height=Inches(height))


def _resolve_image_path(src: str) -> Path | None:
    """Resolve an HTML image src to a local filesystem path."""
    clean_src = str(src or "").strip()
    if not clean_src or clean_src.startswith("data:") or clean_src.startswith("http://") or clean_src.startswith("https://"):
        return None
    parsed = urlparse(clean_src)
    if parsed.scheme == "file":
        candidate = Path(unquote(parsed.path))
    else:
        candidate = Path(clean_src)
        if not candidate.is_absolute():
            try:
                candidate = resolve_workspace_path(clean_src)
            except Exception:
                return None
    if candidate.exists() and candidate.is_file():
        return candidate
    return None


def _fit_image_box(image_path: Path, box: tuple[float, float, float, float]) -> tuple[float, float, float, float]:
    """Fit an image into a target box while preserving aspect ratio."""
    try:
        with Image.open(image_path) as image:
            image_width, image_height = image.size
    except Exception:
        return box
    if image_width <= 0 or image_height <= 0:
        return box
    x, y, width, height = box
    scale = min(width / image_width, height / image_height)
    fitted_width = image_width * scale
    fitted_height = image_height * scale
    return (
        x + (width - fitted_width) / 2,
        y + (height - fitted_height) / 2,
        fitted_width,
        fitted_height,
    )


def _set_slide_background(slide: object, color: RGBColor) -> None:
    """Set a solid slide background color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _parse_css_color(value: str | None) -> RGBColor | None:
    """Parse simple CSS color syntax into python-pptx RGBColor."""
    clean = str(value or "").strip().lower()
    if not clean or clean in {"none", "transparent", "inherit", "initial"} or "gradient" in clean:
        return None
    if clean in _NAMED_COLORS:
        return _NAMED_COLORS[clean]
    hex_match = re.search(r"#([0-9a-f]{3}|[0-9a-f]{6})\b", clean)
    if hex_match:
        raw = hex_match.group(1)
        if len(raw) == 3:
            raw = "".join(ch * 2 for ch in raw)
        return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))
    rgb_match = re.search(r"rgba?\(([^)]+)\)", clean)
    if rgb_match:
        parts = [part.strip() for part in rgb_match.group(1).split(",")[:3]]
        try:
            return RGBColor(*(max(0, min(255, int(float(part)))) for part in parts))
        except (TypeError, ValueError):
            return None
    return None


def _border_color_from_shorthand(border: str) -> str:
    """Extract the first color-looking value from CSS border shorthand."""
    clean = str(border or "")
    hex_match = re.search(r"#[0-9a-fA-F]{3,6}\b", clean)
    if hex_match:
        return hex_match.group(0)
    rgb_match = re.search(r"rgba?\([^)]+\)", clean)
    if rgb_match:
        return rgb_match.group(0)
    for token in clean.split():
        if token.lower() in _NAMED_COLORS:
            return token
    return ""


def _border_width_pt(style: dict[str, str]) -> float:
    """Return border width in points."""
    for key in ("border-width", "border-top-width", "border-left-width", "border"):
        width_px = _css_length_to_px(style.get(key), 1280)
        if width_px is not None:
            return max(0.5, width_px * 0.75)
    return 0.75


def _font_size_pt(style: dict[str, str], *, tag_name: str) -> float:
    """Return a PowerPoint font size in points."""
    explicit = _css_length_to_px(style.get("font-size"), 1280)
    if explicit is not None:
        return max(6.0, min(72.0, explicit * 0.75))
    defaults = {"h1": 42.0, "h2": 34.0, "h3": 28.0, "h4": 22.0, "h5": 18.0, "h6": 16.0, "li": 15.0}
    return defaults.get(tag_name, 16.0)


def _default_text_height_px(tag_name: str) -> float:
    """Return a fallback text box height in pixels."""
    return {"h1": 86.0, "h2": 70.0, "h3": 58.0}.get(tag_name, 44.0)


def _font_family(style: dict[str, str]) -> str:
    """Return a PPT-safe font family name."""
    raw = str(style.get("font-family") or "").strip()
    if not raw:
        return _DEFAULT_FONT
    return raw.split(",")[0].strip().strip("\"'") or _DEFAULT_FONT


def _is_bold(style: dict[str, str], *, tag_name: str) -> bool:
    """Return whether the text should be bold."""
    weight = str(style.get("font-weight") or "").lower()
    if weight in {"bold", "bolder"} or tag_name in {"h1", "h2", "h3", "h4", "h5", "h6", "strong"}:
        return True
    try:
        return int(weight) >= 600
    except ValueError:
        return False


def _paragraph_align(value: str) -> PP_ALIGN | None:
    """Map CSS text-align to PowerPoint paragraph alignment."""
    clean = str(value or "").strip().lower()
    if clean == "center":
        return PP_ALIGN.CENTER
    if clean == "right":
        return PP_ALIGN.RIGHT
    if clean == "justify":
        return PP_ALIGN.JUSTIFY
    if clean == "left":
        return PP_ALIGN.LEFT
    return None


def _safe_int(value: str | None, *, default: int) -> int:
    """Parse an integer without raising."""
    try:
        return int(str(value or "").strip())
    except ValueError:
        return default
