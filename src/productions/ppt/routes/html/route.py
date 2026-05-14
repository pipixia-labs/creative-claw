"""HTML route MVP for PPT generation."""

from __future__ import annotations

import copy
from dataclasses import dataclass, field
import html
import json
import textwrap
from pathlib import Path
from typing import Any

from google.adk.agents import BaseAgent, LlmAgent
from google.adk.apps import App
from google.adk.artifacts import BaseArtifactService, InMemoryArtifactService
from google.adk.memory import InMemoryMemoryService
from google.adk.runners import Runner
from google.adk.sessions import InMemorySessionService
from google.adk.tools.tool_context import ToolContext
from google.genai.types import Content, Part
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt

from conf.llm import build_llm
from src.productions.ppt.schemas import (
    DeckContentPlan,
    DeckPageAsset,
    DeckPagePlan,
    HtmlRouteBuildPackage,
    HtmlTemplatePackage,
)
from src.productions.ppt.routes.html.html_to_pptx import convert_html_pages_to_pptx
from src.productions.ppt.templates.html_registry import load_html_template_package
from src.runtime.tool_context_artifact_service import ToolContextArtifactService
from src.runtime.workspace import resolve_workspace_path, workspace_relative_path

_COLOR_BG = RGBColor(247, 248, 251)
_COLOR_INK = RGBColor(23, 32, 51)
_COLOR_MUTED = RGBColor(102, 112, 133)
_COLOR_ACCENT = RGBColor(36, 87, 214)
_COLOR_ACCENT_2 = RGBColor(67, 166, 255)
_COLOR_LINE = RGBColor(219, 226, 239)
_COLOR_PANEL = RGBColor(255, 255, 255)
_COLOR_WHITE = RGBColor(255, 255, 255)
_FONT_FAMILY = "Aptos"
_FREE_DESIGN_TEMPLATE_ID = "free_design"
PPT_HTML_PAGE_GENERATION_EXPERT_NAME = "PptHtmlPageGenerationExpert"
HTML_PAGE_GENERATION_CONTENT_PLAN_KEY = "ppt_html_route_content_plan"
HTML_PAGE_GENERATION_PAGES_KEY = "ppt_html_route_generated_pages"
HTML_PAGE_GENERATION_AGENT_MESSAGE_KEY = "ppt_html_page_generation_agent_message"
HTML_PAGE_GENERATION_WARNINGS_KEY = "ppt_html_route_page_generation_warnings"
HTML_ROUTE_STAGE_SEQUENCE = (
    "template_preparation",
    "page_generation",
    "pptx_output",
)
HTML_DELIVERY_STAGE = "quality_delivery"


@dataclass(frozen=True)
class HtmlRoutePaths:
    """Filesystem paths used by one HTML route build."""

    output_dir: Path
    preview_dir: Path
    html_path: Path
    pptx_path: Path
    quality_report_path: Path
    build_log_path: Path


@dataclass(frozen=True)
class HtmlTemplatePreparationResult:
    """Result of HTML route template preparation."""

    template: HtmlTemplatePackage
    stage: str = "template_preparation"


@dataclass(frozen=True)
class HtmlPageGenerationResult:
    """Result of HTML route page generation."""

    html_path: Path
    preview_paths: list[Path]
    html_pages: list[str] = field(default_factory=list)
    generation_mode: str = "deterministic_renderer"
    warnings: list[str] = field(default_factory=list)
    stage: str = "page_generation"


@dataclass(frozen=True)
class HtmlPptxOutputResult:
    """Result of HTML route PPTX output."""

    pptx_path: Path
    pptx_strategy: str
    warnings: list[str] = field(default_factory=list)
    conversion_report: dict[str, Any] = field(default_factory=dict)
    stage: str = "pptx_output"


@dataclass(frozen=True)
class HtmlQualityDeliveryResult:
    """Result of HTML route quality and delivery packaging."""

    quality_report_path: Path
    build_log_path: Path
    quality_report: dict[str, Any]
    warnings: list[str]
    stage: str = HTML_DELIVERY_STAGE


def build_html_route(
    *,
    content_plan: DeckContentPlan,
    output_dir: Path,
    aspect_ratio: str = "16:9",
    template_id: str = "",
) -> HtmlRouteBuildPackage:
    """Generate static HTML, PNG previews, and editable PPTX."""
    paths = prepare_html_route_paths(output_dir)
    template_stage = prepare_html_template(template_id=template_id, aspect_ratio=aspect_ratio)
    page_stage = generate_html_pages(
        content_plan=content_plan,
        template=template_stage.template,
        paths=paths,
    )
    pptx_stage = export_html_pptx(
        content_plan=content_plan,
        template=template_stage.template,
        page_generation=page_stage,
        paths=paths,
    )
    quality_stage = deliver_html_route_quality(
        content_plan=content_plan,
        template=template_stage.template,
        page_generation=page_stage,
        pptx_output=pptx_stage,
        paths=paths,
    )

    return HtmlRouteBuildPackage(
        template=template_stage.template,
        html_deck_path=workspace_relative_path(page_stage.html_path),
        preview_paths=[workspace_relative_path(path) for path in page_stage.preview_paths],
        pptx_path=_workspace_relative_existing_file(pptx_stage.pptx_path),
        quality_report_path=workspace_relative_path(quality_stage.quality_report_path),
        build_log_path=workspace_relative_path(quality_stage.build_log_path),
        warnings=quality_stage.warnings,
    )


async def build_html_route_with_agent(
    *,
    content_plan: DeckContentPlan,
    output_dir: Path,
    aspect_ratio: str = "16:9",
    template_id: str = "",
    tool_context: ToolContext | None = None,
    app_name: str = "creative_claw",
    artifact_service: BaseArtifactService | None = None,
    page_generation_agent: BaseAgent | None = None,
) -> HtmlRouteBuildPackage:
    """Generate HTML-route artifacts, using the page-generation agent when available."""
    paths = prepare_html_route_paths(output_dir)
    template_stage = prepare_html_template(template_id=template_id, aspect_ratio=aspect_ratio)
    page_stage = await generate_html_pages_with_agent(
        content_plan=content_plan,
        template=template_stage.template,
        paths=paths,
        tool_context=tool_context,
        app_name=app_name,
        artifact_service=artifact_service,
        page_generation_agent=page_generation_agent,
    )
    pptx_stage = export_html_pptx(
        content_plan=content_plan,
        template=template_stage.template,
        page_generation=page_stage,
        paths=paths,
    )
    quality_stage = deliver_html_route_quality(
        content_plan=content_plan,
        template=template_stage.template,
        page_generation=page_stage,
        pptx_output=pptx_stage,
        paths=paths,
    )

    return HtmlRouteBuildPackage(
        template=template_stage.template,
        html_deck_path=workspace_relative_path(page_stage.html_path),
        preview_paths=[workspace_relative_path(path) for path in page_stage.preview_paths],
        pptx_path=_workspace_relative_existing_file(pptx_stage.pptx_path),
        quality_report_path=workspace_relative_path(quality_stage.quality_report_path),
        build_log_path=workspace_relative_path(quality_stage.build_log_path),
        warnings=[*page_stage.warnings, *quality_stage.warnings],
    )


def prepare_html_route_paths(output_dir: Path) -> HtmlRoutePaths:
    """Prepare output paths for all HTML route stages."""
    output_dir.mkdir(parents=True, exist_ok=True)
    preview_dir = output_dir / "previews"
    preview_dir.mkdir(parents=True, exist_ok=True)
    return HtmlRoutePaths(
        output_dir=output_dir,
        preview_dir=preview_dir,
        html_path=output_dir / "deck.html",
        pptx_path=output_dir / "deck.pptx",
        quality_report_path=output_dir / "quality_report.json",
        build_log_path=output_dir / "build_log.json",
    )


def _workspace_relative_existing_file(path: Path) -> str:
    """Return a workspace-relative path only for an existing file."""
    return workspace_relative_path(path) if path.exists() and path.is_file() else ""


def prepare_html_template(
    *,
    template_id: str,
    aspect_ratio: str,
) -> HtmlTemplatePreparationResult:
    """Prepare the HTML route design package for a route run."""
    if not str(template_id or "").strip():
        return HtmlTemplatePreparationResult(
            template=_build_free_design_template_package(aspect_ratio=aspect_ratio),
        )
    return HtmlTemplatePreparationResult(
        template=load_html_template_package(template_id, aspect_ratio=aspect_ratio),
    )


def _build_free_design_template_package(*, aspect_ratio: str) -> HtmlTemplatePackage:
    """Build the no-system-template package used for free HTML page generation."""
    if aspect_ratio not in {"16:9", "4:3"}:
        raise ValueError("HTML template aspect_ratio must be `16:9` or `4:3`.")
    viewport_width = 1024 if aspect_ratio == "4:3" else 1280
    viewport_height = 768 if aspect_ratio == "4:3" else 720
    return HtmlTemplatePackage(
        template_id=_FREE_DESIGN_TEMPLATE_ID,
        label="Free Design",
        version="0.1.0",
        aspect_ratio=aspect_ratio,
        viewport_width=viewport_width,
        viewport_height=viewport_height,
        page_types={
            "cover": "free-cover",
            "toc": "free-toc",
            "chapter_start": "free-section",
            "chapter_content": "free-content",
            "ending": "free-ending",
        },
        pptx_strategy="native_editable",
        editability_notes=(
            "No system HTML template was selected. The HTML route uses free-design page generation "
            "and exports editable text boxes and vector shapes."
        ),
    )


def generate_html_pages(
    *,
    content_plan: DeckContentPlan,
    template: HtmlTemplatePackage,
    paths: HtmlRoutePaths,
) -> HtmlPageGenerationResult:
    """Generate HTML deck and PNG previews from a content plan."""
    paths.html_path.write_text(_render_html_deck(content_plan, template), encoding="utf-8")
    preview_paths = _render_previews(content_plan, template, paths.preview_dir)
    return HtmlPageGenerationResult(
        html_path=paths.html_path,
        preview_paths=preview_paths,
        html_pages=[_render_html_slide(page, template) for page in content_plan.pages],
    )


async def generate_html_pages_with_agent(
    *,
    content_plan: DeckContentPlan,
    template: HtmlTemplatePackage,
    paths: HtmlRoutePaths,
    tool_context: ToolContext | None,
    app_name: str,
    artifact_service: BaseArtifactService | None,
    page_generation_agent: BaseAgent | None = None,
) -> HtmlPageGenerationResult:
    """Generate per-slide HTML with an ADK page-generation agent, falling back deterministically."""
    if template.template_id != _FREE_DESIGN_TEMPLATE_ID:
        return generate_html_pages(content_plan=content_plan, template=template, paths=paths)
    if tool_context is None or not hasattr(tool_context, "_invocation_context"):
        return generate_html_pages(content_plan=content_plan, template=template, paths=paths)
    if page_generation_agent is None:
        return generate_html_pages(content_plan=content_plan, template=template, paths=paths)

    try:
        page_fragments = await _run_html_page_generation_agent(
            content_plan=content_plan,
            template=template,
            tool_context=tool_context,
            app_name=app_name,
            artifact_service=artifact_service,
            page_generation_agent=page_generation_agent,
        )
        paths.html_path.write_text(
            _render_agent_html_deck(content_plan, template, page_fragments),
            encoding="utf-8",
        )
        preview_paths = await _render_html_deck_screenshots(
            html_path=paths.html_path,
            content_plan=content_plan,
            template=template,
            preview_dir=paths.preview_dir,
        )
        warnings = []
        if not preview_paths:
            warnings.append("HTML page screenshots were unavailable; deterministic preview fallback was used.")
            preview_paths = _render_previews(content_plan, template, paths.preview_dir)
        return HtmlPageGenerationResult(
            html_path=paths.html_path,
            preview_paths=preview_paths,
            html_pages=[page.html_fragment for page in page_fragments],
            generation_mode="llm_agent_html",
            warnings=warnings,
        )
    except Exception as exc:
        agent_name = getattr(page_generation_agent, "name", PPT_HTML_PAGE_GENERATION_EXPERT_NAME)
        warning = f"{agent_name} fallback: {type(exc).__name__}: {exc}"
        _append_html_page_generation_warning(tool_context.state, warning)
        fallback = generate_html_pages(content_plan=content_plan, template=template, paths=paths)
        return HtmlPageGenerationResult(
            html_path=fallback.html_path,
            preview_paths=fallback.preview_paths,
            generation_mode=fallback.generation_mode,
            warnings=[warning, *fallback.warnings],
        )


def export_html_pptx(
    *,
    content_plan: DeckContentPlan,
    template: HtmlTemplatePackage,
    page_generation: HtmlPageGenerationResult,
    paths: HtmlRoutePaths,
) -> HtmlPptxOutputResult:
    """Export PPTX from generated HTML-route pages."""
    warnings: list[str] = []
    if page_generation.generation_mode == "llm_agent_html" or template.pptx_strategy == "html_to_pptx":
        conversion = convert_html_pages_to_pptx(
            html_pages=page_generation.html_pages,
            pptx_path=paths.pptx_path,
            template=template,
        )
        if conversion.ok:
            return HtmlPptxOutputResult(
                pptx_path=paths.pptx_path,
                pptx_strategy=conversion.strategy,
                warnings=conversion.warnings,
                conversion_report=_build_html_to_pptx_conversion_report(
                    conversion,
                    page_count=len(content_plan.pages),
                    fallback_used=False,
                ),
            )
        warnings.extend(
            [
                "HTML-to-PPTX conversion failed; no screenshot PPTX fallback was used because final PPTX must remain editable.",
                *conversion.warnings,
                *conversion.errors,
            ]
        )
        pptx_strategy = "html_to_pptx_failed"
        conversion_report = _build_html_to_pptx_conversion_report(
            conversion,
            page_count=len(content_plan.pages),
            fallback_used=False,
            final_strategy=pptx_strategy,
        )
    elif template.pptx_strategy == "native_editable":
        _export_native_pptx(content_plan, paths.pptx_path, template)
        pptx_strategy = template.pptx_strategy
        conversion_report = _build_non_html_to_pptx_conversion_report(
            engine="native_editable_renderer",
            strategy=pptx_strategy,
            page_count=len(content_plan.pages),
        )
    else:
        _export_previews_to_pptx(
            page_generation.preview_paths,
            paths.pptx_path,
            aspect_ratio=template.aspect_ratio,
        )
        pptx_strategy = template.pptx_strategy
        conversion_report = _build_non_html_to_pptx_conversion_report(
            engine="screenshot_renderer",
            strategy=pptx_strategy,
            page_count=len(content_plan.pages),
        )
    return HtmlPptxOutputResult(
        pptx_path=paths.pptx_path,
        pptx_strategy=pptx_strategy,
        warnings=warnings,
        conversion_report=conversion_report,
    )


def deliver_html_route_quality(
    *,
    content_plan: DeckContentPlan,
    template: HtmlTemplatePackage,
    page_generation: HtmlPageGenerationResult,
    pptx_output: HtmlPptxOutputResult,
    paths: HtmlRoutePaths,
) -> HtmlQualityDeliveryResult:
    """Validate route artifacts and write delivery logs."""
    quality_report = _validate_html_route_output(
        content_plan=content_plan,
        html_path=page_generation.html_path,
        preview_paths=page_generation.preview_paths,
        pptx_path=pptx_output.pptx_path,
        pptx_strategy=pptx_output.pptx_strategy,
    )
    quality_report["route_stages"] = list(HTML_ROUTE_STAGE_SEQUENCE)
    quality_report["delivery_stage"] = HTML_DELIVERY_STAGE
    quality_report["pptx_conversion"] = pptx_output.conversion_report
    paths.quality_report_path.write_text(
        json.dumps(quality_report, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    build_log = {
        "route": "html",
        "workflow_name": "HtmlRouteSequentialAgent",
        "route_stages": list(HTML_ROUTE_STAGE_SEQUENCE),
        "delivery_stage": HTML_DELIVERY_STAGE,
        "template_id": template.template_id,
        "template_source": "none" if template.template_id == _FREE_DESIGN_TEMPLATE_ID else "system",
        "page_generation_mode": _FREE_DESIGN_TEMPLATE_ID
        if template.template_id == _FREE_DESIGN_TEMPLATE_ID
        else "system_template",
        "page_generation_executor": page_generation.generation_mode,
        "pptx_strategy": pptx_output.pptx_strategy,
        "pptx_warnings": pptx_output.warnings,
        "pptx_conversion": pptx_output.conversion_report,
        "slide_count": len(content_plan.pages),
        "html_deck_path": workspace_relative_path(page_generation.html_path),
        "pptx_path": workspace_relative_path(pptx_output.pptx_path),
        "preview_count": len(page_generation.preview_paths),
    }
    paths.build_log_path.write_text(json.dumps(build_log, ensure_ascii=False, indent=2), encoding="utf-8")
    warnings = [*pptx_output.warnings]
    if pptx_output.pptx_strategy == "screenshot" and template.editability_notes:
        warnings.append(template.editability_notes)
    return HtmlQualityDeliveryResult(
        quality_report_path=paths.quality_report_path,
        build_log_path=paths.build_log_path,
        quality_report=quality_report,
        warnings=warnings,
    )


def _build_html_to_pptx_conversion_report(
    conversion: Any,
    *,
    page_count: int,
    fallback_used: bool,
    final_strategy: str | None = None,
) -> dict[str, Any]:
    """Build a user-visible page-level report for HTML-to-PPTX conversion."""
    findings = list((conversion.preflight_report or {}).get("findings") or [])
    findings_by_slide: dict[int, list[dict[str, Any]]] = {}
    for finding in findings:
        try:
            slide_number = int(finding.get("slide_number") or 0)
        except (TypeError, ValueError):
            slide_number = 0
        if slide_number > 0:
            findings_by_slide.setdefault(slide_number, []).append(dict(finding))

    conversion_ok = bool(getattr(conversion, "ok", False)) and not fallback_used
    resolved_final_strategy = final_strategy or (
        "screenshot" if fallback_used else getattr(conversion, "strategy", "html_to_pptx")
    )
    page_status = "screenshot_fallback" if fallback_used else "html_to_pptx" if conversion_ok else "html_to_pptx_failed"
    editable_level = "low" if fallback_used else "high" if conversion_ok else "none"
    pages = []
    for slide_number in range(1, page_count + 1):
        slide_findings = findings_by_slide.get(slide_number, [])
        pages.append(
            {
                "slide_number": slide_number,
                "status": page_status,
                "editable_level": editable_level,
                "warnings": [
                    finding.get("message", "")
                    for finding in slide_findings
                    if finding.get("severity") == "warning"
                ],
                "errors": [
                    finding.get("message", "")
                    for finding in slide_findings
                    if finding.get("severity") == "error"
                ],
            }
        )

    return {
        "engine": getattr(conversion, "engine", "python_structured_html"),
        "requested_strategy": "html_to_pptx",
        "final_strategy": resolved_final_strategy,
        "ok": conversion_ok,
        "fallback_used": fallback_used,
        "editable_level": editable_level,
        "warnings": list(getattr(conversion, "warnings", []) or []),
        "errors": list(getattr(conversion, "errors", []) or []),
        "preflight_report": conversion.preflight_report or {},
        "pages": pages,
    }


def _build_non_html_to_pptx_conversion_report(
    *,
    engine: str,
    strategy: str,
    page_count: int,
) -> dict[str, Any]:
    """Build a conversion report for non-HTML-to-PPTX output strategies."""
    editable_level = "high" if strategy == "native_editable" else "low"
    return {
        "engine": engine,
        "requested_strategy": strategy,
        "final_strategy": strategy,
        "ok": True,
        "fallback_used": False,
        "editable_level": editable_level,
        "warnings": [],
        "errors": [],
        "preflight_report": {},
        "pages": [
            {
                "slide_number": slide_number,
                "status": strategy,
                "editable_level": editable_level,
                "warnings": [],
                "errors": [],
            }
            for slide_number in range(1, page_count + 1)
        ],
    }


def build_html_page_generation_agent() -> LlmAgent:
    """Build the PPT product expert that turns a DeckContentPlan into per-slide HTML."""
    return build_ppt_html_page_generation_expert()


def build_ppt_html_page_generation_expert() -> LlmAgent:
    """Build the PPT product-level expert for editable HTML slide fragments."""
    return LlmAgent(
        name=PPT_HTML_PAGE_GENERATION_EXPERT_NAME,
        model=build_llm(),
        description="Generates editable PPT-friendly HTML fragments from a DeckContentPlan.",
        instruction=(
            "You are Creative Claw's HTML PPT page generation agent.\n"
            "You receive a complete DeckContentPlan JSON and must create one HTML fragment per slide.\n"
            "Do not use a fixed template, fixed card grid, or shared business layout. Choose each page's "
            "composition directly from the page content, audience, page_type, and available assets.\n"
            "Keep visual style and layout language consistent across the whole deck. All `chapter_start` "
            "slides must use a consistent style and layout language with each other.\n"
            "Each fragment must be a single <section> for one slide. Inline CSS is allowed. Do not use "
            "external CSS, external JS, remote images, markdown, JSON, or code fences.\n"
            "For PPTX conversion compatibility, design for a fixed browser viewport and make visible "
            "elements measurable: use explicit width and height for the slide, important text blocks, "
            "images, and decorative shapes. Prefer absolute positioning for final visible elements. "
            "Readable text must be inside h1-h6, p, ul, ol, li, or label tags. Use div/section elements "
            "for backgrounds, cards, borders, and shapes; do not put unwrapped text directly inside divs. "
            "Use real ul/li lists instead of manually typing bullet symbols. Do not put background, border, "
            "or box-shadow directly on text tags; place a separate shape behind the text instead. Avoid CSS "
            "gradients, filters, canvas, svg-only content, remote assets, and background images on divs.\n"
            "Fit each slide into the declared viewport and leave a safe bottom margin. Avoid text overlap, "
            "tiny unreadable text, and image/text collisions.\n"
            "Use ready asset html_src values when useful. Keep image text out of generated images.\n"
            "When all slides are ready, call save_html_route_pages with the full ordered page list."
        ),
        tools=[save_html_route_pages],
        output_key=HTML_PAGE_GENERATION_AGENT_MESSAGE_KEY,
        include_contents="none",
    )


def save_html_route_pages(
    pages: list[dict[str, Any]],
    tool_context: ToolContext,
) -> dict[str, Any]:
    """Save generated per-slide HTML fragments for the HTML route."""
    content_plan_payload = tool_context.state.get(HTML_PAGE_GENERATION_CONTENT_PLAN_KEY) or {}
    content_plan = DeckContentPlan.model_validate(content_plan_payload)
    normalized_pages = _normalize_agent_page_fragments(pages, content_plan=content_plan)
    payload = [page.model_dump(mode="json") for page in normalized_pages]
    tool_context.state[HTML_PAGE_GENERATION_PAGES_KEY] = payload
    tool_context.state["current_output"] = {
        "status": "success",
        "message": f"{PPT_HTML_PAGE_GENERATION_EXPERT_NAME} saved per-slide HTML.",
        "html_pages": payload,
    }
    return {
        "status": "success",
        "message": "HTML route pages saved.",
        "page_count": len(payload),
    }


@dataclass(frozen=True)
class _GeneratedHtmlPage:
    """One generated slide HTML fragment."""

    slide_number: int
    html_fragment: str

    def model_dump(self, *, mode: str = "json") -> dict[str, Any]:
        """Return a JSON-safe representation."""
        return {
            "slide_number": self.slide_number,
            "html": self.html_fragment,
        }


async def _run_html_page_generation_agent(
    *,
    content_plan: DeckContentPlan,
    template: HtmlTemplatePackage,
    tool_context: ToolContext,
    app_name: str,
    artifact_service: BaseArtifactService | None,
    page_generation_agent: BaseAgent,
) -> list[_GeneratedHtmlPage]:
    """Run the child page-generation agent and return saved HTML fragments."""
    invocation_context = tool_context._invocation_context
    child_session_service = InMemorySessionService()
    child_artifact_service = _resolve_child_artifact_service(
        tool_context=tool_context,
        fallback_service=artifact_service or InMemoryArtifactService(),
    )
    child_runner = _build_child_runner(
        agent=page_generation_agent,
        app_name=app_name,
        session_service=child_session_service,
        artifact_service=child_artifact_service,
        invocation_context=invocation_context,
    )
    child_state = _copy_state(tool_context.state)
    child_state[HTML_PAGE_GENERATION_CONTENT_PLAN_KEY] = content_plan.model_dump(mode="json")
    try:
        child_session = await child_session_service.create_session(
            app_name=app_name,
            user_id=invocation_context.user_id,
            state=child_state,
        )
        async for _event in child_runner.run_async(
            user_id=child_session.user_id,
            session_id=child_session.id,
            new_message=Content(
                role="user",
                parts=[
                    Part(
                        text=_build_html_page_generation_user_message(
                            content_plan=content_plan,
                            template=template,
                        )
                    )
                ],
            ),
        ):
            pass
        final_session = await child_session_service.get_session(
            app_name=app_name,
            user_id=child_session.user_id,
            session_id=child_session.id,
        )
        final_state = final_session.state if final_session is not None else child_state
        pages_payload = final_state.get(HTML_PAGE_GENERATION_PAGES_KEY)
        if not pages_payload:
            raise ValueError(
                f"{getattr(page_generation_agent, 'name', PPT_HTML_PAGE_GENERATION_EXPERT_NAME)} "
                "did not save HTML route pages."
            )
        normalized_pages = _normalize_agent_page_fragments(pages_payload, content_plan=content_plan)
        tool_context.state[HTML_PAGE_GENERATION_PAGES_KEY] = [
            page.model_dump(mode="json") for page in normalized_pages
        ]
        if final_state.get(HTML_PAGE_GENERATION_AGENT_MESSAGE_KEY):
            tool_context.state[HTML_PAGE_GENERATION_AGENT_MESSAGE_KEY] = str(
                final_state.get(HTML_PAGE_GENERATION_AGENT_MESSAGE_KEY)
            )
        return normalized_pages
    finally:
        await child_runner.close()


async def run_html_page_generation_expert(
    *,
    content_plan: DeckContentPlan,
    template: HtmlTemplatePackage,
    tool_context: ToolContext,
    app_name: str,
    artifact_service: BaseArtifactService | None,
    page_generation_agent: BaseAgent,
) -> list[dict[str, Any]]:
    """Run the PPT HTML page-generation expert and return JSON-safe page fragments."""
    pages = await _run_html_page_generation_agent(
        content_plan=content_plan,
        template=template,
        tool_context=tool_context,
        app_name=app_name,
        artifact_service=artifact_service,
        page_generation_agent=page_generation_agent,
    )
    return [page.model_dump(mode="json") for page in pages]


def _build_html_page_generation_user_message(
    *,
    content_plan: DeckContentPlan,
    template: HtmlTemplatePackage,
) -> str:
    """Build the explicit user message for HTML page generation."""
    payload = content_plan.model_dump(mode="json")
    page_payloads = []
    for page in payload.get("pages", []):
        prepared_page = dict(page)
        prepared_assets = []
        for asset in prepared_page.get("assets", []) or []:
            prepared_asset = dict(asset)
            asset_path = str(prepared_asset.get("path") or "").strip()
            if asset_path:
                try:
                    prepared_asset["html_src"] = resolve_workspace_path(asset_path).as_uri()
                except Exception:
                    prepared_asset["html_src"] = asset_path
            prepared_assets.append(prepared_asset)
        prepared_page["assets"] = prepared_assets
        page_payloads.append(prepared_page)
    payload["pages"] = page_payloads
    return "\n".join(
        [
            "Generate HTML for each PPT slide from this DeckContentPlan.",
            "The route is free design: do not reuse a fixed layout across slides unless the content naturally asks for it.",
            "Return no prose. Call save_html_route_pages with a list of objects: "
            "`[{\"slide_number\": 1, \"html\": \"<section>...</section>\"}, ...]`.",
            f"Viewport: {template.viewport_width}x{template.viewport_height}.",
            "HTML-to-PPTX rule: every slide must be convertible into editable PowerPoint objects. "
            "Use measurable browser layout, explicit dimensions, semantic text tags, local/file image src values, "
            "and separate div/section shapes behind text.",
            "",
            "DeckContentPlan JSON:",
            "```json",
            json.dumps(payload, ensure_ascii=False, indent=2),
            "```",
        ]
    )


def _normalize_agent_page_fragments(
    pages: list[dict[str, Any]],
    *,
    content_plan: DeckContentPlan,
) -> list[_GeneratedHtmlPage]:
    """Validate and normalize saved page HTML fragments."""
    pages_by_number: dict[int, str] = {}
    for item in pages:
        if not isinstance(item, dict):
            continue
        try:
            slide_number = int(item.get("slide_number") or item.get("number") or 0)
        except (TypeError, ValueError):
            slide_number = 0
        html_fragment = _clean_generated_html_fragment(str(item.get("html") or item.get("html_fragment") or ""))
        if slide_number > 0 and html_fragment:
            pages_by_number[slide_number] = html_fragment

    normalized_pages = []
    for page in content_plan.pages:
        html_fragment = pages_by_number.get(page.slide_number, "")
        if not html_fragment:
            raise ValueError(f"Missing generated HTML for slide {page.slide_number}.")
        normalized_pages.append(
            _GeneratedHtmlPage(
                slide_number=page.slide_number,
                html_fragment=_ensure_slide_section(html_fragment, page=page),
            )
        )
    return normalized_pages


def _clean_generated_html_fragment(fragment: str) -> str:
    """Remove common LLM wrappers around an HTML fragment."""
    clean_fragment = str(fragment or "").strip()
    if clean_fragment.startswith("```"):
        clean_fragment = clean_fragment.strip("`")
        clean_fragment = clean_fragment.removeprefix("html").strip()
    return clean_fragment


def _ensure_slide_section(fragment: str, *, page: DeckPagePlan) -> str:
    """Ensure a generated fragment is wrapped as one slide section."""
    clean_fragment = fragment.strip()
    slide_number = f"{page.slide_number:02d}"
    if "<section" in clean_fragment[:200].lower():
        section = clean_fragment
        if "data-slide-number" not in section[:400]:
            section = section.replace("<section", f'<section data-slide-number="{slide_number}"', 1)
        if "class=" not in section[:400]:
            section = section.replace("<section", '<section class="slide generated-slide"', 1)
        elif "slide" not in section[:400]:
            section = section.replace('class="', 'class="slide generated-slide ', 1)
        return section
    return (
        f'<section class="slide generated-slide" data-slide-number="{slide_number}" '
        f'data-page-type="{html.escape(page.page_type)}">{clean_fragment}</section>'
    )


def _render_agent_html_deck(
    content_plan: DeckContentPlan,
    template: HtmlTemplatePackage,
    pages: list[_GeneratedHtmlPage],
) -> str:
    """Render a full HTML deck from agent-generated slide fragments."""
    slides_html = "\n".join(page.html_fragment for page in pages)
    return f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{html.escape(content_plan.title)}</title>
  <style>
    * {{ box-sizing: border-box; }}
    body {{ margin: 0; background: #20242d; font-family: Inter, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif; }}
    .deck {{ display: grid; gap: 32px; padding: 32px; }}
    .slide {{
      width: {template.viewport_width}px;
      height: {template.viewport_height}px;
      position: relative;
      overflow: hidden;
      background: #fff;
      color: #111827;
    }}
    img {{ max-width: 100%; display: block; }}
  </style>
</head>
<body data-deck-template="{html.escape(template.template_id)}" data-page-generation-mode="llm_agent_html">
  <main class="deck">
{slides_html}
  </main>
</body>
</html>
"""


async def _render_html_deck_screenshots(
    *,
    html_path: Path,
    content_plan: DeckContentPlan,
    template: HtmlTemplatePackage,
    preview_dir: Path,
) -> list[Path]:
    """Render browser screenshots of generated HTML slides."""
    try:
        from playwright.async_api import async_playwright
    except ImportError:
        return []

    preview_paths: list[Path] = []
    try:
        async with async_playwright() as playwright:
            browser = await playwright.chromium.launch()
            try:
                page = await browser.new_page(
                    viewport={
                        "width": template.viewport_width,
                        "height": template.viewport_height,
                    }
                )
                await page.goto(html_path.as_uri(), wait_until="networkidle")
                for slide_plan in content_plan.pages:
                    locator = page.locator(f'[data-slide-number="{slide_plan.slide_number:02d}"]').first
                    preview_path = preview_dir / f"slide_{slide_plan.slide_number:03d}.png"
                    await locator.screenshot(path=str(preview_path))
                    preview_paths.append(preview_path)
            finally:
                await browser.close()
    except Exception:
        return []
    return preview_paths


def _copy_state(state: Any) -> dict[str, Any]:
    """Return a deep copy of an ADK state object or plain dict."""
    if hasattr(state, "to_dict"):
        return copy.deepcopy(state.to_dict())
    return copy.deepcopy(dict(state))


def _resolve_child_artifact_service(
    *,
    tool_context: ToolContext,
    fallback_service: BaseArtifactService,
) -> BaseArtifactService:
    """Pick the artifact service for an internal route runner."""
    required_methods = ("save_artifact", "load_artifact", "list_artifacts")
    if all(hasattr(tool_context, method_name) for method_name in required_methods):
        return ToolContextArtifactService(tool_context)
    return fallback_service


def _build_child_runner(
    *,
    agent: BaseAgent,
    app_name: str,
    session_service: InMemorySessionService,
    artifact_service: BaseArtifactService,
    invocation_context: Any,
) -> Runner:
    """Create a child ADK runner for the HTML page-generation agent."""
    child_plugins = getattr(getattr(invocation_context, "plugin_manager", None), "plugins", None)
    runner_kwargs = {
        "app_name": app_name,
        "session_service": session_service,
        "artifact_service": artifact_service,
        "memory_service": InMemoryMemoryService(),
        "credential_service": getattr(invocation_context, "credential_service", None),
    }
    if child_plugins:
        runner_kwargs["app"] = App(
            name=app_name,
            root_agent=agent,
            plugins=list(child_plugins),
        )
    else:
        runner_kwargs["agent"] = agent
    return Runner(**runner_kwargs)


def _append_html_page_generation_warning(state: Any, warning: str) -> None:
    """Append one HTML page-generation warning to state."""
    clean_warning = str(warning or "").strip()
    if not clean_warning:
        return
    warnings = list(state.get(HTML_PAGE_GENERATION_WARNINGS_KEY) or [])
    warnings.append(clean_warning)
    state[HTML_PAGE_GENERATION_WARNINGS_KEY] = warnings


def _render_html_deck(content_plan: DeckContentPlan, template: HtmlTemplatePackage) -> str:
    """Render one complete static HTML deck."""
    slides_html = "\n".join(_render_html_slide(page, template) for page in content_plan.pages)
    theme = _html_theme_values(content_plan, template)
    return f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{html.escape(content_plan.title)}</title>
  <style>
    :root {{
      --bg: {theme['bg']};
      --ink: {theme['ink']};
      --muted: {theme['muted']};
      --panel: {theme['panel']};
      --accent: {theme['accent']};
      --accent-soft: {theme['accent_soft']};
      --line: {theme['line']};
    }}
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      background: {theme['body_bg']};
      color: var(--ink);
      font-family: Inter, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
    }}
    .deck {{
      display: grid;
      gap: 32px;
      padding: 32px;
    }}
    .slide {{
      width: {template.viewport_width}px;
      height: {template.viewport_height}px;
      background: var(--bg);
      border: {theme['slide_border']};
      border-radius: {theme['slide_radius']};
      overflow: hidden;
      position: relative;
      padding: 58px 72px;
      display: flex;
      flex-direction: column;
      justify-content: space-between;
    }}
    .slide::after {{
      content: attr(data-slide-number);
      position: absolute;
      right: 34px;
      bottom: 24px;
      color: #98a2b3;
      font-size: 18px;
    }}
    .eyebrow {{
      color: var(--accent);
      font-size: 17px;
      font-weight: 700;
      text-transform: uppercase;
      letter-spacing: 0;
      margin-bottom: 18px;
    }}
    h1 {{
      font-size: 56px;
      line-height: 1.06;
      margin: 0;
      max-width: 900px;
    }}
    h2 {{
      font-size: 40px;
      line-height: 1.12;
      margin: 0 0 24px;
      max-width: 960px;
    }}
    .takeaway {{
      color: var(--muted);
      font-size: 24px;
      line-height: 1.35;
      max-width: 900px;
      margin-top: 20px;
    }}
    .content-grid {{
      display: grid;
      grid-template-columns: repeat(2, minmax(0, 1fr));
      gap: 22px;
      margin-top: 24px;
    }}
    .block {{
      background: var(--panel);
      border: 1px solid var(--line);
      border-radius: {theme['block_radius']};
      padding: 22px 24px;
      min-height: 130px;
    }}
    .block-title {{
      font-size: 21px;
      font-weight: 760;
      margin-bottom: 10px;
    }}
    .block-body {{
      font-size: 18px;
      line-height: 1.38;
      color: var(--muted);
    }}
    .toc-list {{
      display: grid;
      gap: 16px;
      margin-top: 30px;
      max-width: 820px;
    }}
    .toc-item {{
      display: grid;
      grid-template-columns: 54px 1fr;
      gap: 16px;
      align-items: start;
      padding: 18px 0;
      border-top: 1px solid var(--line);
    }}
    .toc-index {{
      color: var(--accent);
      font-size: 24px;
      font-weight: 780;
    }}
    .visual-band {{
      height: 92px;
      border-radius: {theme['block_radius']};
      background: {theme['visual_band']};
      margin-top: 34px;
    }}
    .asset-frame {{
      height: 160px;
      border-radius: 8px;
      border: 1px solid var(--line);
      background: var(--panel);
      margin: 26px 0 0;
      overflow: hidden;
    }}
    .asset-frame img {{
      width: 100%;
      height: 100%;
      object-fit: cover;
      display: block;
    }}
  </style>
</head>
<body data-deck-template="{html.escape(template.template_id)}" data-page-generation-mode="{theme['page_generation_mode']}">
  <main class="deck">
{slides_html}
  </main>
</body>
</html>
"""


def _html_theme_values(content_plan: DeckContentPlan, template: HtmlTemplatePackage) -> dict[str, str]:
    """Return CSS values for system-template or no-template page generation."""
    if template.template_id != _FREE_DESIGN_TEMPLATE_ID:
        return {
            "bg": "#f7f8fb",
            "ink": "#172033",
            "muted": "#667085",
            "panel": "#ffffff",
            "accent": "#2457d6",
            "accent_soft": "#e8eefc",
            "line": "#dbe2ef",
            "body_bg": "#d7dce8",
            "slide_border": "1px solid var(--line)",
            "slide_radius": "0",
            "block_radius": "8px",
            "visual_band": "linear-gradient(90deg, var(--accent), #43a6ff)",
            "page_generation_mode": "system_template",
        }
    if _looks_child_friendly_plan(content_plan):
        return {
            "bg": "#fff8e7",
            "ink": "#243049",
            "muted": "#5f6472",
            "panel": "#ffffff",
            "accent": "#ff6b57",
            "accent_soft": "#ffe6d6",
            "line": "#ffd39e",
            "body_bg": "#ffe8f0",
            "slide_border": "0",
            "slide_radius": "22px",
            "block_radius": "18px",
            "visual_band": "linear-gradient(90deg, #ff7a59, #ffd166, #4ecdc4)",
            "page_generation_mode": _FREE_DESIGN_TEMPLATE_ID,
        }
    return {
        "bg": "#fbfbf7",
        "ink": "#191b1f",
        "muted": "#5d6470",
        "panel": "#ffffff",
        "accent": "#0f9f8f",
        "accent_soft": "#dff8f4",
        "line": "#dce7e4",
        "body_bg": "#e9efe9",
        "slide_border": "0",
        "slide_radius": "18px",
        "block_radius": "14px",
        "visual_band": "linear-gradient(90deg, #0f9f8f, #8fd14f)",
        "page_generation_mode": _FREE_DESIGN_TEMPLATE_ID,
    }


def _looks_child_friendly_plan(content_plan: DeckContentPlan) -> bool:
    """Infer whether the free-design deck should use child-friendly styling."""
    text = " ".join(
        [
            content_plan.title,
            content_plan.core_narrative,
            *[page.title for page in content_plan.pages],
            *[page.key_takeaway for page in content_plan.pages],
        ]
    ).lower()
    return any(
        keyword in text
        for keyword in (
            "幼儿园",
            "小朋友",
            "儿童",
            "可爱",
            "apple",
            "cat",
            "dog",
            "ball",
            "sun",
            "kindergarten",
            "kid",
            "child",
        )
    )


def _render_html_slide(page: DeckPagePlan, template: HtmlTemplatePackage) -> str:
    """Render one HTML slide from a page plan."""
    layout = template.page_types.get(page.page_type, "chapter-content")
    blocks = page.content_blocks or [{"title": "Core point", "body": page.key_takeaway}]
    block_html = "\n".join(_render_content_block(block) for block in blocks[:4])
    toc_html = ""
    if page.page_type == "toc":
        toc_html = _render_toc_list(blocks)
    content = toc_html or f'<div class="content-grid">{block_html}</div>'
    visual = _render_page_asset_html(page) or ('<div class="visual-band"></div>' if page.asset_intent else "")
    return f"""    <section class="slide slide-{html.escape(page.page_type)}" data-layout="{html.escape(layout)}" data-slide-number="{page.slide_number:02d}">
      <div>
        <div class="eyebrow">{html.escape(page.page_type.replace("_", " "))}</div>
        <h2>{html.escape(page.title)}</h2>
        <div class="takeaway">{html.escape(page.key_takeaway)}</div>
        {content}
      </div>
      {visual}
    </section>"""


def _render_page_asset_html(page: DeckPagePlan) -> str:
    """Render the first ready slide asset as an HTML image."""
    asset = _first_ready_asset(page)
    if asset is None:
        return ""
    image_path = _resolve_asset_file(asset)
    if image_path is None:
        return ""
    src = image_path.as_uri()
    alt = asset.alt or asset.description or page.title
    return f'<figure class="asset-frame"><img src="{html.escape(src)}" alt="{html.escape(alt)}" /></figure>'


def _render_content_block(block: dict[str, Any]) -> str:
    """Render one content block for a slide."""
    title = str(block.get("title") or block.get("heading") or "Point").strip()
    body = str(block.get("body") or block.get("text") or block.get("description") or "").strip()
    return f"""<div class="block">
  <div class="block-title">{html.escape(title)}</div>
  <div class="block-body">{html.escape(body)}</div>
</div>"""


def _render_toc_list(blocks: list[dict[str, Any]]) -> str:
    """Render a table of contents from planned chapter blocks."""
    items = []
    for index, block in enumerate(blocks[:6], start=1):
        title = str(block.get("title") or block.get("heading") or f"Chapter {index}").strip()
        body = str(block.get("body") or block.get("text") or block.get("description") or "").strip()
        items.append(
            f"""        <div class="toc-item"><div class="toc-index">{index:02d}</div><div><div class="block-title">{html.escape(title)}</div><div class="block-body">{html.escape(body)}</div></div></div>"""
        )
    return "      <div class=\"toc-list\">\n" + "\n".join(items) + "\n      </div>"


def _render_previews(
    content_plan: DeckContentPlan,
    template: HtmlTemplatePackage,
    preview_dir: Path,
) -> list[Path]:
    """Render simple slide preview PNGs for screenshot-style PPTX export."""
    preview_paths: list[Path] = []
    for page in content_plan.pages:
        image = Image.new("RGB", (template.viewport_width, template.viewport_height), "#F7F8FB")
        draw = ImageDraw.Draw(image)
        _draw_slide_preview(image, draw, page, template)
        preview_path = preview_dir / f"slide_{page.slide_number:03d}.png"
        image.save(preview_path)
        preview_paths.append(preview_path)
    return preview_paths


def _draw_slide_preview(
    image: Image.Image,
    draw: ImageDraw.ImageDraw,
    page: DeckPagePlan,
    template: HtmlTemplatePackage,
) -> None:
    """Draw one preview slide with deterministic typography and layout."""
    width = template.viewport_width
    height = template.viewport_height
    margin_x = 72 if template.aspect_ratio == "16:9" else 62
    top = 54
    accent = "#2457D6"
    ink = "#172033"
    muted = "#667085"
    line = "#DBE2EF"
    panel = "#FFFFFF"

    title_font = _load_font(42 if template.aspect_ratio == "16:9" else 36, bold=True)
    body_font = _load_font(21 if template.aspect_ratio == "16:9" else 18)
    label_font = _load_font(16, bold=True)
    small_font = _load_font(15)

    draw.text((margin_x, top), page.page_type.replace("_", " ").upper(), fill=accent, font=label_font)
    title_y = top + 38
    for line_text in _wrap_text(page.title, width=26 if template.aspect_ratio == "16:9" else 22)[:2]:
        draw.text((margin_x, title_y), line_text, fill=ink, font=title_font)
        title_y += 50

    takeaway_y = title_y + 14
    for line_text in _wrap_text(page.key_takeaway, width=42 if template.aspect_ratio == "16:9" else 34)[:3]:
        draw.text((margin_x, takeaway_y), line_text, fill=muted, font=body_font)
        takeaway_y += 30

    blocks = page.content_blocks or [{"title": "Core point", "body": page.key_takeaway}]
    grid_top = min(takeaway_y + 36, height - 300)
    card_gap = 22
    card_width = int((width - margin_x * 2 - card_gap) / 2)
    card_height = 128 if template.aspect_ratio == "16:9" else 112
    for index, block in enumerate(blocks[:4]):
        col = index % 2
        row = index // 2
        x = margin_x + col * (card_width + card_gap)
        y = grid_top + row * (card_height + card_gap)
        draw.rounded_rectangle((x, y, x + card_width, y + card_height), radius=8, fill=panel, outline=line)
        block_title = str(block.get("title") or block.get("heading") or "Point")
        block_body = str(block.get("body") or block.get("text") or block.get("description") or "")
        draw.text((x + 22, y + 18), _truncate(block_title, 30), fill=ink, font=label_font)
        body_y = y + 50
        for line_text in _wrap_text(block_body, width=34)[:2]:
            draw.text((x + 22, body_y), line_text, fill=muted, font=small_font)
            body_y += 23

    band_y = height - 104
    if not _paste_preview_asset(image, draw, page, (margin_x, band_y - 72, width - margin_x, band_y + 58)):
        draw.rounded_rectangle((margin_x, band_y, width - margin_x, band_y + 58), radius=8, fill=accent)
        draw.text((margin_x + 20, band_y + 17), page.asset_intent or "HTML route preview", fill="#FFFFFF", font=small_font)
    draw.text((width - margin_x - 38, height - 40), f"{page.slide_number:02d}", fill="#98A2B3", font=small_font)


def _paste_preview_asset(
    image: Image.Image,
    draw: ImageDraw.ImageDraw,
    page: DeckPagePlan,
    box: tuple[int, int, int, int],
) -> bool:
    """Paste a ready slide asset into the PNG preview."""
    asset = _first_ready_asset(page)
    image_path = _resolve_asset_file(asset) if asset is not None else None
    if image_path is None:
        return False
    try:
        with Image.open(image_path) as asset_image:
            asset_image = asset_image.convert("RGB")
            max_width = max(1, box[2] - box[0])
            max_height = max(1, box[3] - box[1])
            asset_image.thumbnail((max_width, max_height))
            paste_x = box[0] + (max_width - asset_image.width) // 2
            paste_y = box[1] + (max_height - asset_image.height) // 2
            draw.rounded_rectangle(box, radius=8, fill="#FFFFFF", outline="#DBE2EF")
            image.paste(asset_image, (paste_x, paste_y))
        return True
    except Exception:
        return False


def _export_previews_to_pptx(preview_paths: list[Path], pptx_path: Path, *, aspect_ratio: str) -> None:
    """Create a screenshot-style PPTX from preview images."""
    prs = Presentation()
    if aspect_ratio == "4:3":
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    for preview_path in preview_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(preview_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
    prs.save(pptx_path)


def _export_native_pptx(
    content_plan: DeckContentPlan,
    pptx_path: Path,
    template: HtmlTemplatePackage,
) -> None:
    """Create an editable PPTX from the HTML route content plan."""
    prs = Presentation()
    slide_width, slide_height = _configure_slide_size(prs, template.aspect_ratio)
    blank_layout = prs.slide_layouts[6]
    for page in content_plan.pages:
        slide = prs.slides.add_slide(blank_layout)
        _draw_native_slide(slide, page, slide_width=slide_width, slide_height=slide_height)
    prs.save(pptx_path)


def _configure_slide_size(prs: Presentation, aspect_ratio: str) -> tuple[float, float]:
    """Configure slide size and return dimensions in inches."""
    if aspect_ratio == "4:3":
        width, height = 10.0, 7.5
    else:
        width, height = 13.333, 7.5
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    return width, height


def _draw_native_slide(
    slide: Any,
    page: DeckPagePlan,
    *,
    slide_width: float,
    slide_height: float,
) -> None:
    """Draw one editable PPTX slide from a page plan."""
    _set_slide_background(slide)
    margin_x = 0.72 if slide_width > 11 else 0.58
    top = 0.48
    _add_text_box(
        slide,
        page.page_type.replace("_", " ").upper(),
        margin_x,
        top,
        3.2,
        0.26,
        font_size=10,
        bold=True,
        color=_COLOR_ACCENT,
    )

    if page.page_type == "cover":
        _draw_cover_slide(slide, page, slide_width, slide_height, margin_x)
    elif page.page_type == "toc":
        _draw_toc_slide(slide, page, slide_width, margin_x)
    elif page.page_type == "chapter_start":
        _draw_chapter_start_slide(slide, page, slide_width, slide_height, margin_x)
    else:
        _draw_content_slide(slide, page, slide_width, slide_height, margin_x)

    _draw_footer(slide, page, slide_width, slide_height, margin_x)


def _draw_cover_slide(
    slide: Any,
    page: DeckPagePlan,
    slide_width: float,
    slide_height: float,
    margin_x: float,
) -> None:
    """Draw an editable cover slide."""
    title_width = min(9.3, slide_width - margin_x * 2)
    _add_text_box(slide, page.title, margin_x, 1.35, title_width, 1.35, font_size=34, bold=True)
    _add_text_box(slide, page.key_takeaway, margin_x, 3.0, title_width, 0.9, font_size=18, color=_COLOR_MUTED)
    _draw_asset_or_visual_band(slide, page, margin_x, slide_height - 1.55, slide_width - margin_x * 2, 0.82)
    _draw_content_cards(slide, page.content_blocks[:2], margin_x, 4.18, slide_width - margin_x * 2, 1.1)


def _draw_toc_slide(slide: Any, page: DeckPagePlan, slide_width: float, margin_x: float) -> None:
    """Draw an editable table-of-contents slide."""
    _add_text_box(slide, page.title, margin_x, 1.04, slide_width - margin_x * 2, 0.55, font_size=28, bold=True)
    _add_text_box(slide, page.key_takeaway, margin_x, 1.7, slide_width - margin_x * 2, 0.55, font_size=15, color=_COLOR_MUTED)
    y = 2.55
    for index, block in enumerate((page.content_blocks or [])[:5], start=1):
        title = _block_title(block, fallback=f"Chapter {index}")
        body = _block_body(block)
        _add_text_box(slide, f"{index:02d}", margin_x, y + 0.06, 0.45, 0.3, font_size=15, bold=True, color=_COLOR_ACCENT)
        _add_line(slide, margin_x + 0.72, y, slide_width - margin_x, y)
        _add_text_box(slide, title, margin_x + 0.72, y + 0.12, 4.5, 0.32, font_size=16, bold=True)
        _add_text_box(slide, body, margin_x + 5.2, y + 0.12, slide_width - margin_x - 5.2, 0.44, font_size=11, color=_COLOR_MUTED)
        y += 0.72


def _draw_chapter_start_slide(
    slide: Any,
    page: DeckPagePlan,
    slide_width: float,
    slide_height: float,
    margin_x: float,
) -> None:
    """Draw an editable chapter divider slide."""
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(margin_x),
        Inches(1.25),
        Inches(0.18),
        Inches(4.8),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _COLOR_ACCENT
    accent.line.fill.background()
    _add_text_box(slide, page.title, margin_x + 0.45, 1.55, slide_width - margin_x * 2 - 0.45, 0.95, font_size=32, bold=True)
    _add_text_box(slide, page.key_takeaway, margin_x + 0.45, 2.78, slide_width - margin_x * 2 - 0.45, 0.82, font_size=18, color=_COLOR_MUTED)
    _draw_content_cards(slide, page.content_blocks[:2], margin_x + 0.45, 4.1, slide_width - margin_x * 2 - 0.45, 1.25)
    _draw_asset_or_visual_band(slide, page, margin_x, slide_height - 1.45, slide_width - margin_x * 2, 0.72)


def _draw_content_slide(
    slide: Any,
    page: DeckPagePlan,
    slide_width: float,
    slide_height: float,
    margin_x: float,
) -> None:
    """Draw a standard editable content slide."""
    _add_text_box(slide, page.title, margin_x, 0.95, slide_width - margin_x * 2, 0.65, font_size=25, bold=True)
    _add_text_box(slide, page.key_takeaway, margin_x, 1.72, slide_width - margin_x * 2, 0.68, font_size=14, color=_COLOR_MUTED)
    _draw_content_cards(slide, page.content_blocks[:4], margin_x, 2.65, slide_width - margin_x * 2, 2.55)
    if not _draw_asset_or_visual_band(slide, page, margin_x, slide_height - 1.36, slide_width - margin_x * 2, 0.72):
        _add_text_box(slide, page.asset_intent, margin_x + 0.18, slide_height - 1.22, slide_width - margin_x * 2 - 0.36, 0.25, font_size=9, color=_COLOR_WHITE)


def _draw_content_cards(
    slide: Any,
    blocks: list[dict[str, Any]],
    x: float,
    y: float,
    width: float,
    height: float,
) -> None:
    """Draw editable content cards for slide blocks."""
    card_blocks = blocks or [{"title": "Core message", "body": ""}]
    gap = 0.22
    rows = 2 if len(card_blocks) > 2 else 1
    cols = 2 if len(card_blocks) > 1 else 1
    card_width = (width - gap * (cols - 1)) / cols
    card_height = (height - gap * (rows - 1)) / rows
    for index, block in enumerate(card_blocks[:4]):
        col = index % cols
        row = index // cols
        card_x = x + col * (card_width + gap)
        card_y = y + row * (card_height + gap)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(card_x),
            Inches(card_y),
            Inches(card_width),
            Inches(card_height),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _COLOR_PANEL
        card.line.color.rgb = _COLOR_LINE
        _add_text_box(slide, _block_title(block), card_x + 0.18, card_y + 0.15, card_width - 0.36, 0.24, font_size=12, bold=True)
        _add_text_box(slide, _block_body(block), card_x + 0.18, card_y + 0.48, card_width - 0.36, card_height - 0.58, font_size=10, color=_COLOR_MUTED)


def _draw_visual_band(slide: Any, x: float, y: float, width: float, height: float) -> None:
    """Draw a simple editable accent band."""
    band = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = _COLOR_ACCENT
    band.line.fill.background()
    cap = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x + width * 0.72),
        Inches(y),
        Inches(width * 0.28),
        Inches(height),
    )
    cap.fill.solid()
    cap.fill.fore_color.rgb = _COLOR_ACCENT_2
    cap.line.fill.background()


def _draw_asset_or_visual_band(
    slide: Any,
    page: DeckPagePlan,
    x: float,
    y: float,
    width: float,
    height: float,
) -> bool:
    """Draw a ready slide image asset, falling back to the editable visual band."""
    asset = _first_ready_asset(page)
    image_path = _resolve_asset_file(asset) if asset is not None else None
    if image_path is None:
        _draw_visual_band(slide, x, y, width, height)
        return False

    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = _COLOR_PANEL
    frame.line.color.rgb = _COLOR_LINE

    pic_x, pic_y, pic_width, pic_height = _fit_picture_box(image_path, x, y, width, height)
    slide.shapes.add_picture(
        str(image_path),
        Inches(pic_x),
        Inches(pic_y),
        width=Inches(pic_width),
        height=Inches(pic_height),
    )
    return True


def _fit_picture_box(
    image_path: Path,
    x: float,
    y: float,
    width: float,
    height: float,
) -> tuple[float, float, float, float]:
    """Fit an image inside a PowerPoint box while preserving aspect ratio."""
    try:
        with Image.open(image_path) as image:
            image_width, image_height = image.size
    except Exception:
        return x, y, width, height
    if image_width <= 0 or image_height <= 0:
        return x, y, width, height
    scale = min(width / image_width, height / image_height)
    fitted_width = image_width * scale
    fitted_height = image_height * scale
    return (
        x + (width - fitted_width) / 2,
        y + (height - fitted_height) / 2,
        fitted_width,
        fitted_height,
    )


def _draw_footer(
    slide: Any,
    page: DeckPagePlan,
    slide_width: float,
    slide_height: float,
    margin_x: float,
) -> None:
    """Draw a compact editable footer."""
    _add_line(slide, margin_x, slide_height - 0.58, slide_width - margin_x, slide_height - 0.58)
    _add_text_box(slide, page.chapter or page.page_type.replace("_", " "), margin_x, slide_height - 0.45, 4.0, 0.22, font_size=8, color=_COLOR_MUTED)
    _add_text_box(slide, f"{page.slide_number:02d}", slide_width - margin_x - 0.42, slide_height - 0.45, 0.42, 0.22, font_size=8, color=_COLOR_MUTED)


def _set_slide_background(slide: Any) -> None:
    """Apply the route background to one PPTX slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = _COLOR_BG


def _add_text_box(
    slide: Any,
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    font_size: int,
    bold: bool = False,
    color: RGBColor = _COLOR_INK,
) -> Any:
    """Add an editable text box and return the shape."""
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    paragraph = frame.paragraphs[0]
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.05
    run = paragraph.add_run()
    run.text = str(text or "")
    run.font.name = _FONT_FAMILY
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def _add_line(slide: Any, x1: float, y1: float, x2: float, y2: float) -> Any:
    """Add a thin editable divider line."""
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = _COLOR_LINE
    line.line.width = Pt(0.75)
    return line


def _block_title(block: dict[str, Any], *, fallback: str = "Point") -> str:
    """Return the display title for one content block."""
    return str(block.get("title") or block.get("heading") or fallback).strip()


def _block_body(block: dict[str, Any]) -> str:
    """Return the display body for one content block."""
    return str(block.get("body") or block.get("text") or block.get("description") or "").strip()


def _first_ready_asset(page: DeckPagePlan) -> DeckPageAsset | None:
    """Return the first ready image asset for a page."""
    for asset in page.assets:
        validated_asset = DeckPageAsset.model_validate(asset)
        if validated_asset.status == "ready" and validated_asset.path:
            return validated_asset
    return None


def _resolve_asset_file(asset: DeckPageAsset | None) -> Path | None:
    """Resolve one ready asset path into an existing workspace file."""
    if asset is None:
        return None
    try:
        path = resolve_workspace_path(asset.path)
    except Exception:
        return None
    if not path.exists() or not path.is_file():
        return None
    return path


def _validate_html_route_output(
    *,
    content_plan: DeckContentPlan,
    html_path: Path,
    preview_paths: list[Path],
    pptx_path: Path,
    pptx_strategy: str,
) -> dict[str, Any]:
    """Validate the core HTML route artifacts."""
    planned_assets = [
        DeckPageAsset.model_validate(asset)
        for page in content_plan.pages
        for asset in page.assets
    ]
    planned_ready_assets = [asset for asset in planned_assets if asset.status == "ready"]
    ready_asset_count = len(planned_ready_assets)
    asset_paths_exist = all(_resolve_asset_file(asset) is not None for asset in planned_ready_assets)
    checks = {
        "html_exists": html_path.exists() and html_path.stat().st_size > 0,
        "preview_count_matches": len(preview_paths) == len(content_plan.pages),
        "all_previews_exist": all(path.exists() and path.stat().st_size > 0 for path in preview_paths),
        "pptx_exists": pptx_path.exists() and pptx_path.stat().st_size > 0,
        "pptx_slide_count_matches": False,
        "pptx_contains_editable_text": False,
        "pptx_titles_present": False,
        "asset_paths_exist": asset_paths_exist,
        "pptx_ready_assets_rendered": ready_asset_count == 0,
    }
    editable_text_shape_count = 0
    pptx_picture_count = 0
    pptx_text = ""
    if checks["pptx_exists"]:
        prs = Presentation(str(pptx_path))
        checks["pptx_slide_count_matches"] = len(prs.slides) == len(content_plan.pages)
        pptx_picture_count = sum(
            1
            for slide in prs.slides
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        )
        editable_text_shape_count = sum(
            1
            for slide in prs.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and str(getattr(shape, "text", "") or "").strip()
        )
        pptx_text = "\n".join(
            str(getattr(shape, "text", "") or "")
            for slide in prs.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        checks["pptx_contains_editable_text"] = editable_text_shape_count > 0 and pptx_strategy != "screenshot"
        checks["pptx_titles_present"] = pptx_strategy != "screenshot" and all(page.title in pptx_text for page in content_plan.pages)
        checks["pptx_ready_assets_rendered"] = ready_asset_count == 0 or pptx_picture_count >= min(
            ready_asset_count,
            len(content_plan.pages),
        )
    status = "pass" if all(checks.values()) else "failed"
    return {
        "status": status,
        "checks": checks,
        "route": "html",
        "slide_count": len(content_plan.pages),
        "pptx_editable_text_shape_count": editable_text_shape_count,
        "pptx_picture_count": pptx_picture_count,
        "ready_asset_count": ready_asset_count,
    }


def _load_font(size: int, *, bold: bool = False) -> ImageFont.ImageFont:
    """Load a font that can render Chinese and Latin text on common hosts."""
    candidates = [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
        "/Library/Fonts/Arial Unicode.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
    ]
    for font_path in candidates:
        try:
            return ImageFont.truetype(font_path, size=size)
        except OSError:
            continue
    return ImageFont.load_default()


def _wrap_text(text: str, *, width: int) -> list[str]:
    """Wrap Latin text by words and CJK text by character count."""
    clean_text = " ".join(str(text or "").split())
    if not clean_text:
        return [""]
    if " " in clean_text:
        return textwrap.wrap(clean_text, width=width) or [clean_text]
    return [clean_text[index : index + width] for index in range(0, len(clean_text), width)]


def _truncate(text: str, max_chars: int) -> str:
    """Truncate one display string to a compact length."""
    clean_text = str(text or "").strip()
    if len(clean_text) <= max_chars:
        return clean_text
    return f"{clean_text[: max_chars - 1]}..."
